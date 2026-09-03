"""PowerPoint deck export (module split #18, 2026-07-10).

CAVEAT — this function is FUSED to engine runtime state: it reads ~55 engine
globals (24 pipeline artifacts via globals().get + config/derived + helper fns).
It is NOT decoupled; the engine syncs all 55 into this module (_sync_ppt_export)
immediately before each call. Moving it out of the monolith is a line-count win,
not an architectural one — a true decoupling would thread state through a context
object. Validated only by live deck generation (before/after diff).
"""
from __future__ import annotations

import os
import io
import json
import math
from datetime import datetime

import numpy as np
import pandas as pd
import yfinance as yf
from dateutil.relativedelta import relativedelta

import matplotlib.pyplot as plt
import matplotlib.dates as mdates
import matplotlib.ticker as mtick

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN

# All other names (config, runtime pipeline state, helper functions) are injected
# by the engine's _sync_ppt_export() before each call — see _PPT_EXPORT_INJECT.

# Longest run of missing actual-NAV days to bridge when PLOTTING the live line.
# A day with no broker snapshot is NaN once reindexed onto the price panel, and
# matplotlib breaks the line at every NaN. Three days covers a long weekend or
# a missed morning; anything longer is the pipeline being down and must stay
# visible on the chart. Affects the plot only — never stored NAV, never the
# return table.
NAV_GAP_BRIDGE_DAYS = 3


def rebase_to(obj, when):
    """Normalise a Series/DataFrame so its value at `when` is 1.0.

    Every line on the performance slide must share an origin or the chart
    answers a question nobody asked. It did not: Actual NAV was rebased to its
    own first observation (the account's inception) while the strategy and the
    benchmarks were rebased to the chart window's start, a month earlier. The
    strategy therefore carried a head start — +5.19% over 25 May → 23 Jun on
    the live book's weights — before the NAV line began at zero, so the slide
    read as roughly five times more divergence than the drift table's actual
    -1.00% cumulative.

    Falls back to the first row when `when` is absent or unusable, which is the
    pre-inception behaviour and the right answer for a hypothetical run.
    """
    try:
        if when is None or when not in obj.index:
            return obj.div(obj.iloc[0])
        base = obj.ffill().loc[when]
        # A zero or missing base would blow the whole line up rather than
        # rebase it; the window start is a worse origin but a finite one.
        import numpy as _np
        if _np.ndim(base) == 0:
            if not _np.isfinite(base) or base == 0:
                return obj.div(obj.iloc[0])
        else:
            base = base.where(_np.isfinite(base) & (base != 0))
            if base.isna().all():
                return obj.div(obj.iloc[0])
            base = base.fillna(obj.iloc[0])
        return obj.div(base)
    except Exception:
        return obj.div(obj.iloc[0])


def bridge_short_gaps(series, max_days: int = NAV_GAP_BRIDGE_DAYS):
    """Interpolate across NaN runs of `max_days` or fewer. Returns (series, n).

    Fills a hole ENTIRELY or not at all. A plain `limit=` fills the first N days
    of a long outage and stops, drawing a line that marches confidently into
    nothing — a different lie from the one being fixed. `limit_area="inside"`
    keeps it from inventing NAV before the first observation or after the last.
    """
    import pandas as _pd

    if series is None or not isinstance(series, _pd.Series) or series.empty:
        return series, 0
    na = series.isna()
    if not na.any():
        return series, 0
    run = na.groupby((na != na.shift()).cumsum()).transform("sum")
    fillable = na & (run <= int(max_days))
    if not fillable.any():
        return series, 0
    out = series.where(~fillable,
                       series.interpolate(method="time", limit_area="inside"))
    # Count what was ACTUALLY filled, not what we hoped to fill. `fillable`
    # includes short NaN runs at the head or tail, but limit_area="inside"
    # rightly refuses those — so counting `fillable` over-reported by one on
    # every 10:20 run, where the series ends at yesterday's broker snapshot
    # while the panel already has today's row. That one-day trailing gap is
    # why the count read 7 against six explicable days on 08-27 and 09-02, and
    # why it never reproduced at midday or in the evening: by then the
    # snapshot exists and the trailing NaN is gone. The chart was always
    # right; only the number was wrong.
    return out, int((na & out.notna()).sum())


def export_to_ppt(results, trades, charts=None):
    """
    Generates a professional PowerPoint summary based on your custom template.
    """
    # Use the module-level APP_DIR rather than redefining from __file__:
    # under a PyInstaller frozen build, __file__ resolves to the _MEI* temp dir
    # where Portfolio_Optimiser.py is extracted, NOT where the template lives.
    # Template lives in .assets/ (hidden subdir) — falls back to repo root for
    # legacy installs where the move hasn't happened yet.
    _new_template = os.path.join(str(APP_DIR), ".assets", "PowerPoint_Template.pptx")
    _legacy_template = os.path.join(str(APP_DIR), "PowerPoint_Template.pptx")
    template_path = _new_template if os.path.exists(_new_template) else _legacy_template

    # Output path â€” always overwrite this file
    ppt_path = str(EXPORT_DIR / "Portfolio_Report.pptx")

    # Load your custom template
    prs = Presentation(template_path)
   
    # --- SLIDE 1: Title and TimeStamp ---
    slide = prs.slides[0] 
    
    # --- Title text box ---  
    if slide.shapes.title:
        slide.shapes.title.text = "Portfolio Performance Overview"
    
    # --- Timestamp text box ---
    now = datetime.now()
    timestamp = now.strftime("Last updated on the %d %B %Y at %I:%M %p")
    ts_box = slide.shapes.add_textbox(Cm(2.032), Cm(15.24), Cm(22.86), Cm(1.27))
    tf2 = ts_box.text_frame
    tf2.word_wrap = False
    p2 = tf2.add_paragraph()
    p2.text = timestamp
    p2.font.size = Pt(16)
    p2.alignment = PP_ALIGN.LEFT  # uses default colour/font

    # --- SLIDE 2: Trade Plan + Brokerage ---
    slide_layout = prs.slide_layouts[20]  # clean layout from your master
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    if slide.shapes.title:
        slide.shapes.title.text = "Trade Plan and Brokerage Overview"

    # --- Portfolio identity call-out (checkbox + label) ---
    plan_label = str(globals().get("TRADEPLAN_LABEL", "Unknown")).strip()
    
    # Show both options with one checked
    pl = plan_label.lower().replace("_", " ").strip()
    is_with = (pl == "with tilts" or pl == "with_tilts")
    is_ens = (pl == "ensemble")
    is_no = not (is_with or is_ens)
    # Position INSIDE the blue title ribbon. Adjusted 2026-06-20:
    # y=1.55 overlapped the title text; y=2.05 dropped below the ribbon.
    # y=1.85 nests cleanly under the title without overlap.
    box = slide.shapes.add_textbox(Cm(2.15), Cm(1.85), Cm(21.80), Cm(0.45))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = (
        f"{'[x]' if is_ens else '[ ]'} Ensemble    "
        f"{'[x]' if is_with else '[ ]'} With Tilts    "
        f"{'[x]' if is_no else '[ ]'} Optimised (No Tilts)    "
        f"|    Trade plan: {plan_label}"
    )
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    # --- Regime mix annotation (only when ensemble is the active plan) ---
    if is_ens:
        try:
            mix = globals().get("ensemble_mix_live", pd.Series(dtype=float))
            if isinstance(mix, pd.Series) and not mix.empty:
                # Compact mix string: "Def 3% · Modest 15% · Agg 27% · Bold 28% · Max 27%"
                _abbr = {
                    "Modest (SPY+0%)":      "Modest",
                    "Aggressive (SPY+5%)":  "Agg",
                    "Bold (SPY+10%)":       "Bold",
                    "Maximum (SPY+15%)":    "Max",
                    "Stretch (SPY+25%)":    "Stretch",
                }
                parts = []
                for n in ENSEMBLE_SLOT_NAMES:
                    if n in mix.index:
                        parts.append(f"{_abbr.get(n, n)} {float(mix[n])*100:.0f}%")
                mix_str = " · ".join(parts)
                # Sits in the blue ribbon directly under the plan label.
                # Adjusted 2026-06-20: y=2.05 was still too tight on plan_label;
                # y=2.35 gives proper line spacing while staying in the ribbon.
                mix_box = slide.shapes.add_textbox(Cm(2.15), Cm(2.35), Cm(21.80), Cm(0.45))
                tfm = mix_box.text_frame
                tfm.clear()
                tfm.word_wrap = False
                tfm.margin_left = 0
                tfm.margin_top = 0
                pm = tfm.paragraphs[0]
                pm.text = f"Regime mix today (rolling-12M Sortino softmax): {mix_str}"
                pm.font.size = Pt(10)
                pm.font.italic = True
                pm.font.color.rgb = RGBColor(255, 255, 255)
                pm.alignment = PP_ALIGN.LEFT
        except Exception as _e_mix:
            print(f"[pptx] Slide 2 regime mix annotation skipped: {_e_mix}")

    # --- LIVE TLH callout (this-run harvest swaps, baked into rebalance) ---
    # Phase 4: live TLH is injected before the trade plan is built so the
    # rebalance delta in `trades` already includes the harvest swaps. Surface
    # them here so the user can see exactly which lots are being swapped
    # and the loss being crystallised — separately from the backtest
    # cumulative scorecard below.
    try:
        _live_tlh_ppt = globals().get("LIVE_TLH_EVENTS", []) or []
        if _live_tlh_ppt:
            _live_loss = float(sum(e.get("loss_aud", 0.0) for e in _live_tlh_ppt))
            _live_pairs = ", ".join(
                f"{ev.get('ticker_sold', '?')}→{ev.get('ticker_bought', '?')}"
                for ev in _live_tlh_ppt[:5]
            )
            if len(_live_tlh_ppt) > 5:
                _live_pairs += f" (+{len(_live_tlh_ppt)-5} more)"
            _live_text = (
                f"LIVE TLH this run: {len(_live_tlh_ppt)} swap(s)  ·  "
                f"${_live_loss:,.0f} loss to realise  ·  baked into rebalance  ·  "
                f"{_live_pairs}"
            )
            _live_box = slide.shapes.add_textbox(Cm(1.10), Cm(13.50), Cm(23.80), Cm(0.50))
            _ltf = _live_box.text_frame
            _ltf.clear()
            _ltf.word_wrap = False
            _ltf.margin_left = 0
            _ltf.margin_top = 0
            _lp = _ltf.paragraphs[0]
            _lp.text = _live_text
            _lp.font.size = Pt(10)
            _lp.font.bold = True
            _lp.font.color.rgb = RGBColor(176, 0, 0)  # accent red — same as FF cutoff line
            _lp.alignment = PP_ALIGN.CENTER
            print(f"[pptx] Slide 3 live-TLH callout added: {len(_live_tlh_ppt)} swap(s)")
    except Exception as _e_live_tlh_ppt:
        print(f"[pptx] Slide 3 live-TLH callout skipped: {_e_live_tlh_ppt}")

    # --- TLH scorecard line (one-liner with cumulative harvest activity) ---
    # Positioned at the BOTTOM of slide 2 (below the Deferred Tax/Cash
    # callouts at y=14.73cm). Standard widescreen slide height ~19cm, so
    # y=16.7cm sits just above the bottom accent strip. Italic gray so it
    # reads as supporting metadata, not a headline. Earlier it sat between
    # the title ribbon and the trade plan table which crowded the layout —
    # moving it to the bottom matches the cleaned-up layout in the target
    # screenshot from 2026-06-19.
    try:
        _tlh_events_ppt = globals().get("oos_tlh_events", []) or []
        if TLH_ENABLED:
            if _tlh_events_ppt:
                _tlh_loss = float(sum(e.get("loss_aud", 0.0) for e in _tlh_events_ppt))
                _eff_st = _effective_cgt_rate(short_term=True)
                _eff_lt = _effective_cgt_rate(short_term=False)
                _tax_est = _tlh_loss * (_eff_st + _eff_lt) / 2.0
                _oos_rets = globals().get("oos_returns_daily", pd.Series(dtype=float))
                _yrs = max(len(_oos_rets) / ANNUAL_TRADING_DAYS, 1e-6) if isinstance(_oos_rets, pd.Series) else 1.0
                _bps = (_tax_est
                        / float(globals().get("_oos_starting_nav_aud") or 1_000_000.0)
                        / _yrs * 10_000)
                _tlh_text = (f"Tax-loss harvesting: {len(_tlh_events_ppt)} events over backtest  ·  "
                             f"${_tlh_loss:,.0f} loss realised  ·  "
                             f"~${_tax_est:,.0f} gross tax-saved est ({_bps:.0f} bps/yr drag offset; "
                             f"gross — net depends on FY-end netting)")
            else:
                _tlh_text = (f"Tax-loss harvesting: enabled, 0 events triggered  ·  "
                             f"threshold {TLH_MIN_LOSS_PCT*100:+.0f}%, "
                             f"${TLH_MIN_LOSS_AUD:.0f} min, {TLH_COOLDOWN_DAYS}d cooldown")
            tlh_box = slide.shapes.add_textbox(Cm(1.10), Cm(14.10), Cm(23.80), Cm(0.50))
            tft = tlh_box.text_frame
            tft.clear()
            tft.word_wrap = False
            tft.margin_left = 0
            tft.margin_top = 0
            pt_ = tft.paragraphs[0]
            pt_.text = _tlh_text
            pt_.font.size = Pt(10)
            pt_.font.italic = True
            pt_.font.color.rgb = RGBColor(60, 60, 60)
            pt_.alignment = PP_ALIGN.CENTER
            print(f"[pptx] Slide 2 TLH scorecard added (bottom): {len(_tlh_events_ppt)} events")
    except Exception as _e_tlh_ppt:
        print(f"[pptx] Slide 2 TLH scorecard skipped: {_e_tlh_ppt}")

    # --- Draw Trade Plan table ---
    if trades is not None and not trades.empty:
        # Resolve the trade-delta column defensively (handles legacy/mis-encoded headers).
        delta_col = None
        if "_trade_delta_col" in globals():
            try:
                delta_col = _trade_delta_col(trades)
            except Exception:
                delta_col = None
        if not delta_col:
            if "Delta Units" in trades.columns:
                delta_col = "Delta Units"
            else:
                for _c in trades.columns:
                    _cs = str(_c)
                    if _cs.endswith(" Units") and "delta" in _cs.lower():
                        delta_col = _c
                        break

        # Map existing columns to desired display names.
        # IMPORTANT: only choose one brokerage source to avoid duplicate 'Brokerage' columns.
        rename_map = {
            "Curr Units": "Current",
            "Target Units": "Target",
            "Last Px (AUD)": "Last Price",
            "Cash Flow (AUD)": "Cash Flow",
        }
        if "Brokerage" in trades.columns:
            rename_map["Brokerage"] = "Brokerage"
        elif "Brokerage (AUD)" in trades.columns:
            rename_map["Brokerage (AUD)"] = "Brokerage"
        if delta_col:
            rename_map[delta_col] = "Change"
    
        # Select, copy, and rename columns
        cols_needed = ["Security"] + list(rename_map.keys())
        cols_present = [c for c in cols_needed if c in trades.columns]
        df = trades[cols_present].copy()
        if "Security" not in df.columns and trades.index.name == "Security":
            df = trades.reset_index()[cols_present]
        df.rename(columns=rename_map, inplace=True)
        # Reorder to your final order but only for those that exist
        final_order = [c for c in ["Security","Current","Target","Change","Last Price","Brokerage","Cash Flow"] if c in df.columns]
        df = df[final_order]

        # --- Filter: hide universe rows with no position and no trade today ---
        # Show row if it has units (Current > 0 or Target > 0) OR is being
        # traded this run (Change != 0). Hides dormant universe members so the
        # table only shows what's actually in play.
        try:
            _cur_n = pd.to_numeric(df.get("Current"), errors="coerce").fillna(0)
            _tgt_n = pd.to_numeric(df.get("Target"), errors="coerce").fillna(0)
            _chg_n = pd.to_numeric(df.get("Change"), errors="coerce").fillna(0)
            _keep_mask = (_cur_n.abs() > 0) | (_tgt_n.abs() > 0) | (_chg_n.abs() > 0)
            _rows_before = len(df)
            df = df[_keep_mask].reset_index(drop=True)
            print(f"[pptx] Trade plan filtered {_rows_before} -> {len(df)} rows "
                  f"(dormant universe entries hidden)")
        except Exception as _e_filter:
            print(f"[pptx] Trade plan filter skipped: {_e_filter}")

        # --- Clean and format values ---
        df["Security"] = df["Security"].astype(str).str.replace(".AX", "", regex=False)
        for col in ["Last Price", "Cash Flow", "Brokerage"]:
            if col in df.columns:
                df[col] = (
                    pd.to_numeric(df[col], errors="coerce")
                    .round(2)
                    .apply(lambda x: f"-${abs(x):,.2f}" if x < 0 else f"${x:,.2f}")
                )
    
        # --- Determine if we split into two tables ---
        rows, cols = df.shape
        split = rows > 15
        half = math.ceil(rows / 2) if split else rows
        table_sets = [df.iloc[:half]] if not split else [df.iloc[:half], df.iloc[half:]]
        left_positions = [Cm(1.1), Cm(11.5)] if split else [Cm(1.1)]
    
        # --- Draw tables ---
        for idx, subdf in enumerate(table_sets):
            top = Cm(4.0)
        
            table_w = 14.0   # widened from 12.02 so "Current" header + 4-digit unit counts fit on one line
            gap = 0.8
            left_margin = 0.35
            left = Cm(left_margin + idx * (table_w + gap))
            width = Cm(table_w)
            height = Cm(6.94)
        
            table = slide.shapes.add_table(
                rows=subdf.shape[0] + 1,
                cols=subdf.shape[1],
                left=left,
                top=top,
                width=width,
                height=height
            ).table
        
            # Auto-fit widths
            _autofit_table_width(table, subdf, total_width_cm=table_w)
        
            # Disable wrapping
            for cell in table.iter_cells():
                cell.text_frame.word_wrap = False
                cell.text_frame.margin_left = 0
                cell.text_frame.margin_right = 0
                cell.text_frame.margin_top = 0
                cell.text_frame.margin_bottom = 0
        
            # Row height
            for r in range(len(table.rows)):
                table.rows[r].height = Cm(0.584)
        
            # Header
            for j, col_name in enumerate(subdf.columns):
                cell = table.cell(0, j)
                cell.text = col_name
                tf = cell.text_frame
                tf.auto_size = MSO_AUTO_SIZE.NONE
                p = tf.paragraphs[0]
                p.font.bold = True
                p.font.size = Pt(8)
                p.alignment = PP_ALIGN.CENTER
        
            # Data rows
            for i, (_, row) in enumerate(subdf.iterrows(), start=1):
                for j, val in enumerate(row):
                    cell = table.cell(i, j)
                    cell.text = str(val)
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(9)
                    p.font.bold = (j == 0)
                    p.alignment = PP_ALIGN.CENTER
    
        # --- Summary bar across top ---
        left = Cm(2.5)
        top = Cm(2.5)  # just below the title
        width = Cm(20.00)
        height = Cm(1.1)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = False
        tf.clear()
        
        # --- Fetch values ---
        total_portfolio = results.get("total_portfolio_value", 0)
        total_brokerage = results.get("total_brokerage", 0)
        net_invested = results.get("net_invested", 0)
        portfolio_change = results.get("portfolio_change", 0)
        net_invested_change = results.get("net_invested_change", 0)
        
        # --- Main summary line ---
        p = tf.add_paragraph()
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Text with separate runs for coloured numbers
        run1 = p.add_run()
        run1.text = f"Total Portfolio: ${total_portfolio:,.2f}"
        run1.font.size = Pt(14)
        run1.font.bold = True
        _add_change_run(p, portfolio_change)

        run2 = p.add_run()
        run2.text = f"     Total Brokerage: ${total_brokerage:,.2f}     "
        run2.font.size = Pt(14)
        run2.font.bold = True
        run2.font.color.rgb = RGBColor(0, 0, 0)

        run3 = p.add_run()
        run3.text = f"Net Invested: ${net_invested:,.2f}"
        run3.font.size = Pt(14)
        run3.font.bold = True
        _add_change_run(p, net_invested_change)

        # --- Slide 3: Portfolio vs Indices ---
        # --- Cash summary (derived from Trade Plan cash flows) ---
        try:
            cash_balance = 0.0
            if trades is not None and not trades.empty and "Cash Flow (AUD)" in trades.columns:
                cash_balance = float(results.get("cash_balance", 0.0))

            cash_box = slide.shapes.add_textbox(Cm(18.288), Cm(11.70), Cm(6.604), Cm(1.524))
            tfc = cash_box.text_frame
            tfc.clear()
            p = tfc.paragraphs[0]
            p.text = f"Cash: {cash_balance:,.0f} AUD"
            p.font.size = Pt(18)
            p.font.bold = True
            p.alignment = PP_ALIGN.RIGHT
        except Exception:
            pass

        # --- Deferred Tax callout (mirrors Cash). Sign convention:
        # positive = losses incurred → tax saved (carry forward to next FY)
        # negative = gains realised → tax owed this rebalance
        try:
            _cgt_owed = float(results.get("total_cgt", 0.0))
            _cgt_saved = float(results.get("loss_carry_forward_tax_aud", 0.0))
            _net_deferred = _cgt_saved - _cgt_owed
            _tax_box = slide.shapes.add_textbox(Cm(11.684), Cm(11.70), Cm(6.604), Cm(1.524))
            tft = _tax_box.text_frame
            tft.clear()
            pt = tft.paragraphs[0]
            if _net_deferred >= 0:
                pt.text = f"Deferred Tax: +{_net_deferred:,.0f} AUD"
                pt.font.color.rgb = RGBColor(0, 128, 0)  # green: tax saved
            else:
                pt.text = f"Deferred Tax: {_net_deferred:,.0f} AUD"
                pt.font.color.rgb = RGBColor(192, 0, 0)  # red: tax owed
            pt.font.size = Pt(18)
            pt.font.bold = True
            pt.alignment = PP_ALIGN.RIGHT
        except Exception as _e_dt:
            print(f"[pptx] Deferred Tax callout skipped: {_e_dt}")

        # --- Currency footnote (clarifies the Last Price / Cash Flow columns) ---
        # Prices come from px_aud and cash flows are AUD, but the table shows a
        # bare "$" across both .AX (AUD) and US (AUD-converted) names — label it
        # so the mixed magnitudes aren't misread as a currency error.
        try:
            _ccy_box = slide.shapes.add_textbox(Cm(1.10), Cm(11.15), Cm(16.0), Cm(0.5))
            _ctf = _ccy_box.text_frame
            _ctf.clear()
            _ctf.margin_left = 0
            _cp = _ctf.paragraphs[0]
            _cp.text = "All prices and cash flows shown in AUD (foreign holdings converted at the run's FX)."
            _cp.font.size = Pt(9)
            _cp.font.italic = True
            _cp.font.color.rgb = RGBColor(90, 90, 90)
        except Exception as _e_ccy:
            print(f"[pptx] Currency footnote skipped: {_e_ccy}")

        slide_layout = prs.slide_layouts[20]  # clean layout from your master
        slide = prs.slides.add_slide(slide_layout)

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = "Portfolio Performance"

        # --- Get 3-month portfolio + benchmarks ---
        slide3 = slide  # alias for clarity in the callout call below

        # Use the portfolio value series as the date anchor (NOT prices)
        pval_src = globals().get("portfolio_value_series", None)
        if pval_src is None and isinstance(charts, dict):
            pval_src = charts.get("portfolio_value_series", None)
        if pval_src is None:
            raise ValueError("portfolio_value_series is missing. Run Cell 15 first.")
        pval_all = pd.to_numeric(pd.Series(pval_src), errors="coerce").dropna().copy()
        pval_all.index = pd.to_datetime(pval_all.index).tz_localize(None)

        end_dt = pval_all.index.max()
        # IMPORTANT: use the SAME 3-month anchor as the return table below
        # (_period_total_return uses relativedelta(months=3)). A naive
        # timedelta(days=90) gives a different baseline by ~4 trading days,
        # which causes the chart's rightmost y-value to disagree with the
        # table's 3M column (e.g. With Tilts visually below SPY on chart but
        # above in table).
        start_dt = end_dt - relativedelta(months=3)

        # Date callout is deferred until after perf_df is built so it reflects
        # the ACTUAL first/last trading day rendered on the chart — otherwise
        # start_dt can land on a weekend and disagree with the x-axis by 1–2 days.

        pval = pval_all.loc[start_dt:end_dt].copy()
        pval = pval.ffill().bfill()

        benchmarks = ["^AORD", "^GSPC", "^IXIC"]

        # ONE benchmark download — long enough for the 3Y table column. The
        # chart slices this for its 3M window. Using two separate downloads
        # (short for chart, long for table) causes auto_adjust to back-adjust
        # dividends differently across the two windows, which makes the chart's
        # 3M return disagree with the table's 3M column for the same ticker.
        bench_long_start = end_dt - relativedelta(years=3, months=1)
        bench_long = yf.download(
            benchmarks,
            start=bench_long_start,
            end=(end_dt + pd.Timedelta(days=1)),
            progress=False,
            auto_adjust=True,
            threads=False
        )

        # Handle multi-index and clean
        if isinstance(bench_long.columns, pd.MultiIndex):
            bench_long = bench_long["Close"]
        else:
            if "Close" in bench_long.columns:
                bench_long = bench_long["Close"]

        bench_long.index = pd.to_datetime(bench_long.index).tz_localize(None)
        bench_long = bench_long.ffill().bfill()

        # Chart uses the 3M slice of the unified download
        bench = bench_long.reindex(pval.index).ffill().bfill()

        # --- Portfolio line: ACTUAL NAV from lots seed + fills log --------
        # Previously this was (pval / pval.iloc[0]) - 1.0 where pval came
        # from tgt_units_full × prices — a hypothetical of the no-tilts
        # target portfolio. That duplicated the Strategy (Ensemble) line's
        # purpose and hid the fact that the held lots (BEAR / BBUS / HBRD)
        # were actually defensive while the chart pretended they tracked
        # NASDAQ. Now reconstructs real NAV; line stubs in at LIVE_START.
        # Reuse the series the engine already built and validated. Computing
        # it again here is not just wasteful — this call passed neither
        # fx_usdaud nor statement_path, so it silently got the mixed-currency
        # seed-based path, failed validation at 4.97% and fell back to
        # broker-only while the engine's own call passed at 0.67%. The chart
        # was drawing the worse of two answers to the same question.
        _actual_nav_series_local = globals().get("LIVE_NAV_SERIES")
        if not isinstance(_actual_nav_series_local, pd.Series) or _actual_nav_series_local.empty:
            _actual_nav_series_local = pd.Series(dtype=float)
            try:
                # Flex XML if the token has refreshed it, else the manually
                # exported CSV. Imported here rather than added to the ~55
                # engine globals this module already has synced by name.
                import nav as _nav_mod
                _actual_nav_series_local = compute_actual_nav_series_spliced(
                    prices,
                    APP_DIR / "ibkr_fills_log.jsonl",
                    APP_DIR / "lots_seed.json",
                    fx_usdaud=globals().get("fx_usdaud"),
                    statement_path=_nav_mod.statement_path_for(APP_DIR),
                )
            except Exception as _e_nav:
                print(f"[chart] Actual NAV computation failed: {_e_nav}")
        # Stash for the table section below so we don't compute twice.
        actual_nav_for_table = _actual_nav_series_local
        _nav_in_window = _actual_nav_series_local.reindex(pval.index)
        # Bridge SHORT holes so the line reads as a line. A day with no broker
        # NAV snapshot becomes NaN here, and matplotlib breaks the series at
        # every NaN — six such days (07-09, 07-10, 07-14 to 07-16, 08-13) left
        # the live line visibly shredded. The NAV either side of each hole is
        # correct; only the plot was misleading.
        #
        # Bounded deliberately at NAV_GAP_BRIDGE_DAYS. A missed morning is a
        # cosmetic artifact worth hiding; a week of silence is the pipeline
        # being down, which must stay visible on the chart rather than being
        # smoothed into a straight line that implies data we do not have.
        # Interpolation is time-weighted and touches the PLOT only — nothing
        # stored, and the return table below still reads the raw series.
        # Capture WHICH days before filling them — a bare count cannot be
        # reconciled after the fact. Twice the reported number differed by one
        # from what the missing broker snapshots explain (6 known holes vs a
        # count of 7 on 08-27 and again on 09-02), and it could not be
        # reproduced outside the engine: same window, same 67 rows, same
        # universe, same snapshot timing all give 6. The likely culprit is a
        # NaN in the engine's own processed panel — compute_nav_from_statement
        # ends with nav.dropna(), so one missing price drops a whole date from
        # the series — but naming the dates settles it instead of inferring it.
        _gap_days = list(_nav_in_window.index[_nav_in_window.isna()])
        _nav_in_window, _bridged = bridge_short_gaps(
            _nav_in_window, NAV_GAP_BRIDGE_DAYS)
        if _bridged:
            # Count only INTERIOR holes. Everything before the live series
            # begins is NaN too — 22 days here, since the chart opens on
            # 2026-05-25 and the account's NAV starts 2026-06-24 — and calling
            # those an "outage" reports a pipeline failure that never happened.
            _fv = _nav_in_window.first_valid_index()
            _lv = _nav_in_window.last_valid_index()
            _left = (int(_nav_in_window.loc[_fv:_lv].isna().sum())
                     if _fv is not None and _lv is not None else 0)
            _filled = [d for d in _gap_days
                       if d >= _fv and d <= _lv] if (_fv is not None and _lv is not None) else _gap_days
            print(f"[chart] actual-NAV: bridged {_bridged} missing day(s) in "
                  f"the plot ({', '.join(d.strftime('%Y-%m-%d') for d in _filled[:10])}"
                  f"{' ...' if len(_filled) > 10 else ''})"
                  + (f"; {_left} left as visible gaps (outage > "
                     f"{NAV_GAP_BRIDGE_DAYS}d)" if _left else ""))
        _first_valid_nav = _nav_in_window.first_valid_index()
        # Every other line on this chart is rebased here too — see rebase_to.
        # None means no live NAV, so the window start stands.
        _chart_origin = None
        if _first_valid_nav is not None:
            _base = float(_nav_in_window.loc[_first_valid_nav])
            if _base > 0:
                _chart_origin = _first_valid_nav
                portfolio_returns = (_nav_in_window / _base) - 1.0
                portfolio_legend_label = (
                    f"Actual NAV (since {_first_valid_nav.strftime('%d %b')})"
                )
            else:
                portfolio_returns = (pval / pval.iloc[0]) - 1.0
                portfolio_legend_label = "Portfolio (Hypothetical)"
        else:
            # No live NAV history yet — fall back to hypothetical so the
            # slide still renders during dev runs / before seed.
            portfolio_returns = (pval / pval.iloc[0]) - 1.0
            portfolio_legend_label = "Portfolio (Hypothetical)"
            print("[chart] no actual NAV data — falling back to hypothetical")
        benchmark_returns = rebase_to(bench, _chart_origin).subtract(1.0)
                
        # --- Strategy line: synthetic projection of the AUTO-SELECTED plan ---
        # Previously hardcoded to "With Tilts" even when the auto-picker chose
        # Ensemble (~95% of runs), which made the chart legend dishonest about
        # what's driving live returns. Read TRADEPLAN_LABEL/WEIGHTS_SER (both
        # set when the auto-picker runs, ~line 6457) and label the chart line
        # accordingly. Falls back to W_WITH_TILTS_SER if the globals weren't
        # populated for any reason.
        _strategy_label_map = {
            "ensemble":   "Strategy (Ensemble)",
            "with_tilts": "Strategy (With Tilts)",
            "no_tilts":   "Strategy (Optimised)",
        }
        _tp_label_raw = str(globals().get("TRADEPLAN_LABEL", "")).strip().lower()
        strategy_legend_label = _strategy_label_map.get(_tp_label_raw, "Strategy (With Tilts)")

        tilted_returns = None
        try:
            returns_wide_df = globals().get("returns_wide_df", None)
            w_selected = globals().get("TRADEPLAN_WEIGHTS_SER", None)
            if w_selected is None or (isinstance(w_selected, pd.Series) and w_selected.empty):
                # Fallback: use W_WITH_TILTS_SER if the auto-picker didn't publish
                w_selected = globals().get("W_WITH_TILTS_SER", None)
                strategy_legend_label = "Strategy (With Tilts)"

            if isinstance(returns_wide_df, pd.DataFrame) and isinstance(w_selected, (pd.Series, dict)):
                w_ser = pd.Series(w_selected).astype(float)
                common = returns_wide_df.columns.intersection(w_ser.index)
                w_ser = w_ser.reindex(common).fillna(0.0)
                if float(w_ser.sum()) != 0.0:
                    w_ser = w_ser / float(w_ser.sum())

                r_tilt = (returns_wide_df[common].reindex(pval.index).fillna(0.0) @ w_ser).astype(float)
                tilted_curve = (1.0 + r_tilt).cumprod()
                tilted_returns = rebase_to(tilted_curve, _chart_origin) - 1.0
        except Exception:
            tilted_returns = None


        # Friendly labels
        benchmark_returns = benchmark_returns.rename(columns={
            "^AORD": "ASX",
            "^GSPC": "S&P 500",
            "^IXIC": "NASDAQ"
        })
        
        # --- Combine into one DataFrame ---
        series_list = [portfolio_returns.rename(portfolio_legend_label)]
        
        # Add Strategy line (auto-selected plan) if we successfully built it
        if "tilted_returns" in locals() and tilted_returns is not None:
            series_list.append(tilted_returns.rename(strategy_legend_label))
        
        series_list += [benchmark_returns[c].rename(c) for c in benchmark_returns.columns]
        
        perf_df = pd.concat(series_list, axis=1).dropna(how="all")
        
        if perf_df.empty or perf_df.dropna(how="all").empty:
            raise ValueError("Slide 3 perf_df is empty after alignment (portfolio vs benchmarks).")
        
        # Optional: clip extreme outliers (prevents visual spikes)
        perf_df = perf_df.clip(lower=-0.2, upper=0.5)

        # Diagnostic: confirm whether the chart is plotting through the most recent trading day,
        # or whether some series is truncating the join. If max < today by more than a few days,
        # check the benchmark fetch (yfinance lag) or the with-tilts series alignment.
        try:
            print(f"[chart] Slide 3 perf_df range: {perf_df.index.min().date()} -> "
                  f"{perf_df.index.max().date()} ({len(perf_df)} rows)")
        except Exception:
            pass

        # The date callout is deferred until AFTER the live-performance block so
        # the live summary can be appended onto the SAME banner line (the banner
        # only fits one subtitle row — a second line clips its bottom edge).
        _s4_date_start, _s4_date_end = perf_df.index.min(), perf_df.index.max()
        _s4_live_extra = ""

        # --- Live performance callout (real returns across rebalances) ----
        # Single-line summary of the actual NAV's since-inception return,
        # alongside SPY and ASX over the same window, plus rebalance count.
        # Shows the user the real performance figure they care about, so they
        # don't have to read it off the stub line on the chart.
        try:
            _live_nav_clean = (actual_nav_for_table.dropna()
                                if isinstance(actual_nav_for_table, pd.Series)
                                else pd.Series(dtype=float))
            if not _live_nav_clean.empty:
                _v0 = float(_live_nav_clean.iloc[0])
                _v1 = float(_live_nav_clean.iloc[-1])
                if _v0 > 0:
                    _live_ret = (_v1 / _v0) - 1.0
                    _live_start = _live_nav_clean.index[0]
                    _live_end = _live_nav_clean.index[-1]
                    _bench_parts = []
                    for _col, _lbl in [("^GSPC", "SPY"), ("^AORD", "ASX")]:
                        if _col in bench_long.columns:
                            _s = pd.to_numeric(bench_long[_col],
                                                 errors="coerce").dropna()
                            if not _s.empty:
                                _b0_dt = _nearest_on_or_before(_s.index, _live_start)
                                _b1_dt = _nearest_on_or_before(_s.index, _live_end)
                                if _b0_dt is not None and _b1_dt is not None:
                                    _b0 = float(_s.loc[_b0_dt])
                                    _b1 = float(_s.loc[_b1_dt])
                                    if _b0 > 0:
                                        _bench_parts.append(
                                            f"vs {_lbl}: {(_b1/_b0-1.0)*100:+.2f}%"
                                        )
                    _n_rebal = 0
                    try:
                        _fp = APP_DIR / "ibkr_fills_log.jsonl"
                        if _fp.exists():
                            _rebal_dates = set()
                            with open(_fp, "r", encoding="utf-8") as _fh:
                                for _line in _fh:
                                    _line = _line.strip()
                                    if not _line:
                                        continue
                                    try:
                                        _r = json.loads(_line.replace("NaN", "null"))
                                    except Exception:
                                        continue
                                    if float(_r.get("qty_filled") or 0) <= 0:
                                        continue
                                    _ts = (_r.get("exec_timestamp")
                                            or _r.get("rec_log_run_at"))
                                    if _ts:
                                        try:
                                            _rebal_dates.add(
                                                pd.Timestamp(_ts).normalize()
                                            )
                                        except Exception:
                                            pass
                            _n_rebal = len(_rebal_dates)
                    except Exception:
                        pass

                    _abs_pnl = _v1 - _v0
                    _first_part = (
                        f"Live since {_live_start.strftime('%d %b')}: "
                        f"{_live_ret*100:+.2f}%  (${_abs_pnl:+,.0f})"
                    )
                    if _n_rebal:
                        _first_part += (
                            f"  ·  {_n_rebal} rebal"
                            f"{'s' if _n_rebal != 1 else ''}"
                        )
                    _parts = [_first_part] + _bench_parts
                    # Current-FY tax accrual from the ACTUAL-fills ledger, so
                    # the user sees what will settle at lodgement (and should
                    # be contributed back / funded externally to keep live
                    # NAV comparable to the tax-inside-portfolio backtest).
                    try:
                        _fy_led = globals().get("FY_TAX_LEDGER_DF")
                        if isinstance(_fy_led, pd.DataFrame) and not _fy_led.empty:
                            _cur_fy = _fy_led.iloc[-1]
                            _cgt_acc = float(_cur_fy["CGT at Lodgement (AUD)"])
                            _cf_out = float(_cur_fy["Carry-Fwd Out"])
                            if _cgt_acc > 0:
                                _parts.append(
                                    f"{_cur_fy['FY']} CGT accrued: ${_cgt_acc:,.0f}"
                                )
                            elif _cf_out > 0:
                                _parts.append(
                                    f"{_cur_fy['FY']} loss c/f: ${_cf_out:,.0f} "
                                    f"(~${_cf_out * _effective_cgt_rate(short_term=True):,.0f} future tax shield)"
                                )
                    except Exception:
                        pass
                    _live_text = "    ".join(_parts)

                    # Append the live summary onto the banner's date line (added
                    # below) rather than its own textbox — the slide has no free
                    # band (chart + return table fill it), and a header-ribbon
                    # textbox clipped against the banner edge.
                    _s4_live_extra = _live_text
                    print(f"[pptx] Slide 4 live perf callout: {_live_text}")
        except Exception as _e_live:
            print(f"[pptx] Slide 4 live perf callout skipped: {_e_live}")

        # Now emit the banner date callout, with the live summary appended if we
        # built one (one line, clear of the banner's bottom edge).
        _add_date_callout(slide3, _s4_date_start, _s4_date_end,
                          prefix="Data", extra=_s4_live_extra)

        # --- Plot ---
        # Use matplotlib directly (not pandas .plot) so the x-axis is a real
        # date axis. pandas's plotter uses period codes internally, which makes
        # explicit set_xticks([Timestamp,...]) silently fall back to an auto-
        # locator with the wrong range.
        fig, ax = plt.subplots(figsize=(7, 4.5))
        for col in perf_df.columns:
            ax.plot(perf_df.index, perf_df[col].mul(100), linewidth=1.8, label=col)

        # Lock x-axis to actual data range so callout and chart agree.
        ax.set_xlim(perf_df.index.min(), perf_df.index.max())

        # Explicit major ticks: data start + each subsequent month-start within range.
        # Converted to mpl date numbers so FixedLocator places them correctly.
        _start, _end = perf_df.index.min(), perf_df.index.max()
        _month_ticks = pd.date_range(
            start=(_start + pd.offsets.MonthBegin(1)),
            end=_end,
            freq="MS",
        )
        major_ticks = [mdates.date2num(_start)] + [mdates.date2num(t) for t in _month_ticks]
        ax.xaxis.set_major_locator(mtick.FixedLocator(major_ticks))
        ax.xaxis.set_minor_locator(mdates.DayLocator(bymonthday=(15,)))
        ax.xaxis.set_major_formatter(mdates.DateFormatter("%d %b"))
        ax.xaxis.set_minor_formatter(mdates.DateFormatter("%d"))
        ax.tick_params(axis="x", which="major", labelsize=9, rotation=0)
        ax.tick_params(axis="x", which="minor", labelsize=7, rotation=0, colors="#888888")
        fig.subplots_adjust(bottom=0.22)

        # Name the origin. Every line is now rebased to the account's first NAV
        # observation, so all of them cross zero there — without saying so, a
        # reader has to infer why the benchmarks start below zero.
        _origin_note = (f", rebased {_chart_origin:%d %b}"
                        if locals().get("_chart_origin") is not None else "")
        ax.set_title("Portfolio vs ASX, S&P 500, NASDAQ "
                     f"(3-Month Performance{_origin_note})")
        ax.set_ylabel("Return (%)")
        ax.legend(loc="upper left", frameon=False)
        ax.grid(True, linestyle="--", alpha=0.4)

        # Explicit end-date annotation in the chart bottom-right, mirroring the
        # FF-vs-Portfolio slide so the live-data end date is unambiguous.
        try:
            ax.annotate(
                f"End: {pd.Timestamp(perf_df.index.max()).strftime('%d %b %Y')}",
                xy=(0.99, 0.02), xycoords="axes fraction",
                ha="right", va="bottom",
                fontsize=9, color="#404040", style="italic",
                bbox=dict(boxstyle="round,pad=0.3", facecolor="white",
                          edgecolor="#bbbbbb", alpha=0.85),
            )
        except Exception:
            pass

        _perf_buf = io.BytesIO()
        fig.savefig(_perf_buf, format="png", bbox_inches="tight")
        plt.close(fig)
        _perf_buf.seek(0)

        # --- Insert chart in PowerPoint ---
        chart_left, chart_top, chart_w, chart_h = _ppt_anchor(
            slide, slide_layout, "chart_perf",
            fb_left_cm=2.032, fb_top_cm=2.95, fb_w_cm=20.828, fb_h_cm=11.176,
        )
        slide.shapes.add_picture(_perf_buf, chart_left, chart_top, width=chart_w, height=chart_h)

        # --- Performance table (3m / 6m / 12m / 3y) under the chart ---
        try:
            # Portfolio series for table (prefer in-memory value series)
            port_px = globals().get("portfolio_value_series", None)
            if port_px is None:
                port_px = globals().get("portfolio_value_series", None)

            
            # Re-use the single benchmark download from the chart section so the
            # chart's 3M curves and the table's 3M column are computed from
            # IDENTICAL price series (no auto_adjust drift between windows).
            bench_px = {b: bench_long[b] for b in benchmarks if b in bench_long.columns}

        
            end_dt = None
            if port_px is not None and not port_px.dropna().empty:
                end_dt = port_px.dropna().index[-1]
            elif len(bench_px) > 0:
                end_dt = list(bench_px.values())[0].dropna().index[-1]
        
            metrics = ["3M", "6M", "12M", "3Y"]
            rows = {}

            # Portfolio row: actual NAV from lots seed + fills, computed in
            # the chart section above and stashed in actual_nav_for_table.
            # port_px (the hypothetical) stays in scope below as the date
            # anchor for the Strategy synth row — that line is meant to
            # span the full 3M window, while Actual NAV stubs in at
            # LIVE_START. Periods predating LIVE_START → NaN via
            # _period_total_return → empty cells in the table.
            _table_actual_nav = locals().get("actual_nav_for_table", None)
            if (_table_actual_nav is not None
                    and isinstance(_table_actual_nav, pd.Series)
                    and not _table_actual_nav.dropna().empty):
                _nav_clean = _table_actual_nav.dropna()
                _first_valid = (
                    _nav_clean.iloc[0] if not _nav_clean.empty else None
                )
                if _first_valid is not None and float(_first_valid) > 0:
                    _live_label = (
                        f"Actual NAV (since {_nav_clean.index[0].strftime('%d %b')})"
                    )
                else:
                    _live_label = "Actual NAV"
                rows[_live_label] = [
                    _period_total_return(_table_actual_nav, end_dt, months=3),
                    _period_total_return(_table_actual_nav, end_dt, months=6),
                    _period_total_return(_table_actual_nav, end_dt, months=12),
                    _period_total_return(_table_actual_nav, end_dt, years=3),
                ]
            elif port_px is not None:
                # No live NAV history yet — fall back to hypothetical so the
                # table still has a portfolio row (dev runs / pre-seed).
                rows["Portfolio (Hypothetical)"] = [
                    _period_total_return(port_px, end_dt, months=3),
                    _period_total_return(port_px, end_dt, months=6),
                    _period_total_return(port_px, end_dt, months=12),
                    _period_total_return(port_px, end_dt, years=3),
                ]
            # Strategy row: synthetic projection of the AUTO-SELECTED plan, with
            # the legend label matching the chart above (Ensemble / With Tilts /
            # Optimised). Falls back to W_WITH_TILTS_SER if the auto-picker
            # didn't publish for any reason.
            try:
                returns_wide_df = globals().get("returns_wide_df", None)
                _tp_label_tbl = str(globals().get("TRADEPLAN_LABEL", "")).strip().lower()
                _table_label_map = {
                    "ensemble":   "Strategy (Ensemble)",
                    "with_tilts": "Strategy (With Tilts)",
                    "no_tilts":   "Strategy (Optimised)",
                }
                strategy_row_label = _table_label_map.get(_tp_label_tbl, "Strategy (With Tilts)")
                w_selected_tbl = globals().get("TRADEPLAN_WEIGHTS_SER", None)
                if w_selected_tbl is None or (isinstance(w_selected_tbl, pd.Series) and w_selected_tbl.empty):
                    w_selected_tbl = globals().get("W_WITH_TILTS_SER", None)
                    strategy_row_label = "Strategy (With Tilts)"

                if isinstance(returns_wide_df, pd.DataFrame) and isinstance(w_selected_tbl, (pd.Series, dict)) and port_px is not None:
                    w_ser = pd.Series(w_selected_tbl).astype(float)
                    common = returns_wide_df.columns.intersection(w_ser.index)
                    w_ser = w_ser.reindex(common).fillna(0.0)
                    if float(w_ser.sum()) != 0.0:
                        w_ser = w_ser / float(w_ser.sum())

                    # Build a synthetic "price" series over the SAME date index as port_px (so _period_total_return works)
                    r_tilt_tbl = (returns_wide_df[common].reindex(port_px.index).fillna(0.0) @ w_ser).astype(float)
                    px_tilt_tbl = (1.0 + r_tilt_tbl).cumprod()
                    px_tilt_tbl = px_tilt_tbl * float(pd.to_numeric(port_px, errors="coerce").dropna().iloc[0])

                    rows[strategy_row_label] = [
                        _period_total_return(px_tilt_tbl, end_dt, months=3),
                        _period_total_return(px_tilt_tbl, end_dt, months=6),
                        _period_total_return(px_tilt_tbl, end_dt, months=12),
                        _period_total_return(px_tilt_tbl, end_dt, years=3),
                    ]
            except Exception:
                pass

            for b, s in bench_px.items():
                rows[str(b)] = [
                    _period_total_return(s, end_dt, months=3),
                    _period_total_return(s, end_dt, months=6),
                    _period_total_return(s, end_dt, months=12),
                    _period_total_return(s, end_dt, years=3),
                ]
        
            perf_tbl = pd.DataFrame.from_dict(rows, orient="index", columns=metrics)
            name_map = {"^AORD": "ASX", "^GSPC": "S&P 500", "^IXIC": "NASDAQ"}
            perf_tbl = perf_tbl.rename(index=name_map)
            tbl_left, tbl_top, tbl_w, tbl_h = _ppt_anchor(
                slide, slide_layout, "table_perf",
                fb_left_cm=2.032, fb_top_cm=13.85, fb_w_cm=20.828, fb_h_cm=2.40,
            )
            _add_perf_table(
                slide, perf_tbl,
                left=tbl_left, top=tbl_top, width=tbl_w, height=tbl_h,
                title="Return Summary",
            )
        except Exception as e:
            print(f"[pptx] Slide 3 table skipped: {e}")

        
        # --- SLIDE 4: Fama French benchmarks + table (quarterly) ---
        try:
            slide_layout = prs.slide_layouts[20]
            slide4 = prs.slides.add_slide(slide_layout)
            if slide4.shapes.title:
                slide4.shapes.title.text = "Fama French Benchmarks vs Portfolio"
        
            # Pull FF factors (daily) and convert to quarterly returns
            ff = globals().get("ff5_raw", None)
            pxdf = globals().get("prices", None)
        
            if isinstance(ff, pd.DataFrame) and not ff.empty and isinstance(pxdf, pd.DataFrame) and "PortfolioValue" in pxdf.columns:
                port_px = pd.to_numeric(pxdf["PortfolioValue"], errors="coerce").dropna()
        
            # Build daily benchmark series (FF factors are daily; most recent date may lag live markets)
            ff_cols = [c for c in ["Mkt-RF","SMB","HML","RMW","CMA","MOM","RF"] if c in ff.columns]
            ffd = ff[ff_cols].dropna().copy()
            
            # Market total return proxy = (Mkt-RF + RF)
            if ("Mkt-RF" in ffd.columns) and ("RF" in ffd.columns):
                ffd["Market (Mkt-RF)"] = ffd["Mkt-RF"] + ffd["RF"]
            
            # Portfolio daily returns
            port_r = port_px.pct_change().dropna()
            
            # Use the latest common date (FF often lags)
            common_end = min(ffd.index.max(), port_r.index.max())
            # Diagnostic: surface the actual dates being used so we can
            # diagnose any misalignment between the annotation and the
            # chart's visible x-axis tick labels.
            print(f"[ff-chart] ff5_raw.index.max()={ff.index.max()}, "
                  f"ffd (dropna).index.max()={ffd.index.max()}, "
                  f"port_r.index.max()={port_r.index.max()}, "
                  f"common_end={common_end}")
            
            # Table window (up to 3Y of overlap)
            window_start_tbl = common_end - relativedelta(years=3, days=10)
            ffd_tbl = ffd.loc[window_start_tbl:common_end]
            port_r_tbl = port_r.loc[window_start_tbl:common_end]
            
            # Chart window (last ~3 months). Portfolio extends past common_end
            # to TODAY so viewers see live performance; FF series stop at
            # common_end (Ken French publishes ~6 weeks late). A vertical line
            # at common_end marks where FF data ends so the visible gap reads
            # as "FF data not yet available" not "portfolio diverged".
            window_start_chart = common_end - relativedelta(months=3, days=10)
            ffd_chart = ffd.loc[window_start_chart:common_end]
            live_end = port_r.index.max() if not port_r.empty else common_end
            port_r_chart = port_r.loc[window_start_chart:live_end]

            # Date callout under the slide title — Slide 4 anchors to the FF data end
            # (~1mo behind live), so everything on this slide should reference common_end.
            _add_date_callout(slide4, window_start_chart, common_end, prefix="Data (FF-anchored)")
            
            # Choose a small set to chart (readable)
            series_to_show = []
            if "Market (Mkt-RF)" in ffd.columns:
                series_to_show.append("Market (Mkt-RF)")
            for c in ["SMB","HML","RMW","CMA","MOM"]:
                if c in ffd.columns:
                    series_to_show.append(c)
                        
            # Strategy line: project the AUTO-SELECTED plan onto the FF window
            # so the chart shows what the recommended portfolio would have done
            # over the same window, not just the static current holdings.
            _strategy_label_ff_map = {
                "ensemble":   "Strategy (Ensemble)",
                "with_tilts": "Strategy (With Tilts)",
                "no_tilts":   "Strategy (Optimised)",
            }
            _tp_label_ff = str(globals().get("TRADEPLAN_LABEL", "")).strip().lower()
            strategy_label_ff = _strategy_label_ff_map.get(_tp_label_ff, "Strategy")
            strat_r_chart = pd.Series(dtype=float)
            strat_r_tbl = pd.Series(dtype=float)
            try:
                returns_wide_df_ff = globals().get("returns_wide_df", None)
                w_selected_ff = globals().get("TRADEPLAN_WEIGHTS_SER", None)
                if w_selected_ff is None or (isinstance(w_selected_ff, pd.Series) and w_selected_ff.empty):
                    w_selected_ff = globals().get("W_WITH_TILTS_SER", None)
                    strategy_label_ff = "Strategy (With Tilts)"

                if isinstance(returns_wide_df_ff, pd.DataFrame) and isinstance(w_selected_ff, (pd.Series, dict)):
                    w_ser_ff = pd.Series(w_selected_ff).astype(float)
                    common_ff = returns_wide_df_ff.columns.intersection(w_ser_ff.index)
                    w_ser_ff = w_ser_ff.reindex(common_ff).fillna(0.0)
                    if float(w_ser_ff.sum()) != 0.0:
                        w_ser_ff = w_ser_ff / float(w_ser_ff.sum())
                    _strat_returns_full = (
                        returns_wide_df_ff[common_ff].fillna(0.0) @ w_ser_ff
                    ).astype(float)
                    strat_r_chart = _strat_returns_full.loc[window_start_chart:common_end]
                    strat_r_tbl = _strat_returns_full.loc[window_start_tbl:common_end]
            except Exception as _e_strat_ff:
                print(f"[pptx] Slide 4 strategy series skipped: {_e_strat_ff}")

            # Use OUTER join so Portfolio rows past common_end survive (FF
            # cells will be NaN there). Cumret is computed per-column below
            # so NaN propagates correctly — FF lines stop, Portfolio continues.
            # Label says pro-forma: this series backcasts the CURRENT/target
            # book's holdings over the window — it is NOT the fund's live
            # NAV history (that's slide 2's "Actual NAV" line). User flagged
            # the bare "Portfolio" label as misleading 2026-07-06.
            _port_label_ff = "Current Book (pro-forma)"
            chart_df = pd.DataFrame({_port_label_ff: port_r_chart})
            if not strat_r_chart.empty:
                # Extend strategy too so its line matches Portfolio length
                strat_full = globals().get("returns_wide_df", None)
                if isinstance(strat_full, pd.DataFrame):
                    try:
                        _w_full = pd.Series(w_selected_ff).astype(float)
                        _c_full = strat_full.columns.intersection(_w_full.index)
                        _w_full = _w_full.reindex(_c_full).fillna(0.0)
                        if float(_w_full.sum()) != 0.0:
                            _w_full = _w_full / float(_w_full.sum())
                        _sret_full = (strat_full[_c_full].fillna(0.0) @ _w_full).astype(float)
                        chart_df[strategy_label_ff] = _sret_full.loc[
                            window_start_chart:live_end]
                    except Exception:
                        chart_df[strategy_label_ff] = strat_r_chart
                else:
                    chart_df[strategy_label_ff] = strat_r_chart
            chart_df = chart_df.join(ffd_chart[series_to_show], how="outer")
            chart_df = chart_df.loc[window_start_chart:live_end]
            tbl_df = pd.DataFrame({_port_label_ff: port_r_tbl})
            if not strat_r_tbl.empty:
                tbl_df[strategy_label_ff] = strat_r_tbl
            tbl_df = tbl_df.join(ffd_tbl[series_to_show], how="inner")
            
            # Compute cumret per-column so NaN (past FF cutoff) propagates.
            # matplotlib skips NaN, so FF lines naturally end at common_end
            # while Portfolio extends to live_end.
            ret = pd.DataFrame(index=chart_df.index)
            for col in chart_df.columns:
                _s = chart_df[col]
                # Cumprod requires no leading NaN; reindex forward-fill 0 from
                # the first valid index but keep trailing NaN as-is.
                _first = _s.first_valid_index()
                _last = _s.last_valid_index()
                if _first is None or _last is None:
                    continue
                _seg = _s.loc[_first:_last].fillna(0.0)
                _cum = ((1.0 + _seg).cumprod() - 1.0) * 100.0
                ret[col] = _cum.reindex(ret.index)
            fig, ax = plt.subplots(figsize=(7.5, 4.8))
            # x_compat=True forces pandas to use matplotlib's date converter
            # instead of its own. Without this, axvline(common_end) and
            # MonthLocator place ticks/lines in different coordinate systems —
            # symptom: vertical line annotated "30 Apr" lands visually
            # between the 01-Mar and 01-Apr ticks instead of just before 01-May.
            ret.plot(ax=ax, linewidth=1.4, x_compat=True)

            # Make room inside the figure on the right for the legend
            fig.subplots_adjust(right=0.78)
            ax.legend(loc="center left", bbox_to_anchor=(1.01, 0.5), frameon=False, fontsize=9)
            ax.set_title("Portfolio vs Fama French Factors (3-Month Performance)")
            ax.set_ylabel("Return (%)")
            # Pin ticks to the 1st of each month so the x-axis reads cleanly
            # (Feb 01, Mar 01, …) instead of matplotlib's auto-locator picking
            # uneven dates like 14-Jan / 23-Jan / 12-Feb / 06-Mar.
            ax.xaxis.set_major_locator(mdates.MonthLocator(bymonthday=1))
            ax.xaxis.set_major_formatter(mdates.DateFormatter("%d-%b"))
            ax.xaxis.set_minor_locator(mdates.MonthLocator(bymonthday=15))
            ax.grid(True, linestyle="--", alpha=0.4)
            ax.margins(x=0)
            if not ret.empty:
                ax.set_xlim(ret.index.min(), ret.index.max())
                print(f"[ff-chart] xlim set to ({ret.index.min()}, {ret.index.max()}), "
                      f"vertical line at common_end={common_end}")

            # Vertical line at FF cutoff (Ken French publishes ~6 wk late).
            # Lets viewers see at-a-glance that FF lines stop here while
            # Portfolio continues past — gap is "data not published yet",
            # not "Portfolio diverged from factors".
            try:
                _ff_end_ts = pd.Timestamp(common_end)
                ax.axvline(_ff_end_ts, color="#b00000", linestyle="--",
                            linewidth=1.2, alpha=0.65, zorder=1)
                # Label placed near the top of the chart, just right of the line
                _ylim_top = ax.get_ylim()[1]
                ax.annotate(
                    f"FF data ends\n{_ff_end_ts.strftime('%d %b %Y')}",
                    xy=(_ff_end_ts, _ylim_top),
                    xytext=(4, -4), textcoords="offset points",
                    ha="left", va="top",
                    fontsize=8, color="#b00000", style="italic",
                )
            except Exception:
                pass

            # Live-end annotation bottom-right.
            try:
                ax.annotate(
                    f"Live end: {pd.Timestamp(live_end).strftime('%d %b %Y')}",
                    xy=(0.99, 0.02), xycoords="axes fraction",
                    ha="right", va="bottom",
                    fontsize=9, color="#404040", style="italic",
                    bbox=dict(boxstyle="round,pad=0.3", facecolor="white",
                              edgecolor="#bbbbbb", alpha=0.85),
                )
            except Exception:
                pass
            _ff_buf = io.BytesIO()
            fig.savefig(_ff_buf, format="png", bbox_inches="tight")
            plt.close(fig)
            _ff_buf.seek(0)

            chart_left, chart_top, chart_w, chart_h = _ppt_anchor(
                slide4, slide_layout, "chart_ff",
                fb_left_cm=2.032, fb_top_cm=3.05, fb_w_cm=20.32, fb_h_cm=8.65,
            )
            slide4.shapes.add_picture(_ff_buf, chart_left, chart_top, width=chart_w, height=chart_h)
            
            # Table: 3M/6M/12M/3Y (compounded) using available daily points.
            # ALL rows — including Portfolio — anchor to the FF data end so the
            # table is internally consistent with the FF-anchored chart above.
            # Previously the Portfolio row used the live end date and the factor
            # rows used the FF end, which made values incomparable (e.g. 9.43%
            # Portfolio 3M against 3.09% Mkt-RF 3M because they covered different
            # 90-day windows). Slide 3 still reports live-end performance.
            end_dt_tbl = tbl_df.index.max()
            # Factor and Strategy rows query the FULL upstream series (ffd
            # spans 1963→present; _strat_returns_full spans the 10Y OOS
            # window) rather than the inner-joined tbl_df, which is capped
            # at Portfolio's ~2Y price window. Without this, 3Y / 12M cells
            # would be NaN for factor rows even though Ken French data goes
            # back decades. Portfolio row stays on port_px (so 3Y is NaN
            # while the live price window is <3Y — honest, populates as
            # history accrues).
            _strat_full_local = locals().get("_strat_returns_full", None)
            rows = {}
            for name in tbl_df.columns:
                if name == _port_label_ff:
                    rows[name] = [
                        _period_total_return(port_px, end_dt_tbl, months=3),
                        _period_total_return(port_px, end_dt_tbl, months=6),
                        _period_total_return(port_px, end_dt_tbl, months=12),
                        _period_total_return(port_px, end_dt_tbl, years=3),
                    ]
                elif (name == strategy_label_ff
                        and isinstance(_strat_full_local, pd.Series)
                        and not _strat_full_local.empty):
                    rr = _strat_full_local
                    rows[name] = [
                        _window_compound_total(rr, end_dt_tbl, months=3),
                        _window_compound_total(rr, end_dt_tbl, months=6),
                        _window_compound_total(rr, end_dt_tbl, months=12),
                        _window_compound_total(rr, end_dt_tbl, years=3),
                    ]
                else:
                    rr = ffd[name] if name in ffd.columns else tbl_df[name]
                    rows[name] = [
                        _window_compound_total(rr, end_dt_tbl, months=3),
                        _window_compound_total(rr, end_dt_tbl, months=6),
                        _window_compound_total(rr, end_dt_tbl, months=12),
                        _window_compound_total(rr, end_dt_tbl, years=3),
                    ]

            ff_tbl = pd.DataFrame.from_dict(rows, orient="index", columns=["3M", "6M", "12M", "3Y"])
            # Position + size tuned 2026-06-21 to user's measured target:
            # left=2.03cm, top=11.55cm, width=20.32cm, height=6.48cm.
            # Previous height 2.794cm was too tight (text overflowed).
            tbl_left, tbl_top, tbl_w, tbl_h = _ppt_anchor(
                slide4, slide_layout, "table_ff",
                fb_left_cm=2.03, fb_top_cm=11.55, fb_w_cm=20.32, fb_h_cm=6.48,
            )
            _add_perf_table(
                slide4, ff_tbl,
                left=tbl_left, top=tbl_top, width=tbl_w, height=tbl_h,
                title="Return Summary",
            )
        except Exception as e:
            print(f"[pptx] Slide 4 skipped: {e}")

        # --- SLIDE 5: Efficient Frontier (chart + points table) ---
        try:
            if isinstance(charts, dict):
                print("[pptx] Slide 5 chart keys:", list(charts.keys()))
                print("[pptx] tilts_comparison_rows present:", "tilts_comparison_rows" in charts)
                if "tilts_comparison_rows" in charts:
                    print("[pptx] tilts_comparison_rows sample:", charts["tilts_comparison_rows"][:3])
        except Exception:
            pass

        try:
            slide_layout = prs.slide_layouts[20]
            slide5 = prs.slides.add_slide(slide_layout)
            if slide5.shapes.title:
                slide5.shapes.title.text = "Efficient Frontier"
        
            eff_image = None
            if isinstance(charts, dict):
                eff_image = charts.get("efficient_frontier_image", None)

            # Always build rows for the points table (even if the chart image is missing)
            pts = {}
            if isinstance(charts, dict):
                pts = charts.get("frontier_points", {}) or {}

            rows = []
            # New marker set as of 2026-06-22: 5 ensemble slots + the blended
            # Ensemble portfolio (replaced Current/Previous/Optimised/With
            # Tilts). Falls back to legacy keys if an older charts dict is
            # passed in so the slide doesn't blank out on stale state.
            # Today's Weight column (added 2026-06-22): pulls each slot's
            # softmax weight from ensemble_mix_live so the regime mix is
            # visible right next to each slot's vol/return.
            _slot_keys = list(ENSEMBLE_SLOT_NAMES) + ["Ensemble"]
            _legacy_keys = ["Current", "Previous", "Optimised", "With Tilts", "Ensemble"]
            _ordered_keys = _slot_keys if any(k in pts for k in _slot_keys) else _legacy_keys
            _ens_mix = globals().get("ensemble_mix_live", pd.Series(dtype=float))
            for k in _ordered_keys:
                v = pts.get(k, None)
                if v is None:
                    continue
                try:
                    vol, ret = float(v[0]), float(v[1])
                    if np.isfinite(vol) and np.isfinite(ret):
                        _label = k.split(" (")[0] if "(" in k else k
                        if k == "Ensemble":
                            _wt = float("nan")
                        elif isinstance(_ens_mix, pd.Series) and k in _ens_mix.index:
                            _wt = float(_ens_mix[k])
                        else:
                            _wt = float("nan")
                        rows.append({"Point": _label, "Vol (ann.)": vol,
                                     "Return (ann.)": ret, "Weight": _wt})
                except Exception:
                    pass

            # Chart is OPTIONAL: only add if a buffer is present.
            if eff_image is not None:
                try:
                    eff_image.seek(0)  # defensive — in case anyone read from it earlier
                except Exception:
                    pass
                chart_left, chart_top, chart_w, chart_h = _ppt_anchor(
                    slide5, slide_layout, "chart_frontier",
                    fb_left_cm=1.52, fb_top_cm=3.56, fb_w_cm=14.50, fb_h_cm=11.50,
                )
                slide5.shapes.add_picture(
                    eff_image, chart_left, chart_top, width=chart_w, height=chart_h,
                )
            
            # Points table (always add if we have data)
            if rows:
                df_pts = pd.DataFrame(rows).set_index("Point").rename(
                    columns={"Vol (ann.)": "Vol.", "Return (ann.)": "Return"}
                )
                # Widened from 7.72cm to 8.70cm so the new Today's Weight
                # column fits without crowding the existing 3. Left edge
                # unchanged (16.50cm) so chart layout is untouched.
                tbl_left, tbl_top, tbl_w, tbl_h = _ppt_anchor(
                    slide5, slide_layout, "table_frontier_points",
                    fb_left_cm=16.50, fb_top_cm=4.06, fb_w_cm=8.70, fb_h_cm=4.32,
                )
                _add_perf_table(
                    slide5, df_pts,
                    left=tbl_left, top=tbl_top, width=tbl_w, height=tbl_h,
                    title="Portfolio",
                )

            # Tilts table dropped from Slide 5 (2026-06-22): the With Tilts
            # vs Without Tilts contrast no longer appears anywhere on the
            # slide (markers replaced by per-slot frontier points), so the
            # table sat orphaned with identical Achieved/Target columns.
            # The data is still stored in charts["tilts_comparison_rows"]
            # if any other consumer wants it.

        except Exception as e:
            print(f"[pptx] Slide 5 skipped: {e}")

        # ---- ROADSHOW SLIDE (Phase 3): inserted at position 2 after build. ----
        try:
            # Reporting-truncated series so the chart ends where the metrics
            # table does. oos_returns_daily stays FULL for the drift tracker.
            oos_rets = globals().get("oos_returns_report", pd.Series(dtype=float))
            if not isinstance(oos_rets, pd.Series) or oos_rets.empty:
                oos_rets = globals().get("oos_returns_daily", pd.Series(dtype=float))
            oos_mtx = globals().get("oos_metrics_table", pd.DataFrame())
            # Reporting-lockboxed frame so the CHART ends on the same date as
            # the metrics table; falls back to the full panel if absent.
            oos_px_long = globals().get("oos_prices_report", pd.DataFrame())
            if not isinstance(oos_px_long, pd.DataFrame) or oos_px_long.empty:
                oos_px_long = globals().get("oos_prices_aud_long", pd.DataFrame())

            if (isinstance(oos_rets, pd.Series) and not oos_rets.empty and
                isinstance(oos_mtx, pd.DataFrame) and not oos_mtx.empty):

                slide_layout = prs.slide_layouts[20]
                road = prs.slides.add_slide(slide_layout)
                if road.shapes.title:
                    road.shapes.title.text = "Fund Performance vs Benchmarks"

                # Aligned daily returns for chart.
                end_dt_rs = oos_rets.index.max()
                start_dt_rs = end_dt_rs - pd.DateOffset(years=10)
                rs_strat = oos_rets[(oos_rets.index >= start_dt_rs) &
                                    (oos_rets.index <= end_dt_rs)].copy()

                # Benchmark daily returns from the unified long download.
                def _bench_rets(col):
                    if col not in oos_px_long.columns:
                        return pd.Series(dtype=float)
                    px = pd.to_numeric(oos_px_long[col], errors="coerce").dropna()
                    return px.pct_change().dropna()
                # AU benchmark: total-return series (VAS.AX) not the price-only
                # ^AORD, which excludes dividends and understated the AU
                # alternative by ~4.2%/yr on this very chart.
                _au_tkr = globals().get("AU_BENCH_TICKER", "VAS.AX")
                _au_lbl = globals().get("AU_BENCH_LABEL", "AU equities (TR)")
                if _au_tkr not in oos_px_long.columns:
                    _au_tkr = globals().get("AU_BENCH_FALLBACK", "^AORD")
                    _au_lbl = "^AORD (price only)"
                spy_rs = _bench_rets("SPY").reindex(rs_strat.index).fillna(0.0)
                aord_rs = _bench_rets(_au_tkr).reindex(rs_strat.index).fillna(0.0)

                # Normalized cumulative-return curves (all series start at
                # 100%). This is the apples-to-apples view requested by
                # the user 2026-06-26: comparing Fund @ $X to SPY @ $Y in
                # absolute dollars is a category error (different starting
                # points), so plot everything as % cumulative return and
                # let the slope/final speak. The actual starting NAV is
                # still preserved in `_oos_starting_nav_aud` for the
                # title + the metrics table.
                #
                # Guard against the 2026-06-26 dual-mode bug where a
                # silently-corrupted SPY benchmark series produced flat
                # 0% lines and degenerate metrics: refuse to render the
                # chart if any series collapses to zero variance or is
                # all-zero — surface the issue instead of shipping a
                # broken slide.
                def _is_degenerate(s, label):
                    if not isinstance(s, pd.Series) or s.empty:
                        return f"{label}: series empty"
                    if not np.isfinite(s).any():
                        return f"{label}: no finite values"
                    if float(s.abs().sum()) < 1e-9:
                        return f"{label}: all-zero (likely silent fetch failure)"
                    return None

                _series_health = [
                    _is_degenerate(rs_strat, "Fund returns"),
                    _is_degenerate(spy_rs, "SPY returns"),
                    _is_degenerate(aord_rs, f"{_au_lbl} returns"),
                ]
                _bad = [m for m in _series_health if m]
                if _bad:
                    print(f"[pptx] Roadshow chart degenerate; refusing to render: "
                          f"{'; '.join(_bad)}")
                    raise RuntimeError(f"Roadshow chart degenerate: {_bad}")

                w_strat = (1.0 + rs_strat).cumprod()
                w_spy = (1.0 + spy_rs).cumprod()
                w_aord = (1.0 + aord_rs).cumprod()

                # Starting NAVs preserved for title + metric-table labels.
                _user_nav = float(globals().get(
                    "_oos_starting_nav_aud", 1_000_000.0))

                _rs_rets_rs = globals().get("oos_returns_daily_roadshow",
                                            pd.Series(dtype=float))
                _has_rs_strat = (isinstance(_rs_rets_rs, pd.Series)
                                 and not _rs_rets_rs.empty)
                if _has_rs_strat:
                    rs_strat_rs = _rs_rets_rs[
                        (_rs_rets_rs.index >= start_dt_rs) &
                        (_rs_rets_rs.index <= end_dt_rs)
                    ].copy()
                    if _is_degenerate(rs_strat_rs, "Fund@RS returns"):
                        print("[pptx] Roadshow second backtest degenerate; "
                              "dropping the @RS line but keeping the rest")
                        _has_rs_strat = False
                        w_strat_rs = pd.Series(dtype=float)
                        _rs_base = None
                    else:
                        _rs_base = float(globals().get(
                            "_roadshow_nav_aud", 1_000_000.0))
                        w_strat_rs = (1.0 + rs_strat_rs).cumprod()
                else:
                    w_strat_rs = pd.Series(dtype=float)
                    _rs_base = None

                # `base` retained for downstream code that expects it.
                # Now represents the normalised start (1.0 == 100%).
                base = 1.0

                # Date callout under title (using actual rendered range).
                _add_date_callout(road, w_strat.index.min(), w_strat.index.max(),
                                  prefix="Backtest")

                # Cumulative wealth chart + ensemble regime evolution.
                # Two stacked subplots sharing the x-axis: top = wealth curves,
                # bottom = softmax-weighted regime mix over time.
                oos_soft = globals().get("oos_softmax_history", pd.DataFrame())
                has_softmax = isinstance(oos_soft, pd.DataFrame) and not oos_soft.empty
                if has_softmax:
                    fig, (ax, ax_mix) = plt.subplots(
                        2, 1, figsize=(11.5, 5.5),
                        gridspec_kw={"height_ratios": [3.5, 1.0]},
                        sharex=True,
                    )
                else:
                    fig, ax = plt.subplots(figsize=(11.5, 4.5))
                    ax_mix = None

                def _nav_label(nav):
                    return (f"${nav/1_000_000:,.2f}M" if nav >= 1_000_000
                            else f"${nav/1000:,.0f}k")

                _strat_main_label = (
                    f"Fund @ {_nav_label(_user_nav)}" if _has_rs_strat
                    else "Fund (Strategy)"
                )
                ax.plot(w_strat.index, w_strat.values, linewidth=2.2,
                        label=_strat_main_label, color="#1f4e8a")
                if _has_rs_strat and not w_strat_rs.empty:
                    ax.plot(w_strat_rs.index, w_strat_rs.values, linewidth=2.2,
                            label=f"Fund @ {_nav_label(_rs_base)}",
                            color="#1f4e8a", linestyle="--")
                ax.plot(w_spy.index, w_spy.values, linewidth=1.6,
                        label="SPY (AUD)", color="#c53030", alpha=0.85)
                ax.plot(w_aord.index, w_aord.values, linewidth=1.6,
                        label=_au_lbl, color="#2f855a", alpha=0.85)
                ax.set_xlim(w_strat.index.min(), w_strat.index.max())
                ax.xaxis.set_major_locator(mdates.YearLocator())
                ax.xaxis.set_major_formatter(mdates.DateFormatter("%Y"))
                # Y-axis: cumulative return as %. 1.0 → 0%, 2.0 → +100%, etc.
                ax.yaxis.set_major_formatter(mtick.FuncFormatter(
                    lambda x, _p: f"{(x-1.0)*100:+.0f}%"))
                ax.axhline(1.0, color="#888888", linewidth=0.8, linestyle=":", alpha=0.6)
                if _has_rs_strat and _rs_base is not None:
                    _scale_note = (f"Fund @ {_nav_label(_user_nav)} + "
                                   f"Fund @ {_nav_label(_rs_base)}")
                else:
                    _scale_note = f"Fund @ {_nav_label(_user_nav)}"
                # State the risk-free rate. It is the CURRENT RBA cash rate
                # applied flat across the whole window, which is harsh on the
                # near-zero-rate years (2020-21 cash was 0.10%) — a reader
                # cannot judge a Sharpe without knowing which rf produced it.
                _rf_note = ""
                try:
                    _rfs = globals().get("rf_series")
                    if isinstance(_rfs, pd.Series) and len(_rfs):
                        _w = _rfs[_rfs.index >= (end_dt_rs - pd.DateOffset(years=10))]
                        _rf_note = (f"; Sharpe/Sortino net of the RBA cash rate "
                                    f"as it stood each day (10Y avg "
                                    f"{_w.mean()*100:.2f}%)")
                    else:
                        _rf_v = float(globals().get("rf_annual", 0.0) or 0.0)
                        _rf_note = (f"; Sharpe/Sortino at rf={_rf_v*100:.2f}% "
                                    f"(current RBA cash rate, applied flat)")
                except Exception:
                    pass
                ax.set_title(
                    f"SIMULATED BACKTEST — not a live track record    "
                    f"({_scale_note}, net of {BROKER_CONFIG['name']} "
                    f"brokerage + AU CGT [{ACTIVE_CGT_PROFILE}]; "
                    f"benchmarks GROSS of tax and costs{_rf_note})",
                    fontsize=9,
                )
                ax.set_ylabel("Cumulative return")
                ax.legend(loc="upper left", frameon=False)
                ax.grid(True, linestyle="--", alpha=0.4)

                # Standing disclosures. These are the three things a reader
                # cannot infer from the curves and would otherwise misread:
                # what the backtest is NOT, how concentrated the live book is,
                # and that a third of the universe post-dates the start.
                try:
                    _n_late = int(globals().get("_universe_post_start_count", 0) or 0)
                    _n_univ = int(globals().get("_universe_total_count", 0) or 0)
                    _live_top = globals().get("TRADEPLAN_WEIGHTS_SER", None)
                    _conc = ""
                    if _live_top is not None and len(_live_top):
                        _t2 = _live_top.sort_values(ascending=False).head(2)
                        _conc = ("  •  Live book is concentrated: "
                                 + " + ".join(f"{k} {v*100:.0f}%" for k, v in _t2.items())
                                 + f" = {_t2.sum()*100:.0f}% in two positions.")
                    _uni = (f"  •  {_n_late} of {_n_univ} universe tickers began trading "
                            f"AFTER the backtest start, so early years ran on a smaller "
                            f"universe." if _n_late and _n_univ else "")
                    ax.text(
                        0.0, -0.16,
                        "Simulated results. Past simulated performance does not predict "
                        "future returns.  •  Live/paper trading began 2026-06-22."
                        + _conc + _uni,
                        transform=ax.transAxes, fontsize=7, color="#555555",
                        va="top", ha="left", wrap=True,
                    )
                except Exception as _e:
                    print(f"[pptx] disclosure footnote skipped: {_e}")

                # Terminal-value annotations on right edge: show final %.
                _annotations = [(w_strat, "Fund", "#1f4e8a")]
                if _has_rs_strat and not w_strat_rs.empty:
                    _annotations.append((w_strat_rs, "Fund(RS)", "#1f4e8a"))
                _annotations.extend([
                    (w_spy, "SPY", "#c53030"),
                    (w_aord, _au_lbl, "#2f855a"),
                ])
                for s, lbl, col in _annotations:
                    if not s.empty:
                        ax.annotate(f"  {(s.iloc[-1]-1.0)*100:+.0f}%",
                                    xy=(s.index[-1], s.iloc[-1]),
                                    xytext=(4, 0), textcoords="offset points",
                                    va="center", fontsize=9,
                                    fontweight="bold", color=col)

                # Bottom subplot: ensemble regime mix stacked area.
                if ax_mix is not None:
                    regime_colors = {
                        "Modest (SPY+0%)":      "#5b9bd5",  # blue - lowest aggression
                        "Aggressive (SPY+5%)":  "#70ad47",  # green
                        "Bold (SPY+10%)":       "#ffc000",  # yellow
                        "Maximum (SPY+15%)":    "#ed7d31",  # orange
                        "Stretch (SPY+25%)":    "#c00000",  # dark red - top aggression
                    }
                    cols_in_order = [n for n in ENSEMBLE_SLOT_NAMES if n in oos_soft.columns]
                    if cols_in_order:
                        soft_plot = oos_soft[cols_in_order].copy()
                        # Reindex to a regular forward-fill so the area chart
                        # interpolates between monthly rebalance dates.
                        idx_daily = w_strat.index
                        soft_plot = soft_plot.reindex(idx_daily, method="ffill").fillna(0.0)
                        ax_mix.stackplot(
                            soft_plot.index,
                            *[soft_plot[c].values for c in cols_in_order],
                            labels=[c.split(" ")[0] for c in cols_in_order],
                            colors=[regime_colors.get(c, "#888888") for c in cols_in_order],
                            alpha=0.85,
                        )
                        ax_mix.set_ylim(0, 1)
                        ax_mix.set_ylabel("Regime", fontsize=9)
                        ax_mix.yaxis.set_major_formatter(mtick.PercentFormatter(1.0, 0))
                        ax_mix.tick_params(axis="y", labelsize=8)
                        ax_mix.tick_params(axis="x", labelsize=9)
                        # Legend below the strip with clearly-sized labels —
                        # 5 entries in one row.
                        ax_mix.legend(loc="upper center", ncol=5, fontsize=9,
                                       frameon=False, bbox_to_anchor=(0.5, -0.30),
                                       handlelength=1.4, handleheight=1.0,
                                       columnspacing=1.5, borderpad=0.2)
                        ax_mix.grid(True, axis="y", linestyle="--", alpha=0.3)

                fig.tight_layout()
                _rs_buf = io.BytesIO()
                fig.savefig(_rs_buf, format="png", bbox_inches="tight")
                plt.close(fig)
                _rs_buf.seek(0)

                # Chart + table dimensions tuned 2026-06-20/21 to match user's
                # measured target layout (Format Picture dialog values):
                # chart 8.61 H × 22.6 W cm at (1.4, 2.78), table 6.48 H ×
                # 20.32 W at (2.03, 11.55). Both centred horizontally.
                road.shapes.add_picture(_rs_buf, Cm(1.4), Cm(2.78),
                                        width=Cm(22.6), height=Cm(8.61))

                # ---- Metrics table (3Y / 5Y / 10Y) ----
                # Restructure for display: rows = (horizon, series), cols = metric.
                # FF5 alpha dropped from the slide (still in the Excel sheet) to
                # keep the table compact enough to fit alongside the bigger chart.
                display_metrics = ["Annualised Return", "Annualised Volatility",
                                   "Sharpe Ratio", "Sortino Ratio",
                                   "Max Drawdown"]
                # Optional roadshow-NAV strategy gets its own per-horizon
                # row pulled from oos_metrics_table_roadshow (Strategy
                # column only — benchmarks are identical).
                _oos_mtx_rs = globals().get(
                    "oos_metrics_table_roadshow", pd.DataFrame())
                _has_rs_mtx = (isinstance(_oos_mtx_rs, pd.DataFrame)
                               and not _oos_mtx_rs.empty)
                _rs_nav_for_label = float(globals().get(
                    "_roadshow_nav_aud", 1_000_000.0))
                _rs_label_suffix = (
                    f"${_rs_nav_for_label/1_000_000:,.2f}M"
                    if _rs_nav_for_label >= 1_000_000
                    else f"${_rs_nav_for_label/1000:,.0f}k"
                )
                _user_nav_for_label = float(globals().get(
                    "_oos_starting_nav_aud", 100_000.0))
                _user_label_suffix = (
                    f"${_user_nav_for_label/1_000_000:,.2f}M"
                    if _user_nav_for_label >= 1_000_000
                    else f"${_user_nav_for_label/1000:,.0f}k"
                )

                rows = []
                row_labels = []
                for h in ("3Y", "5Y", "10Y"):
                    for series_name in ("Strategy", "SPY (AUD)", _au_lbl):
                        col_key = (h, series_name)
                        if col_key not in oos_mtx.columns:
                            continue
                        row = []
                        for m in display_metrics:
                            v = oos_mtx.at[m, col_key] if m in oos_mtx.index else np.nan
                            row.append(v)
                        rows.append(row)
                        if series_name == "Strategy" and _has_rs_mtx:
                            row_labels.append(
                                f"{h} — Strategy @ {_user_label_suffix}")
                        else:
                            row_labels.append(f"{h} — {series_name}")

                    if _has_rs_mtx:
                        col_key_rs = (h, "Strategy")
                        if col_key_rs in _oos_mtx_rs.columns:
                            row_rs = []
                            for m in display_metrics:
                                v = (_oos_mtx_rs.at[m, col_key_rs]
                                     if m in _oos_mtx_rs.index else np.nan)
                                row_rs.append(v)
                            rows.append(row_rs)
                            row_labels.append(
                                f"{h} — Strategy @ {_rs_label_suffix}")

                if rows:
                    n_rows = len(rows) + 1  # +1 header
                    n_cols = len(display_metrics) + 1  # +1 row label
                    # Table position pinned flush against chart bottom: top=11.39
                    # = chart bottom (2.78 + 8.61 = 11.39), no gap (user spec
                    # 2026-06-21). Width 22.6, height 6.77 — bottom ends at
                    # 18.16cm which fits within the 19.05cm slide height.
                    tbl_shape = road.shapes.add_table(
                        n_rows, n_cols,
                        Cm(1.4), Cm(11.39), Cm(22.6), Cm(6.77)
                    )
                    tbl = tbl_shape.table
                    # Header row
                    tbl.cell(0, 0).text = "Horizon — Series"
                    for j, m in enumerate(display_metrics, start=1):
                        tbl.cell(0, j).text = m
                    # Data rows
                    pct_metrics = {"Annualised Return", "Annualised Volatility",
                                   "Max Drawdown", "Alpha vs FF5 (ann)"}
                    for i, (label, row) in enumerate(zip(row_labels, rows), start=1):
                        tbl.cell(i, 0).text = label
                        for j, (m, v) in enumerate(zip(display_metrics, row), start=1):
                            if v is None or (isinstance(v, float) and not np.isfinite(v)):
                                txt = ""
                            elif m in pct_metrics:
                                txt = f"{v*100:+.2f}%"
                            else:
                                txt = f"{v:.2f}"
                            tbl.cell(i, j).text = txt
                    # Format: bold + colour Strategy rows
                    for rr in range(n_rows):
                        for cc in range(n_cols):
                            cell = tbl.cell(rr, cc)
                            cell.text_frame.word_wrap = False
                            for p in cell.text_frame.paragraphs:
                                p.font.size = Pt(10)
                                p.alignment = PP_ALIGN.CENTER
                                if rr == 0 or (rr > 0 and "Strategy" in row_labels[rr-1]):
                                    p.font.bold = True

                _rs_terminal_msg = (
                    f", Fund@RS({_nav_label(_rs_base)}) 10y = "
                    f"{(w_strat_rs.iloc[-1]-1.0)*100:+.1f}%"
                    if _has_rs_strat and not w_strat_rs.empty
                    else ""
                )
                print(f"[pptx] Roadshow slide built (normalised view) — "
                      f"Fund @ {_nav_label(_user_nav)} 10y = "
                      f"{(w_strat.iloc[-1]-1.0)*100:+.1f}%"
                      f"{_rs_terminal_msg}, "
                      f"SPY = {(w_spy.iloc[-1]-1.0)*100:+.1f}%")
        except Exception as e:
            print(f"[pptx] Roadshow slide skipped: {e}")

        # Reorder: move the roadshow slide (last added) into position 2.
        try:
            _xml_slides = prs.slides._sldIdLst
            _sl = list(_xml_slides)
            if len(_sl) >= 3:
                _last = _sl[-1]
                _xml_slides.remove(_last)
                _xml_slides.insert(1, _last)
        except Exception as _e_reorder:
            print(f"[pptx] Roadshow reorder skipped: {_e_reorder}")

        # ---- SCALE-SENSITIVITY SLIDE ----
        # Strategy at N NAVs (default $100k / $250k / $500k / $1M) on
        # a single normalised chart + per-NAV metrics table. Gated by
        # SCALE_SENSITIVITY=1 because the extra backtests cost ~100s
        # each. Inserted after Roadshow in the deck (position 3).
        _scale_results_pptx = globals().get("oos_scale_results", {})
        _scale_metrics_pptx = globals().get("oos_scale_metrics", {})
        if _scale_results_pptx and _scale_metrics_pptx:
            try:
                scale_layout = prs.slide_layouts[20]
                scale_slide = prs.slides.add_slide(scale_layout)
                if scale_slide.shapes.title:
                    scale_slide.shapes.title.text = "Performance by Portfolio Scale"

                # Common 10Y window aligned to the longest available series.
                _scale_end = max(p["returns"].index.max()
                                 for p in _scale_results_pptx.values()
                                 if isinstance(p["returns"], pd.Series)
                                 and not p["returns"].empty)
                _scale_start = _scale_end - pd.DateOffset(years=10)
                _add_date_callout(scale_slide, _scale_start, _scale_end,
                                  prefix="Backtest")

                # SPY benchmark — common reference across scales.
                def _bench_pct(col):
                    if col not in oos_prices_aud_long.columns:
                        return pd.Series(dtype=float)
                    px = pd.to_numeric(oos_prices_aud_long[col], errors="coerce").dropna()
                    return px.pct_change().dropna()
                _spy_full = _bench_pct("SPY")

                fig_sc, (ax_sc, ax_sc_mix) = plt.subplots(
                    2, 1, figsize=(11.5, 5.5),
                    gridspec_kw={"height_ratios": [3.5, 1.0]},
                    sharex=True,
                )

                # Blue palette ramping from light (small NAV) to dark
                # (large NAV) so the friction-tax visual is intuitive.
                _navs_sorted = sorted(_scale_results_pptx.keys())
                _palette = ["#a8c5e6", "#6c9bd2", "#3f73b8", "#1f4e8a", "#0a2f5c"]
                _line_colors = {nav: _palette[min(i, len(_palette)-1)]
                                for i, nav in enumerate(_navs_sorted)}

                _scale_terminals: list[tuple[str, float, str]] = []
                _common_idx = None
                for _nav in _navs_sorted:
                    _rets = _scale_results_pptx[_nav]["returns"]
                    _slice = _rets[(_rets.index >= _scale_start) &
                                   (_rets.index <= _scale_end)]
                    if _slice.empty:
                        continue
                    _w = (1.0 + _slice).cumprod()
                    if _common_idx is None:
                        _common_idx = _w.index
                    _lbl = (f"Fund @ ${_nav/1_000_000:,.2f}M"
                            if _nav >= 1_000_000
                            else f"Fund @ ${_nav/1000:,.0f}k")
                    ax_sc.plot(_w.index, _w.values, linewidth=2.0,
                               label=_lbl, color=_line_colors[_nav])
                    _scale_terminals.append((_lbl, float(_w.iloc[-1]),
                                              _line_colors[_nav]))

                if _common_idx is not None and not _spy_full.empty:
                    _spy_slice = _spy_full.reindex(_common_idx).fillna(0.0)
                    _w_spy_sc = (1.0 + _spy_slice).cumprod()
                    ax_sc.plot(_w_spy_sc.index, _w_spy_sc.values,
                               linewidth=1.6, label="SPY (AUD)",
                               color="#c53030", alpha=0.85)
                    _scale_terminals.append(("SPY", float(_w_spy_sc.iloc[-1]),
                                              "#c53030"))

                ax_sc.axhline(1.0, color="#888888", linewidth=0.8,
                              linestyle=":", alpha=0.6)
                ax_sc.set_xlim(_scale_start, _scale_end)
                ax_sc.xaxis.set_major_locator(mdates.YearLocator())
                ax_sc.xaxis.set_major_formatter(mdates.DateFormatter("%Y"))
                ax_sc.yaxis.set_major_formatter(mtick.FuncFormatter(
                    lambda x, _p: f"{(x-1.0)*100:+.0f}%"))
                ax_sc.set_title(
                    f"Cumulative return by portfolio scale — normalised to "
                    f"0% start    (net of {BROKER_CONFIG['name']} brokerage "
                    f"+ AU CGT [{ACTIVE_CGT_PROFILE}])",
                    fontsize=10,
                )
                ax_sc.set_ylabel("Cumulative return")
                ax_sc.legend(loc="upper left", frameon=False, fontsize=8)
                ax_sc.grid(True, linestyle="--", alpha=0.4)

                # Auto-spread overlapping terminal annotations. With 4
                # Fund lines stacked within a 5pp band (e.g. $100k +297%,
                # $250k +319%, $500k +321%, $1M +333%) the per-line text
                # labels overprint at the right edge. Sort by y-value,
                # walk top-down, and nudge each label up only if the
                # previous one is closer than `_min_gap_data` units in
                # data-y space. Anchor points (the line endpoints) stay
                # exactly on the curves — only the text positions move.
                _y_min, _y_max = ax_sc.get_ylim()
                _min_gap_data = 0.030 * (_y_max - _y_min)
                _terminals_sorted = sorted(_scale_terminals,
                                            key=lambda t: t[1])
                _adj_y = [t[1] for t in _terminals_sorted]
                for _i in range(1, len(_adj_y)):
                    if _adj_y[_i] - _adj_y[_i-1] < _min_gap_data:
                        _adj_y[_i] = _adj_y[_i-1] + _min_gap_data
                for (lbl, val, col), _ty in zip(_terminals_sorted, _adj_y):
                    # Leader from line endpoint to nudged y if displaced.
                    if abs(_ty - val) > 1e-9:
                        ax_sc.annotate(
                            "",
                            xy=(_scale_end, val),
                            xytext=(_scale_end, _ty),
                            arrowprops=dict(arrowstyle="-",
                                             color=col,
                                             linewidth=0.6,
                                             alpha=0.5),
                            annotation_clip=False,
                        )
                    ax_sc.annotate(
                        f"  {(val-1.0)*100:+.0f}%",
                        xy=(_scale_end, _ty),
                        xytext=(8, 0),
                        textcoords="offset points",
                        va="center", fontsize=8,
                        fontweight="bold", color=col,
                        annotation_clip=False,
                    )

                # Bottom strip: ensemble regime mix from the primary
                # backtest (regime mix is NAV-independent — same engine,
                # same softmax).
                _soft_sc = globals().get("oos_softmax_history", pd.DataFrame())
                if (isinstance(_soft_sc, pd.DataFrame) and not _soft_sc.empty
                        and _common_idx is not None):
                    _regime_colors_sc = {
                        "Modest (SPY+0%)":      "#5b9bd5",
                        "Aggressive (SPY+5%)":  "#70ad47",
                        "Bold (SPY+10%)":       "#ffc000",
                        "Maximum (SPY+15%)":    "#ed7d31",
                        "Stretch (SPY+25%)":    "#c00000",
                    }
                    _cols_in_order_sc = [n for n in ENSEMBLE_SLOT_NAMES
                                         if n in _soft_sc.columns]
                    if _cols_in_order_sc:
                        _soft_plot_sc = _soft_sc[_cols_in_order_sc].reindex(
                            _common_idx, method="ffill").fillna(0.0)
                        ax_sc_mix.stackplot(
                            _soft_plot_sc.index,
                            *[_soft_plot_sc[c].values for c in _cols_in_order_sc],
                            labels=[c.split(" ")[0] for c in _cols_in_order_sc],
                            colors=[_regime_colors_sc.get(c, "#888888")
                                    for c in _cols_in_order_sc],
                            alpha=0.85,
                        )
                        ax_sc_mix.set_xlim(_scale_start, _scale_end)
                        ax_sc_mix.set_ylim(0, 1)
                        ax_sc_mix.yaxis.set_major_formatter(mtick.FuncFormatter(
                            lambda y, _p: f"{y*100:.0f}%"))
                        ax_sc_mix.set_ylabel("Regime")
                        ax_sc_mix.legend(loc="lower center",
                                         bbox_to_anchor=(0.5, -0.45),
                                         ncol=5, frameon=False, fontsize=7)

                import io as _io_sc
                _sc_buf = _io_sc.BytesIO()
                fig_sc.tight_layout()
                fig_sc.savefig(_sc_buf, format="png", dpi=180,
                               bbox_inches="tight")
                plt.close(fig_sc)
                _sc_buf.seek(0)

                from pptx.util import Cm as _CmSc
                scale_slide.shapes.add_picture(_sc_buf, _CmSc(1.4), _CmSc(2.78),
                                               width=_CmSc(22.6),
                                               height=_CmSc(8.61))

                # ---- Metrics table: NAV-as-row × horizon-as-col ----
                _horizons = ("3Y", "5Y", "10Y")
                _nav_label_row = lambda n: (
                    f"Strategy @ ${n/1_000_000:,.2f}M"
                    if n >= 1_000_000 else f"Strategy @ ${n/1000:,.0f}k"
                )

                # Row labels: one per NAV + SPY benchmark row.
                _row_labels_sc = [_nav_label_row(n) for n in _navs_sorted]
                _row_labels_sc.append("SPY (AUD)")

                # Header: Series | 3Y Ret/Shr/Sort/DD | 5Y... | 10Y...
                _header_cells = ["NAV / Series"]
                for h in _horizons:
                    _header_cells.extend([f"{h} Return", f"{h} Sharpe",
                                          f"{h} Sortino", f"{h} MaxDD"])
                _ncols_sc = len(_header_cells)
                _nrows_sc = len(_row_labels_sc) + 1
                # User-specified table geometry (2026-06-27):
                # position (left, top) = (0.59cm, 11.7cm)
                # size (width, height) = (24.39cm, 3.56cm)
                # Font: 8pt bold for every cell.
                _tbl_shape_sc = scale_slide.shapes.add_table(
                    _nrows_sc, _ncols_sc,
                    _CmSc(0.59), _CmSc(11.7),
                    _CmSc(24.39), _CmSc(3.56),
                )
                _tbl_sc = _tbl_shape_sc.table
                for j, h in enumerate(_header_cells):
                    _tbl_sc.cell(0, j).text = h

                def _fmt_pct(v):
                    if v is None or (isinstance(v, float) and not np.isfinite(v)):
                        return ""
                    return f"{v*100:+.2f}%"
                def _fmt_ratio(v):
                    if v is None or (isinstance(v, float) and not np.isfinite(v)):
                        return ""
                    return f"{v:.2f}"

                def _fill_horizon_block(row_idx, mtx, series_name):
                    """Write Return / Sharpe / Sortino / MaxDD for each
                    horizon into 4 cells starting at column 1. Returns
                    the next free column index after all horizons."""
                    j = 1
                    for h in _horizons:
                        _col = (h, series_name)
                        if _col not in mtx.columns:
                            j += 4
                            continue
                        _ret = (mtx.at["Annualised Return", _col]
                                if "Annualised Return" in mtx.index else None)
                        _shr = (mtx.at["Sharpe Ratio", _col]
                                if "Sharpe Ratio" in mtx.index else None)
                        _sor = (mtx.at["Sortino Ratio", _col]
                                if "Sortino Ratio" in mtx.index else None)
                        _dd = (mtx.at["Max Drawdown", _col]
                               if "Max Drawdown" in mtx.index else None)
                        _tbl_sc.cell(row_idx, j).text = _fmt_pct(_ret)
                        _tbl_sc.cell(row_idx, j+1).text = _fmt_ratio(_shr)
                        _tbl_sc.cell(row_idx, j+2).text = _fmt_ratio(_sor)
                        _tbl_sc.cell(row_idx, j+3).text = _fmt_pct(_dd)
                        j += 4

                # Per-NAV rows
                for i, _nav in enumerate(_navs_sorted, start=1):
                    _mtx_for_nav = _scale_metrics_pptx.get(_nav, pd.DataFrame())
                    _tbl_sc.cell(i, 0).text = _nav_label_row(_nav)
                    if not _mtx_for_nav.empty:
                        _fill_horizon_block(i, _mtx_for_nav, "Strategy")

                # SPY row — pull from the primary metrics table (SPY is
                # NAV-invariant: same daily returns regardless of strategy
                # NAV, so any of the existing per-NAV metric tables work).
                _spy_row_idx = len(_navs_sorted) + 1
                _tbl_sc.cell(_spy_row_idx, 0).text = "SPY (AUD)"
                _any_mtx = next(iter(_scale_metrics_pptx.values()),
                                pd.DataFrame())
                if not _any_mtx.empty:
                    _fill_horizon_block(_spy_row_idx, _any_mtx, "SPY (AUD)")

                # User-specified formatting (2026-06-27): every cell 8pt
                # bold. Iterate all rows × cols × paragraphs × runs so
                # both header and data rows pick up the styling.
                for _ri in range(_nrows_sc):
                    for _ci in range(_ncols_sc):
                        _cell = _tbl_sc.cell(_ri, _ci)
                        for _para in _cell.text_frame.paragraphs:
                            for _run in _para.runs:
                                _run.font.size = Pt(8)
                                _run.font.bold = True

                # Summary print: cumulative return at end of period for
                # each NAV. Series indexed by date, so we need iloc[-1]
                # AFTER cumprod — the raw 'returns' field is daily
                # percentage returns, not cumulative.
                def _nav_summary(n):
                    _r = _scale_results_pptx[n]["returns"]
                    if not isinstance(_r, pd.Series) or _r.empty:
                        return ("$1.00M" if n >= 1_000_000
                                else f"${n/1000:.0f}k") + " (no data)"
                    _w = (1.0 + _r).cumprod()
                    _final = float(_w.iloc[-1])
                    _lbl = (f"${n/1_000_000:.2f}M" if n >= 1_000_000
                            else f"${n/1000:.0f}k")
                    return f"@{_lbl} {(_final-1.0)*100:+.1f}%"

                print(f"[pptx] Scale slide built — "
                      f"{len(_navs_sorted)} NAVs: "
                      + ", ".join(_nav_summary(n) for n in _navs_sorted))

                # Reorder: insert immediately after Roadshow (position 3).
                try:
                    _xml_slides_sc = prs.slides._sldIdLst
                    _sl_sc = list(_xml_slides_sc)
                    if len(_sl_sc) >= 4:
                        _last_sc = _sl_sc[-1]
                        _xml_slides_sc.remove(_last_sc)
                        _xml_slides_sc.insert(2, _last_sc)
                except Exception as _e_reorder_sc:
                    print(f"[pptx] Scale slide reorder skipped: {_e_reorder_sc}")
            except Exception as _e_scale_slide:
                print(f"[pptx] Scale slide skipped: {_e_scale_slide}")

        # --- FINAL SLIDE: Headline Metrics Dashboard (dump slide) ---
        # All the key numbers in one place: build/config stamp, 10Y headline
        # metrics vs SPY, engine totals (TLH, CGT, brokerage), live regime
        # mix + top positions. Added 2026-06-19 per user request — useful as
        # a quick-reference tear sheet without flipping through the deck.
        try:
            dump_layout = prs.slide_layouts[20]
            dump_slide = prs.slides.add_slide(dump_layout)
            if dump_slide.shapes.title:
                dump_slide.shapes.title.text = "Engine Metrics Dashboard"

            # Build stamp + production config line (italic, in title ribbon)
            _prod_slot = (next(iter(PRODUCTION_SLOT_OVERRIDE.keys()))
                           if PRODUCTION_SLOT_OVERRIDE else "5-slot ensemble blend")
            _hedge_lbl = "ON" if PRODUCTION_CRASH_HEDGE else "off"
            _stamp_text = (f"Build {_BUILD_GIT_SHA} at {_BUILD_TIME}    |    "
                            f"Production: {_prod_slot}  ·  Crash hedge: {_hedge_lbl}  ·  "
                            f"CGT: {int(CGT_CONFIG['marginal_tax_rate']*100)}% MTR")
            _stamp_box = dump_slide.shapes.add_textbox(Cm(2.15), Cm(1.85), Cm(21.80), Cm(0.45))
            _stf = _stamp_box.text_frame
            _stf.clear()
            _stf.word_wrap = False
            _stf.margin_left = 0
            _stp = _stf.paragraphs[0]
            _stp.text = _stamp_text
            _stp.font.size = Pt(10)
            _stp.font.italic = True
            _stp.font.color.rgb = RGBColor(255, 255, 255)

            # === HEADLINE METRICS — Strategy vs SPY two-column comparison ===
            _mt = globals().get("oos_metrics_table", None)
            def _mt_get(metric, horizon, series):
                try:
                    return float(_mt.loc[metric, (horizon, series)])
                except Exception:
                    return None

            # --- Styled-table helper (colored header, zebra rows) so the dashboard
            #     reads as a tear-sheet of tables rather than a wall of text. ---
            def _styled_table(left_cm, top_cm, width_cm, header_cells, data_rows,
                              col_fracs=None, font_pt=11, row_h_cm=0.64,
                              header_fill=(31, 78, 161)):
                _nr, _nc = len(data_rows) + 1, len(header_cells)
                _shp = dump_slide.shapes.add_table(
                    _nr, _nc, Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(row_h_cm * _nr))
                _t = _shp.table
                try:
                    _t.first_row = False
                    _t.horz_banding = False
                except Exception:
                    pass
                for _j, _h in enumerate(header_cells):
                    _c = _t.cell(0, _j)
                    _c.text = str(_h)
                    _c.fill.solid(); _c.fill.fore_color.rgb = RGBColor(*header_fill)
                    _p = _c.text_frame.paragraphs[0]
                    _p.font.bold = True; _p.font.size = Pt(font_pt)
                    _p.font.color.rgb = RGBColor(255, 255, 255)
                    _p.alignment = PP_ALIGN.LEFT if _j == 0 else PP_ALIGN.CENTER
                for _i, _row_vals in enumerate(data_rows, start=1):
                    _shade = (234, 240, 248) if (_i % 2 == 1) else (255, 255, 255)
                    for _j, _val in enumerate(_row_vals):
                        _c = _t.cell(_i, _j)
                        _c.text = str(_val)
                        _c.fill.solid(); _c.fill.fore_color.rgb = RGBColor(*_shade)
                        _p = _c.text_frame.paragraphs[0]
                        _p.font.size = Pt(font_pt)
                        _p.font.bold = (_j == 0)
                        _p.font.color.rgb = RGBColor(40, 40, 40)
                        _p.alignment = PP_ALIGN.LEFT if _j == 0 else PP_ALIGN.CENTER
                for _r in range(_nr):
                    _t.rows[_r].height = Cm(row_h_cm)
                if col_fracs:
                    for _j, _fr in enumerate(col_fracs):
                        _t.columns[_j].width = int(Cm(width_cm) * _fr)
                for _c in _t.iter_cells():
                    _c.text_frame.word_wrap = False
                    _c.margin_left = Cm(0.15); _c.margin_right = Cm(0.1)
                    _c.margin_top = Cm(0.02); _c.margin_bottom = Cm(0.02)
                return _t

            def _section_label(left_cm, top_cm, text):
                _b = dump_slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(11.0), Cm(0.55))
                _tf = _b.text_frame; _tf.clear(); _tf.margin_left = 0; _tf.margin_top = 0
                _p = _tf.paragraphs[0]
                _p.text = text; _p.font.size = Pt(12); _p.font.bold = True
                _p.font.color.rgb = RGBColor(31, 78, 161)

            def _f_pct(v):
                return f"{v*100:+.2f}%" if v is not None else "n/a"
            def _f_ratio(v):
                return f"{v:.2f}" if v is not None else "n/a"

            _LX, _RX = 1.4, 13.2           # left / right column x-anchors
            _LW, _RW = 11.2, 10.8          # column widths

            # === (1) 10-Year headline vs SPY (left) ===
            _alpha_10y = _mt_get("Alpha vs SPY (ann)", "10Y", "Strategy")
            _section_label(_LX, 3.35, "10-Year Headline (since OOS start)")
            _hl_rows = [
                ["Annualised Return", _f_pct(_mt_get("Annualised Return", "10Y", "Strategy")),
                 _f_pct(_mt_get("Annualised Return", "10Y", "SPY (AUD)"))],
                ["Sharpe Ratio", _f_ratio(_mt_get("Sharpe Ratio", "10Y", "Strategy")),
                 _f_ratio(_mt_get("Sharpe Ratio", "10Y", "SPY (AUD)"))],
                ["Sortino Ratio", _f_ratio(_mt_get("Sortino Ratio", "10Y", "Strategy")),
                 _f_ratio(_mt_get("Sortino Ratio", "10Y", "SPY (AUD)"))],
                ["Max Drawdown", _f_pct(_mt_get("Max Drawdown", "10Y", "Strategy")),
                 _f_pct(_mt_get("Max Drawdown", "10Y", "SPY (AUD)"))],
                ["Alpha vs SPY (ann.)",
                 (f"{_alpha_10y*100:+.2f}%/yr" if _alpha_10y is not None else "n/a"), "—"],
            ]
            _styled_table(_LX, 3.95, _LW, ["Metric", "Strategy", "SPY (AUD)"],
                          _hl_rows, col_fracs=[0.46, 0.27, 0.27])

            # === (2) By-horizon (right) ===
            _section_label(_RX, 3.35, "Strategy by Horizon")
            _bh_rows = [
                ["Ann. Return", _f_pct(_mt_get("Annualised Return", "3Y", "Strategy")),
                 _f_pct(_mt_get("Annualised Return", "5Y", "Strategy")),
                 _f_pct(_mt_get("Annualised Return", "10Y", "Strategy"))],
                ["Sharpe", _f_ratio(_mt_get("Sharpe Ratio", "3Y", "Strategy")),
                 _f_ratio(_mt_get("Sharpe Ratio", "5Y", "Strategy")),
                 _f_ratio(_mt_get("Sharpe Ratio", "10Y", "Strategy"))],
            ]
            _styled_table(_RX, 3.95, _RW, ["Horizon", "3Y", "5Y", "10Y"],
                          _bh_rows, col_fracs=[0.40, 0.20, 0.20, 0.20])

            # === (3) Engine totals (left) ===
            _tlh_events_ds = globals().get("oos_tlh_events", []) or []
            _tlh_n = len(_tlh_events_ds)
            _tlh_loss = float(sum(e.get("loss_aud", 0.0) for e in _tlh_events_ds))
            _eff_st = _effective_cgt_rate(short_term=True)
            _eff_lt = _effective_cgt_rate(short_term=False)
            _tax_saved_est = _tlh_loss * (_eff_st + _eff_lt) / 2.0
            _cost_ser = globals().get("oos_rebalance_costs", pd.Series(dtype=float))
            _tax_ser = globals().get("oos_rebalance_taxes", pd.Series(dtype=float))
            _oos_rets = globals().get("oos_returns_daily", pd.Series(dtype=float))
            _years = max(len(_oos_rets) / ANNUAL_TRADING_DAYS, 1e-6) if isinstance(_oos_rets, pd.Series) else 1.0
            _brk_bps = (float(_cost_ser.sum()) / _years * 10_000) if not _cost_ser.empty else 0.0
            _cgt_bps = (float(_tax_ser.sum()) / _years * 10_000) if not _tax_ser.empty else 0.0
            _section_label(_LX, 8.7, "Engine Totals (10Y OOS window)")
            _et_rows = [
                ["TLH events", f"{_tlh_n}"],
                ["Loss realised", f"${_tlh_loss:,.0f}"],
                ["Tax saved (gross, est.)", f"${_tax_saved_est:,.0f}"],
                ["Brokerage drag", f"{_brk_bps:.0f} bps/yr"],
                ["CGT drag", f"{_cgt_bps:.0f} bps/yr"],
                ["Total cost", f"{_brk_bps + _cgt_bps:.0f} bps/yr"],
            ]
            _styled_table(_LX, 9.3, _LW, ["Metric", "Value"], _et_rows,
                          col_fracs=[0.62, 0.38])

            # === (4) Today — regime mix (right) ===
            _mix_live = globals().get("ensemble_mix_live", pd.Series(dtype=float))
            _abbr2 = {"Modest (SPY+0%)": "Modest", "Aggressive (SPY+5%)": "Aggressive",
                       "Bold (SPY+10%)": "Bold", "Maximum (SPY+15%)": "Maximum",
                       "Stretch (SPY+25%)": "Stretch"}
            if isinstance(_mix_live, pd.Series) and not _mix_live.empty:
                _section_label(_RX, 6.6, "Today — Regime Mix")
                _rm_rows = [[_abbr2.get(n, n), f"{float(_mix_live.get(n, 0))*100:.0f}%"]
                            for n in ENSEMBLE_SLOT_NAMES if n in _mix_live.index]
                _styled_table(_RX, 7.2, _RW, ["Slot", "Weight"], _rm_rows,
                              col_fracs=[0.62, 0.38])

            # === (5) Today — top positions (right) ===
            _w_live = globals().get("W_ENSEMBLE_SER", pd.Series(dtype=float))
            if isinstance(_w_live, pd.Series) and not _w_live.empty:
                _section_label(_RX, 11.6, "Today — Top 5 Positions")
                _tp_rows = [[str(k).replace(".AX", ""), f"{float(v)*100:.0f}%"]
                            for k, v in _w_live.nlargest(5).items()]
                _styled_table(_RX, 12.2, _RW, ["Ticker", "Weight"], _tp_rows,
                              col_fracs=[0.62, 0.38])

            # === Footer timestamp ===
            _ftext = (f"Generated {pd.Timestamp.now().isoformat(timespec='seconds')}    ·    "
                       f"Metrics from {len(_oos_rets)}-day OOS walk-forward")
            _foot_box = dump_slide.shapes.add_textbox(Cm(1.5), Cm(17.3), Cm(22.5), Cm(0.6))
            _fft = _foot_box.text_frame
            _fft.clear()
            _ffp = _fft.paragraphs[0]
            _ffp.text = _ftext
            _ffp.font.size = Pt(9)
            _ffp.font.italic = True
            _ffp.font.color.rgb = RGBColor(120, 120, 120)
            _ffp.alignment = PP_ALIGN.CENTER

            print(f"[pptx] Engine Metrics Dashboard slide added (last position)")
        except Exception as _e_dump:
            print(f"[pptx] Dashboard slide skipped: {_e_dump}")

        # --- FINAL FINAL SLIDE: PDS-style disclosure (placeholder) ---
        # Standard Product Disclosure Statement structure: fund identity,
        # strategy, target, benchmark, fees modelled, key risks, legal
        # disclosures. Conservative placeholder text — user can edit the
        # strings here when actual fund details are finalised.
        try:
            pds_layout = prs.slide_layouts[20]
            pds_slide = prs.slides.add_slide(pds_layout)
            if pds_slide.shapes.title:
                pds_slide.shapes.title.text = "Disclosure & Product Statement"

            # Fund identity sub-line in the title ribbon
            _pds_subtitle = ("Guina Family Managed Investments  ·  Active Asset Allocation "
                              "(multi-region equity ensemble)  ·  Benchmark: SPY (AUD)")
            _pds_sub_box = pds_slide.shapes.add_textbox(Cm(2.15), Cm(1.85), Cm(21.80), Cm(0.45))
            _psf = _pds_sub_box.text_frame
            _psf.clear()
            _psf.word_wrap = False
            _psf.margin_left = 0
            _psp = _psf.paragraphs[0]
            _psp.text = _pds_subtitle
            _psp.font.size = Pt(10)
            _psp.font.italic = True
            _psp.font.color.rgb = RGBColor(255, 255, 255)

            # Body sections — two-column layout
            def _section_box(left_cm, top_cm, width_cm, height_cm, header, body_lines,
                              header_size=12, body_size=10):
                box = pds_slide.shapes.add_textbox(
                    Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))
                tf = box.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.margin_left = Cm(0.2)
                tf.margin_top = Cm(0.1)
                # header
                ph = tf.paragraphs[0]
                ph.text = header
                ph.font.size = Pt(header_size)
                ph.font.bold = True
                ph.font.color.rgb = RGBColor(31, 78, 161)
                for line in body_lines:
                    p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(body_size)
                    p.font.color.rgb = RGBColor(40, 40, 40)

            # --- Live metrics for the disclosure (consistency with slides 2 & 7) ---
            # Read from the SAME globals the performance/dashboard slides use so the
            # PDS can never drift from the headline again (the stale 0.97/0.83 bug:
            # the disclosure carried June's numbers while slides 2/7 showed 1.03/0.81).
            # Every value falls back to conservative generic wording if unavailable.
            _pds_mt = globals().get("oos_metrics_table", None)
            def _pds_get(metric, horizon="10Y", series="Strategy"):
                try:
                    return float(_pds_mt.loc[metric, (horizon, series)])
                except Exception:
                    return None
            _pds_shr_s = _pds_get("Sharpe Ratio", "10Y", "Strategy")
            _pds_shr_b = _pds_get("Sharpe Ratio", "10Y", "SPY (AUD)")
            _pds_mdd_s = _pds_get("Max Drawdown", "10Y", "Strategy")
            _pds_cost = globals().get("oos_rebalance_costs", pd.Series(dtype=float))
            _pds_tax  = globals().get("oos_rebalance_taxes", pd.Series(dtype=float))
            _pds_rets = globals().get("oos_returns_daily", pd.Series(dtype=float))
            _pds_years = (max(len(_pds_rets) / ANNUAL_TRADING_DAYS, 1e-6)
                          if isinstance(_pds_rets, pd.Series) and len(_pds_rets) else 1.0)
            _pds_brk_bps = (float(_pds_cost.sum()) / _pds_years * 10_000
                            if isinstance(_pds_cost, pd.Series) and not _pds_cost.empty else None)
            _pds_cgt_bps = (float(_pds_tax.sum()) / _pds_years * 10_000
                            if isinstance(_pds_tax, pd.Series) and not _pds_tax.empty else None)
            _pds_sharpe_line = (
                f"  Backtest 10Y Sharpe ≈ {_pds_shr_s:.2f} vs SPY {_pds_shr_b:.2f}."
                if (_pds_shr_s is not None and _pds_shr_b is not None)
                else "  Backtest 10Y Sharpe ≈ 1.0 vs SPY ~0.8.")
            _pds_brk_line = (
                f"  · Modelled cost ~{_pds_brk_bps:.0f} bps/yr at current NAV."
                if _pds_brk_bps is not None
                else "  · Modelled cost varies with NAV (fixed minima).")
            _pds_cgt_line = (
                f"  · Modelled drag ~{_pds_cgt_bps:.0f} bps/yr at current NAV."
                if _pds_cgt_bps is not None
                else "  · Modelled drag varies with turnover + FY netting.")
            _pds_mdd_line = (
                f"• MaxDD: backtest {_pds_mdd_s*100:.0f}% modern, -25% GFC stress test."
                if _pds_mdd_s is not None
                else "• MaxDD: backtest ~-21% modern, -25% GFC stress test.")

            # --- Left column: Fund summary ---
            # Top blocks: body 9pt + height 8.0cm so the longer post-2026-06-22
            # content (trustee + investor lines) doesn't overflow into the
            # bottom blocks. Bottom blocks pushed to top=12.0 for clearance.
            _section_box(
                left_cm=1.0, top_cm=3.7, width_cm=11.5, height_cm=8.0,
                header="FUND SUMMARY", body_size=9,
                body_lines=[
                    "Fund:  Guina Family Managed Investments (managed fund).",
                    "Trustee / Responsible Entity:  Fionn Guina (AFSL pending).",
                    "Target Investor:  Wholesale only — HNW + Sophisticated",
                    "  Investors (Corporations Act s708 / s761G).",
                    "",
                    "Strategy:  5-slot regime-aware ensemble (Modest, Aggressive,",
                    "  Bold, Maximum, Stretch). Softmax-blended via rolling 12-month",
                    "  Sortino + forward SPY-regime signal. Rebalanced ~9×/year.",
                    "",
                    "Asset Class:  Multi-region equity ETFs (US + Australia +",
                    "  Europe + Japan + Emerging Markets), plus defensive sleeve",
                    "  (gold, cash equivalents, long bonds, inverse equity).",
                    "",
                    "Universe:  ~46 ETFs; ~45 pass FF5 universe validation each run.",
                    "",
                    "Target Return:  SPY (AUD)-comparable, risk-adjusted (Sharpe).",
                    _pds_sharpe_line,
                    "",
                    "Investment Horizon:  5+ years recommended.",
                ],
            )

            # --- Right column: Fees + Costs ---
            # Build fee block dynamically from the FUND_FEES_ACTIVE flag so
            # that flipping fees on later updates this slide automatically.
            _fee_status = "ACTIVE" if FUND_FEES_ACTIVE else "currently waived"
            _fee_lines = [
                f"Management Fee:  {MANAGEMENT_FEE_PCT_ANN*100:.1f}% per annum",
                f"  ({_fee_status}). Accrued daily on NAV.",
                f"Performance Fee:  {PERFORMANCE_FEE_PCT*100:.0f}% over high-water mark",
                f"  ({_fee_status}). Crystallised {PERFORMANCE_FEE_CRYSTALLISE_FREQ}-ly.",
                ("  Hurdle: none."
                 if PERFORMANCE_FEE_HURDLE_ANN <= 0 else
                 f"  Hurdle: {PERFORMANCE_FEE_HURDLE_ANN*100:.1f}% p.a."),
                "",
                "Brokerage:  Interactive Brokers Pro AU schedule",
                "  · AU min $5.00 + 0.080% · US min $1.50 + 0.020%",
                _pds_brk_line,
                "",
                "CGT:  Australian personal MTR 30% + Medicare 2%",
                "  · 50% LT discount on holdings ≥365 days",
                "  · FY netting + carry-forward losses honoured.",
                _pds_cgt_line,
                "",
                "FX:  AUD-denominated. USD assets unhedged.",
                "Liquidity:  Redemptions at quarterly rebalance only.",
                "Custody:  Interactive Brokers Australia (IBKR Pty Ltd).",
                "Distribution:  Direct payment only — no platforms.",
                "Auditor:  TBC.",
            ]
            _section_box(
                left_cm=13.0, top_cm=3.7, width_cm=11.5, height_cm=8.0,
                header="FEES & TERMS", body_size=9,
                body_lines=_fee_lines,
            )

            # --- Bottom-left: Key risks ---
            _section_box(
                left_cm=1.0, top_cm=12.0, width_cm=11.5, height_cm=5.0,
                header="KEY RISKS",
                body_lines=[
                    "• Equity-like volatility (backtest annualised ~14%).",
                    _pds_mdd_line,
                    "• Tail risk: 2008-class crash modelled to -25% drawdown;",
                    "  uncovered tail regimes (e.g. stagflation) untested.",
                    "• Concentration: individual ticker weight up to 5% (capped).",
                    "• Regime risk: design assumes regime patterns persist.",
                    "• FX risk: USD holdings unhedged.",
                    "• Liquidity: ETF universe; assumes daily fills at close.",
                ],
                body_size=9,
            )

            # --- Bottom-right: Legal disclosures ---
            _section_box(
                left_cm=13.0, top_cm=12.0, width_cm=11.5, height_cm=5.0,
                header="DISCLOSURES",
                body_lines=[
                    "• Wholesale-only offer under Corporations Act s708/s761G.",
                    "  Not available to retail investors.",
                    "• AFSL: pending application. Trustee operates under",
                    "  s911A exemptions until issued. Not licensed advice.",
                    "• Past performance ≠ future returns. Forecasts are model",
                    "  outputs, not promises.",
                    "• Backtest is OOS walk-forward but design choices were",
                    "  made on 2016-2026 window — pseudo-overfit risk remains.",
                    "• Backtest excludes: tracking error, slippage gates,",
                    "  liquidity constraints, dividend reinvestment timing.",
                    "• Tax: modelled at AU 30% MTR; investor's own MTR applies.",
                    "• Operational risk: see RUNBOOK.md. Single-operator key",
                    "  person risk acknowledged.",
                ],
                body_size=9,
            )

            # Footer
            _pds_foot_text = (
                f"Disclosure as at {pd.Timestamp.now().strftime('%d %B %Y')}    ·    "
                f"Build {_BUILD_GIT_SHA}    ·    Wholesale-only — Corporations Act s708/s761G    ·    "
                f"DRAFT — pending AFSL + legal review"
            )
            _pds_foot_box = pds_slide.shapes.add_textbox(Cm(1.5), Cm(17.3), Cm(22.5), Cm(0.6))
            _pft = _pds_foot_box.text_frame
            _pft.clear()
            _pfp = _pft.paragraphs[0]
            _pfp.text = _pds_foot_text
            _pfp.font.size = Pt(9)
            _pfp.font.italic = True
            _pfp.font.color.rgb = RGBColor(120, 120, 120)
            _pfp.alignment = PP_ALIGN.CENTER

            print(f"[pptx] PDS disclosure slide added (final position)")
        except Exception as _e_pds:
            print(f"[pptx] PDS slide skipped: {_e_pds}")

        # Atomic-ish save: write beside the target, then swap in.
        #
        # The swap FAILS with WinError 5 whenever the destination is open in
        # PowerPoint — and the engine opens the deck itself after saving
        # (OPEN_PPT_AFTER_SAVE), so any run whose deck is still on screen locks
        # the NEXT run out of its own output. Previously that raised straight
        # out of here: the caller swallowed it, the run reported "PPT generated:
        # FAILED / not saved", a ~12MB .__tmp__.pptx was orphaned, and the deck
        # was simply lost — while the wrapper still emailed the RUN verdict. On
        # an unattended run (Mon 09:30, deck left open over the weekend) that is
        # a silent loss of the one artefact the user reviews before trading.
        #
        # So: never lose the deck. Fall back to a timestamped sibling, which the
        # user can open and which leaves no debris.
        tmp_path = ppt_path.replace(".pptx", ".__tmp__.pptx")
        prs.save(tmp_path)
        try:
            os.replace(tmp_path, ppt_path)
        except OSError as _e_swap:
            _stamp = datetime.now().strftime("%Y-%m-%d_%H%M%S")
            _fallback = ppt_path.replace(".pptx", f".{_stamp}.pptx")
            try:
                os.replace(tmp_path, _fallback)
                print(f"[ppt][WARN] could not overwrite {os.path.basename(ppt_path)} "
                      f"({_e_swap.__class__.__name__}: {_e_swap}). It is most likely "
                      f"still open in PowerPoint. Deck saved instead to: {_fallback}")
                return _fallback
            except Exception as _e_fb:
                # Both paths failed — clean up so we don't strand a 12MB temp.
                try:
                    if os.path.exists(tmp_path):
                        os.remove(tmp_path)
                except Exception:
                    pass
                print(f"[ppt][ERROR] deck could not be saved: swap failed "
                      f"({_e_swap}), fallback failed ({_e_fb}). Temp cleaned up.")
                raise
        print(f"[ppt] Report saved to: {ppt_path}")
        return ppt_path
