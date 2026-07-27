"""PPT/report formatting + date-window primitives (module split #18, 2026-07-10).

Pure helpers used by export_to_ppt (which stays in the engine and imports these
back). No engine config — just pandas/numpy/pptx. Split out so the report layer's
reusable primitives are testable in isolation.
  _nearest_on_or_before / _period_total_return / _window_compound_total  date-window returns.
  _ppt_anchor / _format_perf_value / _add_date_callout / _add_perf_table /
  _add_change_run / add_header_footer                                    slide builders.
"""
from __future__ import annotations

import numpy as np
import pandas as pd
from dateutil.relativedelta import relativedelta

from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def _nearest_on_or_before(idx, dt):
    """Return the largest value in `idx` that is <= dt, or None if no such value exists.

    Previously returned `idx[0]` when dt fell before the series start, which
    violated the contract and caused _period_total_return to compute returns
    from idx[0] instead of returning NaN. Symptom: a 1-week-old Actual NAV
    row showed identical -0.83% for 3M, 6M, 12M, 3Y (all four lookups
    silently fell back to idx[0]). Both callers (_period_total_return and
    _window_compound_total) already guard `if start_dt is None: return NaN`,
    so returning None here makes those guards live again.
    """
    if len(idx) == 0:
        return None
    dt = pd.to_datetime(dt)
    pos = idx.searchsorted(dt, side="right") - 1
    if pos < 0:
        return None
    return idx[min(pos, len(idx) - 1)]


def _period_total_return(px, end_dt, months=None, years=None):
    """Price-based total return over the lookback window ending at end_dt."""
    s = pd.to_numeric(pd.Series(px), errors="coerce").dropna()
    if s.empty:
        return np.nan
    end_dt = _nearest_on_or_before(s.index, end_dt)
    if end_dt is None:
        return np.nan
    start_target = pd.to_datetime(end_dt)
    if years:
        start_target = start_target - relativedelta(years=int(years))
    if months:
        start_target = start_target - relativedelta(months=int(months))
    start_dt = _nearest_on_or_before(s.index, start_target)
    if start_dt is None:
        return np.nan
    v0 = float(s.loc[start_dt])
    v1 = float(s.loc[end_dt])
    if not np.isfinite(v0) or not np.isfinite(v1) or v0 == 0:
        return np.nan
    return (v1 / v0) - 1.0


def _window_compound_total(r, end_dt, months=None, years=None):
    """Returns-based compounded total return over the lookback window ending at end_dt."""
    r = pd.to_numeric(pd.Series(r), errors="coerce").dropna()
    if r.empty:
        return np.nan
    start_target = end_dt
    if years:
        start_target = start_target - relativedelta(years=years)
    if months:
        start_target = start_target - relativedelta(months=months)
    start_dt = _nearest_on_or_before(r.index, start_target)
    end_dt2 = _nearest_on_or_before(r.index, end_dt)
    if start_dt is None or end_dt2 is None or start_dt >= end_dt2:
        return np.nan
    rr = r.loc[start_dt:end_dt2]
    if rr.empty:
        return np.nan
    return float((1.0 + rr).prod() - 1.0)


def _ppt_anchor(slide, layout, name, fb_left_cm, fb_top_cm, fb_w_cm, fb_h_cm):
    """
    Look up a named shape on the slide or its layout and return its
    (left, top, width, height) in EMU. Falls back to the hardcoded
    cm values when no shape with that name exists.

    To customise positions: in PowerPoint, add a no-fill placeholder shape
    on the slide layout (e.g. layout 20), open Selection Pane, and rename
    the shape to match `name` (e.g. "chart_main", "table_perf").
    """
    for src in (slide, layout):
        if src is None:
            continue
        try:
            for shp in src.shapes:
                if getattr(shp, "name", "") == name:
                    return (shp.left, shp.top, shp.width, shp.height)
        except Exception:
            continue
    return (Cm(fb_left_cm), Cm(fb_top_cm), Cm(fb_w_cm), Cm(fb_h_cm))


def _format_perf_value(v, fmt="pct2"):
    """Format a numeric value for a perf-style table cell. Uncomputable values
    (NaN / non-finite — e.g. a 3Y column when live history is shorter than 3Y)
    render as 'n/a' rather than an empty cell, so a blank never reads as a bug."""
    if pd.isna(v):
        return "n/a"
    try:
        fv = float(v)
    except (TypeError, ValueError):
        return "n/a"
    if not np.isfinite(fv):
        return "n/a"
    if fmt == "pct2":
        return f"{fv*100:.2f}%"
    if fmt == "dec3":
        return f"{fv:.3f}"
    return str(fv)


def _add_date_callout(slide, start_dt, end_dt, prefix: str = "Data", extra: str = ""):
    """Add a left-aligned white callout under the slide title showing the data window.
    Makes it impossible to mis-read the chart's date range — Slide 3 anchors to live
    portfolio end, Slide 4 anchors to FF data end, and these may differ by ~1 month.

    `extra` appends a second clause on the SAME line (e.g. the live-performance
    summary). The banner is only tall enough for one subtitle line, so this keeps
    everything on one row (10pt when extra is present) rather than a second line
    that would clip against the banner's bottom edge."""
    try:
        tb = slide.shapes.add_textbox(Cm(2.032), Cm(1.92), Cm(22.8), Cm(0.7))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = False
        p = tf.paragraphs[0]
        _txt = (
            f"{prefix}: {pd.Timestamp(start_dt).strftime('%d %b %Y')}"
            f"  →  {pd.Timestamp(end_dt).strftime('%d %b %Y')}"
        )
        if extra:
            _txt += f"      ·      {extra}"
        p.text = _txt
        p.font.size = Pt(10) if extra else Pt(11)
        p.font.italic = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT
    except Exception as _e:
        print(f"[pptx] Date callout skipped: {_e}")


def _add_perf_table(slide, df_metrics, left, top, width, height,
                    title=None, value_fmt="pct2", font_pt=11):
    """Add a formatted PPT table from a DataFrame. Values formatted per `value_fmt`."""
    rows = df_metrics.shape[0] + 1
    cols = df_metrics.shape[1] + 1  # include row label column
    shp = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = shp.table
    tbl.cell(0, 0).text = str(title) if title else ""
    for j, c in enumerate(df_metrics.columns, start=1):
        tbl.cell(0, j).text = str(c)
    for i, (idx, row) in enumerate(df_metrics.iterrows(), start=1):
        tbl.cell(i, 0).text = str(idx)
        for j, c in enumerate(df_metrics.columns, start=1):
            tbl.cell(i, j).text = _format_perf_value(row[c], fmt=value_fmt)
    for r in range(rows):
        for c in range(cols):
            for p in tbl.cell(r, c).text_frame.paragraphs:
                p.font.size = Pt(font_pt)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
    # Give the row-label column ~34% of the width so long labels ("Actual NAV
    # (since 08 Jul)", "Current Book (pro-forma)", "Aggressive") don't wrap; split
    # the rest evenly across the value columns.
    try:
        _label_w = int(width * 0.34)
        _data_w = int((width - _label_w) / max(cols - 1, 1))
        tbl.columns[0].width = _label_w
        for _cj in range(1, cols):
            tbl.columns[_cj].width = _data_w
    except Exception:
        pass
    return tbl


def _add_change_run(paragraph, val, font_pt=14):
    """Add a coloured (+/-) change run after a number in a summary line."""
    run = paragraph.add_run()
    if val == 0:
        run.text = ""
        return
    sign = "+" if val > 0 else ""
    run.text = f" ({sign}{val:,.2f})"
    run.font.size = Pt(font_pt)
    if val > 0:
        run.font.color.rgb = RGBColor(0, 128, 0)
    elif val < 0:
        run.font.color.rgb = RGBColor(192, 0, 0)
    else:
        run.font.color.rgb = RGBColor(80, 80, 80)


def add_header_footer(slide, title_text: str, footer_text: str = ""):
    """Adds a consistent header and footer banner with text."""
    # Header banner
    header = slide.shapes.add_shape(
        1,  # mso_shape.rectangle
        Cm(0), Cm(0),
        Cm(25.4), Cm(2.54)
    )
    fill = header.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)  # dark navy
    header.line.fill.background()  # no border

    # Header text
    tf = header.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = 1  # centre

    # Footer banner
    footer = slide.shapes.add_shape(
        1, Cm(0), Cm(17.78), Cm(25.4), Cm(1.016)
    )
    fill = footer.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(230, 230, 230)
    footer.line.fill.background()

    # Footer text
    tf = footer.text_frame
    tf.text = footer_text or "Generated by Portfolio Optimiser"
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = 1
