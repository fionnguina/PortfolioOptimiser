"""
Grabbing Data From Yahoo Finance For Stock Build last updated 19/03/2026

Portfolio Optimiser - single-file application.

Pipeline (executes top to bottom):
   1. Imports
   2. Config, paths, trade-plan + validation helpers
   3. Data download: prices, Fama-French 5 + momentum factors, FX, benchmarks
   4. Holdings / tilts input dialog
   5. Optimisation engine: factor betas, tilt recommendations, efficient frontier
   6. Transaction costs & CGT
   7. Excel workbook writer
   8. Excel launcher
   9. PowerPoint report generator
  10. PowerPoint launcher

Run top-to-bottom as __main__ (via Main.py); bundled into the .exe by build_helper.py.
"""


# =====================================================================
# BLOCK 1 imports
# =====================================================================
# Standard library imports
import datetime
import hashlib
import io
import json
import math
import multiprocessing as mp
import os
import pathlib
import re
import shutil
import sys
import time
import zipfile
from datetime import datetime
from pathlib import Path

# Third-party imports
import cvxpy as cp
import matplotlib.dates as mdates
import matplotlib.pyplot as plt
import matplotlib.ticker as mtick
import numpy as np
import numpy.linalg as la
import openpyxl
from openpyxl import Workbook, load_workbook
from openpyxl.utils.dataframe import dataframe_to_rows
import pandas as pd
import pptx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Cm, Inches, Pt
import requests
import scipy.optimize
from scipy.optimize import minimize
from numpy.linalg import pinv
import statsmodels.api as sm
import xlwings as xw
import yfinance as yf
from dateutil.relativedelta import relativedelta

# Optional Windows-specific import
try:
    import win32com.client as win32
    HAS_WIN32COM = True
except ImportError:
    HAS_WIN32COM = False

# Project-local imports (extracted modules — Phase 4 module split, 2026-06-29)
from jsonl_logs import (
    append_trade_recommendation_log,
    _load_recommendation_log,
    append_live_nav_history,
    _load_live_nav_series,
    append_cash_ledger,
)
from metrics import (
    _annualized_sharpe,
    _series_metrics,
    _ir_vs_bench,
    _capm_alpha_beta,
    _ff5_alpha,
)
from ensemble import softmax_ensemble_weights
from cgt import (
    TRADE_DELTA_CANDIDATES,
    _trade_delta_col,
    _security_from_row,
    _is_long_term_au,
    _allocate_sale_to_lots,
    compute_cgt_tax,
    CGT_PROFILES,
    ACTIVE_CGT_PROFILE,
    CGT_CONFIG,
    LotBook,
    _effective_cgt_rate,
    compute_cgt_for_rebalance,
    au_financial_year,
    compute_fy_tax_ledger,
)
from drift import (
    DRIFT_MONTHLY_THRESH,
    DRIFT_CUMULATIVE_THRESH,
    DRIFT_DD_ALERT_THRESH,
    DRIFT_SLIPPAGE_BPS_THRESH,
    DRIFT_FEE_MULTIPLIER,
    _match_fill_to_recommendation,
    compute_fill_drift,
    compute_live_max_drawdown,
    compute_monthly_nav_drift,
    _print_drift_warnings,
)
from brokerage import (
    BROKER_PROFILES,
    ACTIVE_BROKER_PROFILE,
    BROKER_CONFIG,
    BROKERAGE,
    ASX_MIN_MARKETABLE_PARCEL_AUD,
    MIN_TRADE_VALUE,
    _market_of,
    suppress_small_trades_by_value,
    compute_brokerage,
    IBKR_DIVERGENCE_WARN_BPS,
    _ibkr_pick_price,
    apply_ibkr_price_override,
)
from tlh import (
    TLH_ENABLED,
    TLH_MIN_LOSS_PCT,
    TLH_COOLDOWN_DAYS,
    TLH_MIN_LOSS_AUD,
    _load_tlh_pairs,
    _run_tlh_pass,
    _build_lot_book_from_df,
    _load_tlh_cooldown_state,
    _save_tlh_cooldown_state,
)

# MV optimisation core (module split #18, 2026-07-09). Constants (caps + slots)
# stay canonical in this engine; they are synced into `solvers` once after both
# are defined (see the `_solvers.` assignments near the ENSEMBLE_SLOTS block).
import solvers as _solvers
from solvers import (
    max_sharpe_long_only,
    solve_frontier_point_cvxpy,
    solve_frontier_point_cvxpy_with_tilts,
    solve_candidate_portfolios,
    _ledoit_wolf_cc,
)

# Lot-book construction + expansion (module split #18, 2026-07-09). Pure
# functions, no engine globals; re-exported here for backward compat.
from lots import (
    _read_lots_from_path,
    _build_lots_from_fills_log,
    _build_lots_from_holdings,
    expand_with_lots,
)

# Actual (broker-truth) NAV series (module split #18, 2026-07-09). Engine syncs
# nav.APP_DIR after APP_DIR is defined (default broker-nav-log path).
import nav as _nav
from nav import (
    compute_actual_nav_series,
    _load_broker_nav_series,
    compute_actual_nav_series_spliced,
)

# Excel/PPT sheet writers + formatting utils (module split #18, 2026-07-09).
# Engine syncs excel_sheets.TARGET_PORTFOLIO_VALUE_AUD after it is defined.
import excel_sheets as _excel_sheets
from excel_sheets import (
    get_or_clear_sheet,
    set_truefalse_validation,
    set_number_formats,
    _ensure_actual_fills_sheet,
    _read_actual_fills,
    _write_drift_sheets,
    _write_cash_ledger_sheet,
    _write_tilts_sheet,
    _write_holdings_sheet,
    _autofit_table_width,
)

# USD/AUD FX conversion (module split #18, 2026-07-10). Engine syncs fx.fx_usdaud
# after it builds the live FX series (~line 2650). Re-exported for backward compat.
import fx as _fx
from fx import (
    _last_numeric,
    get_usd_aud_fx,
    fx_to_aud_for_tickers,
)

# PPT/report formatting + date-window primitives (module split #18, 2026-07-10).
# Pure helpers used by export_to_ppt (which stays here). Re-exported for compat.
from ppt_utils import (
    _nearest_on_or_before,
    _period_total_return,
    _window_compound_total,
    _ppt_anchor,
    _format_perf_value,
    _add_date_callout,
    _add_perf_table,
    _add_change_run,
    add_header_footer,
)

# PowerPoint deck export (module split #18, 2026-07-10). export_to_ppt is FUSED
# to engine runtime state — it reads ~55 engine globals (24 pipeline artifacts via
# globals().get + config/derived + helper fns). Rather than decouple, the engine
# syncs all 55 into ppt_export via _sync_ppt_export() immediately before each call
# (see ppt_export.py caveat). Names are AST-derived from the function body.
import ppt_export as _ppt_export
_PPT_EXPORT_INJECT = (
    "ACTIVE_CGT_PROFILE", "ANNUAL_TRADING_DAYS", "APP_DIR", "BROKER_CONFIG",
    "CGT_CONFIG", "ENSEMBLE_SLOT_NAMES", "EXPORT_DIR", "FUND_FEES_ACTIVE",
    "FY_TAX_LEDGER_DF", "LIVE_TLH_EVENTS", "MANAGEMENT_FEE_PCT_ANN",
    "PERFORMANCE_FEE_CRYSTALLISE_FREQ", "PERFORMANCE_FEE_HURDLE_ANN",
    "PERFORMANCE_FEE_PCT", "PRODUCTION_CRASH_HEDGE", "PRODUCTION_SLOT_OVERRIDE",
    "TLH_COOLDOWN_DAYS", "TLH_ENABLED", "TLH_MIN_LOSS_AUD", "TLH_MIN_LOSS_PCT",
    "TRADEPLAN_LABEL", "TRADEPLAN_WEIGHTS_SER", "W_ENSEMBLE_SER", "W_WITH_TILTS_SER",
    "_BUILD_GIT_SHA", "_BUILD_TIME", "_add_change_run", "_add_date_callout",
    "_add_perf_table", "_autofit_table_width", "_effective_cgt_rate",
    "_nearest_on_or_before", "_oos_starting_nav_aud", "_period_total_return",
    "_ppt_anchor", "_roadshow_nav_aud", "_trade_delta_col", "_window_compound_total",
    "compute_actual_nav_series_spliced", "ensemble_mix_live", "ff5_raw",
    "oos_metrics_table", "oos_metrics_table_roadshow", "oos_prices_aud_long",
    "oos_rebalance_costs", "oos_rebalance_taxes", "oos_returns_daily",
    "oos_returns_daily_roadshow", "oos_scale_metrics", "oos_scale_results",
    "oos_softmax_history", "oos_tlh_events", "portfolio_value_series", "prices",
    "returns_wide_df",
)


def _sync_ppt_export():
    """Push the engine names export_to_ppt reads (config + runtime pipeline state
    + helper fns) into the ppt_export module namespace so its globals() lookups
    resolve. Lenient .get() matches export_to_ppt's own tolerance for optional
    runtime artifacts (roadshow/scale paths). Call immediately before export."""
    _g = globals()
    for _n in _PPT_EXPORT_INJECT:
        setattr(_ppt_export, _n, _g.get(_n))


def export_to_ppt(results, trades, charts=None):
    """Back-compat shim: sync engine state into ppt_export, then delegate."""
    _sync_ppt_export()
    return _ppt_export.export_to_ppt(results, trades, charts)


# OOS ensemble walk-forward backtest engine (module split #18, 2026-07-10). The
# analytics core (run_oos + 7 helpers + _is_us_ticker) — a CLOSED move-set that
# calls only each other + already-extracted modules + libs. The engine syncs ~25
# config values into it via _sync_oos_engine() once after config is defined
# (called after TLH_PAIRS ~line 1005; config is static post-load). Re-exported
# for the 15 run_oos call sites + the live-pipeline helper calls.
import oos_engine as _oos_engine
from oos_engine import (
    run_oos_ensemble_walk_forward,
    _apply_crash_hedge,
    _check_crash_trigger,
    _apply_mu_shrinkage,
    _compute_trend_sleeve,
    blend_ensemble_signals,
    compute_forward_regime_signal,
    estimate_rebalance_cost_fraction,
    _is_us_ticker,
)
_OOS_ENGINE_INJECT = (
    "COV_SHRINKAGE", "CRASH_HEDGE_BASKET", "CRASH_HEDGE_DD_RELEASE",
    "CRASH_HEDGE_DD_TRIGGER", "CRASH_HEDGE_LOOKBACK_DAYS", "CRISIS_HEDGE_BAND_SD",
    "CRISIS_HEDGE_MA_DAYS", "CRISIS_HEDGE_TICKER", "CRISIS_HEDGE_WEIGHT",
    "EARLY_TRIGGER_DD_DEEPEN", "EARLY_TRIGGER_MIN_DAYS", "ENSEMBLE_SLOT_NAMES",
    "LT_DEFER_DD_CONDITIONAL", "LT_DEFER_RELEASE_DD", "LT_DEFER_WINDOW_DAYS",
    "MU_SHRINKAGE_LAMBDA", "PER_ASSET_WEIGHT_CAPS", "RETURN_OUTLIER_THRESHOLD",
    "SKIP_REBAL_DELTA", "SKIP_REBAL_DELTA_CALM", "STRETCH_FLOOR_CALM",
    "STRETCH_FLOOR_PREDICTIVE", "TLH_PAIRS", "TREND_SLEEVE_WEIGHT", "VOL_TARGET_ANNUAL",
)


def _sync_oos_engine():
    """Push the ~25 engine config values the OOS engine reads into oos_engine.
    Config is static after load, so one call (after config is defined) suffices;
    covers both run_oos and the live-pipeline's direct helper calls."""
    _g = globals()
    for _n in _OOS_ENGINE_INJECT:
        setattr(_oos_engine, _n, _g.get(_n))


# Research/diagnostic CLI drivers (module split #18, 2026-07-10). The 12 --flag
# research modes, extracted as leaf drivers. Inventory symtable-verified gap-free
# (incl. the conditional-def _apply_data_lockbox). The engine injects 6 SHARED
# helper fns + 9 config/runtime values via _sync_research_modes() before dispatch
# (all defined before the first `if flag:` block). Re-exported for the dispatch.
import research_modes as _research_modes
from research_modes import (
    _run_gfc_stress_test,
    _run_scale_analysis,
    _run_dev_validation,
    _run_rebal_skip_sweep,
    _run_turnover_penalty_sweep,
    _run_walk_forward_cv,
    _run_attribution,
    _run_crash_hedge_test,
    _run_crash_hedge_release_sweep,
    _run_stretch_only_test,
    _run_stretch_hedge_sweep,
    _run_tilted_ensemble_test,
)
_RESEARCH_INJECT = (
    "_apply_data_lockbox", "_evaluate_sweep_result", "_normalize_yfinance_close",
    "_print_sweep_verdict", "compute_oos_metrics",
    "ANNUAL_TRADING_DAYS", "APP_DIR", "CRASH_HEDGE_BASKET", "CRASH_HEDGE_DD_RELEASE",
    "CRASH_HEDGE_DD_TRIGGER", "CRASH_HEDGE_LOOKBACK_DAYS", "ENSEMBLE_SLOT_NAMES",
    "REBALANCE_FREQ", "prices",
)


def _sync_research_modes():
    """Push the 6 shared helper fns + 9 config/runtime values the research drivers
    read into research_modes. Call once before the first dispatch — all 15 names
    are defined by then (config + the helpers, incl. conditional _apply_data_lockbox)."""
    _g = globals()
    for _n in _RESEARCH_INJECT:
        setattr(_research_modes, _n, _g.get(_n))

# Debug: Print Python executable path
print(sys.executable)

# === Build stamp (F7/L15) ====================================================
# _version.py is auto-generated by build_helper.py before PyInstaller runs and
# bundled into the .exe. When running from source it may not exist — falls
# back to a "dev" marker.
try:
    from _version import GIT_SHA as _BUILD_GIT_SHA, BUILD_TIME as _BUILD_TIME
except Exception:
    _BUILD_GIT_SHA = "dev"
    _BUILD_TIME = "n/a"
print(f"[build] version: GIT_SHA={_BUILD_GIT_SHA}  BUILD_TIME={_BUILD_TIME}")

# Track script start time for the [health] runtime line at the end of every
# live run (see _print_run_health_summary at bottom of file).
import time as _time_for_health
_SCRIPT_START_TIME = _time_for_health.perf_counter()


# === Swallowed-exception helper (F1/F3/L18) ==================================
# Use inside `except: pass` blocks where pass is the desired runtime behaviour
# but you want diagnostics available. Default: print exc type + a one-line
# summary so swallows aren't truly invisible. With DEBUG_SWALLOWED=1 env var:
# full traceback. Centralised so we have one switch for all of them.
def _log_swallowed(stage: str = "") -> None:
    import traceback as _tb
    exc_type = sys.exc_info()[0]
    exc_msg = sys.exc_info()[1]
    exc_name = exc_type.__name__ if exc_type else "Unknown"
    if os.environ.get("DEBUG_SWALLOWED", "0") == "1":
        print(f"[swallowed][{stage}]\n{_tb.format_exc()}")
    else:
        print(f"[swallowed][{stage}] {exc_name}: {exc_msg} "
              f"(set DEBUG_SWALLOWED=1 for traceback)")

# === CLI flags ===
# `--stress-test`: skip live dialog + PPT pipeline, jump to GFC stress runner
# (defined later, just before "# === Rebuild core analytics ===").
_STRESS_TEST_MODE = "--stress-test" in sys.argv
if _STRESS_TEST_MODE:
    print("[stress] --stress-test detected; will skip dialog + live pipeline")
# `--scale-analysis`: skip dialog + run OOS walk-forward at 6 NAV scales to
# produce a fee-drag-vs-AUM evidence pack. See _run_scale_analysis().
_SCALE_ANALYSIS_MODE = "--scale-analysis" in sys.argv
if _SCALE_ANALYSIS_MODE:
    print("[scale] --scale-analysis detected; will skip dialog + live pipeline")
# `--dev-validation`: skip dialog + run OOS twice (dev window vs locked-box
# validation window) to expose meta-parameter overfitting. See
# _run_dev_validation(). Lock box: 2020-02-20 → today; dev: 2015-01-01 →
# 2020-02-19 (SPY pre-COVID ATH).
_DEV_VALIDATION_MODE = "--dev-validation" in sys.argv
if _DEV_VALIDATION_MODE:
    print("[devval] --dev-validation detected; will skip dialog + live pipeline")
# `--rebal-skip-sweep`: dev-only sweep of SKIP_REBAL_DELTA at 3/4/5/6/7%
# to find the highest-Sharpe knob value, then open VALIDATION lock-box ONCE
# on the winner. Used to test if higher skip-delta improves net Sharpe via
# more 12-month-LT-discount qualifications.
_REBAL_SKIP_SWEEP_MODE = "--rebal-skip-sweep" in sys.argv
if _REBAL_SKIP_SWEEP_MODE:
    print("[skip-sweep] --rebal-skip-sweep detected; will skip dialog + live pipeline")
# `--turnover-penalty-sweep`: dev-only sweep of the cost-aware solver
# penalty (γ_cgt × ||w - w_prev||_1 inside cvxpy). Tests if penalising
# regime-switch turnover at the optimiser level reduces CGT drag enough
# to lift net Sharpe. Picks DEV winner, then opens VAL lock-box once.
_TURNOVER_SWEEP_MODE = "--turnover-penalty-sweep" in sys.argv
if _TURNOVER_SWEEP_MODE:
    print("[turnover-sweep] --turnover-penalty-sweep detected; will skip dialog + live pipeline")
# `--walk-forward-cv`: run engine once on full history, slice OOS returns
# into non-overlapping calendar-year folds, report per-fold Sharpe + mean ±
# std. Gives N independent OOS observations of the current engine config
# instead of the single point estimate from dev/validation split. Used as
# the preferred statistical hygiene tool for parameter selection — does NOT
# burn the validation lock-box budget.
_WALK_FORWARD_CV_MODE = "--walk-forward-cv" in sys.argv
if _WALK_FORWARD_CV_MODE:
    print("[wf-cv] --walk-forward-cv detected; will skip dialog + live pipeline")
# `--attribution`: decompose OOS returns by slot, asset, and regime to
# answer "where does the engine earn its money / lose to SPY?" Purely
# descriptive — no parameters to tune, no validation budget consumed.
_ATTRIBUTION_MODE = "--attribution" in sys.argv
if _ATTRIBUTION_MODE:
    print("[attr] --attribution detected; will skip dialog + live pipeline")
# `--crash-hedge-test`: walk-forward CV with crash hedge ON vs OFF.
# Reports per-fold uplift to test the asymmetric defensive overlay.
_CRASH_HEDGE_TEST_MODE = "--crash-hedge-test" in sys.argv
if _CRASH_HEDGE_TEST_MODE:
    print("[hedge] --crash-hedge-test detected; will skip dialog + live pipeline")
# `--crash-hedge-release-sweep`: walk-forward CV across release thresholds
# {-3%, -5%, -8%, -10%, -12%} with trigger fixed at -15%. Used to find
# the optimum release that balances slow-bear capture vs V-shape recovery.
_CRASH_HEDGE_RELEASE_SWEEP_MODE = "--crash-hedge-release-sweep" in sys.argv
if _CRASH_HEDGE_RELEASE_SWEEP_MODE:
    print("[hedge-rel] --crash-hedge-release-sweep detected; will skip dialog + live pipeline")
# `--stretch-only-test`: walk-forward CV comparing 5-slot ensemble vs
# Stretch-only (slot_weights_override forcing 100% on Stretch slot).
# Tests whether the defensive layer is worth its alpha cost.
_STRETCH_ONLY_TEST_MODE = "--stretch-only-test" in sys.argv
if _STRETCH_ONLY_TEST_MODE:
    print("[stretch] --stretch-only-test detected; will skip dialog + live pipeline")
# `--stretch-hedge-sweep`: Stretch-only + crash hedge across release values.
# The synthesis test — bull alpha from Stretch + tail protection from hedge.
_STRETCH_HEDGE_SWEEP_MODE = "--stretch-hedge-sweep" in sys.argv
if _STRETCH_HEDGE_SWEEP_MODE:
    print("[stretch-hedge] --stretch-hedge-sweep detected; will skip dialog + live pipeline")
# `--show-metrics-history`: print the metrics_history.jsonl as a comparison
# table so we can spot regressions across versions at a glance.
_SHOW_METRICS_HISTORY_MODE = "--show-metrics-history" in sys.argv
if _SHOW_METRICS_HISTORY_MODE:
    print("[metrics-history] --show-metrics-history detected; will print history then exit")
# `--factor-recs`: preview the trailing-3M factor-momentum auto-recommendation.
# Pure preview — no engine impact. Use to see what tilts WOULD be applied if
# auto-tilts were turned on as a baseline.
_FACTOR_RECS_MODE = "--factor-recs" in sys.argv
if _FACTOR_RECS_MODE:
    print("[factor-recs] --factor-recs detected; will print recommendation table then exit")
# `--preflight`: fast system-check mode. Validates yfinance reachability,
# Excel COM, workbook lock state, IBKR TWS (if enabled), config files,
# disk space — and exits in seconds with a PASS/FAIL list. Catches setup
# issues BEFORE the user spends 5 minutes on a heavy pipeline that crashes.
_PREFLIGHT_MODE = "--preflight" in sys.argv
if _PREFLIGHT_MODE:
    print("[preflight] --preflight detected; will run system checks then exit")
# `--auto-pipeline`: non-interactive run for the scheduled daily wrapper.
# Skips the dialog (forces TRADE_PLAN_MODE='ensemble') and surfaces a
# [rebal-trigger] verdict line in run.log so a wrapper script can decide
# whether to notify the user / kick off Phase 3 execution.
_AUTO_PIPELINE_MODE = "--auto-pipeline" in sys.argv
if _AUTO_PIPELINE_MODE:
    print("[auto-pipeline] --auto-pipeline detected; non-interactive run, "
          "ensemble mode, rebal-trigger verdict written to run.log")
# `--tilted-ensemble-test`: walk-forward CV with auto-factor-tilts ON vs OFF.
# Threads dynamic tilt targets (recomputed per rebal from trailing-3M factor
# Sharpes) through solve_candidate_portfolios so each ensemble slot expresses
# the factor view. Compares to baseline (no tilts) on the modern OOS window.
_TILTED_ENSEMBLE_TEST_MODE = "--tilted-ensemble-test" in sys.argv
if _TILTED_ENSEMBLE_TEST_MODE:
    print("[tilted-ens] --tilted-ensemble-test detected; will skip dialog + live pipeline")
# OOS kernel mode (Phase 3b, 2026-06-29) — workers spawned by the parallel
# scale-sensitivity loop set OOS_KERNEL_MODE=1 so this re-exec only loads
# imports + function defs, then sys.exits at the sentinel near the OOS
# function. Folded into _SKIP_LIVE_PIPELINE so workers bypass the holdings
# dialog (otherwise each worker pops its own Tk window before reaching the
# sentinel — see screenshot from 2026-06-29). Full definition + comment
# block lives below near line 2386.
OOS_KERNEL_MODE = bool(os.environ.get("OOS_KERNEL_MODE", "").strip())
# All diagnostic modes follow the same skip-everything-heavy path.
_SKIP_LIVE_PIPELINE = (_STRESS_TEST_MODE or _SCALE_ANALYSIS_MODE
                       or _DEV_VALIDATION_MODE or _REBAL_SKIP_SWEEP_MODE
                       or _TURNOVER_SWEEP_MODE or _WALK_FORWARD_CV_MODE
                       or _ATTRIBUTION_MODE or _CRASH_HEDGE_TEST_MODE
                       or _CRASH_HEDGE_RELEASE_SWEEP_MODE
                       or _STRETCH_ONLY_TEST_MODE
                       or _STRETCH_HEDGE_SWEEP_MODE
                       or _SHOW_METRICS_HISTORY_MODE
                       or _FACTOR_RECS_MODE
                       or _TILTED_ENSEMBLE_TEST_MODE
                       or _PREFLIGHT_MODE
                       or OOS_KERNEL_MODE)


# =====================================================================
# BLOCK 2 Global codes and Data Retrieval from the web
# =====================================================================
# ---------------------------------------------------------------------
# Central base directory
# ---------------------------------------------------------------------
def _app_dir() -> Path:
    """
    Determine the application directory dynamically:
      - When frozen (PyInstaller): use the exe folder
      - When run as a script: use the script's folder
      - When interactive (Jupyter/IPython): use cwd
    """
    if getattr(sys, "frozen", False):
        return Path(sys.executable).parent
    if "__file__" in globals():
        return Path(__file__).resolve().parent
    return Path(os.getcwd())


# Absolute path to your central config root (for dev use)
_DEV_BASE = Path.home() / "Portfolio_Optimiser"

# Use the dev folder if it exists, otherwise fall back to dynamic app dir
APP_DIR = _DEV_BASE if _DEV_BASE.exists() else _app_dir()
# Sync the resolved APP_DIR into the nav module (its only engine coupling — the
# default broker-nav-log path). Same pattern as factors.REGIONS_JSON_PATH.
_nav.APP_DIR = APP_DIR

# ---------------------------------------------------------------------
# Config file and Excel workbook paths
# ---------------------------------------------------------------------
def _default_excel_path() -> str:
    """Return full path to the default Excel workbook."""
    app_dir = Path(APP_DIR)  # APP_DIR might be a string in the notebook
    return str((app_dir / "Stock Analysis.xlsm").resolve())

CONFIG_PATH = APP_DIR / "config.json"

# ---------------------------------------------------------------------
# Export directory (for generated reports)
# ---------------------------------------------------------------------
EXPORT_DIR = APP_DIR / "Reports"
EXPORT_DIR.mkdir(exist_ok=True)

# ---------------------------------------------------------------------
# Default configuration values
# ---------------------------------------------------------------------
_DEFAULTS = {
    "excel_path": _default_excel_path(),
    "marginal_tax_rate": 0.37,
    "carry_forward_losses": 0.0,
    "lot_match_method": "HIFO",
    "open_after_save": True,
    "use_xlwings": True,
    "open_excel_after_save": True,
    "open_ppt_after_save": True,
    # OOS validation: walk-forward backtest + roadshow slide. Default ON so
    # every build refreshes the OOS_Validation sheet and the executive-summary
    # slide. Set to False in config.json to skip (saves ~5–10s per run).
    "oos_validation": True,
}

# ---------------------------------------------------------------------
# Config loader
# ---------------------------------------------------------------------
def load_config() -> dict:
    """Load configuration from config.json, falling back to defaults."""
    cfg = _DEFAULTS.copy()
    try:
        if CONFIG_PATH.exists():
            with CONFIG_PATH.open("r", encoding="utf-8") as f:
                user_cfg = json.load(f)
            cfg.update({k: v for k, v in user_cfg.items() if k in cfg})
    except Exception as e:
        print(f"[config] Using defaults (error reading config.json): {e}")

    # Ensure workbook directory exists
    try:
        Path(cfg["excel_path"]).parent.mkdir(parents=True, exist_ok=True)
    except Exception:
        pass

    return cfg

# ---------------------------------------------------------------------
# PPTX handling
# ---------------------------------------------------------------------
CFG = load_config()

def open_ppt_if_enabled(pptx_path: str) -> None:
    """Open PPTX file if enabled in config."""
    if not CFG.get("open_ppt_after_save", True):
        return

    pptx_path = os.path.abspath(pptx_path)

    # Kill any orphaned PowerPoint processes from prior notebook runs (Windows only)
    _kill_orphan_powerpoint()

    # Always use OS open to launch the user's normal PowerPoint UI instance
    _os_open(pptx_path)


def _os_open(path: str) -> None:
    """Open file with default OS application."""
    try:
        os.startfile(path)  # Windows
    except AttributeError:
        import subprocess
        if sys.platform == "darwin":
            subprocess.run(["open", path])
        else:
            subprocess.run(["xdg-open", path])


def _kill_orphan_powerpoint() -> None:
    """Kill orphaned PowerPoint processes (Windows only)."""
    try:
        import subprocess
        subprocess.run(
            ["taskkill", "/F", "/IM", "POWERPNT.EXE"],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
    except Exception:
        pass


user_opts = {}

# ---------------------------------------------------------------------
# Validation Code
# ---------------------------------------------------------------------

# Constants
TRADE_PLAN_MODE = "ask"  # Options: "ask", "auto", "no_tilts", "with_tilts", "ensemble"
# Force non-interactive ensemble mode when scheduled wrapper invokes the engine.
if _AUTO_PIPELINE_MODE:
    TRADE_PLAN_MODE = "ensemble"
    print(f"[auto-pipeline] TRADE_PLAN_MODE forced to 'ensemble'")
VALIDATION_LOOKBACK_DAYS = 252  # 1 year of daily data
ANNUAL_TRADING_DAYS = 252  # For Sharpe calculation


# ============================================================================
# BROKER TRANSACTION COST MODEL
# ----------------------------------------------------------------------------
# Each profile carries BOTH the OOS-backtest cost model (flat fee + spread bps
# + FX bps) AND the live trade-plan brokerage model (min_fee + % rate). Switch
# brokers by changing ACTIVE_BROKER_PROFILE — backtest and live stay in sync.
# All bps values are in basis points (5 = 5 bps = 0.05% = 0.0005 fraction).
# All AUD fee values are in dollars.
# ============================================================================
# NOTE: BROKER_PROFILES + ACTIVE_BROKER_PROFILE + BROKER_CONFIG canonical
# definitions moved to brokerage.py (Phase 4 split, 2026-06-29). Imported at
# top of file. The original 100+ line block lived here.
# ============================================================================
# CAPITAL GAINS TAX MODEL (Australian rules)
# ----------------------------------------------------------------------------
# Used by the OOS backtest to produce NET-of-tax returns. When the user's
# marginal tax bracket changes, or AU budget passes new CGT legislation,
# update THIS BLOCK only — everything downstream picks it up.
#
# Current AU rules applied:
#   • Marginal tax rate (MTR) applies to net taxable capital gain
#   • 50% discount on gains from assets held >= 365 days
#   • Within-rebalance loss offset (ST losses offset ST gains first, then
#     spill over to LT gains, and vice versa)
#   • Carry-forward losses NOT modelled here (within-rebalance only) — TODO
#
# Future budget changes can be modelled by editing fields below (e.g. if the
# 50% discount is reduced to 40%, set lt_discount_rate = 0.40).
# ============================================================================
# CGT_PROFILES + ACTIVE_CGT_PROFILE + CGT_CONFIG canonical definitions moved
# to cgt.py (Phase 4 split, 2026-06-29). Imported at top of file.


# ============================================================================
# DRIFT TRACKER (Tier-1 #3)
# ----------------------------------------------------------------------------
# Drives live-vs-backtest comparison + execution-quality monitoring. Set
# LIVE_TRADING_START_DATE to the first day of live trading (e.g. "2026-08-01")
# once IBKR is connected. Until then drift NAV tracking stays dormant but the
# infrastructure (recommendation log, fill sheet, NAV history) keeps running.
# ============================================================================
LIVE_TRADING_START_DATE: str | None = "2026-06-22"  # paper trading commenced this date; drift tracker active from here. Update to real-money start once AFSL issues + first live fill.
# DRIFT_* threshold canonical definitions moved to drift.py (Phase 4 split).
# Imported at top of file.
# "Drift vs target" in the cash ledger anchors to the account's ACTUAL NAV
# (portfolio_state.json), not a $1M nominal — otherwise it always reads ~-$750k
# for a ~$250k account, drowning any real signal. Falls back to $1M if the
# state file is missing (e.g. a fresh checkout).
TARGET_PORTFOLIO_VALUE_AUD = 1_000_000.0
try:
    _tgt_state = APP_DIR / "portfolio_state.json"
    if _tgt_state.exists():
        _tgt_pv = float(json.loads(_tgt_state.read_text(encoding="utf-8"))
                        .get("portfolio_value", 0) or 0)
        if np.isfinite(_tgt_pv) and _tgt_pv > 0:
            TARGET_PORTFOLIO_VALUE_AUD = _tgt_pv
except Exception:
    pass
# Sync into excel_sheets (its only engine coupling — the cash-ledger anchor).
_excel_sheets.TARGET_PORTFOLIO_VALUE_AUD = TARGET_PORTFOLIO_VALUE_AUD

# === Hybrid pricing (Tier-1 #1) ============================================
# When True + TWS paper connection works, the live trade plan uses IBKR's
# delayed last-price (free, ~15min lag) instead of yfinance's close. Fixes
# the after-hours/corporate-action gaps that yfinance has. 10y OOS backtest
# still uses yfinance for history depth + bulk download. Silent fallback to
# yfinance if TWS isn't running.
USE_IBKR_LIVE_PRICES   = True
IBKR_HOST              = "127.0.0.1"
IBKR_PORT              = 7497    # paper TWS (7496 = live; never use here)
IBKR_CLIENT_ID         = 10      # distinct from helper scripts (7/8/9)
IBKR_CONNECT_TIMEOUT   = 8
IBKR_SNAPSHOT_WAIT_SEC = 6
# IBKR_DIVERGENCE_WARN_BPS canonical definition moved to ibkr.py (Phase 4 split).

# Cash-fit sizing: size the live trade plan to (holdings + available cash -
# reserve) so net buys never exceed the cash on hand (see make_trade_plan +
# _get_available_cash_aud). Reserve = max($ floor, pct of investable), covering
# brokerage + FX spread + unit rounding so orders don't reject on funds.
CASH_RESERVE_PCT       = 0.005   # 0.5% of investable
CASH_RESERVE_MIN_AUD   = 300.0   # never reserve less than this
# Imported at top of file.


# ============================================================================
# REBALANCE FREQUENCY
# ----------------------------------------------------------------------------
# Drives how often the OOS engine (and live recommendation cadence) rebalances.
# Use pandas freq aliases:
#   "MS"  = month start  → 12 rebalances/year (default)
#   "QS"  = quarter start →  4 rebalances/year (more LT discount eligibility)
#   "6W"  = every 6 weeks →  ~8.7 rebalances/year (split-the-difference)
#   "YS"  = year start   →  1 rebalance/year (most tax efficient, least responsive)
# Lower frequency = lower brokerage + more LT-eligible CGT discount, but slower
# regime adaptation (misses fast events like 2020 COVID).
# ============================================================================
REBALANCE_FREQ = "6W"
REBALANCES_PER_YEAR = {"MS": 12, "QS": 4, "6W": 8.67, "YS": 1}.get(REBALANCE_FREQ, 12)


# ============================================================================
# PER-ASSET WEIGHT CAPS
# ----------------------------------------------------------------------------
# Some assets (leveraged ETFs, volatility products) have structural problems
# for mean-variance optimisation: their trailing μ is inflated by gamma drag
# and decay effects that don't appear in the covariance matrix. The solver
# would over-allocate to them based on backward-looking returns that don't
# predict forward performance.
#
# Capping their max weight per candidate portfolio keeps them OPTIONAL (the
# solver can still use them when their μ is good) without letting them
# dominate the allocation. 5% is a defensible default — leveraged exposure
# of that size adds meaningful return tail without crippling drawdowns.
#
# To disable a cap: remove the entry. To exclude a ticker entirely: set the
# cap to 0.0 (forces weight = 0).
# ============================================================================
PER_ASSET_WEIGHT_CAPS: dict[str, float] = {
    "SOXL":  0.05,   # 3x daily leveraged semis — high decay
    "TQQQ":  0.05,   # 3x daily leveraged NASDAQ — high decay
    "SVIX":  0.05,   # short-VIX 1x — looks high-μ in calm regimes, blows up in vol spikes
    "UVIX":  0.00,   # 2x leveraged long-VIX — structural ~-30%/yr decay, never useful long-only
    "VXX":   0.00,   # 1x long-VIX — structural decay, never useful long-only
    # --- Thematic ETFs: REVERTED to solver-excluded 2026-07-02 ------------
    # The expansion A/B/C (wf_cv_X_46t_equiv.log / walk_forward_cv_expanded
    # .log / wf_cv_Z_pruned.log) measured it at best Sharpe-neutral with
    # -1.3%/yr full-period return cost, and -0.07 Sharpe / -2.75% CAGR in
    # the production frame. Mean-variance error-maximisation: more noisy-μ
    # candidates → solver selects on estimation noise. Cap 0 = excluded
    # from allocation but still priced, which keeps SOXX / SEMI.AX usable
    # as TLH substitutes for SMH (tlh_pairs.json). Re-enable only if the
    # alpha-shrinkage experiment (task #13) validates through the full
    # gate (production frame, full-period, delta > 2×SE).
    "BOTZ":     0.00,   # AI / robotics
    "CIBR":     0.00,   # cybersecurity
    "ITA":      0.00,   # aerospace & defense
    "MAGS":     0.00,   # Magnificent-7 concentrated mega-cap
    "PAVE.AX":  0.00,   # infrastructure build-out
    "SEMI.AX":  0.00,   # semiconductors (AU-listed) — TLH sub for SMH
    "SKYY":     0.00,   # cloud / SaaS
    "SOXX":     0.00,   # semiconductors (US-listed) — TLH sub for SMH
    "XBI":      0.00,   # biotech (equal-weighted)
    # --- TLH-substitute-only tickers (cap 0 = solver never allocates) -----
    # In the universe purely so the TLH pass can price and buy them as
    # swap substitutes (see tlh_pairs.json). Not solver candidates.
    "PMGOLD.AX": 0.00,  # Perth Mint gold — TLH substitute for GOLD.AX
}


# ============================================================================
# SECTOR / THEME GROUP CAPS
# ----------------------------------------------------------------------------
# Per-asset caps don't stop the solver from stacking one bet across several
# wrappers: SMH + SOXX + SEMI.AX + SOXL are ~95% correlated, so 4 × 8% caps
# still allows a 32%+ semis cluster. Factor-implied Σ sees the correlation
# but a spurious trailing α on any one wrapper drags the whole cluster in.
#
# Each group: summed weight of member tickers ≤ cap, applied in BOTH
# solve_frontier_point_cvxpy and max_sharpe_long_only (mirrors the
# per-asset cap plumbing). Tickers absent from the current universe are
# ignored. Ungrouped tickers are unconstrained at group level.
#
# Groups deliberately cover only genuine same-bet clusters — XBI / ITA /
# PAVE.AX etc. are diversifying themes and stay ungrouped (their 8%
# per-asset caps are sufficient).
# ============================================================================
# EMPTY since the 2026-07-02 thematic revert. NOTE: do NOT re-add groups
# containing SMH/SOXL while running the 46-ticker production config — the
# group cap binds on members that remain in the universe, silently
# re-introducing the semis haircut that cost -4.2%/yr of 3Y return.
# Previous definitions (for the alpha-shrinkage re-test, task #13):
#   semiconductors: cap 0.20, [SMH, SOXX, SEMI.AX, SOXL]
#   tech_thematic:  cap 0.20, [SKYY, BOTZ, CIBR, MAGS]
SECTOR_GROUP_CAPS: dict[str, dict] = {}

# Sweep support: JSON env override for A/B cap experiments — lets a driver
# run multiple cap configs through --walk-forward-cv without code edits
# between runs (code edits would also confuse the git-sha cache key).
#   PORTOPT_CAP_OVERRIDES='{"SOXX": 0.0, "MAGS": 0.0}'  → merged into
#     PER_ASSET_WEIGHT_CAPS (cap 0 = solver exclusion, data still loaded).
#   PORTOPT_SECTOR_CAPS_DISABLE=1  → empties SECTOR_GROUP_CAPS.
_cap_env = os.environ.get("PORTOPT_CAP_OVERRIDES")
if _cap_env:
    try:
        _cap_ov = {str(k): float(v) for k, v in json.loads(_cap_env).items()}
        PER_ASSET_WEIGHT_CAPS.update(_cap_ov)
        print(f"[config-override] PER_ASSET_WEIGHT_CAPS += {_cap_ov}")
    except Exception as _e_cap_ov:
        print(f"[config-override] bad PORTOPT_CAP_OVERRIDES ignored: {_e_cap_ov}")
if os.environ.get("PORTOPT_SECTOR_CAPS_DISABLE") == "1":
    SECTOR_GROUP_CAPS = {}
    print("[config-override] SECTOR_GROUP_CAPS disabled via env")


# ============================================================================
# μ SHRINKAGE (James-Stein toward cross-sectional median)
# ----------------------------------------------------------------------------
# The candidate solvers consume trailing 24-month mean returns — the noisiest
# estimator available. Mean-variance is an error maximiser: it allocates
# hardest to the asset whose estimation NOISE is most positive, and the
# expected magnitude of that selection error grows with universe size
# (measured 2026-07-02: 9 extra free-to-ignore candidates cost -1.3%/yr —
# see wf_cv_X_46t_equiv.log vs walk_forward_cv_expanded.log).
#
# Shrinking each asset's μ toward the cross-sectional MEDIAN pulls extreme
# (mostly lucky) estimates in, making the solver skeptical-by-default.
# Median not mean: the universe contains structurally-negative-μ
# instruments (BEAR, BBUS, UVIX, VXX) that would drag a mean prior down.
#
#   μ_shrunk = (1 - λ)·μ + λ·median(μ)     λ=0 → off (production default)
#
# Applied identically in the OOS walk-forward (per training window) and the
# live ensemble path so backtest and live use the same estimator.
# Sweep via env: PORTOPT_MU_SHRINKAGE=0.5 (see task #13 experiment).
# ============================================================================
MU_SHRINKAGE_LAMBDA = 0.0

_mu_shrink_env = os.environ.get("PORTOPT_MU_SHRINKAGE")
if _mu_shrink_env:
    try:
        MU_SHRINKAGE_LAMBDA = max(0.0, min(1.0, float(_mu_shrink_env)))
        print(f"[config-override] MU_SHRINKAGE_LAMBDA={MU_SHRINKAGE_LAMBDA} via env")
    except Exception as _e_mu_env:
        print(f"[config-override] bad PORTOPT_MU_SHRINKAGE ignored: {_e_mu_env}")


# ============================================================================
# LT-DISCOUNT-AWARE SELL DEFERRAL (tax-code arbitrage, task 2026-07-02)
# ----------------------------------------------------------------------------
# At 6W rebalancing, most realised gains are SHORT-TERM (taxed at full MTR
# ~32%). Held past 12 months, the same gain gets the 50% CGT discount
# (~16%). The 12-month boundary is a tax-code fact, not an estimated edge:
# deferring the sale of a gain lot that is within LT_DEFER_WINDOW_DAYS of
# eligibility halves the tax on that gain, at the cost of holding the lot
# (tracking drift) for up to that window.
#
# Modeled HONESTLY in the walk-forward: a deferred sell means the cash was
# never raised, so the rebalance's BUYS are reduced by the deferred value
# (pro-rata) — the return stream carries the drift cost, and the lot book
# skips the shielded lots explicitly (LotBook.sell protect=). Loss lots are
# never deferred (TLH wants them realised); already-LT lots sell freely.
#
# 0 = OFF (production default). Sweep via PORTOPT_LT_DEFER_DAYS; natural
# scale is rebalance periods (42d = one 6W cycle).
# ============================================================================
LT_DEFER_WINDOW_DAYS = 0

_lt_defer_env = os.environ.get("PORTOPT_LT_DEFER_DAYS")
if _lt_defer_env:
    try:
        LT_DEFER_WINDOW_DAYS = max(0, int(float(_lt_defer_env)))
        print(f"[config-override] LT_DEFER_WINDOW_DAYS={LT_DEFER_WINDOW_DAYS} via env")
    except Exception as _e_defer_env:
        print(f"[config-override] bad PORTOPT_LT_DEFER_DAYS ignored: {_e_defer_env}")

# DD-conditional variant (parked 2026-07-02, formulated after the unconditional
# shield failed dev/validation): the shield's measured failure mode is slowing
# de-risking at violent turns — shielded lots are recent winners, so protecting
# them fights the solver exactly when the engine earns its edge (COVID/2022).
# When True, the shield is RELEASED (sells run unshielded) at early-triggered
# rebalances — the SPY DD-deepen insertions — keeping the tax deferral only on
# calm scheduled rebalances. No effect unless LT_DEFER_WINDOW_DAYS > 0.
LT_DEFER_DD_CONDITIONAL = False

_lt_defer_cond_env = os.environ.get("PORTOPT_LT_DEFER_DD_COND")
if _lt_defer_cond_env is not None:
    LT_DEFER_DD_CONDITIONAL = _lt_defer_cond_env.strip() not in ("", "0", "false", "False")
    print(f"[config-override] LT_DEFER_DD_CONDITIONAL={LT_DEFER_DD_CONDITIONAL} via env")

# Sharper release condition (variant B, 2026-07-03): instead of releasing at
# early-triggered rebalances (which fire in ~30% of ALL 6W gaps — ordinary
# chop, not stress), release when the benchmark's trailing-252d drawdown at
# the rebalance date is at or below this threshold (e.g. -0.10). Targets true
# bear states (COVID, 2022, late-2018) and keeps the deferral running through
# noise. 0.0 = off. When set, this REPLACES the early-trigger release rule.
LT_DEFER_RELEASE_DD = 0.0

_lt_defer_reldd_env = os.environ.get("PORTOPT_LT_DEFER_RELEASE_DD")
if _lt_defer_reldd_env:
    try:
        LT_DEFER_RELEASE_DD = min(0.0, float(_lt_defer_reldd_env))
        print(f"[config-override] LT_DEFER_RELEASE_DD={LT_DEFER_RELEASE_DD} via env")
    except Exception as _e_reldd_env:
        print(f"[config-override] bad PORTOPT_LT_DEFER_RELEASE_DD ignored: {_e_reldd_env}")


# ============================================================================
# CONDITIONAL REBALANCING (Q2)
# ----------------------------------------------------------------------------
# Layered on top of the scheduled REBALANCE_FREQ cadence to reduce turnover-
# driven CGT drag (currently ~250 bps/yr at 6W cadence) without sacrificing
# regime responsiveness.
#
#   SKIP_REBAL_DELTA          Skip a scheduled rebalance if the target weight
#                             change |Δw|.sum() < this threshold. Saves the
#                             cost of tiny re-trims that don't move the needle.
#                             Set to 0.0 to disable.
#
#   EARLY_TRIGGER_DD_DEEPEN   Force an early rebalance between scheduled dates
#                             when SPY drawdown deepens by more than this from
#                             the last executed rebalance. Set to 0.0 to disable.
#
#   EARLY_TRIGGER_MIN_DAYS    Minimum gap (calendar days) between an executed
#                             rebalance and an early-trigger insertion — avoids
#                             noise-driven trigger stacking.
# ============================================================================
SKIP_REBAL_DELTA          = 0.03   # 3% summed |Δw|

# Calm-conditional skip threshold — PRE-REGISTERED experiment 2026-07-08
# (memory: reference-asym-rebal-experiment; run exactly as written there).
# On CALM scheduled rebalances (NOT an early-trigger insertion AND benchmark
# trailing-252d DD > -5%) the skip threshold widens to this value, letting
# winners drift instead of realising ST gains on no-conviction re-trims.
# Stress rebalances keep SKIP_REBAL_DELTA; early-trigger machinery untouched.
# 0 = OFF (production default). Sweep via PORTOPT_SKIP_DELTA_CALM.
SKIP_REBAL_DELTA_CALM     = 0.0

_skip_calm_env = os.environ.get("PORTOPT_SKIP_DELTA_CALM")
if _skip_calm_env:
    try:
        SKIP_REBAL_DELTA_CALM = max(0.0, float(_skip_calm_env))
        print(f"[config-override] SKIP_REBAL_DELTA_CALM={SKIP_REBAL_DELTA_CALM} via env")
    except Exception as _e_calm_env:
        print(f"[config-override] bad PORTOPT_SKIP_DELTA_CALM ignored: {_e_calm_env}")

# Insurance-premium experiment (PRE-REGISTERED 2026-07-09; memory:
# reference-insurance-premium-experiment). On CALM rebalances (not early-
# triggered AND benchmark 252d DD > -5%) floor the top/Stretch slot weight at
# this value, redistributing from the defensive slots — captures more bull
# return, releases instantly on stress. 0 = OFF (production). This LOOSENS the
# drawdown defense; gate hardest (full-period MaxDD must not deepen > 1pp).
STRETCH_FLOOR_CALM = 0.0

_stretch_floor_env = os.environ.get("PORTOPT_STRETCH_FLOOR_CALM")
if _stretch_floor_env:
    try:
        STRETCH_FLOOR_CALM = min(1.0, max(0.0, float(_stretch_floor_env)))
        print(f"[config-override] STRETCH_FLOOR_CALM={STRETCH_FLOOR_CALM} via env")
    except Exception as _e_sf_env:
        print(f"[config-override] bad PORTOPT_STRETCH_FLOOR_CALM ignored: {_e_sf_env}")

# Predictive-release variant (2026-07-09, 2nd formulation): also require the
# benchmark to be above its 200d MA for the floor to engage — a trend filter
# that stands the floor down through sustained bears (all of 2022) rather than
# waiting for the reactive -5% DD. No effect unless STRETCH_FLOOR_CALM > 0.
STRETCH_FLOOR_PREDICTIVE = False
_sf_pred_env = os.environ.get("PORTOPT_STRETCH_FLOOR_PREDICTIVE")
if _sf_pred_env is not None:
    STRETCH_FLOOR_PREDICTIVE = _sf_pred_env.strip() not in ("", "0", "false", "False")
    if STRETCH_FLOOR_PREDICTIVE:
        print("[config-override] STRETCH_FLOOR_PREDICTIVE=1 (SPY>200dMA gate) via env")

# Trend-following sleeve (PRE-REGISTERED 2026-07-09; memory:
# reference-trend-sleeve-experiment). Core-satellite: blend a long-only
# inverse-vol TSMOM sleeve with the ensemble at this weight. First experiment
# adding a NEW return SOURCE (crisis-alpha diversifier). 0 = OFF (production).
TREND_SLEEVE_WEIGHT = 0.0

_trend_sleeve_env = os.environ.get("PORTOPT_TREND_SLEEVE_WEIGHT")
if _trend_sleeve_env:
    try:
        TREND_SLEEVE_WEIGHT = min(0.6, max(0.0, float(_trend_sleeve_env)))
        print(f"[config-override] TREND_SLEEVE_WEIGHT={TREND_SLEEVE_WEIGHT} via env")
    except Exception as _e_ts_env:
        print(f"[config-override] bad PORTOPT_TREND_SLEEVE_WEIGHT ignored: {_e_ts_env}")

# Ledoit-Wolf covariance shrinkage (memory: reference-cov-shrinkage-experiment).
# Σ-side robustness: shrink the noisy sample covariance toward the constant-
# correlation target for better-conditioned + more STABLE MV weights.
# Parameter-free. SHIPPED 2026-07-09 as production default (True) after CV +
# 2× dev/val: +0.03 full-period Sharpe, +0.06 dev Sharpe, crisis cost ~0.5pp
# within noise, textbook-correct estimator, more stable than sample cov.
COV_SHRINKAGE = True
_cov_shrink_env = os.environ.get("PORTOPT_COV_SHRINKAGE")
if _cov_shrink_env is not None:
    COV_SHRINKAGE = _cov_shrink_env.strip() not in ("", "0", "false", "False")
    if COV_SHRINKAGE:
        print("[config-override] COV_SHRINKAGE=1 (Ledoit-Wolf) via env")

# Volatility targeting (memory: reference-vol-targeting-experiment). Cap
# ex-ante portfolio vol at this annual target by scaling the blend toward cash
# (long-only, de-risk only). SHIPPED 2026-07-09 at 0.16 after CV + dev/val:
# validation MaxDD -3pp (crisis window) at flat Sharpe, bull window no-op,
# ~47bps/yr cost — the clearest win of the search, on the fund's core axis.
# 0 = OFF. Env PORTOPT_VOL_TARGET_ANNUAL, fingerprint-aware.
VOL_TARGET_ANNUAL = 0.16
_vol_target_env = os.environ.get("PORTOPT_VOL_TARGET_ANNUAL")
if _vol_target_env:
    try:
        VOL_TARGET_ANNUAL = max(0.0, float(_vol_target_env))
        print(f"[config-override] VOL_TARGET_ANNUAL={VOL_TARGET_ANNUAL} via env")
    except Exception as _e_vt_env:
        print(f"[config-override] bad PORTOPT_VOL_TARGET_ANNUAL ignored: {_e_vt_env}")

# Market-timed inverse-ETF crisis hedge (PRE-REGISTERED 2026-07-09; memory:
# reference-inverse-hedge-experiment). When SPY is below its 200d SMA (trend
# down), carve this fraction of the book into BEAR.AX (-1x AU inverse) and
# scale the long book to (1-w). First lever with GENUINE NEGATIVE beta — the
# axis every prior defensive tool (cash/bonds/gold/defensive slots) lacked;
# they all co-fell in 2022. Trend-triggered (not DD-triggered like the dormant
# HBRD/GOLD crash hedge, which engaged too late). 0 = OFF (production default).
# Env PORTOPT_CRISIS_HEDGE_WEIGHT, fingerprint-aware. BEAR.AX chosen over BBUS
# (which is -2x + has corrupted yfinance history: 517% ann vol).
CRISIS_HEDGE_WEIGHT = 0.0
CRISIS_HEDGE_TICKER = "BEAR.AX"
# Signal = lower Bollinger band: hedge ON iff SPY_close < SMA(MA_DAYS) - BAND_SD·σ,
# where σ = rolling std of SPY over the same MA_DAYS window. BAND_SD=0 → plain
# SMA crossover (the formulation that FAILED 2026-07-09: fired ~43% of rebals,
# mostly chop, deepened MaxDD). BAND_SD>0 = selective dislocation trigger meant
# to fire ONLY on sharp selloffs (exploratory variant, user-proposed).
CRISIS_HEDGE_MA_DAYS = 200
CRISIS_HEDGE_BAND_SD = 0.0
_crisis_hedge_env = os.environ.get("PORTOPT_CRISIS_HEDGE_WEIGHT")
if _crisis_hedge_env:
    try:
        CRISIS_HEDGE_WEIGHT = min(0.5, max(0.0, float(_crisis_hedge_env)))
        print(f"[config-override] CRISIS_HEDGE_WEIGHT={CRISIS_HEDGE_WEIGHT} via env")
    except Exception as _e_ch_env:
        print(f"[config-override] bad PORTOPT_CRISIS_HEDGE_WEIGHT ignored: {_e_ch_env}")
_crisis_ma_env = os.environ.get("PORTOPT_CRISIS_HEDGE_MA_DAYS")
if _crisis_ma_env:
    try:
        CRISIS_HEDGE_MA_DAYS = max(5, int(_crisis_ma_env))
        print(f"[config-override] CRISIS_HEDGE_MA_DAYS={CRISIS_HEDGE_MA_DAYS} via env")
    except Exception as _e_cma_env:
        print(f"[config-override] bad PORTOPT_CRISIS_HEDGE_MA_DAYS ignored: {_e_cma_env}")
_crisis_band_env = os.environ.get("PORTOPT_CRISIS_HEDGE_BAND_SD")
if _crisis_band_env:
    try:
        CRISIS_HEDGE_BAND_SD = max(0.0, float(_crisis_band_env))
        print(f"[config-override] CRISIS_HEDGE_BAND_SD={CRISIS_HEDGE_BAND_SD} via env")
    except Exception as _e_cband_env:
        print(f"[config-override] bad PORTOPT_CRISIS_HEDGE_BAND_SD ignored: {_e_cband_env}")

EARLY_TRIGGER_DD_DEEPEN   = 0.05   # 5% SPY DD deepen since last rebal
EARLY_TRIGGER_MIN_DAYS    = 10     # min days from prior rebal before re-trigger

# ============================================================================
# TAX-LOSS HARVESTING (TLH)
# ----------------------------------------------------------------------------
# At each rebalance, scan current lots for unrealised losses ≥ TLH_MIN_LOSS_PCT.
# For each loss lot whose ticker has a substitute defined in tlh_pairs.json:
# sell the loss lot (realises the loss into the FY bucket → offsets gains),
# buy equivalent dollar value of the substitute (same economic exposure).
# Cooldown prevents immediate swap-back (ATO anti-avoidance under TR 2008/1).
#
#   TLH_ENABLED            Master switch; False fully disables the pass.
#   TLH_MIN_LOSS_PCT       Threshold for a lot to be swap-eligible. Conservative
#                          default avoids harvesting noise-level losses where
#                          the brokerage cost would exceed the tax benefit.
#   TLH_COOLDOWN_DAYS      Min days between a ticker being swapped OUT and being
#                          swapped back IN. AU has no formal wash-sale rule but
#                          this satisfies the anti-avoidance test in TR 2008/1.
#   TLH_PAIRS              Substitute mapping. Loaded from tlh_pairs.json at
#                          startup; falls back to a baked-in defensive default.
# ============================================================================
# NOTE: TLH_* constants + _load_tlh_pairs moved to tlh.py (Phase 4 split,
# 2026-06-29). Constants imported at top of file. Pairs loaded from APP_DIR
# here so the path resolves to the right place at engine-load time.
TLH_PAIRS = _load_tlh_pairs(APP_DIR / "tlh_pairs.json")

# ============================================================================
# CRASH HEDGE (asymmetric defensive overlay)
# ----------------------------------------------------------------------------
# Attribution showed the engine's defensive slots and assets bleed in bull
# regimes without providing real crash protection — the "vol-managed beta"
# framing was actually two different concepts (low-vol smoothing vs true
# tail hedge) conflated. The crash hedge is the *true tail hedge*: stays
# OFF in normal regimes (zero carry cost) and engages only when a SPY peak-
# to-current drawdown trigger fires.
#
#   CRASH_HEDGE_ENABLED       Master switch.
#   CRASH_HEDGE_DD_TRIGGER    DD level at which hedge engages (negative).
#   CRASH_HEDGE_DD_RELEASE    DD level at which hedge releases (negative,
#                             less negative than trigger → hysteresis).
#   CRASH_HEDGE_LOOKBACK_DAYS Rolling window for peak detection.
#   CRASH_HEDGE_BASKET        Target weights when active. Must sum to 1.
#                             Default: 60% cash (HBRD.AX) + 40% gold
#                             (GOLD.AX). Deliberately avoids inverse ETFs
#                             (BEAR/BBUS) — first-cut conservatism.
#
# Hysteresis avoids whipsaw around the threshold: enter at -15%, exit when
# recovery to -5%. In 2016-2025, this triggers on COVID-2020 + 2022 bear +
# minor 2018 + minor 2025-Apr — designed to fire well before bottoms.
# ============================================================================
CRASH_HEDGE_ENABLED      = False  # off by default; turned on via --crash-hedge-test or live config
CRASH_HEDGE_DD_TRIGGER   = -0.15  # engage hedge at SPY peak-to-current DD ≤ -15%
CRASH_HEDGE_DD_RELEASE   = -0.05  # release hedge when DD recovers above -5%
CRASH_HEDGE_LOOKBACK_DAYS = 252   # rolling 1y peak for DD calculation
CRASH_HEDGE_BASKET       = {"HBRD.AX": 0.60, "GOLD.AX": 0.40}

# ============================================================================
# PRODUCTION CONFIG (THE shipped engine setup)
# ----------------------------------------------------------------------------
# Empirically chosen from session-long testing (commits a4f0053 → 34195f3):
# attribution showed defensive slots are alpha tax (Modest α = -5.5%/yr,
# Stretch α = +4.8%); 10-fold walk-forward CV confirmed Stretch-only beats
# 5-slot blend by +3.6% modern alpha; GFC stress test showed crash hedge
# overlay closes the 8% tail-drawdown gap (Stretch+hedge GFC MaxDD -26%
# matches 5-slot blend -25%). The synthesis ships.
#
#   PRODUCTION_SLOT_OVERRIDE  Forces all softmax weight onto a single slot.
#                             None → full 5-slot ensemble (legacy default).
#                             {"Stretch (SPY+25%)": 1.0} → ship Stretch-only.
#   PRODUCTION_CRASH_HEDGE    Enables asymmetric crash hedge overlay in the
#                             live OOS run and trade-plan generation.
#
# To revert to the 5-slot blend without hedge, set both to None / False.
# ============================================================================
# REVERTED 2026-06-19: empirical live run showed the Stretch+hedge config
# made Sharpe materially WORSE (0.99 → 0.83) and full-window MaxDD MUCH
# deeper (-20.5% → -34.4%) for negligible return uplift (+0.3%/yr, ~+$11k
# over 10y on $100k). The sweeps that justified Stretch+hedge used FOLD-
# MEAN MaxDD which understates multi-year peak-to-trough drawdowns. The
# GFC stress test also used a GOLD-only basket (HBRD.AX didn't list till
# 2017) which inflated hedge performance. Both pieces of evidence were
# measurement artifacts — the live full-period numbers are the truth.
#
# Reverted to legacy 5-slot blend with no hedge. To re-enable, set:
#   PRODUCTION_SLOT_OVERRIDE = {"Stretch (SPY+25%)": 1.0}
#   PRODUCTION_CRASH_HEDGE   = True
# (and validate with FULL-PERIOD peak-to-trough MaxDD, not fold-mean).
PRODUCTION_SLOT_OVERRIDE = None
PRODUCTION_CRASH_HEDGE   = False


# === Fund economics (currently DISABLED — wired for future activation) =======
# Set FUND_FEES_ACTIVE = True to start accruing in the live engine + backtest.
# Until then these constants are surfaced on the PDS slide but do NOT touch
# any return / NAV calculation anywhere in the pipeline. Designed so that
# flipping fees on later is a config change, not a code change.
#
# Management fee: % of NAV per annum, accrued daily on (NAV / 252).
# Performance fee: % of return above the high-water mark, crystallised at
# `PERFORMANCE_FEE_CRYSTALLISE_FREQ`. HWM uses peak-NAV-since-inception logic
# (no rolling reset). Optional hurdle rate undercuts the performance fee
# only on excess above hurdle_ann * mgmt_period — set to 0 for unhurdled.
FUND_FEES_ACTIVE                  = False
MANAGEMENT_FEE_PCT_ANN            = 0.02   # 2% per annum
PERFORMANCE_FEE_PCT               = 0.20   # 20% of NAV above HWM
PERFORMANCE_FEE_HWM               = True   # enforce high-water mark
PERFORMANCE_FEE_HURDLE_ANN        = 0.00   # 0 = unhurdled
PERFORMANCE_FEE_CRYSTALLISE_FREQ  = "Q"    # quarterly crystallisation


def compute_fund_fees(
    nav_series: "pd.Series",
    *,
    mgmt_pct_ann: float = MANAGEMENT_FEE_PCT_ANN,
    perf_pct: float = PERFORMANCE_FEE_PCT,
    use_hwm: bool = PERFORMANCE_FEE_HWM,
    hurdle_ann: float = PERFORMANCE_FEE_HURDLE_ANN,
    crystallise_freq: str = PERFORMANCE_FEE_CRYSTALLISE_FREQ,
) -> "pd.DataFrame":
    """Compute daily management + performance fee accrual on a NAV series.

    INACTIVE in the live pipeline (gated by FUND_FEES_ACTIVE). Wired so the
    fund can flip fees on without retrofitting the accounting layer.

    Returns a DataFrame indexed by date with columns:
      - mgmt_fee_daily: daily accrued management fee (AUD)
      - mgmt_fee_ytd:   cumulative management fee since last crystallisation
      - perf_fee_daily: zero except on crystallisation dates
      - hwm:            high-water mark as-of each row
      - nav_net_of_fees: NAV after subtracting accrued (but not crystallised)
                        management fee + crystallised performance fee

    HWM convention: peak NAV (net of management fee, before performance fee)
    since inception, never reset. Performance fee crystallises on the last
    business day of each crystallisation period (Q = quarterly: Mar/Jun/Sep/Dec).
    Hurdle: only the portion above (hwm * (1 + hurdle_ann * period_days/252))
    is fee-eligible.

    Per-investor accounting (HWM that resets on redemption-and-rebuy etc) is
    NOT modelled here — this is fund-level. Real per-investor subledger is
    a future build when the fund actually opens to outside capital.
    """
    import pandas as _pd
    import numpy as _np

    if nav_series is None or len(nav_series) == 0:
        return _pd.DataFrame()

    nav = _pd.Series(nav_series, dtype=float).sort_index()
    idx = nav.index

    mgmt_daily_rate = float(mgmt_pct_ann) / 252.0

    mgmt_fee_daily = nav * mgmt_daily_rate
    mgmt_fee_cum_period = _pd.Series(0.0, index=idx)
    perf_fee_daily = _pd.Series(0.0, index=idx)
    hwm_series = _pd.Series(0.0, index=idx)
    nav_net = _pd.Series(0.0, index=idx)

    # Crystallisation: last trading day in each period
    period_label = _pd.PeriodIndex(idx, freq=crystallise_freq)
    is_period_end = period_label != period_label.shift(-1, fill_value=period_label[-1])

    hwm = float(nav.iloc[0])
    running_mgmt_cum = 0.0
    running_period_start_nav = float(nav.iloc[0])
    for i, (dt, nav_t) in enumerate(nav.items()):
        nav_t = float(nav_t)
        running_mgmt_cum += float(mgmt_fee_daily.iloc[i])
        mgmt_fee_cum_period.iloc[i] = running_mgmt_cum

        nav_pre_perf = nav_t - running_mgmt_cum
        # Hurdle compounding over the current crystallisation period.
        if use_hwm:
            period_days = max(1, i - 0 + 1)  # simplified; real fund uses actual days
            hurdle_mult = (1.0 + float(hurdle_ann) * period_days / 252.0)
            fee_eligible_nav = max(0.0, nav_pre_perf - hwm * hurdle_mult)
        else:
            fee_eligible_nav = max(0.0, nav_pre_perf - running_period_start_nav)

        hwm_series.iloc[i] = hwm

        if is_period_end[i]:
            perf_fee = float(perf_pct) * fee_eligible_nav
            perf_fee_daily.iloc[i] = perf_fee
            nav_net.iloc[i] = nav_pre_perf - perf_fee
            # Update HWM only if the new NAV (net of perf fee) exceeds prior HWM
            if nav_net.iloc[i] > hwm:
                hwm = float(nav_net.iloc[i])
            running_mgmt_cum = 0.0
            running_period_start_nav = float(nav_net.iloc[i])
        else:
            nav_net.iloc[i] = nav_pre_perf

    return _pd.DataFrame({
        "mgmt_fee_daily":  mgmt_fee_daily,
        "mgmt_fee_ytd":    mgmt_fee_cum_period,
        "perf_fee_daily":  perf_fee_daily,
        "hwm":             hwm_series,
        "nav_net_of_fees": nav_net,
    }, index=idx)


# === Config snapshot (L16) ===================================================
# Dump every operationally-meaningful knob at startup so any run.log is
# self-describing — no need to grep the code to know what version of the
# config produced a given output.
def _log_config_snapshot() -> None:
    print("=" * 80)
    print("[config] === SNAPSHOT ===")
    print(f"[config] BUILD                 GIT_SHA={_BUILD_GIT_SHA}  BUILD_TIME={_BUILD_TIME}")
    print(f"[config] BROKER                profile={ACTIVE_BROKER_PROFILE} ({BROKER_CONFIG['name']})")
    print(f"[config]                       AU live: min ${BROKER_CONFIG['live_asx_min_fee']:.2f} + {BROKER_CONFIG['live_asx_rate']*100:.3f}%")
    print(f"[config]                       US live: min ${BROKER_CONFIG['live_us_min_fee']:.2f} + {BROKER_CONFIG['live_us_rate']*100:.3f}%")
    print(f"[config] CGT                   profile={ACTIVE_CGT_PROFILE}  MTR={CGT_CONFIG['marginal_tax_rate']*100:.0f}%  "
          f"LT discount={float(CGT_CONFIG['lt_discount_rate'])*100:.0f}%  "
          f"medicare={'on' if CGT_CONFIG.get('include_medicare', True) else 'off'}")
    print(f"[config] REBAL                 freq={REBALANCE_FREQ} (~{REBALANCES_PER_YEAR:.1f}/yr)  "
          f"skip<{SKIP_REBAL_DELTA*100:.0f}%Δw  early-trigger>{EARLY_TRIGGER_DD_DEEPEN*100:.0f}% DD "
          f"(min {EARLY_TRIGGER_MIN_DAYS}d)")
    if PER_ASSET_WEIGHT_CAPS:
        _caps_str = ", ".join(f"{k}={v*100:.0f}%" for k, v in PER_ASSET_WEIGHT_CAPS.items())
        print(f"[config] WEIGHT CAPS           {_caps_str}")
    if SECTOR_GROUP_CAPS:
        _gcaps_str = ", ".join(
            f"{k}≤{float(v.get('cap', 1.0))*100:.0f}% ({len(v.get('tickers', []))} tickers)"
            for k, v in SECTOR_GROUP_CAPS.items()
        )
        print(f"[config] SECTOR CAPS           {_gcaps_str}")
    if MU_SHRINKAGE_LAMBDA > 0:
        print(f"[config] MU SHRINKAGE          λ={MU_SHRINKAGE_LAMBDA:.2f} toward cross-sectional median")
    if LT_DEFER_WINDOW_DAYS > 0:
        if LT_DEFER_RELEASE_DD < 0:
            _cond_str = f"released when SPY 252d DD ≤ {LT_DEFER_RELEASE_DD*100:.0f}%"
        elif LT_DEFER_DD_CONDITIONAL:
            _cond_str = "released at early-triggered rebals"
        else:
            _cond_str = "unconditional"
        print(f"[config] LT-DEFER              gain lots within {LT_DEFER_WINDOW_DAYS}d of "
              f"12mo discount are shielded from sells ({_cond_str})")
    print(f"[config] DRIFT                 LIVE_START={LIVE_TRADING_START_DATE or '(not set)'}  "
          f"monthly={DRIFT_MONTHLY_THRESH*100:.0f}%  cumulative={DRIFT_CUMULATIVE_THRESH*100:.0f}%  "
          f"DD={DRIFT_DD_ALERT_THRESH*100:+.0f}%  slip={DRIFT_SLIPPAGE_BPS_THRESH:.0f}bps")
    print(f"[config] TLH                   enabled={TLH_ENABLED}  "
          f"min_loss={TLH_MIN_LOSS_PCT*100:+.0f}%  min_$=${TLH_MIN_LOSS_AUD:.0f}  "
          f"cooldown={TLH_COOLDOWN_DAYS}d  pairs={len(TLH_PAIRS)}")
    print(f"[config] CRASH HEDGE           enabled={CRASH_HEDGE_ENABLED}  "
          f"trigger={CRASH_HEDGE_DD_TRIGGER*100:+.0f}%DD  "
          f"release={CRASH_HEDGE_DD_RELEASE*100:+.0f}%DD  "
          f"lookback={CRASH_HEDGE_LOOKBACK_DAYS}d  "
          f"basket={CRASH_HEDGE_BASKET}")
    print(f"[config] CRISIS HEDGE          weight={CRISIS_HEDGE_WEIGHT*100:.0f}%  "
          f"ticker={CRISIS_HEDGE_TICKER}  "
          f"trigger=SPY<{CRISIS_HEDGE_MA_DAYS}dMA"
          f"{'' if CRISIS_HEDGE_BAND_SD<=0 else f'-{CRISIS_HEDGE_BAND_SD:g}σ'}")
    print(f"[config] FUND FEES             active={FUND_FEES_ACTIVE}  "
          f"mgmt={MANAGEMENT_FEE_PCT_ANN*100:.1f}%/yr  "
          f"perf={PERFORMANCE_FEE_PCT*100:.0f}%>HWM  "
          f"hwm={PERFORMANCE_FEE_HWM}  "
          f"hurdle={PERFORMANCE_FEE_HURDLE_ANN*100:.1f}%/yr  "
          f"crystallise={PERFORMANCE_FEE_CRYSTALLISE_FREQ}")
    _prod_slot_str = (next(iter(PRODUCTION_SLOT_OVERRIDE.keys()))
                       if PRODUCTION_SLOT_OVERRIDE else "5-slot blend")
    print(f"[config] PRODUCTION ENGINE     slot={_prod_slot_str}  "
          f"crash_hedge={PRODUCTION_CRASH_HEDGE}")
    print(f"[config] TARGET PORTFOLIO      ${TARGET_PORTFOLIO_VALUE_AUD:,.0f} AUD")
    print(f"[config] IBKR LIVE PRICES      enabled={USE_IBKR_LIVE_PRICES}  "
          f"host={IBKR_HOST}:{IBKR_PORT}  client_id={IBKR_CLIENT_ID}  "
          f"warn>{IBKR_DIVERGENCE_WARN_BPS:.0f}bps")
    print(f"[config] APP_DIR               {APP_DIR}")
    print(f"[config] STRESS TEST MODE      {_STRESS_TEST_MODE}")
    print("=" * 80)


_log_config_snapshot()


def _show_metrics_history(max_rows: int = 12) -> int:
    """Print the last N entries from metrics_history.jsonl as a table.

    Designed for quick visual diff across builds: shows git_sha, timestamp,
    production config flags, and 10Y Strategy Sharpe / MaxDD / Alpha. The
    most recent run is the bottom row so the eye reads time-forward.
    """
    path = APP_DIR / "metrics_history.jsonl"
    if not path.exists():
        print(f"[metrics-history] no log yet at {path} — run the live pipeline once to populate.")
        return 0
    try:
        with path.open("r", encoding="utf-8") as f:
            entries = [json.loads(ln) for ln in f if ln.strip()]
    except Exception as e:
        print(f"[metrics-history] failed to read {path}: {e}")
        return 1
    if not entries:
        print(f"[metrics-history] {path} is empty.")
        return 0
    entries = entries[-max_rows:]
    print("\n" + "=" * 116)
    print(f"METRICS HISTORY — last {len(entries)} runs (most recent at bottom)")
    print("=" * 116)
    print(f"  {'When':<19} {'Git SHA':<14} {'Slot':<24} {'Hedge':<6} "
          f"{'10Y Sh':>7} {'10Y MaxDD':>10} {'10Y α':>8} {'10Y Ret':>8} {'TLH':>4}")
    print(f"  {'-'*19} {'-'*14} {'-'*24} {'-'*6} {'-'*7} {'-'*10} {'-'*8} {'-'*8} {'-'*4}")
    for e in entries:
        ts = str(e.get("timestamp", ""))[:19]
        sha = str(e.get("git_sha", ""))[:14]
        cfg = e.get("config", {}) or {}
        slot = cfg.get("production_slot_override")
        slot_label = (next(iter(slot.keys())) if isinstance(slot, dict) and slot
                      else "5-slot blend")
        slot_label = slot_label[:24]
        hedge = "ON" if cfg.get("production_crash_hedge") else "off"
        # Find 10Y horizon
        h10 = next((h for h in (e.get("horizons") or [])
                     if h.get("horizon") == "10Y"), None)
        if h10:
            sh = h10.get("strategy_sharpe")
            dd = h10.get("strategy_max_drawdown")
            al = h10.get("strategy_alpha_vs_spy")
            rt = h10.get("strategy_ann_return")
            sh_s = f"{sh:+.2f}" if sh is not None else "  ?  "
            dd_s = f"{dd*100:+.2f}%" if dd is not None else "    ?    "
            al_s = f"{al*100:+.2f}%" if al is not None else "   ?   "
            rt_s = f"{rt*100:+.2f}%" if rt is not None else "   ?   "
        else:
            sh_s = dd_s = al_s = rt_s = "  -  "
        tlh = e.get("tlh_events", 0)
        print(f"  {ts:<19} {sha:<14} {slot_label:<24} {hedge:<6} "
              f"{sh_s:>7} {dd_s:>10} {al_s:>8} {rt_s:>8} {tlh:>4}")
    print()
    # Latest vs prior delta
    if len(entries) >= 2:
        prev = entries[-2]
        cur = entries[-1]
        prv_10 = next((h for h in (prev.get("horizons") or []) if h.get("horizon") == "10Y"), None)
        cur_10 = next((h for h in (cur.get("horizons") or []) if h.get("horizon") == "10Y"), None)
        if prv_10 and cur_10:
            d_sh = (cur_10.get("strategy_sharpe") or 0) - (prv_10.get("strategy_sharpe") or 0)
            d_dd = (cur_10.get("strategy_max_drawdown") or 0) - (prv_10.get("strategy_max_drawdown") or 0)
            d_al = (cur_10.get("strategy_alpha_vs_spy") or 0) - (prv_10.get("strategy_alpha_vs_spy") or 0)
            print(f"  Latest vs prior:  ΔSharpe {d_sh:+.3f}   "
                  f"ΔMaxDD {d_dd*100:+.2f}%   "
                  f"Δα {d_al*100:+.2f}%")
    print(f"\n  Full log: {path}")
    return 0


if _SHOW_METRICS_HISTORY_MODE:
    _exit_code = _show_metrics_history()
    sys.exit(_exit_code)


def _run_preflight() -> int:
    """Fast system-check before the heavy live pipeline.

    Returns 0 if all checks pass, 1 if any FAIL (blocking) issue found.
    WARN-level issues don't fail the check but are surfaced.
    """
    import os as _os
    import socket as _socket
    import shutil as _shutil

    print("\n" + "=" * 88)
    print("PREFLIGHT — system checks before heavy pipeline")
    print("=" * 88)

    pass_count = 0
    warn_count = 0
    fail_count = 0

    def _check(label: str, status: str, detail: str = ""):
        nonlocal pass_count, warn_count, fail_count
        if status == "PASS":
            pass_count += 1
            print(f"  ✓ PASS  {label:<32}  {detail}")
        elif status == "WARN":
            warn_count += 1
            print(f"  ⚠ WARN  {label:<32}  {detail}")
        else:
            fail_count += 1
            print(f"  ✗ FAIL  {label:<32}  {detail}")

    # 1) Build stamp
    if _BUILD_GIT_SHA not in ("dev", "unknown"):
        _check("Build stamp", "PASS", f"{_BUILD_GIT_SHA} at {_BUILD_TIME}")
    else:
        _check("Build stamp", "WARN",
                f"running from source (git={_BUILD_GIT_SHA}) — fine for dev")

    # 2) Required runtime config files
    for fname in ("regions.json", "tlh_pairs.json"):
        fpath = APP_DIR / fname
        if fpath.exists():
            _check(f"Config: {fname}", "PASS", f"size {fpath.stat().st_size} bytes")
        else:
            _check(f"Config: {fname}", "WARN", "file not found — engine uses defaults")

    # 3) Excel workbook present and not locked
    try:
        xl_path = CFG.get("excel_path") or _default_excel_path()
        if not _os.path.exists(xl_path):
            _check("Excel workbook", "FAIL", f"NOT FOUND at {xl_path}")
        else:
            # Try opening for read+write to detect lock
            try:
                with open(xl_path, "r+b") as _fh:
                    pass
                _check("Excel workbook", "PASS", f"writable at {xl_path}")
            except PermissionError:
                _check("Excel workbook", "FAIL",
                        "LOCKED — close Excel and any open instance of the file")
            except Exception as _e:
                _check("Excel workbook", "WARN", f"open check failed: {_e}")
    except Exception as _e:
        _check("Excel workbook", "FAIL", str(_e))

    # 4) Disk space (need at least 500 MB for reports + logs)
    try:
        usage = _shutil.disk_usage(str(APP_DIR))
        free_mb = usage.free / 1024 / 1024
        if free_mb < 200:
            _check("Disk space", "FAIL", f"only {free_mb:.0f} MB free in {APP_DIR}")
        elif free_mb < 500:
            _check("Disk space", "WARN", f"{free_mb:.0f} MB free — tight")
        else:
            _check("Disk space", "PASS", f"{free_mb:.0f} MB free")
    except Exception as _e:
        _check("Disk space", "WARN", f"check failed: {_e}")

    # 5) yfinance reachability (small probe download)
    try:
        _probe = yf.download("SPY", period="5d", interval="1d",
                              auto_adjust=True, threads=False, progress=False)
        if _probe is None or (hasattr(_probe, "empty") and _probe.empty):
            _check("yfinance reachable", "FAIL", "SPY probe returned empty")
        else:
            _n = len(_probe) if hasattr(_probe, "__len__") else 0
            _check("yfinance reachable", "PASS", f"SPY 5-day probe = {_n} rows")
    except Exception as _e:
        _check("yfinance reachable", "FAIL", f"{type(_e).__name__}: {_e}")

    # 6) IBKR TWS port reachable (if enabled)
    if USE_IBKR_LIVE_PRICES:
        try:
            sock = _socket.socket(_socket.AF_INET, _socket.SOCK_STREAM)
            sock.settimeout(1.5)
            result = sock.connect_ex((IBKR_HOST, IBKR_PORT))
            sock.close()
            if result == 0:
                _check("IBKR TWS port", "PASS",
                        f"port {IBKR_HOST}:{IBKR_PORT} open")
            else:
                _check("IBKR TWS port", "WARN",
                        f"port {IBKR_HOST}:{IBKR_PORT} closed — will fall back to yfinance")
        except Exception as _e:
            _check("IBKR TWS port", "WARN", f"check failed: {_e}")
    else:
        _check("IBKR TWS", "PASS", "disabled in config")

    # 7) Excel COM probe (xlwings) — only attempt on Windows
    if USE_XLWINGS:
        try:
            import xlwings as _xw
            # Don't actually open a workbook here; just test that the COM object is reachable.
            _app = _xw.App(visible=False, add_book=False)
            _app.quit()
            _check("Excel COM (xlwings)", "PASS", "COM accessible")
        except Exception as _e:
            _check("Excel COM (xlwings)", "FAIL", f"{type(_e).__name__}: {_e}")
    else:
        _check("Excel COM (xlwings)", "PASS", "xlwings disabled in config")

    # 8) Production config snapshot — visible at preflight too
    _prod_slot = (next(iter(PRODUCTION_SLOT_OVERRIDE.keys()))
                   if PRODUCTION_SLOT_OVERRIDE else "5-slot blend")
    _check("Production config", "PASS",
            f"slot={_prod_slot}, hedge={'ON' if PRODUCTION_CRASH_HEDGE else 'off'}")

    # 9) Past metrics history (regression tripwire)
    try:
        _hist = APP_DIR / "metrics_history.jsonl"
        if _hist.exists():
            with _hist.open(encoding="utf-8") as f:
                _n_hist = sum(1 for line in f if line.strip())
            _check("Metrics history", "PASS",
                    f"{_n_hist} run snapshots in metrics_history.jsonl")
        else:
            _check("Metrics history", "WARN",
                    "no history yet — first live run will create it")
    except Exception as _e:
        _check("Metrics history", "WARN", str(_e))

    # Summary
    print()
    print("=" * 88)
    if fail_count == 0:
        if warn_count == 0:
            print(f"PREFLIGHT RESULT: ALL CLEAR  ({pass_count} passed)")
            print(f"  Safe to run the live pipeline.")
        else:
            print(f"PREFLIGHT RESULT: PASS WITH WARNINGS  "
                  f"({pass_count} passed, {warn_count} warnings, 0 failed)")
            print(f"  Live pipeline will likely succeed but check warnings above.")
        rc = 0
    else:
        print(f"PREFLIGHT RESULT: BLOCKED  "
              f"({pass_count} passed, {warn_count} warnings, {fail_count} FAILED)")
        print(f"  Fix the FAILED items before running the live pipeline.")
        rc = 1
    print("=" * 88)
    return rc


# Preflight dispatch DEFERRED — see after USE_XLWINGS / IBKR config
# constants are defined (search for: --preflight dispatch).


def _run_factor_recs() -> int:
    """Preview the auto-recommended factor tilts at multiple lookback windows.

    Downloads US FF5+MOM factor data, computes trailing 3M/6M/12M Sharpe per
    factor, and prints both the raw stats and the recommended tilt targets at
    each window. Lets the user eyeball whether the recommender is picking the
    right horses before any deeper integration into the engine.
    """
    print("\n" + "=" * 88)
    print("FACTOR TILT RECOMMENDATIONS — trailing-window auto-scorer")
    print("=" * 88)
    print("  Score: each factor's annualised Sharpe over the trailing window")
    print(f"  Tilt magnitude: Sharpe × {FACTOR_TILT_SHARPE_TO_MAG:.2f}, "
          f"clipped to ±{FACTOR_TILT_MAX_MAGNITUDE:.2f}")
    print(f"  Region: US (most relevant given universe concentration)")
    print()
    try:
        # Lockboxed: tilt recommendations must not read post-lockbox factor
        # returns (moot while Ken French publishes ~2mo behind, but principled).
        ff = _apply_data_lockbox(get_ff5_mom_daily(region="US"))
    except Exception as e:
        print(f"[factor-recs] failed to load FF5+MOM data: {e}")
        return 1
    if ff is None or ff.empty:
        print("[factor-recs] no FF5 data returned; aborting.")
        return 1
    last_date = ff.index.max()
    print(f"  FF5+MOM data: {ff.index.min().date()} → {last_date.date()}  "
          f"({len(ff)} daily obs)")
    print()

    LOOKBACK_WINDOWS = [
        ("3M", 63),
        ("6M", 126),
        ("12M", 252),
    ]
    # Build per-window stats
    print(f"  {'Factor':<8} " + "  ".join(
        f"{lbl + ' AnnRet':>11} {lbl + ' AnnVol':>11} {lbl + ' Sharpe':>11}"
        for lbl, _ in LOOKBACK_WINDOWS))
    print(f"  {'-'*8} " + "  ".join(
        f"{'-'*11} {'-'*11} {'-'*11}" for _ in LOOKBACK_WINDOWS))
    per_window_stats = {lbl: compute_factor_recent_stats(ff, lookback_days=days)
                         for lbl, days in LOOKBACK_WINDOWS}
    for f in FACTOR_NAMES:
        row_parts = []
        for lbl, _ in LOOKBACK_WINDOWS:
            s = per_window_stats[lbl]
            if f in s.index:
                r = float(s.loc[f, "ann_return"])
                v = float(s.loc[f, "ann_vol"])
                sh = float(s.loc[f, "sharpe"])
                row_parts.append(f"{r*100:>+10.2f}%  {v*100:>+10.2f}%  {sh:>+11.2f}")
            else:
                row_parts.append(f"{'?':>11} {'?':>11} {'?':>11}")
        print(f"  {f:<8} " + "  ".join(row_parts))
    print()

    print("=" * 88)
    print("RECOMMENDED TILT TARGETS (Sharpe-scaled, clipped)")
    print("=" * 88)
    print(f"  {'Factor':<8} {'3M tilt':>11} {'6M tilt':>11} {'12M tilt':>11}  Reading")
    print(f"  {'-'*8} {'-'*11} {'-'*11} {'-'*11}  {'-'*40}")
    recs_3m = auto_recommend_factor_tilts(ff, lookback_days=63)
    recs_6m = auto_recommend_factor_tilts(ff, lookback_days=126)
    recs_12m = auto_recommend_factor_tilts(ff, lookback_days=252)
    for f in FACTOR_NAMES:
        t3 = recs_3m.get(f, 0.0)
        t6 = recs_6m.get(f, 0.0)
        t12 = recs_12m.get(f, 0.0)
        # Reading: directional verdict on the 3M view
        if t3 >= 0.20:
            reading = "STRONG long tilt — factor running hot"
        elif t3 >= 0.10:
            reading = "moderate long tilt"
        elif t3 >= -0.05:
            reading = "neutral"
        elif t3 >= -0.15:
            reading = "moderate short tilt"
        else:
            reading = "STRONG short tilt — factor underperforming"
        print(f"  {f:<8} {t3:>+10.3f}  {t6:>+10.3f}  {t12:>+10.3f}   {reading}")
    print()

    print("=" * 88)
    print("INTERPRETATION")
    print("=" * 88)
    # Rank factors by 3M Sharpe
    s3 = per_window_stats.get("3M", pd.DataFrame())
    if not s3.empty:
        ranked = s3["sharpe"].sort_values(ascending=False)
        top = ranked.head(2)
        bottom = ranked.tail(2)
        print(f"  Trailing 3M factor leaders:  "
              + "  ·  ".join(f"{f} ({sh:+.2f})" for f, sh in top.items()))
        print(f"  Trailing 3M factor laggards: "
              + "  ·  ".join(f"{f} ({sh:+.2f})" for f, sh in bottom.items()))
    print()
    print("  Next step (if numbers look right):")
    print("    1) Plumb auto_recommend_factor_tilts() into the dialog as the")
    print("       default Auto Recommend behaviour.")
    print("    2) Thread tilt_targets through solve_candidate_portfolios so")
    print("       each ensemble slot picks up the factor view.")
    print("    3) Validate with walk-forward CV BEFORE making it production.")
    print()
    return 0


# Dispatch deferred until after FACTOR_TILT_* constants and
# auto_recommend_factor_tilts() are defined (~line 2155).


# === Metrics history (run-over-run regression tracking) ====================
# Append one JSON line per live run to `metrics_history.jsonl` capturing the
# build/config/key metrics so we can spot regressions across versions.
# Added 2026-06-19 after a Stretch+hedge ship was reverted because the live
# full-period MaxDD blew out from -20.5% to -34.4% without my noticing — the
# sweep-mean MaxDD looked fine in isolation but full-period peak-to-trough
# was much worse. Without a persistent log, that kind of regression is
# invisible run-to-run.
def _evaluate_sweep_result(
    baseline: dict,
    treatment: dict,
    *,
    sharpe_threshold: float = -0.10,
    maxdd_threshold: float = -0.05,
    alpha_threshold: float = -0.01,
    return_threshold: float = -0.01,
    label_baseline: str = "baseline",
    label_treatment: str = "treatment",
) -> dict:
    """Honest verdict for any A/B sweep: treatment vs baseline.

    Each dict should contain: sharpe, max_drawdown, alpha_vs_spy, ann_return.
    Missing keys are treated as 'unmeasured' (skipped). The verdict requires
    treatment NOT to regress materially on ANY of the four dimensions —
    fixes the 2026-06-19 bug where sweep verdicts said PROCEED based on
    Sharpe improvement while silently ignoring a -5.5%/yr return collapse.

    Thresholds are *worst tolerable delta* (negative = regression). Defaults:
      sharpe   ≥ -0.10  (Sharpe must not drop by more than 0.10)
      maxdd    ≥ -0.05  (MaxDD must not deepen by more than 5pp)
      alpha    ≥ -0.01  (alpha vs SPY must not drop by more than 1pp)
      return   ≥ -0.01  (annualised return must not drop by more than 1pp)

    Returns dict with:
      warnings   list[str]  — one line per regressing dimension
      improvements list[str] — one line per materially improving dimension
      verdict    str        — 'SHIP' / 'NEUTRAL' / 'REVERT'
      summary    str        — one-line headline
    """
    def _d(key, neg_is_bad=False):
        b = baseline.get(key); t = treatment.get(key)
        if b is None or t is None:
            return None, None, None
        delta = t - b
        return b, t, delta

    warnings = []
    improvements = []

    b_sh, t_sh, d_sh = _d("sharpe")
    if d_sh is not None:
        if d_sh <= sharpe_threshold:
            warnings.append(f"Sharpe regressed {d_sh:+.3f} ({b_sh:.2f} → {t_sh:.2f})")
        elif d_sh >= -sharpe_threshold:
            improvements.append(f"Sharpe improved {d_sh:+.3f} ({b_sh:.2f} → {t_sh:.2f})")

    b_dd, t_dd, d_dd = _d("max_drawdown")
    if d_dd is not None:
        # max_drawdown is negative; "regressed" = MORE negative (deeper)
        if d_dd <= maxdd_threshold:
            warnings.append(f"MaxDD deepened by {d_dd*100:+.2f}pp "
                              f"({b_dd*100:+.2f}% → {t_dd*100:+.2f}%)")
        elif d_dd >= -maxdd_threshold:
            improvements.append(f"MaxDD improved by {d_dd*100:+.2f}pp "
                                  f"({b_dd*100:+.2f}% → {t_dd*100:+.2f}%)")

    b_al, t_al, d_al = _d("alpha_vs_spy")
    if d_al is not None:
        if d_al <= alpha_threshold:
            warnings.append(f"Alpha vs SPY worsened {d_al*100:+.2f}pp "
                              f"({b_al*100:+.2f}% → {t_al*100:+.2f}%)")
        elif d_al >= -alpha_threshold:
            improvements.append(f"Alpha improved {d_al*100:+.2f}pp "
                                  f"({b_al*100:+.2f}% → {t_al*100:+.2f}%)")

    b_rt, t_rt, d_rt = _d("ann_return")
    if d_rt is not None:
        if d_rt <= return_threshold:
            warnings.append(f"Ann return dropped {d_rt*100:+.2f}pp "
                              f"({b_rt*100:+.2f}% → {t_rt*100:+.2f}%)")
        elif d_rt >= -return_threshold:
            improvements.append(f"Ann return rose {d_rt*100:+.2f}pp "
                                  f"({b_rt*100:+.2f}% → {t_rt*100:+.2f}%)")

    # Overall verdict
    if warnings:
        verdict = "REVERT"
        summary = (f"{label_treatment} regresses on {len(warnings)} dimension(s) "
                    f"vs {label_baseline} — DO NOT SHIP without addressing the warnings")
    elif improvements:
        verdict = "SHIP"
        summary = (f"{label_treatment} improves on {len(improvements)} dimension(s) "
                    f"vs {label_baseline} with no material regression — candidate for shipping")
    else:
        verdict = "NEUTRAL"
        summary = (f"{label_treatment} is within noise of {label_baseline} on all "
                    f"dimensions — no material change, keep current default")

    return {
        "warnings": warnings,
        "improvements": improvements,
        "verdict": verdict,
        "summary": summary,
    }


def _print_sweep_verdict(eval_result: dict) -> None:
    """Pretty-print the result of _evaluate_sweep_result() to console + run.log.
    Always prints warnings FIRST so the eye catches regressions immediately,
    then improvements, then the overall verdict header."""
    print("=" * 88)
    print(f"VERDICT — {eval_result['verdict']}")
    print("=" * 88)
    print(f"  {eval_result['summary']}")
    if eval_result["warnings"]:
        print(f"\n  REGRESSIONS (worst-tolerable thresholds breached):")
        for w in eval_result["warnings"]:
            print(f"    ✗  {w}")
    if eval_result["improvements"]:
        print(f"\n  Improvements (beyond noise):")
        for im in eval_result["improvements"]:
            print(f"    ✓  {im}")
    if not eval_result["warnings"] and not eval_result["improvements"]:
        print(f"  All four dimensions (Sharpe, MaxDD, Alpha, Ann Return) within noise.")


def _append_metrics_snapshot(metrics_table, ensemble_mix_live, w_ensemble_live,
                              tlh_events_n: int = 0, tlh_loss_aud: float = 0.0,
                              n_executed: int = 0, n_skipped: int = 0,
                              scale_metrics: dict | None = None) -> None:
    """Append a metrics snapshot to metrics_history.jsonl + warn on regressions.

    Captures git_sha + build_time + PRODUCTION_* config + 3Y/5Y/10Y Strategy
    Sharpe/MaxDD/Alpha/AnnReturn (and SPY benchmark). Compares to the prior
    snapshot and prints [metrics-warn] lines if 10Y Sharpe regressed by
    ≥0.10, MaxDD deepened by ≥5%, or alpha worsened by ≥1%. Non-fatal —
    purely diagnostic.

    Optional `scale_metrics`: dict keyed by NAV_aud (float) → DataFrame
    of OOS metrics at that NAV (output of compute_oos_metrics on the
    scale-sensitivity backtests). When provided, the per-NAV Strategy
    rows for 3Y/5Y/10Y get nested under a `per_nav_horizons` field so
    metrics_history.jsonl accumulates a continuous track at every scale
    in parallel — used for the wholesale-fund pitch evidence narrative
    (user can answer "what was the Sharpe at $1M six months ago?"
    without re-running the backtest).
    """
    try:
        if metrics_table is None or metrics_table.empty:
            print("[metrics] snapshot skipped — empty metrics_table")
            return
        # Resolve git SHA / build time from the build stamp module if frozen,
        # else fall back to git command on the source tree.
        gsha, btime = "unknown", "unknown"
        try:
            gsha = str(_BUILD_GIT_SHA)
            btime = str(_BUILD_TIME)
        except Exception:
            pass

        # Extract per-horizon Strategy + SPY metrics. The metrics_table is a
        # DataFrame with MultiIndex columns (horizon, series). The series we
        # care about: Strategy and SPY (AUD).
        def _get(metric_name, horizon_label, series_label):
            try:
                return float(metrics_table.loc[metric_name, (horizon_label, series_label)])
            except Exception:
                return None

        horizons = []
        # The horizons depend on what was passed to compute_oos_metrics; cover the
        # standard 3Y / 5Y / 10Y. Skip any horizon that produced no data.
        for hz in ["3Y", "5Y", "10Y"]:
            row = {
                "horizon": hz,
                "strategy_ann_return": _get("Annualised Return", hz, "Strategy"),
                "strategy_sharpe":     _get("Sharpe Ratio",      hz, "Strategy"),
                "strategy_sortino":    _get("Sortino Ratio",     hz, "Strategy"),
                "strategy_max_drawdown": _get("Max Drawdown",    hz, "Strategy"),
                "strategy_alpha_vs_spy": _get("Alpha vs SPY (ann)", hz, "Strategy"),
                "spy_ann_return":     _get("Annualised Return", hz, "SPY (AUD)"),
                "spy_sharpe":         _get("Sharpe Ratio",      hz, "SPY (AUD)"),
                "spy_max_drawdown":   _get("Max Drawdown",      hz, "SPY (AUD)"),
            }
            # Skip horizon if all None (e.g. <3y of data)
            if all(v is None for k, v in row.items() if k != "horizon"):
                continue
            horizons.append(row)

        # Live recommendation snapshot (top 5)
        live_top5 = {}
        try:
            if isinstance(w_ensemble_live, pd.Series) and not w_ensemble_live.empty:
                live_top5 = {str(k): float(v) for k, v in
                              w_ensemble_live.nlargest(5).items()}
        except Exception:
            pass

        # Regime mix snapshot
        live_mix = {}
        try:
            if isinstance(ensemble_mix_live, pd.Series) and not ensemble_mix_live.empty:
                live_mix = {str(k): float(v) for k, v in ensemble_mix_live.items()}
        except Exception:
            pass

        # Optional per-NAV evidence track. Populated when SCALE_SENSITIVITY
        # is on and the scale-sensitivity sweep ran. Strategy-only (SPY
        # is NAV-invariant, already captured in horizons[*].spy_*).
        per_nav_horizons: list[dict] = []
        if scale_metrics:
            for _nav, _mtx in scale_metrics.items():
                if not isinstance(_mtx, pd.DataFrame) or _mtx.empty:
                    continue
                _nav_block: dict = {"nav_aud": float(_nav), "horizons": []}
                for hz in ["3Y", "5Y", "10Y"]:
                    def _gm(metric_name, hz=hz, _mtx=_mtx):
                        try:
                            return float(_mtx.loc[metric_name, (hz, "Strategy")])
                        except Exception:
                            return None
                    _row = {
                        "horizon": hz,
                        "strategy_ann_return": _gm("Annualised Return"),
                        "strategy_sharpe":     _gm("Sharpe Ratio"),
                        "strategy_sortino":    _gm("Sortino Ratio"),
                        "strategy_max_drawdown": _gm("Max Drawdown"),
                        "strategy_alpha_vs_spy": _gm("Alpha vs SPY (ann)"),
                    }
                    if not all(v is None for k, v in _row.items() if k != "horizon"):
                        _nav_block["horizons"].append(_row)
                if _nav_block["horizons"]:
                    per_nav_horizons.append(_nav_block)
            per_nav_horizons.sort(key=lambda b: b["nav_aud"])

        snapshot = {
            "timestamp": pd.Timestamp.now().isoformat(timespec="seconds"),
            "git_sha": gsha,
            "build_time": btime,
            "config": {
                "production_slot_override": (
                    PRODUCTION_SLOT_OVERRIDE if PRODUCTION_SLOT_OVERRIDE else None
                ),
                "production_crash_hedge": bool(PRODUCTION_CRASH_HEDGE),
                "rebalance_freq": REBALANCE_FREQ,
                "skip_rebal_delta": SKIP_REBAL_DELTA,
                "tlh_enabled": TLH_ENABLED,
                "crash_hedge_trigger_dd": CRASH_HEDGE_DD_TRIGGER,
                "crash_hedge_release_dd": CRASH_HEDGE_DD_RELEASE,
                "broker_profile": ACTIVE_BROKER_PROFILE,
                "cgt_profile": ACTIVE_CGT_PROFILE,
            },
            "horizons": horizons,
            "per_nav_horizons": per_nav_horizons,
            "tlh_events": int(tlh_events_n),
            "tlh_loss_realised_aud": float(tlh_loss_aud),
            "n_rebal_executed": int(n_executed),
            "n_rebal_skipped": int(n_skipped),
            "live_regime_mix": live_mix,
            "live_top5_positions": live_top5,
        }

        path = APP_DIR / "metrics_history.jsonl"
        # Read prior snapshot for regression check BEFORE appending.
        prior = None
        try:
            if path.exists():
                with path.open("r", encoding="utf-8") as f:
                    lines = [ln for ln in f if ln.strip()]
                if lines:
                    prior = json.loads(lines[-1])
        except Exception:
            prior = None

        # Append the new snapshot.
        with path.open("a", encoding="utf-8") as f:
            f.write(json.dumps(snapshot) + "\n")
        print(f"[metrics] snapshot → {path}")

        # Regression warnings (10Y horizon is the canonical pitch number).
        if prior:
            cur_10y = next((h for h in horizons if h["horizon"] == "10Y"), None)
            prv_10y = next((h for h in (prior.get("horizons") or [])
                             if h.get("horizon") == "10Y"), None)
            if cur_10y and prv_10y:
                warns = []
                def _d(key, sign=1):
                    cv = cur_10y.get(key); pv = prv_10y.get(key)
                    if cv is None or pv is None:
                        return None
                    return (cv - pv) * sign
                d_sh = _d("strategy_sharpe")
                d_dd = _d("strategy_max_drawdown")  # more negative = worse
                d_al = _d("strategy_alpha_vs_spy")
                if d_sh is not None and d_sh <= -0.10:
                    warns.append(f"10Y Sharpe regressed {d_sh:+.3f} "
                                  f"({prv_10y['strategy_sharpe']:.2f} → {cur_10y['strategy_sharpe']:.2f})")
                if d_dd is not None and d_dd <= -0.05:
                    warns.append(f"10Y MaxDD deepened by {d_dd*100:+.1f}% "
                                  f"({prv_10y['strategy_max_drawdown']*100:+.1f}% → {cur_10y['strategy_max_drawdown']*100:+.1f}%)")
                if d_al is not None and d_al <= -0.01:
                    warns.append(f"10Y α vs SPY worsened by {d_al*100:+.2f}% "
                                  f"({prv_10y['strategy_alpha_vs_spy']*100:+.2f}% → {cur_10y['strategy_alpha_vs_spy']*100:+.2f}%)")
                if warns:
                    print(f"[metrics-warn] REGRESSION vs prior run (git {prior.get('git_sha','?')[:8]} at {prior.get('timestamp','?')[:19]}):")
                    for w in warns:
                        print(f"[metrics-warn]   • {w}")
                else:
                    print(f"[metrics] no material regression vs prior run "
                          f"(git {prior.get('git_sha','?')[:8]} at {prior.get('timestamp','?')[:10]})")
    except Exception as e:
        print(f"[metrics] snapshot write failed: {e}")


# NOTE: _effective_cgt_rate + LotBook + compute_cgt_for_rebalance moved to
# cgt.py (Phase 4 split, 2026-06-29). Imported at top of file.

# NOTE: _run_tlh_pass, _build_lot_book_from_df, _load_tlh_cooldown_state,
# _save_tlh_cooldown_state moved to tlh.py (Phase 4 split, 2026-06-29).
# Imported at top of file.


def ask_tradeplan_portfolio_choice() -> str:
    """
    Prompt user to choose between 'ensemble', 'with_tilts' or 'no_tilts'.

    Returns:
        str: "ensemble", "with_tilts" or "no_tilts"
    """
    try:
        import tkinter as tk
        root = tk.Tk()
        root.title("Trade Plan Portfolio")
        # Centre and size the dialog
        root.geometry("420x230")
        root.resizable(False, False)
        choice = {"value": "ensemble"}

        tk.Label(root,
                 text="Which portfolio drives the trade plan?",
                 font=("Arial", 11, "bold")).pack(pady=(18, 8))
        tk.Label(root,
                 text="Ensemble — regime-adaptive blend (recommended)\n"
                      "With Tilts — factor-target portfolio\n"
                      "Optimised — pure max-Sharpe tangency",
                 font=("Arial", 9), justify="left").pack(pady=(0, 12))

        btn = tk.Frame(root)
        btn.pack(pady=4)

        def _pick(val):
            choice["value"] = val
            root.quit()
            root.destroy()

        tk.Button(btn, text="Ensemble", width=12, default="active",
                  command=lambda: _pick("ensemble")).pack(side="left", padx=4)
        tk.Button(btn, text="With Tilts", width=12,
                  command=lambda: _pick("with_tilts")).pack(side="left", padx=4)
        tk.Button(btn, text="Optimised", width=12,
                  command=lambda: _pick("no_tilts")).pack(side="left", padx=4)

        root.protocol("WM_DELETE_WINDOW", lambda: _pick("ensemble"))
        root.mainloop()
        return choice["value"]
    except Exception:
        return "ensemble"  # Safe fallback — ensemble is the recommended live engine


def build_trade_plan_from_units(
    current_units: pd.Series,
    last_prices: pd.Series,
    target_weights: pd.Series,
    cash_buffer: float = 0.0,
    min_trade_aud: float = 0.0,
    round_to_whole_units: bool = True,
    portfolio_value_override: float | None = None
) -> pd.DataFrame:
    """
    Build a rebalance trade list from current units to target weights.

    Args:
        current_units: Current holdings (ticker-indexed).
        last_prices: Latest prices (ticker-indexed).
        target_weights: Target weights (ticker-indexed, sum to 1).
        cash_buffer: Fraction to hold as cash (0-1).
        min_trade_aud: Minimum trade size in AUD.
        round_to_whole_units: Whether to round to whole units.
        portfolio_value_override: Override portfolio value.

    Returns:
        DataFrame with trade plan details.
    """
    # Validate inputs
    if len(target_weights.index) and isinstance(target_weights.index[0], (int, np.integer)):
        raise ValueError("target_weights must be ticker-indexed Series.")

    # Align and clean data
    universe = sorted(set(current_units.index) | set(last_prices.index) | set(target_weights.index))
    u = current_units.reindex(universe).fillna(0.0).astype(float)
    p = last_prices.reindex(universe).fillna(np.nan).astype(float)

    # Filter tradable assets
    tradable = p.notna() & (p > 0) & np.isfinite(p)
    u = u.loc[tradable]
    p = p.loc[tradable]

    w_tgt = target_weights.reindex(u.index).fillna(0.0)

    # Normalize target weights
    w_sum = w_tgt.sum()
    if w_sum > 0:
        w_tgt /= w_sum
    w_tgt *= (1.0 - cash_buffer)

    curr_value = u * p
    port_value = curr_value.sum()

    if portfolio_value_override is not None and np.isfinite(portfolio_value_override) and portfolio_value_override > 0:
        port_value = portfolio_value_override

    if port_value <= 0:
        return pd.DataFrame({
            "Security": u.index,
            "Price": p.values,
            "CurrUnits": u.values,
            "CurrValue": curr_value.values,
            "CurrW": np.zeros(len(u)),
            "TgtW": w_tgt.values,
            "TgtValue": np.zeros(len(u)),
            "DeltaValue": np.zeros(len(u)),
            "TradeUnits": np.zeros(len(u)),
            "Side": ["HOLD"] * len(u),
        })

    curr_w = curr_value / port_value
    tgt_value = w_tgt * port_value
    delta_value = tgt_value - curr_value
    trade_units = delta_value / p

    if round_to_whole_units:
        trade_units = np.where(trade_units >= 0, np.floor(trade_units), np.ceil(trade_units))

    # Apply min trade filter
    notional = np.abs(trade_units * p)
    trade_units = np.where(notional >= min_trade_aud, trade_units, 0.0)

    side = np.where(trade_units > 0, "BUY", np.where(trade_units < 0, "SELL", "HOLD"))

    out = pd.DataFrame({
        "Security": u.index,
        "Price": p.values,
        "CurrUnits": u.values,
        "CurrValue": curr_value.values,
        "CurrW": curr_w.values,
        "TgtW": w_tgt.values,
        "TgtValue": tgt_value.values,
        "DeltaValue": delta_value.values,
        "TradeUnits": trade_units,
        "Side": side
    })

    # Sort: sells first, then buys
    side_rank = {"SELL": 0, "HOLD": 1, "BUY": 2}
    out["SideRank"] = out["Side"].map(side_rank).fillna(1).astype(int)
    out = out.sort_values(["SideRank", "Security"]).drop(columns=["SideRank"]).reset_index(drop=True)

    return out


# NOTE: _annualized_sharpe moved to metrics.py (Phase 4 split).


def choose_portfolio_for_tradeplan(
    returns_df: pd.DataFrame,
    w_no_tilts: pd.Series,
    w_with_tilts: pd.Series,
    rf_annual: float,
    lookback_days: int = VALIDATION_LOOKBACK_DAYS,
    w_ensemble: pd.Series | None = None,
) -> tuple[str, pd.Series, dict]:
    """
    Choose portfolio based on Sharpe ratio over lookback period.
    With w_ensemble provided, compares THREE candidates and picks the highest
    Sharpe — typically the ensemble wins because it's regime-adaptive.

    Returns:
        (choice, weights, diagnostics)
    """
    r = returns_df.tail(lookback_days).replace([np.inf, -np.inf], np.nan).dropna(how="all")

    w0 = w_no_tilts.reindex(r.columns).fillna(0.0)
    w1 = w_with_tilts.reindex(r.columns).fillna(0.0)
    we = (w_ensemble.reindex(r.columns).fillna(0.0)
          if isinstance(w_ensemble, pd.Series) and not w_ensemble.empty
          else None)

    # Normalize weights
    for w in [w0, w1]:
        w_sum = w.sum()
        if w_sum != 0:
            w /= w_sum
    if we is not None:
        we_sum = we.sum()
        if we_sum != 0:
            we /= we_sum

    p0 = (r @ w0).dropna()
    p1 = (r @ w1).dropna()
    pe = (r @ we).dropna() if we is not None else pd.Series(dtype=float)

    sh0 = _annualized_sharpe(p0, rf_annual)
    sh1 = _annualized_sharpe(p1, rf_annual)
    she = _annualized_sharpe(pe, rf_annual) if not pe.empty else np.nan

    diag = {"sharpe_no_tilts": sh0, "sharpe_with_tilts": sh1, "sharpe_ensemble": she}

    candidates = []
    if np.isfinite(she) and we is not None:
        candidates.append(("ensemble", we, she))
    if np.isfinite(sh1):
        candidates.append(("with_tilts", w1, sh1))
    if np.isfinite(sh0):
        candidates.append(("no_tilts", w0, sh0))

    if not candidates:
        return "no_tilts", w0, diag
    # Stable sort: highest Sharpe wins, ties broken by entry order (ensemble first).
    candidates.sort(key=lambda x: x[2], reverse=True)
    choice, w_chosen, _ = candidates[0]
    return choice, w_chosen, diag

# ---------------------------------------------------------------------
# Main Configuration Binding
# ---------------------------------------------------------------------
TILT_FACTORS = ["Mkt-RF", "SMB", "HML", "RMW", "CMA", "MOM"]

# Bind config to globals
filename = CFG["excel_path"]
MARGINAL_TAX_RATE = CFG["marginal_tax_rate"]
CAPITAL_LOSS_CARRY_FWD = CFG["carry_forward_losses"]
LOT_MATCH_METHOD = CFG["lot_match_method"].upper()
# How the Lots sheet is reconstructed at the end of each run.
#   'fills'    — read ibkr_fills_log.jsonl (only confirmed fills count).
#                Default and what production runs use.
#   'holdings' — single-lot-per-ticker at today's last AUD price. Loses
#                cost-basis history; intended as a fallback for brokers
#                without a per-fill export.
# The legacy 'recommendations' path (applying the engine's trade plan as
# if every recommended order filled) was the source of the 2026-06 lot
# inflation bug — it's no longer reachable from the engine.
LOTS_REBUILD_MODE = str(CFG.get("lots_rebuild_mode", os.environ.get("LOTS_REBUILD_MODE", "fills"))).lower()
# Dual-NAV Roadshow chart. When ROADSHOW_DUAL_NAV is truthy (env var or
# config), the engine runs a second OOS backtest at ROADSHOW_BASE_NAV in
# addition to the existing one driven by portfolio_state.json. The
# Roadshow PPT slide then plots both Strategy lines and the metrics
# table gets three extra rows (one per horizon) for the larger-NAV run.
# Default OFF — daily ops keeps single-backtest speed.
ROADSHOW_DUAL_NAV = bool(int(os.environ.get("ROADSHOW_DUAL_NAV", "0") or "0"))
try:
    ROADSHOW_BASE_NAV = float(os.environ.get("ROADSHOW_BASE_NAV") or CFG.get("roadshow_base_nav", 1_000_000.0))
except Exception:
    ROADSHOW_BASE_NAV = 1_000_000.0
# Scale-sensitivity slide. When SCALE_SENSITIVITY truthy, the engine
# runs N additional OOS backtests (one per NAV in SCALE_SENSITIVITY_NAVS)
# and adds a dedicated slide after Roadshow showing how friction scales
# from $100k to $1M. Off by default because each extra backtest adds
# ~100s — only useful when generating PDS / Roadshow PPTs.
SCALE_SENSITIVITY = bool(int(os.environ.get("SCALE_SENSITIVITY", "0") or "0"))
try:
    _scale_str = os.environ.get("SCALE_SENSITIVITY_NAVS") or CFG.get(
        "scale_sensitivity_navs", "100000,250000,500000,1000000")
    SCALE_SENSITIVITY_NAVS = sorted({float(x.strip()) for x in str(_scale_str).split(",")
                                     if x.strip() and float(x.strip()) > 0})
except Exception:
    SCALE_SENSITIVITY_NAVS = [100_000.0, 250_000.0, 500_000.0, 1_000_000.0]
OPEN_AFTER_SAVE = CFG.get("open_after_save", True)
USE_XLWINGS = CFG.get("use_xlwings", True)
# OOS kernel mode (Phase 3b, 2026-06-29) — when subprocess workers want
# to import the engine PURELY to grab run_oos_ensemble_walk_forward and
# its dependencies, set OOS_KERNEL_MODE=1. The script will exec all
# imports + config + function definitions, then sys.exit(0) right after
# the OOS function is defined. Workers catch the SystemExit, get the
# fully-built module namespace, and call the function. Avoids the
# full pipeline (dialog, OOS execution, PPT/Excel write).
# NOTE: the flag itself is read early (near _SKIP_LIVE_PIPELINE around
# line 215) so workers bypass the holdings dialog. Re-reading here would
# clobber the early read if env changed — left as a no-op assert instead.
assert OOS_KERNEL_MODE == bool(os.environ.get("OOS_KERNEL_MODE", "").strip())

# === STRICT DATA LOCKBOX (2026-06-27 — user directive) ====================
# All decision-making data must end at DATA_LOCKBOX_DATE so the engine is
# truly blind to the future. From the lockbox date onward, NEW market data
# (prices, FF5 factors, FX) gets clamped at every ingestion point. Forward
# walk-through validation lives in paper_simulator.py (Phase 2), which
# applies the engine's frozen logic to post-lockbox market data WITHOUT
# the engine itself seeing it.
#
# SCOPE (2026-07-06, user directive after lockbox refresh #1): the lockbox
# governs RESEARCH honesty — backtests used to select parameters must not
# see post-boundary data. It does NOT freeze the live engine: once today
# moved past the boundary, lockboxed live runs were trading week-old
# regimes (Stretch 68% solved pre-semis-selloff). Doctrine now:
#   research CLI modes  → truncated at the boundary (default below)
#   live pipeline / --auto-pipeline / diagnostics → full current data
# The forward validation window stays honest through the PEEK BUDGET
# (don't select knobs on post-boundary backtests), not by blinding the
# live engine. See LOCKBOX.md.
#
# User can override via env var DATA_LOCKBOX_DATE=YYYY-MM-DD (applies to
# any mode). Set DATA_LOCKBOX_DATE="" (empty) to fully disable — DON'T DO
# THIS for research without a deliberate methodology justification.
#
# Implementation: monkey-patch yfinance + wrap FF5 loader. Any DataFrame
# or Series with a DatetimeIndex returned by yf.download gets truncated.
# Verification check at engine startup raises if any data source returns
# rows past the lockbox — defends against bypass.
_DATA_LOCKBOX_RESEARCH_MODE = (_STRESS_TEST_MODE or _SCALE_ANALYSIS_MODE
                               or _DEV_VALIDATION_MODE or _REBAL_SKIP_SWEEP_MODE
                               or _TURNOVER_SWEEP_MODE or _WALK_FORWARD_CV_MODE
                               or _ATTRIBUTION_MODE or _CRASH_HEDGE_TEST_MODE
                               or _CRASH_HEDGE_RELEASE_SWEEP_MODE
                               or _STRETCH_ONLY_TEST_MODE
                               or _STRETCH_HEDGE_SWEEP_MODE
                               or _TILTED_ENSEMBLE_TEST_MODE)
_lockbox_env = os.environ.get("DATA_LOCKBOX_DATE")
if _lockbox_env is None:
    if _DATA_LOCKBOX_RESEARCH_MODE:
        DATA_LOCKBOX_DATE = pd.Timestamp("2026-06-30")
    else:
        DATA_LOCKBOX_DATE = None
        print("[lockbox] live/diagnostic run — full current data "
              "(lockbox scoped to research modes; directive 2026-07-06)")
elif _lockbox_env.strip() == "":
    DATA_LOCKBOX_DATE = None
    print("[lockbox] DISABLED (DATA_LOCKBOX_DATE='') — engine sees all available data")
else:
    try:
        DATA_LOCKBOX_DATE = pd.Timestamp(_lockbox_env.strip())
    except Exception:
        print(f"[lockbox] env var DATA_LOCKBOX_DATE={_lockbox_env!r} unparseable; "
              f"falling back to 2026-06-30")
        DATA_LOCKBOX_DATE = pd.Timestamp("2026-06-30")

# Propagate the RESOLVED state to child processes: scale-sensitivity kernel
# workers re-exec this script and must inherit the parent's lockbox view,
# not re-decide from their own CLI flags (which say kernel, not research).
os.environ["DATA_LOCKBOX_DATE"] = ("" if DATA_LOCKBOX_DATE is None
                                   else DATA_LOCKBOX_DATE.date().isoformat())

if DATA_LOCKBOX_DATE is not None:
    _orig_yf_download = yf.download

    def _yf_download_lockbox(*args, **kwargs):
        """Yfinance wrapper that truncates results at DATA_LOCKBOX_DATE so
        the engine cannot see post-lockbox market data. Applied globally
        via monkey-patch — every existing yf.download call site
        automatically picks this up without code changes."""
        df = _orig_yf_download(*args, **kwargs)
        if df is None or (hasattr(df, "empty") and df.empty):
            return df
        try:
            if isinstance(df.index, pd.DatetimeIndex):
                return df[df.index <= DATA_LOCKBOX_DATE].copy()
        except Exception:
            pass
        return df
    yf.download = _yf_download_lockbox

    def _apply_data_lockbox(df_or_series):
        """Helper for non-yfinance data sources (FF5 CSV, custom loaders).
        Truncates any DatetimeIndex-keyed object at DATA_LOCKBOX_DATE.
        Returns the input unchanged if not a date-indexed pandas object."""
        if df_or_series is None:
            return df_or_series
        try:
            if hasattr(df_or_series, "index") and isinstance(
                    df_or_series.index, pd.DatetimeIndex):
                return df_or_series[df_or_series.index <= DATA_LOCKBOX_DATE].copy()
        except Exception:
            pass
        return df_or_series

    print(f"[lockbox] ENABLED — all market data truncated at "
          f"{DATA_LOCKBOX_DATE.date().isoformat()} "
          f"(engine cannot see future)")
else:
    # No-op helper so call sites work whether or not lockbox is active.
    def _apply_data_lockbox(df_or_series):
        return df_or_series

# ---------------------------------------------------------------------
# Factor-data layer — extracted to factors.py (module split, 2026-07-03).
# Constants + loaders canonical there; imported back for engine-wide use.
# ---------------------------------------------------------------------
import factors as _factors
from factors import (
    get_rba_cash_rate_target_current,
    FF5_REGION_URLS, EUROPEAN_EXCHANGE_SUFFIXES, FF5_DAILY_ZIP, MOM_DAILY_ZIP,
    TICKER_REGION_OVERRIDES, USER_REGION_OVERRIDES, region_for_ticker,
    _load_regions_json, _save_regions_json,
    get_mom_daily, get_ff5_daily, get_ff5_mom_daily,
    FACTOR_TILT_LOOKBACK_DAYS, FACTOR_TILT_MAX_MAGNITUDE,
    FACTOR_TILT_SHARPE_TO_MAG, FACTOR_NAMES,
    compute_factor_recent_stats, auto_recommend_factor_tilts,
    compute_ff5_betas, compute_ff5_betas_multi_region,
)

# regions.json lives beside the workbook/exe (APP_DIR handles frozen case),
# not wherever the module happens to load from.
_factors.REGIONS_JSON_PATH = APP_DIR / "regions.json"

# Dispatch --factor-recs HERE (constants + auto_recommend_factor_tilts now defined).
# Runs before heavy price download / FF5 universe build that follows.
if _FACTOR_RECS_MODE:
    _exit_code = _run_factor_recs()
    sys.exit(_exit_code)

# --preflight dispatch — placed here so USE_XLWINGS, IBKR config, etc are
# all defined, but BEFORE the heavy yfinance bulk download at line ~2290.
if _PREFLIGHT_MODE:
    _exit_code = _run_preflight()
    sys.exit(_exit_code)


# ---------------------------------------------------------------------
# Foreign Exchange
# ---------------------------------------------------------------------
# Runtime overrides to enforce 3-month tilt recommendation horizon in notebook runs.
TILT_RECOMMENDATION_LOOKBACK_DAYS = 63
FF5_LOOKBACK_DAYS = 63

# In notebook mode we default to validation-based trade-plan choice.
if str(globals().get("TRADE_PLAN_MODE", "ask")).lower().strip() == "ask":
    TRADE_PLAN_MODE = "auto"

print("[override] TILT_RECOMMENDATION_LOOKBACK_DAYS=63, TRADE_PLAN_MODE=", TRADE_PLAN_MODE)


# =====================================================================
# BLOCK 3: Data Download — Prices, Factors, FX, Benchmarks
# =====================================================================

# Constants for data processing
PRICE_DOWNLOAD_PERIOD = "2y"
FF5_BETA_WINDOW_DAYS = 504  # ~2 years of business days
FX_CACHE_PERIOD = "5y"
BENCHMARK_PERIOD = "6y"
BENCHMARK_INDICES = ["^AORD", "^GSPC", "^IXIC"]  # ASX, S&P500, NASDAQ
STATIC_STARTERS = ["^AORD"]
EXCLUDE_FROM_OPT = {"^AORD"}

# Initialize data storage
data_dict = {}

# =====================================================================
# 1) Configuration & Universe Setup
# =====================================================================
_XL_PATH = globals().get("filename", _default_excel_path())
rf_annual = get_rba_cash_rate_target_current()
rf_label = f"{rf_annual * 100:.2f}%"


def _extract_tickers_from_holdings(xl_path: str, sheet: str = "Holdings") -> list[str]:
    """Extract unique tickers from Holdings sheet Security column."""
    try:
        df = pd.read_excel(xl_path, sheet_name=sheet)
        if not isinstance(df, pd.DataFrame) or df.empty or "Security" not in df.columns:
            return []
        tickers = df["Security"].dropna().astype(str).str.strip()
        return list(dict.fromkeys([t for t in tickers if t]))
    except Exception:
        return []


def _build_ticker_universe(sheet_tickers: list[str], starters: list[str]) -> list[str]:
    """Build deduped ticker universe with mandatory benchmark."""
    universe = list(dict.fromkeys(sheet_tickers + starters))
    if "^AORD" not in universe:
        universe.insert(0, "^AORD")
    return universe


# Build universe
tickers_from_sheet = _extract_tickers_from_holdings(_XL_PATH, sheet="Holdings")
tickers = _build_ticker_universe(tickers_from_sheet, STATIC_STARTERS)
print(f"XL_PATH = {_XL_PATH}")
print(f"Tickers loaded from sheet: {tickers_from_sheet}")

# Load user region overrides from regions.json. Populated entries take
# precedence over TICKER_REGION_OVERRIDES and the suffix heuristic. Users
# pick a region per-ticker via the add-ticker dialog dropdown; the JSON
# is the persistence layer for those choices.
USER_REGION_OVERRIDES.update(_load_regions_json())
if USER_REGION_OVERRIDES:
    print(f"[region] User overrides from regions.json: {USER_REGION_OVERRIDES}")

# =====================================================================
# 2) Download Prices
# =====================================================================
def _normalize_yfinance_close(dl) -> pd.DataFrame:
    """Handle yfinance output (may be Series or DataFrame) and return DataFrame."""
    if isinstance(dl, pd.DataFrame) and "Close" in dl.columns:
        return dl["Close"]
    # Single ticker returns Series
    if isinstance(dl, pd.Series):
        return dl.to_frame()
    return pd.DataFrame()


dl = yf.download(
    tickers,
    period=PRICE_DOWNLOAD_PERIOD,
    auto_adjust=True,
    threads=False,
    progress=False
)
prices = _normalize_yfinance_close(dl)

# Clean: ensure datetime index, fill gaps, dedupe columns
prices.index = pd.to_datetime(prices.index)
idx = pd.date_range(start=prices.index.min(), end=prices.index.max(), freq="B")
prices = prices.reindex(idx).ffill().bfill()
prices.index.name = "Date"
prices = prices.loc[:, ~prices.columns.duplicated()]

# =====================================================================
# 3) Fama-French 5 Factors + Momentum (multi-region)
# =====================================================================
# Load US (canonical, used everywhere downstream for factor moments / FF5F
# sheet / FX-adjusted aggregate) plus AP-ex-Japan and Japan (used only for
# per-security beta regressions via the regional dispatch — see Task #6).
expected_cols = ["Mkt-RF", "SMB", "HML", "RMW", "CMA", "MOM", "RF"]

def _safe_load_region(region: str) -> pd.DataFrame:
    """Load a region's FF5+MOM with friendly fallback. On failure logs a warning
    and returns the US series so downstream regressions degrade gracefully
    rather than crashing the pipeline. Applies the data lockbox to the
    loaded factor frame so post-lockbox factor returns can't leak into
    the engine's regressions."""
    try:
        df = get_ff5_mom_daily(region=region)
        df = df.loc[:, ~df.columns.duplicated()].reindex(columns=expected_cols).copy()
        return _apply_data_lockbox(df)
    except Exception as e:
        print(f"[ff5] {region} factor download failed ({e}); falling back to US factors for this region")
        df = get_ff5_mom_daily(region="US")
        df = df.loc[:, ~df.columns.duplicated()].reindex(columns=expected_cols).copy()
        return _apply_data_lockbox(df)

ff5_raw = _safe_load_region("US")
if ff5_raw.empty:
    # Empty US factors cascade into an empty optimiser universe and crash
    # ~9k lines later at the cov build with an opaque pandas error
    # (2026-07-06: poisoned empty MOM cache). Say it HERE, loudly.
    print("=" * 72)
    print("[ff5][WARN] US factor frame is EMPTY — reference factor set missing.")
    print("[ff5][WARN] Live pipeline WILL fail downstream. Likely a failed Ken")
    print("[ff5][WARN] French download; re-run shortly. Cache no longer stores")
    print("[ff5][WARN] empty results, so the next run self-heals.")
    print("=" * 72)
ff5_win_for_betas = ff5_raw.tail(FF5_BETA_WINDOW_DAYS)

# Conditional download: only fetch regions actually used by the current ticker
# universe. US is always loaded (it's the reference for factor standardisation),
# regardless of whether any US tickers exist. Saves ~1 HTTP round-trip per
# unused region — small per-run but cumulative across daily runs.
_used_regions = {"US"} | {region_for_ticker(t) for t in tickers}
_used_regions = _used_regions & set(FF5_REGION_URLS.keys())
print(f"[ff5] regions required by universe: {sorted(_used_regions)}")

ff5_regional_raw = {"US": ff5_raw}
for _region in sorted(_used_regions - {"US"}):
    ff5_regional_raw[_region] = _safe_load_region(_region)
ff5_regional_windows = {r: df.tail(FF5_BETA_WINDOW_DAYS) for r, df in ff5_regional_raw.items()}
print(
    "[ff5] region windows: "
    + ", ".join(
        f"{r}={ff5_regional_windows[r].shape[0]}d ({ff5_regional_windows[r].index.min().date()} → {ff5_regional_windows[r].index.max().date()})"
        for r in ff5_regional_windows
        if not ff5_regional_windows[r].empty
    )
)

# =====================================================================
# 4) FX Rates (AUD/USD for factor adjustment & USD/AUD for holdings)
# =====================================================================
def _download_fx_series(ticker: str, period: str = FX_CACHE_PERIOD) -> pd.Series:
    """Download FX rate and return as Series."""
    try:
        dl = yf.download(
            ticker,
            period=period,
            interval="1d",
            auto_adjust=True,
            progress=False,
            threads=False
        )
        fx = dl["Close"] if isinstance(dl, pd.DataFrame) else dl
        # Newer yfinance returns Close as a 1-column DataFrame (MultiIndex columns) — squeeze to Series
        # so pd.to_numeric doesn't reject it with "arg must be a list, tuple, 1-d array, or Series".
        if isinstance(fx, pd.DataFrame):
            fx = fx.squeeze("columns")
        return pd.to_numeric(fx, errors="coerce")
    except Exception as e:
        print(f"Warning: Failed to download {ticker}: {e}")
        return pd.Series(dtype=float)


# AUD/USD for factor conversion (align to factor index)
fx_audusd = _download_fx_series("AUDUSD=X", period=BENCHMARK_PERIOD)
if fx_audusd.empty:
    print("Warning: AUDUSD=X unavailable; using flat FX=1.0")
    fx_audusd = pd.Series(1.0, index=ff5_raw.index)
else:
    fx_audusd = fx_audusd.reindex(ff5_raw.index).ffill()

fx_ret = fx_audusd.pct_change().fillna(0.0)

# USD/AUD for holdings conversion (align to price index)
fx_usdaud = _download_fx_series("USDAUD=X", period=FX_CACHE_PERIOD)
if fx_usdaud.empty:
    fx_usdaud = pd.Series(1.5, index=prices.index)  # fallback
else:
    fx_usdaud = fx_usdaud.reindex(prices.index).ffill()
# Sync the built FX series into the fx module (get_usd_aud_fx reads it there).
_fx.fx_usdaud = fx_usdaud

# =====================================================================
# 5) AUD-Adjusted Returns
# =====================================================================
# Identify USD-priced tickers (no .AX suffix, not an index)
usd_tickers = [str(c) for c in prices.columns 
               if not str(c).endswith(".AX") and not str(c).startswith("^")]

# Convert USD prices to AUD for return calculation
prices_aud = prices.copy()
if usd_tickers:
    prices_aud[usd_tickers] = prices[usd_tickers].mul(fx_usdaud, axis=0)

# Defensive filter: ETF daily returns >|30%| are almost always yfinance data
# errors (missed split/consolidation adjustments — e.g. BBUS.AX on 2025-12-01
# where price jumped $2.84 -> $28.70 because the consolidation factor was not
# back-applied to history). Drop the bad row from the return inputs so a
# single bogus print can't dominate the annualised mean / covariance.
RETURN_OUTLIER_THRESHOLD = 0.30

def _drop_return_outliers(d: pd.DataFrame, *, verbose: bool) -> pd.DataFrame:
    """NaN out rows where |Return| > RETURN_OUTLIER_THRESHOLD. Mutates and returns d."""
    mask = d["Return"].abs() > RETURN_OUTLIER_THRESHOLD
    if mask.any():
        if verbose:
            print(f"[data] Dropped {int(mask.sum())} return outlier(s) (|r| > {RETURN_OUTLIER_THRESHOLD:.0%}):")
            for _, row in d.loc[mask].iterrows():
                print(f"  {row['Security']}  {pd.Timestamp(row['Date']).date()}  ret={float(row['Return']):+.4f}")
        d.loc[mask, "Return"] = np.nan
    return d

# Compute returns
df_returns = (
    prices_aud.reset_index()
    .melt(id_vars="Date", var_name="Security", value_name="Close")
    .sort_values(["Security", "Date"])
)
df_returns["Return"] = df_returns.groupby("Security", sort=False)["Close"].pct_change()
df_returns = _drop_return_outliers(df_returns, verbose=True)
df_returns = df_returns.dropna()

# FX map for holdings sheet last-price conversion
usd_aud = get_usd_aud_fx()
fx_map_all = fx_to_aud_for_tickers(prices.columns, usd_aud)

# =====================================================================
# 6) Download Benchmark Data
# =====================================================================
def _download_benchmarks(
    tickers: list[str],
    start_date: pd.Timestamp,
    end_date: pd.Timestamp
) -> pd.DataFrame:
    """Download benchmark data and return as DataFrame."""
    try:
        data = yf.download(
            tickers,
            start=start_date,
            end=end_date,
            progress=False,
            auto_adjust=True,
            threads=False
        )
        # Extract Close prices
        if isinstance(data.columns, pd.MultiIndex):
            data = data["Close"]
        return data.ffill().bfill()
    except Exception as e:
        print(f"[data] Benchmark download failed: {e}")
        return pd.DataFrame()


benchmark_data = _download_benchmarks(
    BENCHMARK_INDICES,
    prices.index[0],
    prices.index[-1]
)
if not benchmark_data.empty:
    print(f"[data] Benchmarks downloaded: {list(benchmark_data.columns)}")

data_dict["benchmark_data"] = benchmark_data

# =====================================================================
# 7) Helper: Fetch New Tickers
# =====================================================================
def _fetch_prices_for_new_tickers(
    new_tickers: list[str],
    base_prices: pd.DataFrame,
    period: str = FX_CACHE_PERIOD
) -> pd.DataFrame:
    """Download prices for new tickers and merge into base_prices."""
    if base_prices is None or not isinstance(base_prices, pd.DataFrame):
        base_prices = pd.DataFrame()

    # Filter to genuinely new tickers
    new_only = [str(t) for t in new_tickers if str(t) not in base_prices.columns]
    if not new_only:
        return base_prices

    # Download
    dl = yf.download(
        new_only,
        period=period,
        auto_adjust=True,
        threads=False,
        progress=False
    )
    new_px = _normalize_yfinance_close(dl)

    if new_px is None or new_px.empty:
        return base_prices

    # Align to existing index if present
    new_px.index = pd.to_datetime(new_px.index).sort_values()
    new_px = new_px.loc[:, ~new_px.columns.duplicated()]
    
    if not base_prices.empty:
        new_px = new_px.reindex(index=base_prices.index).ffill()

    # Merge
    if base_prices.empty:
        return new_px

    combined = base_prices.copy()
    for col in new_px.columns:
        if col not in combined.columns:
            combined[col] = new_px[col]
    return combined


# =====================================================================
# BLOCK 4 Holdings dialog — extracted to dialogs.py (module split, 2026-07-08)
# =====================================================================
import dialogs as _dialogs
from dialogs import (
    HAS_CTK,
    _to_bool_flag,
    _to_float,
    _read_holdings_seed_from_path,
    _read_tilts_seed_from_path,
    edit_holdings_and_tilts_dialog,
)


def _wire_dialogs_module() -> None:
    """Inject engine callbacks + config mirrors into dialogs.py.

    Called immediately before the dialog opens (NOT at import time —
    recommended_tilts_for_universe is defined after this point in the
    file). After the dialog returns, the engine syncs back the three
    globals the Save handlers write: TRADE_PLAN_MODE,
    OPEN_EXCEL_AFTER_SAVE, OPEN_PPT_AFTER_SAVE.
    """
    _dialogs._fetch_prices_for_new_tickers = _fetch_prices_for_new_tickers
    _dialogs.ask_tradeplan_portfolio_choice = ask_tradeplan_portfolio_choice
    _dialogs.recommended_tilts_for_universe = recommended_tilts_for_universe
    _dialogs.OPEN_EXCEL_AFTER_SAVE = OPEN_EXCEL_AFTER_SAVE
    _dialogs.OPEN_PPT_AFTER_SAVE = OPEN_PPT_AFTER_SAVE
    _dialogs.TILT_FACTORS = TILT_FACTORS
    _dialogs.TRADE_PLAN_MODE = TRADE_PLAN_MODE


# =====================================================================
# BLOCK 5 Creating the Covariance Matrix and the Rest of the OPT
# =====================================================================
# === Analytics helpers (moved from Block 4) ===================================
gamma_cgt = 1.0
rf_label = f"{rf_annual*100:.2f}%"
chart_title = "Efficient Frontier"
TRADING_DAYS = 252


def holdings_portfolio_returns(prices: pd.DataFrame, units: pd.Series) -> pd.Series:
    units = pd.Series(units).reindex(prices.columns).fillna(0.0)
    if units.abs().sum() == 0:
        return pd.Series(dtype=float)
    px = prices.reindex(columns=units.index).ffill()
    port_val = (px * units.values).sum(axis=1)
    return port_val.pct_change(fill_method=None).dropna()


def current_holdings_weights(
    units: pd.Series,
    last_prices: pd.Series,
    investable: list[str],
    fx_to_aud: pd.Series | float | None = None,
) -> pd.Series:
    if isinstance(fx_to_aud, pd.Series):
        fx = fx_to_aud.reindex(units.index).fillna(1.0)
    else:
        fx = float(fx_to_aud) if isinstance(fx_to_aud, (int, float)) else 1.0

    mv = pd.Series(units, dtype=float) * pd.Series(last_prices, dtype=float) * fx
    mv = mv.reindex(investable).fillna(0.0)
    den = float(mv.sum())
    return (mv / den) if den > 0 else mv


# ------------------------------------------------------------
# 3) COVARIANCE MATRIX (daily)
# ------------------------------------------------------------
# Backward-compatible aliases (Block 7 now emits df_returns/prices_aud).
if "df_melt" not in globals():
    if "df_returns" in globals():
        df_melt = df_returns.copy()
    else:
        raise NameError("df_melt (or df_returns) is not defined; run Block 3 first.")

if "prices_aud_for_returns" not in globals():
    if "prices_aud" in globals():
        prices_aud_for_returns = prices_aud.copy()
    elif "prices" in globals():
        prices_aud_for_returns = prices.copy()
    else:
        raise NameError("prices_aud_for_returns (or prices_aud/prices) is not defined; run Block 3 first.")

df_cov_wide = (
    df_melt[["Date", "Security", "Return"]]
    .pivot(index="Date", columns="Security", values="Return")
)
Sigma_daily = df_cov_wide.cov()

# Optional sanity check that Sigma was built from AUD-converted prices.
# Apply the same return-outlier filter as the canonical df_returns build so
# the cross-check compares apples-to-apples; otherwise the diagnostic falsely
# reports "Using FX-adjusted returns?: False" whenever an outlier was dropped
# upstream.
try:
    _sigma_check = (
        pd.melt(
            prices_aud_for_returns.reset_index(),
            id_vars="Date",
            var_name="Security",
            value_name="Close",
        )
        .sort_values(["Security", "Date"])
    )
    _sigma_check["Return"] = _sigma_check.groupby("Security")["Close"].pct_change(fill_method=None)
    _sigma_check = _drop_return_outliers(_sigma_check, verbose=False)
    Sigma_from_aud = (
        _sigma_check.pivot(index="Date", columns="Security", values="Return").cov()
    ).reindex(index=Sigma_daily.index, columns=Sigma_daily.columns)

    diff = (Sigma_daily - Sigma_from_aud).abs().to_numpy()
    max_abs_diff = float(np.nanmax(diff)) if diff.size else np.nan
    using_fx = np.allclose(
        Sigma_daily.to_numpy(),
        Sigma_from_aud.to_numpy(),
        rtol=0,
        atol=1e-12,
        equal_nan=True,
    )
    print(f"Using FX-adjusted returns for Sigma?: {using_fx} (max |diff|={max_abs_diff:.2e})")
except Exception as e:
    print(f"[sigma-check] skipped: {e}")


# ------------------------------------------------------------
# 4) GEOMETRIC EXPECTED RETURNS (annual)
# ------------------------------------------------------------
df_melt["LogRet"] = np.log1p(df_melt["Return"])
mu_log_ann = df_melt.groupby("Security")["LogRet"].mean() * TRADING_DAYS
mu_ann_geo = np.expm1(mu_log_ann)

securities_all = [s for s in Sigma_daily.columns if s != "PortfolioValue"]
Sigma_daily = Sigma_daily.loc[securities_all, securities_all]
mu_vec_all = mu_ann_geo.reindex(securities_all)

valid_all = [
    s
    for s in securities_all
    if pd.notna(mu_vec_all.get(s, np.nan)) and pd.notna(Sigma_daily.loc[s, s])
]
Sigma_daily = Sigma_daily.loc[valid_all, valid_all]
mu_vec_all = mu_vec_all.reindex(valid_all)


# ------------------------------------------------------------
# 5) FF5 helper functions
# ------------------------------------------------------------
def compute_factor_feasible_ranges(
    B: pd.DataFrame,
    include_flags: dict,
    factor_order: list[str] | None = None,
) -> pd.DataFrame:
    """
    Under long-only and sum(w)=1, feasible factor beta range is [min(beta_i), max(beta_i)]
    over included securities.
    """
    if B is None or B.empty:
        return pd.DataFrame(columns=["Min beta", "Max beta"])

    tickers = [t for t in B.index if include_flags.get(t, False)]
    if not tickers:
        return pd.DataFrame(columns=["Min beta", "Max beta"])

    B_sub = B.loc[tickers]
    factors = factor_order if factor_order else list(B_sub.columns)

    out = pd.DataFrame(index=factors, columns=["Min beta", "Max beta"], dtype=float)
    for f in factors:
        if f in B_sub.columns:
            col = pd.to_numeric(B_sub[f], errors="coerce")
            out.loc[f, "Min beta"] = float(col.min())
            out.loc[f, "Max beta"] = float(col.max())
    return out


def get_ff5_mom_aud(ff_factors: pd.DataFrame, fx_ret_series: pd.Series) -> pd.DataFrame:
    ff = ff_factors.copy()
    ff = ff.loc[:, ~ff.columns.duplicated()]
    fx_series = pd.to_numeric(fx_ret_series.reindex(ff.index), errors="coerce").fillna(0.0)

    for col in ff.columns:
        if col != "RF":
            ff[col] = pd.to_numeric(ff[col], errors="coerce").fillna(0.0) + fx_series
    return ff


def recommend_factor_tilts(f_mean_ann: pd.Series, Fcov_daily: pd.DataFrame, normalise: bool = True) -> pd.Series:
    """Theoretical unconstrained tilt recommendation: t = Sigma^-1 * mu."""
    fac = list(f_mean_ann.index)
    mu = f_mean_ann.to_numpy(dtype=float)
    Sigma = Fcov_daily.loc[fac, fac].to_numpy(dtype=float) * TRADING_DAYS
    t_opt = np.linalg.pinv(Sigma) @ mu

    if normalise and "Mkt-RF" in fac:
        i_mkt = fac.index("Mkt-RF")
        if abs(t_opt[i_mkt]) > 1e-12:
            # Scale by the MAGNITUDE of the market tilt so every factor's sign is
            # preserved. Dividing by the signed value flips all signs whenever the
            # market tilt is negative (e.g. a down-market estimate window).
            t_opt = t_opt / abs(t_opt[i_mkt])

    return pd.Series(t_opt, index=fac, name="Recommended beta")


def recommend_factor_tilts_achievable(B: pd.DataFrame, f_mean_ann: pd.Series, Fcov_daily: pd.DataFrame):
    """
    Project theoretical factor tilts into the achievable space from investable betas B.
    Returns: (achievable_tilts, theoretical_tilts)
    """
    fac = list(f_mean_ann.index)
    theoretical = recommend_factor_tilts(f_mean_ann, Fcov_daily, normalise=False)

    if B is None or B.empty:
        return theoretical.copy(), theoretical

    B_use = B.reindex(columns=fac).dropna(how="any")
    if B_use.empty:
        return theoretical.copy(), theoretical

    Bmat = B_use.to_numpy(dtype=float)
    n = len(B_use)

    def obj(w):
        return float(np.sum((Bmat.T @ w - theoretical.values) ** 2))

    cons = ({"type": "eq", "fun": lambda w: np.sum(w) - 1.0},)
    bnds = [(0.0, 1.0)] * n
    w0 = np.full(n, 1.0 / n)

    try:
        sol = minimize(obj, w0, method="SLSQP", bounds=bnds, constraints=cons)
        if sol.success:
            w = np.asarray(sol.x, dtype=float)
            achievable = pd.Series(Bmat.T @ w, index=fac)
        else:
            achievable = theoretical.copy()
    except Exception:
        achievable = theoretical.copy()

    try:
        mins = pd.to_numeric(B_use[fac].min(axis=0), errors="coerce")
        maxs = pd.to_numeric(B_use[fac].max(axis=0), errors="coerce")
        achievable = achievable.clip(lower=mins, upper=maxs)
    except Exception:
        pass

    return achievable, theoretical


def compute_achieved_tilts(
    B: pd.DataFrame,
    w: pd.Series,
    factors=None,
    renormalise_missing: bool = True,
) -> pd.Series:
    """Factor betas of a portfolio: B.T @ w (long-only weights renormalised by default)."""
    if B is None or B.empty:
        return pd.Series(dtype=float)

    w_all = pd.Series(w).reindex(B.index).fillna(0.0)
    if renormalise_missing and float(w_all.sum()) > 0:
        w_use = w_all / float(w_all.sum())
    else:
        w_use = w_all

    out = (B.T @ w_use).rename("Achieved beta")
    return out.reindex(factors) if factors is not None else out


def optimal_portfolio_tilts(B, mu_assets, Sigma_assets, factor_index, included=None, rf: float = 0.0) -> pd.Series:
    """Recommended factor tilts = factor betas of the long-only max-Sharpe portfolio.

    Risk-aware and achievable by construction, so it never collapses to the extreme
    feasible corners the way the old direction-matching projection did. Filters out
    securities with NaN factor betas (e.g. Dimson regression failed for them) so
    they can't masquerade as "risk-free safe assets" to the solver.
    """
    zero = pd.Series(0.0, index=list(factor_index), dtype=float)
    try:
        if B is None or getattr(B, "empty", True):
            return zero
        factors = list(factor_index)
        usable_factor_cols = [f for f in factors if f in B.columns]
        if not usable_factor_cols:
            return zero
        B_clean = B.dropna(subset=usable_factor_cols)
        if B_clean.empty:
            return zero
        mu_raw = pd.to_numeric(pd.Series(mu_assets), errors="coerce").dropna()
        Sig_raw = pd.DataFrame(Sigma_assets)
        candidates = list(B_clean.index)
        if included:
            candidates = [t for t in included if t in candidates]
        candidates = [t for t in candidates
                      if t in mu_raw.index and t in Sig_raw.index and t in Sig_raw.columns]
        if not candidates:
            return zero
        Sig_sub = Sig_raw.loc[candidates, candidates]
        keep = [t for t in candidates if not Sig_sub.loc[t].isna().any()]
        if not keep:
            return zero
        mu_s = mu_raw.reindex(keep)
        Sig = Sig_sub.loc[keep, keep]
        w = max_sharpe_long_only(mu_s, Sig, rf=rf)
        if w is None or w.empty:
            return zero
        return compute_achieved_tilts(B_clean, w, factors=factors).reindex(factors).fillna(0.0)
    except Exception as e:
        print(f"[optimal_portfolio_tilts] {type(e).__name__}: {e}")
        return zero


def recommended_tilts_for_universe(included_tickers, factor_index):
    """Auto-recommend handler for the dialog: factor betas of the risk-optimal
    (max-Sharpe) long-only portfolio for the currently included tickers.
    """
    zero = pd.Series(0.0, index=list(factor_index), dtype=float)
    try:
        B_g = globals().get("B")
        mu_g = globals().get("mu_vec_opt")
        Sig_g = globals().get("Sigma_opt")
        if B_g is None or mu_g is None or Sig_g is None or getattr(B_g, "empty", True):
            return zero
        rf = float(globals().get("rf_annual", 0.0))
        return optimal_portfolio_tilts(
            B_g, mu_g, Sig_g, factor_index, included=included_tickers, rf=rf
        )
    except Exception:
        return zero


# ------------------------------------------------------------
# 6) Build optimiser moments (FF5 model or sample moments)
# ------------------------------------------------------------
WINDOW = 504

# Blended factor expected returns: a long-run premium anchor plus a small, capped
# recent tilt. This stops one quarter's regime (e.g. a market dip) from dictating
# the factor tilts, while still letting recent data nudge them. Short-term momentum
# is captured separately by the MOM factor.
FACTOR_MU_LONG_DAYS = 252 * 5      # ~5yr long-run premium anchor
FACTOR_MU_RECENT_DAYS = 252        # ~1yr recent window
FACTOR_MU_RECENT_WEIGHT = 0.20     # weight on the recent window (0..1)

USE_FF5 = True

# ---------------------------------------------------------------------------
# DIALOG-FIRST ARCHITECTURE: FF5 build deferred until AFTER the dialog.
# Previously, FF5 regression + Sigma_opt build ran here at module load on the
# sheet's tickers. Dialog-added tickers then needed a second-pass refresh.
# Now we initialise placeholders only — `_run_ff5_and_frontier_setup()` is
# called post-dialog as the single source of truth for the OPT universe.
# ---------------------------------------------------------------------------
B = pd.DataFrame()
alpha_daily = pd.Series(dtype=float)
resid_var = pd.Series(dtype=float)
ff5_regression_stats = pd.DataFrame()
f_mean_ann = pd.Series(dtype=float)
Fcov_daily = pd.DataFrame()
ff_aud = pd.DataFrame()
B_aligned = pd.DataFrame()
Sigma_ff_daily = pd.DataFrame()
securities_opt: list = []
Sigma_opt = pd.DataFrame()
mu_vec_opt = pd.Series(dtype=float)
exp_ret_label = "Expected Return (annual, FF5 AUD-adjusted)"
Sigma_frontier = pd.DataFrame()
mu_frontier = pd.Series(dtype=float)
mu_plus = pd.Series(dtype=float)
cov_plus = pd.DataFrame()
exp_ret_df = pd.DataFrame()
tilt_reco_achievable = pd.Series(dtype=float)
w_tilt = None
tilt_reco = pd.Series(dtype=float)
# Frontier outputs — also deferred.
W = pd.DataFrame()
stats_df = pd.DataFrame()
tan_ret = float("nan")
tan_vol = float("nan")


def _run_ff5_and_frontier_setup(new_prices: pd.DataFrame) -> None:
    """Canonical FF5 + frontier setup. Runs on the dialog's final universe.

    Architectural anchor for the dialog-first flow: this function IS the OPT
    universe build. It runs once post-dialog (whether the dialog was saved or
    cancelled) and is the single source of truth for Sigma_opt / mu_vec_opt /
    B / Frontier / cov_plus / exp_ret_df. The initial module-level FF5 block
    is intentionally skipped — every downstream consumer reads from globals
    that THIS function writes.

    Builds in order:
      1. df_cov_wide from new_prices (FX-adjusted)
      2. Regional FF5 windows (auto-load any new regions the universe needs)
      3. FF5 regression → B, alpha_daily, resid_var, ff5_regression_stats
      4. Factor-implied Σ and μ → Sigma_opt, mu_vec_opt
      5. cov_plus + exp_ret_df
      6. Frontier (W, stats_df, tan_ret, tan_vol)
      7. Diagnostic prints (mu top/bottom 5, tilt recommendation, Sigma diag)

    All outputs published to globals(). Previously named
    `_refresh_ff5_universe_after_dialog` and had a no-new-tickers early-return —
    that early-return removed: this function is now the ONLY FF5 setup path.
    """
    g = globals()
    px_cols = [c for c in new_prices.columns if c != "PortfolioValue"]
    existing_B = g.get("B", pd.DataFrame())
    existing_idx = set(existing_B.index) if isinstance(existing_B, pd.DataFrame) else set()
    new_tickers = sorted(set(px_cols) - existing_idx)

    if new_tickers:
        print(f"[ff5-setup] Building FF5 universe (new tickers: {new_tickers}).")
    else:
        print(f"[ff5-setup] Building FF5 universe ({len(px_cols)} tickers from dialog).")

    # 1) Rebuild df_cov_wide on the expanded prices (FX-adjust USD tickers to AUD).
    # USE fx_usdaud (USDAUD=X, ~1.41) — the price of USD in AUD. Multiplying
    # USD prices by this gives AUD prices. The earlier draft of this function
    # mistakenly used fx_audusd (AUDUSD=X, ~0.71 — the reciprocal), which
    # systematically inflated USD-ticker μ in Sigma_opt.
    px = new_prices.copy()
    px = px.drop(columns=[c for c in ["PortfolioValue"] if c in px.columns], errors="ignore")
    _fx_usdaud = g.get("fx_usdaud")
    if isinstance(_fx_usdaud, pd.Series) and not _fx_usdaud.empty:
        usd_cols = [c for c in px.columns
                    if not str(c).endswith(".AX") and not str(c).startswith("^")]
        fx_reidx = _fx_usdaud.reindex(px.index).ffill()
        if usd_cols:
            px.update(px.loc[:, usd_cols].mul(fx_reidx, axis=0))
    px = px.ffill().bfill()
    # Drop columns that are entirely NaN (failed yfinance fetch — e.g. ticker
    # was delisted, renamed, or doesn't exist on Yahoo). Without this, the
    # downstream `dropna(how="any")` would drop EVERY row that has any NaN
    # in the failed column, which is every row → 0 observations → FF5
    # regression yields no betas → universe build fails silently.
    _all_nan_cols = [c for c in px.columns if px[c].isna().all()]
    if _all_nan_cols:
        print(f"[ff5-setup] Dropping {len(_all_nan_cols)} ticker(s) with no price "
              f"data (failed yfinance fetch): {sorted(_all_nan_cols)}")
        px = px.drop(columns=_all_nan_cols, errors="ignore")
    rets_w = px.pct_change()
    rets_w = rets_w.where(rets_w.abs() <= RETURN_OUTLIER_THRESHOLD).dropna(how="any")
    df_cov_wide_new = rets_w

    # 2) Ensure ff5_regional_windows has data for every region required.
    ff5_rw = dict(g.get("ff5_regional_windows", {}))
    for t in new_tickers:
        r = region_for_ticker(t)
        if r and r not in ff5_rw and r in FF5_REGION_URLS:
            try:
                raw_r = _safe_load_region(r)
                ff5_rw[r] = raw_r.tail(FF5_BETA_WINDOW_DAYS)
                print(f"[ff5-setup] Loaded factor window for region '{r}'.")
            except Exception as _e:
                print(f"[ff5-setup] Could not load region '{r}': {_e}")
    g["ff5_regional_windows"] = ff5_rw

    # Surface which region each security got regressed against (was in the old
    # initial FF5 block — moved here so post-dialog runs still get the audit).
    # Note: the regression itself happens below; we print this after.

    # 3) Re-run FF5 regression on the expanded universe.
    B_new, alpha_daily_new, resid_var_new, ff5_stats_new = compute_ff5_betas_multi_region(
        df_cov_wide_new,
        regional_factors=ff5_rw,
        region_map=region_for_ticker,
        min_obs=120,
        return_stats=True,
    )
    if B_new is None or B_new.empty:
        print("[ff5-setup] FF5 regression yielded no betas — cannot build OPT universe.")
        return

    # Audit: surface regional beta assignments (was in old initial block).
    if not B_new.empty:
        _reg_summary: dict[str, list[str]] = {}
        for sec in B_new.index:
            _reg_summary.setdefault(region_for_ticker(sec), []).append(sec)
        print("[ff5] regional beta assignment:")
        for r, secs in _reg_summary.items():
            print(f"  {r}: {len(secs)} securities -> {secs}")

    # 4) Rebuild factor-implied Σ + μ — mirrors the module-level FF5 block recipe.
    ff_aud_new = get_ff5_mom_aud(g.get("ff5_raw"), g.get("fx_ret"))
    ff5_win = ff_aud_new.tail(WINDOW)
    fac_cols = [c for c in ff5_win.columns if c != "RF"]

    Fcov_daily_new = ff5_win[fac_cols].cov()
    _mu_long = ff_aud_new[fac_cols].tail(FACTOR_MU_LONG_DAYS).mean() * TRADING_DAYS
    _mu_recent = ff_aud_new[fac_cols].tail(FACTOR_MU_RECENT_DAYS).mean() * TRADING_DAYS
    f_mean_ann_new = (1.0 - FACTOR_MU_RECENT_WEIGHT) * _mu_long + FACTOR_MU_RECENT_WEIGHT * _mu_recent

    alpha_ann_new = pd.to_numeric(alpha_daily_new, errors="coerce").fillna(0.0) * TRADING_DAYS
    B_aligned_new = B_new.reindex(columns=fac_cols)
    mu_ff_ann_new = (
        alpha_ann_new.reindex(B_aligned_new.index).fillna(0.0)
        + (B_aligned_new @ f_mean_ann_new).fillna(0.0)
        + float(rf_annual)
    )
    securities_opt_new = [t for t in B_aligned_new.index if t not in EXCLUDE_FROM_OPT]

    F_np = Fcov_daily_new.to_numpy(dtype=float)
    Bmat_np = B_aligned_new.fillna(0.0).to_numpy(dtype=float)
    resid_diag_np = np.diag(
        pd.to_numeric(resid_var_new.reindex(B_aligned_new.index), errors="coerce")
          .clip(lower=0).fillna(0.0).to_numpy(dtype=float)
    )
    Sigma_ff_np = Bmat_np @ F_np @ Bmat_np.T + resid_diag_np
    Sigma_ff_daily_new = pd.DataFrame(Sigma_ff_np,
                                      index=B_aligned_new.index,
                                      columns=B_aligned_new.index)

    Sigma_opt_new = Sigma_ff_daily_new.loc[securities_opt_new, securities_opt_new].copy()
    mu_vec_opt_new = mu_ff_ann_new.reindex(securities_opt_new).copy()

    if "PortfolioValue" in Sigma_opt_new.index:
        Sigma_opt_new = Sigma_opt_new.drop(index="PortfolioValue",
                                           columns="PortfolioValue", errors="ignore")
    if "PortfolioValue" in mu_vec_opt_new.index:
        mu_vec_opt_new = mu_vec_opt_new.drop(index="PortfolioValue", errors="ignore")

    # 5) Rebuild cov_plus + exp_ret_df.
    n_opt_new = len(Sigma_opt_new.index)
    cov_plus_new = pd.DataFrame(0.0,
                                index=list(Sigma_opt_new.index) + ["w"],
                                columns=list(Sigma_opt_new.index) + ["w"])
    cov_plus_new.iloc[:n_opt_new, :n_opt_new] = Sigma_opt_new.values
    exp_ret_df_new = mu_vec_opt_new.rename(
        g.get("exp_ret_label", "Expected Return (annual, FF5 AUD-adjusted)")
    ).to_frame()

    # 6) Publish FF5 outputs to globals.
    g["prices"] = new_prices
    g["df_cov_wide"] = df_cov_wide_new
    g["B"] = B_new
    g["alpha_daily"] = alpha_daily_new
    g["resid_var"] = resid_var_new
    g["ff5_regression_stats"] = ff5_stats_new
    g["ff_aud"] = ff_aud_new
    g["Fcov_daily"] = Fcov_daily_new
    g["f_mean_ann"] = f_mean_ann_new
    g["Sigma_ff_daily"] = Sigma_ff_daily_new
    g["securities_opt"] = securities_opt_new
    g["Sigma_opt"] = Sigma_opt_new
    g["mu_vec_opt"] = mu_vec_opt_new
    g["Sigma_frontier"] = Sigma_opt_new.copy()
    g["mu_frontier"] = mu_vec_opt_new.copy()
    g["mu_plus"] = mu_vec_opt_new.copy()
    g["cov_plus"] = cov_plus_new
    g["exp_ret_df"] = exp_ret_df_new
    g["exp_ret_label"] = g.get("exp_ret_label", "Expected Return (annual, FF5 AUD-adjusted)")

    # 7) mu top/bottom 5 diagnostic (was in initial block).
    try:
        _mu_sorted = pd.to_numeric(mu_vec_opt_new, errors="coerce").dropna().sort_values(ascending=False)
        print(f"[diag] mu top 5 (annualized): {_mu_sorted.head(5).to_dict()}")
        print(f"[diag] mu bottom 5 (annualized): {_mu_sorted.tail(5).to_dict()}")
    except Exception as _e_mu_diag:
        print(f"[diag] mu summary skipped: {_e_mu_diag}")

    # 8) Tilt recommendation (was in initial block).
    g["tilt_reco_achievable"] = pd.Series(dtype=float)
    g["tilt_reco"] = pd.Series(dtype=float)
    g["w_tilt"] = None
    if not f_mean_ann_new.empty and not Fcov_daily_new.empty:
        try:
            g["tilt_reco_achievable"] = optimal_portfolio_tilts(
                B_new, mu_vec_opt_new, Sigma_opt_new, TILT_FACTORS, rf=rf_annual
            )
            g["tilt_reco"] = recommend_factor_tilts(f_mean_ann_new, Fcov_daily_new)
            print("\nRecommended factor tilts (betas of the risk-optimal long-only portfolio):")
            print(g["tilt_reco_achievable"].round(3))
        except Exception as _e_tilt_reco:
            print(f"[tilts] recommendation skipped: {_e_tilt_reco}")

    # 9) Frontier build (was at line 3561 in the initial block — moved here so
    # the frontier is built on the dialog's final universe).
    try:
        W_new, stats_df_new, tan_ret_new, tan_vol_new = _build_frontier(
            mu_vec_opt_new.copy(),
            Sigma_opt_new.copy(),
            target_returns=None,
            n_points=24,
        )
        g["W"] = W_new
        g["stats_df"] = stats_df_new
        g["tan_ret"] = tan_ret_new
        g["tan_vol"] = tan_vol_new
    except Exception as _e_frontier:
        print(f"[frontier] build skipped: {_e_frontier}")
        g["W"] = pd.DataFrame()
        g["stats_df"] = pd.DataFrame()
        g["tan_ret"] = float("nan")
        g["tan_vol"] = float("nan")

    # 10) Sigma_opt / mu_vec_opt sanity diagnostic (was at line 3691).
    try:
        print("\n--- DEBUG CHECK: Sigma_opt / mu_vec_opt ---")
        print("Any NaN in Sigma_opt:", bool(Sigma_opt_new.isna().any().any()))
        print("Any NaN in mu_vec_opt:", bool(mu_vec_opt_new.isna().any()))
        if len(Sigma_opt_new) > 0:
            print("Min variance:", float(np.nanmin(np.diag(Sigma_opt_new))))
        print("Number of assets:", len(Sigma_opt_new))
        if not Sigma_opt_new.empty:
            print(Sigma_opt_new.head())
        if not mu_vec_opt_new.empty:
            print(mu_vec_opt_new.head())
        print("OPT TICKERS:", list(securities_opt_new))
        print("mu:", mu_vec_opt_new.describe())
        if len(Sigma_opt_new) > 0:
            print("Sigma diag min/max:",
                  float(np.nanmin(Sigma_opt_new.values.diagonal())),
                  float(np.nanmax(Sigma_opt_new.values.diagonal())))
        if not f_mean_ann_new.empty:
            print(f_mean_ann_new)
    except Exception as _e_diag:
        print(f"[diag] Sigma_opt summary skipped: {_e_diag}")

    print(f"[ff5-setup] Universe built → {len(securities_opt_new)} securities. "
          f"FF5 covers: {sorted(B_new.index.tolist())}")


# ------------------------------------------------------------
# 8) OPTIMISATION UTILITIES (unconstrained + tilt-constrained)
# ------------------------------------------------------------
def optimise_unconstrained_analytic(mu, Sigma, target_return):
    mu = np.asarray(mu, dtype=float)
    Sigma = np.asarray(Sigma, dtype=float)

    n = len(mu)
    ones = np.ones(n)
    Sigma_inv = np.linalg.pinv(Sigma)

    A = ones @ Sigma_inv @ ones
    Bv = ones @ Sigma_inv @ mu
    C = mu @ Sigma_inv @ mu

    M = np.array([[A, Bv], [Bv, C]])
    rhs = np.array([1.0, float(target_return)])

    try:
        alpha, beta = np.linalg.solve(M, rhs)
        w = Sigma_inv @ (alpha * ones + beta * mu)
        return w, "Analytic solution."
    except np.linalg.LinAlgError:
        return np.full(n, np.nan), "Analytic solver failed (singular)."


def optimise_long_only_with_tilts(mu, Sigma, target_return, B, tilt_targets, tilt_bands, use_mask):
    mu_arr = np.asarray(mu, dtype=float)
    Sigma_arr = np.asarray(Sigma, dtype=float)
    n = len(mu_arr)

    def obj(w):
        return float(w @ Sigma_arr @ w)

    constraints = [
        {"type": "eq", "fun": lambda w: np.sum(w) - 1.0},
        {"type": "eq", "fun": lambda w: float(mu_arr @ w) - float(target_return)},
    ]

    factors = list(tilt_targets.keys()) if hasattr(tilt_targets, "keys") else []
    for f in factors:
        if not use_mask.get(f, True):
            continue
        if hasattr(B, "columns") and f not in B.columns:
            continue

        t = float(tilt_targets.get(f, 0.0))
        b = float(tilt_bands.get(f, 0.05))
        v = np.asarray(pd.Series(B[f]).reindex(range(n)).fillna(0.0), dtype=float)

        constraints.append({"type": "ineq", "fun": lambda w, v=v, t=t, b=b: (t + b) - float(v @ w)})
        constraints.append({"type": "ineq", "fun": lambda w, v=v, t=t, b=b: float(v @ w) - (t - b)})

    x0 = np.full(n, 1.0 / n)
    bounds = [(0.0, 1.0)] * n

    try:
        sol = minimize(obj, x0, method="SLSQP", bounds=bounds, constraints=constraints)
        if not sol.success:
            return np.full(n, np.nan), f"SLSQP failed: {sol.message}"
        return np.asarray(sol.x, dtype=float), "SLSQP success"
    except Exception as e:
        return np.full(n, np.nan), f"SLSQP error: {e}"


def _build_frontier(
    mu_vec_opt: pd.Series,
    Sigma_opt: pd.DataFrame,
    target_returns: list[float] | None = None,
    *,
    n_points: int = 24,
) -> tuple[pd.DataFrame, pd.DataFrame, float, float]:
    """
    Build long-only efficient frontier on realistic target returns.
    """
    mu = pd.to_numeric(mu_vec_opt, errors="coerce").reindex(Sigma_opt.index)
    keep = mu.index[mu.notna()]

    Sigma_clean = Sigma_opt.loc[keep, keep].copy()
    good = ~(Sigma_clean.isna().any(axis=1) | Sigma_clean.isna().any(axis=0))
    Sigma_clean = Sigma_clean.loc[good, good]
    mu_clean = mu.reindex(Sigma_clean.index).astype(float)

    if len(mu_clean) == 0:
        raise ValueError("No valid assets after cleaning mu/Sigma")

    assets = list(Sigma_clean.index)
    S = Sigma_clean.to_numpy(dtype=float)
    S = S + 1e-10 * np.eye(len(S))
    mu_arr = mu_clean.to_numpy(dtype=float)

    # Global minimum-variance portfolio
    n = len(assets)
    w_var = cp.Variable(n)
    prob_mvp = cp.Problem(cp.Minimize(cp.quad_form(w_var, S)), [cp.sum(w_var) == 1, w_var >= 0])
    try:
        prob_mvp.solve(solver=cp.OSQP, verbose=False)
        if w_var.value is None:
            prob_mvp.solve(solver=cp.ECOS, verbose=False)
    except Exception:
        pass

    if w_var.value is None:
        w_mvp = np.full(n, 1.0 / n)
    else:
        w_mvp = np.asarray(w_var.value).reshape(-1)

    R_mvp_ann = float(w_mvp @ mu_arr)
    mu_max = float(np.nanmax(mu_arr))

    # Size the frontier range to cover the upper portion of the asset return distribution.
    # MAD-based robust SD is immune to a single outlier asset (e.g. a noisy 450% mu blowing up
    # raw stdev); 1.4826 is the scaling factor that makes MAD comparable to SD under normality.
    # We extend to whichever is LARGER: MVP + 3 robust SDs (covers the dispersion-driven range)
    # or the 90th percentile of asset μ (covers the high-return regime). Capped at mu_max
    # because a long-only sum=1 portfolio cannot exceed any single asset's return.
    mu_finite = mu_arr[np.isfinite(mu_arr)]
    if mu_finite.size:
        mu_median = float(np.median(mu_finite))
        mad = float(np.median(np.abs(mu_finite - mu_median)))
        robust_sd = 1.4826 * mad if mad > 0 else float(np.std(mu_finite, ddof=0))
        mu_p90 = float(np.percentile(mu_finite, 90))
    else:
        robust_sd = 0.05  # last-resort fallback
        mu_p90 = 0.20

    low = R_mvp_ann
    # Extend `high` to encompass concentrated portfolios (e.g. the regime-
    # adaptive Ensemble at SMH×17% + VLUE×52%, return ~25% in current data).
    # Previously capped at p90/3*robust_sd which left the Ensemble marker off
    # the right edge of the visible frontier. Push high halfway from p90 to
    # mu_max so the curve covers high-concentration single-asset-leaning
    # portfolios without producing a degenerate single-point frontier corner.
    mu_p90_extended = mu_p90 + 0.5 * (mu_max - mu_p90)
    high = min(mu_max - 0.005, max(R_mvp_ann + 3.0 * robust_sd, mu_p90_extended))

    if high <= low + 0.01:
        high = low + 0.06

    if target_returns is None:
        target_returns = np.linspace(low, high, n_points).tolist()

    print(f"[frontier] target_returns from {target_returns[0]:.4%} to {target_returns[-1]:.4%}")

    weights_dict = {}
    stats_rows = []

    for R in target_returns:
        w_full, ok, note = solve_frontier_point_cvxpy(
            mu_clean,
            Sigma_clean,
            R,
            use_inequality=True,
        )

        # w_full returned on Sigma_clean.index length
        if len(w_full) != len(assets):
            w_series = pd.Series(0.0, index=assets)
        else:
            w_series = pd.Series(w_full, index=assets)

        weights_dict[R] = w_series.to_numpy(dtype=float)

        if ok and np.isfinite(w_series.to_numpy(dtype=float)).all():
            wv = w_series.to_numpy(dtype=float)
            vol_ann = float(np.sqrt(max(wv @ S @ wv, 0.0)) * np.sqrt(TRADING_DAYS))
            achieved = float(mu_arr @ wv)
        else:
            vol_ann = np.nan
            achieved = np.nan

        sharpe = (
            (achieved - float(rf_annual)) / vol_ann
            if (pd.notna(vol_ann) and vol_ann > 0 and pd.notna(achieved))
            else np.nan
        )

        stats_rows.append(
            {
                "Target Return": float(R),
                "Achieved Return": achieved,
                "Volatility (ann.)": vol_ann,
                "Sharpe": sharpe,
                "Method": "Frontier CVXPY",
                "Note": note,
            }
        )

    target_returns = [r for r in target_returns if np.isfinite(r)]
    cols = [
        f"{r*100:.1f}%" if (r * 100) % 1 != 0 else f"{int(r*100):d}%"
        for r in target_returns
    ]

    W = pd.DataFrame(
        {c: weights_dict[R] for c, R in zip(cols, target_returns)},
        index=assets,
    )

    stats_df = pd.DataFrame(stats_rows)
    stats_df.insert(0, "Target (%)", cols)
    stats_df = stats_df.drop(columns=["Target Return"])

    sh = pd.to_numeric(stats_df["Sharpe"], errors="coerce")
    if sh.notna().any():
        best_idx = int(sh.idxmax())
    else:
        vol_series = pd.to_numeric(stats_df["Volatility (ann.)"], errors="coerce")
        best_idx = int(vol_series.idxmin()) if vol_series.notna().any() else 0

    tan_ret = float(pd.to_numeric(stats_df.loc[best_idx, "Achieved Return"], errors="coerce"))
    tan_vol = float(pd.to_numeric(stats_df.loc[best_idx, "Volatility (ann.)"], errors="coerce"))

    if not np.isfinite(tan_ret) or not np.isfinite(tan_vol):
        tan_ret, tan_vol = np.nan, np.nan

    print(f"[frontier] tangency ret ~ {tan_ret:.4f}, vol ~ {tan_vol:.4f}")
    return W, stats_df, tan_ret, tan_vol


# ------------------------------------------------------------
# 9) FRONTIER: MVP-centred long-only
# ------------------------------------------------------------
# DIALOG-FIRST: frontier build moved into _run_ff5_and_frontier_setup()
# which runs post-dialog. Module-level placeholders defined in block 6.


# ------------------------------------------------------------
# 10) PREPARE A TRADE PLAN
# ------------------------------------------------------------
# DIALOG-FIRST: cov_plus + exp_ret_df also built inside
# _run_ff5_and_frontier_setup(). No module-level execution here.


def make_trade_plan(
    units_cur,
    last_px,
    fx_map,
    w_target,
    include_flags,
    include_zero_lines: bool = False,
    portfolio_value_override=None,
    available_cash_aud=None,
):
    """
    Return (trade_df, residual_cash) to move from current units to target weights (AUD).

    Cash-fit sizing: when `available_cash_aud` is given, the target book is sized
    to (current holdings + available_cash - reserve) so the net buys can never
    exceed the cash on hand. The reserve (max of CASH_RESERVE_MIN_AUD and
    CASH_RESERVE_PCT of investable) covers brokerage + FX spread + unit rounding.
    Falls back to `portfolio_value_override` (NAV sizing) when cash is unknown.
    """
    tickers = pd.Index(w_target.index, name="Security")

    lp = pd.to_numeric(last_px, errors="coerce").reindex(tickers).fillna(0.0)
    fx = pd.Series(1.0, index=tickers)
    if isinstance(fx_map, (dict, pd.Series)):
        fx = pd.to_numeric(pd.Series(fx_map), errors="coerce").reindex(tickers).fillna(1.0)

    px_aud = (lp * fx).replace([np.inf, -np.inf], np.nan).fillna(0.0)
    cur_units = pd.to_numeric(units_cur, errors="coerce").reindex(tickers).fillna(0).astype(int)

    cur_val = float((cur_units * px_aud).sum())  # actual current holdings value (AUD)
    if available_cash_aud is not None and np.isfinite(available_cash_aud):
        # Cash-fit: target = holdings + (cash - reserve). Net buys (= target -
        # holdings) then never exceed the cash on hand, so the plan is fundable.
        _reserve = max(CASH_RESERVE_MIN_AUD,
                       CASH_RESERVE_PCT * (cur_val + float(available_cash_aud)))
        cur_val = cur_val + max(0.0, float(available_cash_aud) - _reserve)
    elif portfolio_value_override is not None and np.isfinite(portfolio_value_override) and portfolio_value_override > 0:
        cur_val = float(portfolio_value_override)

    tgt_val = pd.to_numeric(w_target, errors="coerce").reindex(tickers).fillna(0.0) * cur_val
    tgt_units = (tgt_val / px_aud.replace(0.0, np.nan)).fillna(0.0).round().astype(int)

    if isinstance(include_flags, dict):
        inc = pd.Series(include_flags).reindex(tickers).fillna(True).astype(bool)
        tgt_units.loc[~inc] = cur_units.loc[~inc]

    # ASX minimum marketable parcel: a BUY that would ESTABLISH a position
    # under ~$500 AUD is rejected at the exchange (adds to existing holdings
    # are exempt) — IBKR cancelled a 4-unit VAE.AX buy for this (2026-07-06).
    # This is THE live plan builder (all three trade-plan modes), so the rule
    # lives here; compute_target_units_for_holdings carries it too.
    for _t in tickers:
        _cu = int(cur_units.get(_t, 0))
        _tu = int(tgt_units.get(_t, 0))
        _pv = _tu * float(px_aud.get(_t, 0.0))
        if (str(_t).upper().endswith(".AX") and _cu == 0 and _tu > 0
                and _pv < ASX_MIN_MARKETABLE_PARCEL_AUD):
            print(f"[trade-plan] {_t}: new-position target {_tu}u ≈ ${_pv:,.0f} AUD "
                  f"< ${ASX_MIN_MARKETABLE_PARCEL_AUD:,.0f} ASX min marketable parcel — "
                  f"dropped (accumulates at a future rebalance)")
            tgt_units.loc[_t] = 0

    delta = (tgt_units - cur_units).astype(int)
    cash_flow = (-delta * px_aud).astype(float)

    df = pd.DataFrame(
        {
            "Security": tickers,
            "Curr Units": cur_units.values,
            "Target Units": tgt_units.values,
            "Delta Units": delta.values,
            "Last Px (AUD)": px_aud.values,
            "Cash Flow (AUD)": cash_flow.values,
        }
    ).set_index("Security")

    if not include_zero_lines:
        df = df.loc[df["Delta Units"] != 0]

    residual_cash = float(df["Cash Flow (AUD)"].sum())
    return df, residual_cash


class SanityViolation(Exception):
    """Raised when a trade plan fails structural sanity checks. Halts
    the engine BEFORE PPT generation, exec log writes, or state file
    updates — i.e. before any side effects that the corrupted plan
    could leak into. The violations list is attached for forensic
    triage in sanity_alerts.jsonl.
    """
    def __init__(self, violations: list[dict]):
        self.violations = violations
        super().__init__(
            f"{len(violations)} sanity violation(s) — see sanity_alerts.jsonl"
        )


def _validate_trade_plan_sanity(
    trade_rec: "pd.DataFrame",
    portfolio_value_aud: float,
    *,
    max_turnover: float = 2.0,
    max_single_trade_pct: float = 0.80,
    max_position_multiple: float = 5.0,
    max_total_volume_multiple: float = 3.0,
    max_single_trade_aud: float = 200_000.0,
    max_total_volume_aud: float = 600_000.0,
) -> None:
    """Halt the engine on structurally absurd trade plans.

    Designed to catch silent state-corruption bugs like the 2026-06-26
    SMH→SOXX phantom-lots incident: the lot book had accumulated 3.4M
    units of SMH from broken `_update_lots_after_trades` logic,
    producing a trade plan with $6.3B turnover on a $1M portfolio.
    The wrapper's [rebal-trigger] verdict was RUN; only TWS being
    down prevented submission. In a live account that's bankruptcy.

    Four checks, each tunable but defaulting to thresholds that
    comfortably accommodate the engine's legitimate behavior under
    PRODUCTION_SLOT_OVERRIDE (typical Σ|Δw| ≈ 1.5-2.0 on first run
    after triage; <0.5 on subsequent runs) while flagging anything
    above:

      1. Σ|Δw| ≤ 2.0  — gross turnover bounded
      2. Single trade ≤ 20% of NAV — no one position dominates
      3. Current position value ≤ 5× NAV — catches state corruption
      4. Total trade volume ≤ 3× NAV — caps round-trip churn

    On violation: writes a structured record to sanity_alerts.jsonl,
    prints a prominent block to stdout, and raises SanityViolation.
    The caller MUST not catch this exception silently — that defeats
    the purpose. Let it propagate to the top-level run loop where it
    aborts the run with a clear error message.
    """
    if portfolio_value_aud is None or not np.isfinite(portfolio_value_aud) or portfolio_value_aud <= 0:
        # Without NAV we cannot validate ratios. Refuse to skip silently —
        # NAV being missing is itself a bug.
        print("[sanity] WARNING: portfolio_value_aud invalid — skipping sanity check (this is itself a bug)")
        return
    if trade_rec is None or trade_rec.empty:
        return

    violations: list[dict] = []

    _delta_col = (_trade_delta_col(trade_rec) if "_trade_delta_col" in globals()
                  else ("Delta Units" if "Delta Units" in trade_rec.columns else None))
    _last_px_col = "Last Px (AUD)" if "Last Px (AUD)" in trade_rec.columns else None
    _curr_units_col = "Curr Units" if "Curr Units" in trade_rec.columns else None

    if _delta_col is None or _last_px_col is None:
        # Can't compute the checks without these. Don't claim sanity if we
        # haven't actually checked anything.
        print("[sanity] WARNING: trade_rec missing required columns "
              f"(delta_col={_delta_col}, last_px_col={_last_px_col}) — skipping check")
        return

    _delta_units = pd.to_numeric(trade_rec[_delta_col], errors="coerce").fillna(0.0)
    _last_px = pd.to_numeric(trade_rec[_last_px_col], errors="coerce").fillna(0.0)
    _delta_value_aud = (_delta_units * _last_px).abs()
    _trade_pcts = _delta_value_aud / portfolio_value_aud

    # Ticker resolver — when trade_rec has Security as a column (engine's
    # actual layout, with a default RangeIndex), `series.idxmax()` returns
    # the positional integer instead of the ticker name. Caused
    # 2026-06-27 violation to read "Trade in 7" instead of "Trade in HBRD".
    # When Security is the index, both paths return the same value.
    def _ticker_label(idx_val) -> str:
        try:
            if "Security" in trade_rec.columns:
                return str(trade_rec["Security"].iloc[int(idx_val)]
                          if isinstance(idx_val, (int, np.integer))
                          else trade_rec.loc[idx_val, "Security"])
        except Exception:
            pass
        return str(idx_val)

    # Check 1: Σ|Δw|
    sum_abs_dw = float(_trade_pcts.sum())
    if sum_abs_dw > max_turnover:
        violations.append({
            "check": "turnover_too_high",
            "actual": sum_abs_dw,
            "limit": max_turnover,
            "msg": (f"Σ|Δw|={sum_abs_dw:.2f} > {max_turnover} — trade plan would rebalance "
                    f"{sum_abs_dw*100:.0f}% of NAV in one run (limit {max_turnover*100:.0f}%)")
        })

    # Check 2: any single trade > max_single_trade_pct of NAV
    worst_trade_pct = float(_trade_pcts.max()) if not _trade_pcts.empty else 0.0
    worst_trade_dv_global = float(_delta_value_aud.max()) if not _delta_value_aud.empty else 0.0
    if worst_trade_pct > max_single_trade_pct:
        _wt_idx = _trade_pcts.idxmax()
        worst_trade_ticker = _ticker_label(_wt_idx)
        worst_trade_dv = float(_delta_value_aud.loc[_wt_idx])
        violations.append({
            "check": "single_trade_too_big",
            "actual_pct": worst_trade_pct,
            "limit_pct": max_single_trade_pct,
            "ticker": worst_trade_ticker,
            "delta_value_aud": worst_trade_dv,
            "msg": (f"Trade in {worst_trade_ticker} = ${worst_trade_dv:,.0f} "
                    f"({worst_trade_pct*100:.1f}% of NAV) > {max_single_trade_pct*100:.0f}% limit")
        })

    # Check 2b: any single trade > max_single_trade_aud absolute cap.
    # Backstop for large NAVs where pct check leaves too much headroom
    # (was added 2026-06-28 after simulator Phase 2c proved $1M accounts
    # were underprotected by % thresholds alone). Catches corruption at
    # any NAV: SMH at $1.6B trips here regardless of NAV ratio.
    if worst_trade_dv_global > max_single_trade_aud:
        _wt_idx = _delta_value_aud.idxmax()
        worst_trade_ticker = _ticker_label(_wt_idx)
        violations.append({
            "check": "single_trade_abs_too_big",
            "actual_aud": worst_trade_dv_global,
            "limit_aud": max_single_trade_aud,
            "ticker": worst_trade_ticker,
            "msg": (f"Trade in {worst_trade_ticker} = ${worst_trade_dv_global:,.0f} "
                    f"> ${max_single_trade_aud:,.0f} absolute cap")
        })

    # Check 3: any current position quantity × price > max_position_multiple × NAV
    if _curr_units_col is not None:
        _curr_units = pd.to_numeric(trade_rec[_curr_units_col], errors="coerce").fillna(0.0)
        _curr_value_abs = (_curr_units.abs() * _last_px)
        max_pos_value = float(_curr_value_abs.max()) if not _curr_value_abs.empty else 0.0
        max_pos_limit = max_position_multiple * portfolio_value_aud
        if max_pos_value > max_pos_limit:
            _wp_idx = _curr_value_abs.idxmax()
            worst_pos_ticker = _ticker_label(_wp_idx)
            worst_pos_units = float(_curr_units.loc[_wp_idx])
            violations.append({
                "check": "position_absurd",
                "actual_value_aud": max_pos_value,
                "limit_value_aud": max_pos_limit,
                "ticker": worst_pos_ticker,
                "units": worst_pos_units,
                "portfolio_value_aud": float(portfolio_value_aud),
                "msg": (f"Current position in {worst_pos_ticker} = {worst_pos_units:,.0f} units "
                        f"(${max_pos_value:,.0f}) > {max_position_multiple}× NAV — "
                        f"almost certainly state corruption (lot book / holdings / fills log)")
            })

    # Check 4: total trade volume > max_total_volume_multiple × NAV
    total_volume = float(_delta_value_aud.sum())
    total_volume_limit = max_total_volume_multiple * portfolio_value_aud
    if total_volume > total_volume_limit:
        violations.append({
            "check": "total_volume_too_high",
            "actual_aud": total_volume,
            "limit_aud": total_volume_limit,
            "msg": (f"Total trade volume ${total_volume:,.0f} > "
                    f"{max_total_volume_multiple}× NAV (${total_volume_limit:,.0f})")
        })

    # Check 4b: total trade volume > max_total_volume_aud absolute cap.
    # Same scale-protection rationale as Check 2b.
    if total_volume > max_total_volume_aud:
        violations.append({
            "check": "total_volume_abs_too_high",
            "actual_aud": total_volume,
            "limit_aud": max_total_volume_aud,
            "msg": (f"Total trade volume ${total_volume:,.0f} > "
                    f"${max_total_volume_aud:,.0f} absolute cap")
        })

    if not violations:
        return

    # Persist forensic record before raising — survival of the alert
    # matters more than survival of the run.
    try:
        _alert_path = APP_DIR / "sanity_alerts.jsonl"
        _alert = {
            "timestamp": pd.Timestamp.now().isoformat(timespec="seconds"),
            "portfolio_value_aud": float(portfolio_value_aud),
            "n_trades": int((_delta_units != 0).sum()),
            "violations": violations,
        }
        with _alert_path.open("a", encoding="utf-8") as f:
            f.write(json.dumps(_alert) + "\n")
    except Exception as _e_alert:
        print(f"[sanity] failed to write alert log: {_e_alert}")

    # Loud, structured stdout output so the user sees the problem
    # immediately and the daily_auto toast surfaces it.
    _bar = "=" * 88
    print()
    print(_bar)
    print("[SANITY VIOLATION] Trade plan rejected — engine halting before any side effects")
    print(_bar)
    for v in violations:
        print(f"  • {v['msg']}")
    print(_bar)
    print(f"  NAV: ${portfolio_value_aud:,.0f}  ·  Alerts logged to: sanity_alerts.jsonl")
    print(f"  Most likely cause: state corruption in lot book, Holdings, or ibkr_fills_log.")
    print(f"  Do NOT execute this trade plan. Investigate state files before re-running.")
    print(_bar)
    print()

    raise SanityViolation(violations)


# === Drift tracker v2/v3 helpers (Tier-1 #3) ================================
# NOTE: append_trade_recommendation_log moved to jsonl_logs.py (Phase 4 split).
# Fill comparison + monthly NAV drift + live MaxDD alert. Read by the live
# pipeline after the recommendation log is written. All Excel writes are
# guarded; nothing here can break a live run.

# NOTE: _load_recommendation_log moved to jsonl_logs.py (Phase 4 split).


# NOTE: _match_fill_to_recommendation, compute_fill_drift,
# compute_live_max_drawdown, compute_monthly_nav_drift moved to drift.py
# (Phase 4 split, 2026-06-29).


# NOTE: _print_drift_warnings moved to drift.py (Phase 4 split).


# === Persistent cash ledger (Tier-1 #3 add-on) ===============================
# One JSONL entry per run. Surfaces cumulative brokerage + CGT + drift vs the
# anchor target portfolio value, so the user can see where money is going
# (currently the engine reports brokerage but never subtracts it from NAV).

# NOTE: append_cash_ledger moved to jsonl_logs.py (Phase 4 split).


def _load_cash_ledger(ledger_path) -> pd.DataFrame:
    """Load cash_ledger.jsonl into a DataFrame, computing cumulative columns
    + drift vs first record + drift vs TARGET_PORTFOLIO_VALUE_AUD anchor."""
    p = Path(ledger_path)
    if not p.exists():
        return pd.DataFrame()
    rows: list[dict] = []
    try:
        with open(p, "r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                try:
                    rows.append(json.loads(line))
                except json.JSONDecodeError:
                    continue
    except Exception:
        return pd.DataFrame()
    if not rows:
        return pd.DataFrame()
    df = pd.DataFrame(rows)
    df["date"] = pd.to_datetime(df.get("date"), errors="coerce")
    df = df.sort_values("run_at").reset_index(drop=True)
    # Cumulative + drift columns
    df["cum_brokerage_aud"] = pd.to_numeric(df["brokerage_this_run_aud"],
                                            errors="coerce").fillna(0).cumsum().round(2)
    df["cum_cgt_aud"] = pd.to_numeric(df["cgt_this_run_aud"],
                                       errors="coerce").fillna(0).cumsum().round(2)
    first_val = float(df["portfolio_value_aud"].iloc[0]) if not df.empty else 0.0
    df["delta_vs_prev_aud"] = pd.to_numeric(df["portfolio_value_aud"],
                                            errors="coerce").diff().round(2)
    df["drift_vs_start_aud"] = (
        pd.to_numeric(df["portfolio_value_aud"], errors="coerce") - first_val
    ).round(2)
    df["drift_vs_target_aud"] = (
        pd.to_numeric(df["portfolio_value_aud"], errors="coerce") - TARGET_PORTFOLIO_VALUE_AUD
    ).round(2)
    # Reconciliation: portfolio change between runs SHOULD equal
    # (market_move) - (brokerage_paid) - (cgt_paid). So
    # market_move = Δ_portfolio + brokerage + cgt. If market_move is wildly
    # different from what underlying prices actually did, that's unexplained
    # drift (slippage, FX, math bug). Blank on the first row — there's no
    # prior NAV to diff against.
    df["unexplained_delta_aud"] = (
        df["delta_vs_prev_aud"]
        + pd.to_numeric(df["brokerage_this_run_aud"], errors="coerce").fillna(0)
        + pd.to_numeric(df["cgt_this_run_aud"], errors="coerce").fillna(0)
    ).round(2)
    return df


# === IBKR hybrid live-price fetch (Tier-1 #1) ================================
# Pulls delayed last-prices from IBKR for the live trade plan. Read-only;
# requires TWS or IB Gateway running on PAPER (DUQ... account). Returns {}
# on any failure so the engine falls back to yfinance cleanly.

# NOTE: _ibkr_pick_price moved to ibkr.py (Phase 4 split, 2026-06-29).


def fetch_ibkr_live_prices_native(tickers: list[str]) -> dict[str, float]:
    """Return native-currency last prices for `tickers` keyed by engine ticker.

    Engine tickers ending with '.AX' are mapped to ASX/AUD; '^*' (benchmarks)
    are skipped; everything else is treated as SMART/USD. Empty dict on any
    failure (TWS not running, non-paper account, contract qualification fail).
    """
    try:
        from ib_insync import IB, Stock
    except ImportError:
        print("[ibkr-price] ib_insync not installed; skipping IBKR price fetch")
        return {}
    ib = IB()
    out: dict[str, float] = {}
    try:
        ib.connect(IBKR_HOST, IBKR_PORT, clientId=IBKR_CLIENT_ID,
                   timeout=IBKR_CONNECT_TIMEOUT)
    except Exception as e:
        print(f"[ibkr-price] connection skipped ({type(e).__name__}); "
              f"using yfinance prices")
        return {}
    try:
        managed = ib.managedAccounts() or []
        if not managed or not str(managed[0]).startswith("DU"):
            print(f"[ibkr-price] non-paper account {managed} — refusing; "
                  f"using yfinance")
            return {}
        ib.reqMarketDataType(3)  # delayed (free)
        plan: list[tuple[str, object]] = []
        for t in tickers:
            ts = str(t).strip()
            if ts.startswith("^"):
                continue
            if ts.endswith(".AX"):
                c = Stock(ts[:-3], exchange="SMART", currency="AUD",
                          primaryExchange="ASX")
            else:
                c = Stock(ts, exchange="SMART", currency="USD")
            plan.append((ts, c))
        if not plan:
            return {}
        ib.qualifyContracts(*[c for _, c in plan])
        plan = [(t, c) for t, c in plan if getattr(c, "conId", 0)]
        ticks = [(t, ib.reqMktData(c, "", snapshot=True,
                                    regulatorySnapshot=False))
                 for t, c in plan]
        deadline = time.time() + IBKR_SNAPSHOT_WAIT_SEC
        while time.time() < deadline:
            ready = sum(1 for _, tk in ticks if _ibkr_pick_price(tk) is not None)
            if ready == len(ticks):
                break
            ib.sleep(0.4)
        for ticker, tk in ticks:
            p = _ibkr_pick_price(tk)
            if p is not None:
                out[ticker] = float(p)
        return out
    except Exception as e:
        print(f"[ibkr-price] fetch failed ({type(e).__name__}): {e}")
        return {}
    finally:
        if ib.isConnected():
            ib.disconnect()


def _fetch_ibkr_cash_aud() -> "float | None":
    """Live IBKR TotalCashValue in the account's base currency (AUD), paper
    account only. Returns None on any failure (Gateway down, non-paper, etc.)
    so the caller falls back to the last broker snapshot. Read-only."""
    try:
        from ib_insync import IB
    except ImportError:
        return None
    ib = IB()
    try:
        ib.connect(IBKR_HOST, IBKR_PORT, clientId=IBKR_CLIENT_ID + 3,
                   timeout=IBKR_CONNECT_TIMEOUT)
    except Exception as e:
        print(f"[cash-fit] IBKR cash query skipped ({type(e).__name__}); "
              f"will try last broker snapshot")
        return None
    try:
        managed = ib.managedAccounts() or []
        if not managed or not str(managed[0]).startswith("DU"):
            print(f"[cash-fit] non-paper account {managed} — refusing cash query")
            return None
        ib.sleep(1.5)  # let the account-update subscription populate
        cash = None
        for v in ib.accountValues(managed[0]):
            if v.tag == "TotalCashValue" and v.currency in ("AUD", "BASE"):
                cash = float(v.value)
                break
        return cash
    except Exception as e:
        print(f"[cash-fit] IBKR cash query failed ({type(e).__name__}): {e}")
        return None
    finally:
        if ib.isConnected():
            ib.disconnect()


def _get_available_cash_aud() -> "float | None":
    """Best available cash figure for sizing the live plan to fit:
      1. live IBKR TotalCashValue (freshest — Gateway up), else
      2. cash_aud from the latest ibkr_nav_log.jsonl snapshot (offline), else
      3. None (caller keeps NAV-based sizing = current behavior).
    """
    cash = _fetch_ibkr_cash_aud()
    if cash is not None and np.isfinite(cash):
        print(f"[cash-fit] live IBKR TotalCashValue = ${cash:,.2f} AUD")
        return float(cash)
    try:
        p = APP_DIR / "ibkr_nav_log.jsonl"
        if p.exists():
            last = None
            for line in p.read_text(encoding="utf-8").splitlines():
                line = line.strip()
                if not line:
                    continue
                try:
                    last = json.loads(line)
                except Exception:
                    continue
            if last is not None and last.get("cash_aud") is not None:
                c = float(last["cash_aud"])
                if np.isfinite(c):
                    print(f"[cash-fit] using last broker snapshot cash_aud = "
                          f"${c:,.2f} AUD (ts {last.get('ts', '?')})")
                    return c
    except Exception as e:
        print(f"[cash-fit] snapshot cash read failed: {e}")
    return None


# NOTE: apply_ibkr_price_override moved to ibkr.py (Phase 4 split).


def compute_target_units_for_holdings(
    units_cur,
    last_px,
    fx_map,
    w_target,
    include_flags,
    portfolio_value_override=None,
    available_cash_aud=None,
):
    # Cash-fit sizing (mirrors make_trade_plan): when available_cash_aud is
    # given, size the target book to (holdings + cash - reserve) so the implied
    # net buys never exceed cash on hand. Falls back to NAV sizing otherwise.
    tickers = list(pd.Index(w_target.index))

    inc = pd.Series(include_flags).reindex(tickers).fillna(True).astype(bool)
    tickers = [t for t in tickers if inc.get(t, True)]

    lp_aud = (
        pd.Series(last_px).reindex(tickers).astype(float)
        * pd.Series(fx_map).reindex(tickers).fillna(1.0).astype(float)
    )
    cur_units = pd.Series(units_cur).reindex(tickers).fillna(0.0).astype(float)

    cur_val = float((cur_units * lp_aud).sum())  # actual current holdings value (AUD)
    if available_cash_aud is not None and np.isfinite(available_cash_aud):
        _reserve = max(CASH_RESERVE_MIN_AUD,
                       CASH_RESERVE_PCT * (cur_val + float(available_cash_aud)))
        cur_val = cur_val + max(0.0, float(available_cash_aud) - _reserve)
    elif portfolio_value_override is not None and np.isfinite(portfolio_value_override) and portfolio_value_override > 0:
        cur_val = float(portfolio_value_override)

    if cur_val <= 0:
        return pd.Series(0, index=w_target.index, dtype=int)

    tgt_val = pd.Series(w_target).reindex(tickers).fillna(0.0) * cur_val
    tgt_units_float = (tgt_val / lp_aud).replace([np.inf, -np.inf], np.nan).fillna(0.0)
    tgt_units_int = tgt_units_float.round().astype(int)
    # ASX minimum marketable parcel: a BUY that would ESTABLISH a position
    # under ~$500 AUD gets rejected at the exchange (adds to existing
    # holdings are exempt). Drop those targets so the plan never emits an
    # unexecutable order (IBKR cancelled a 4-unit VAE.AX buy, 2026-07-06).
    for t in tickers:
        _tu = int(tgt_units_int.get(t, 0))
        _pv = _tu * float(lp_aud.get(t, 0.0))
        if (str(t).upper().endswith(".AX")
                and float(cur_units.get(t, 0.0)) == 0.0
                and _tu > 0
                and _pv < ASX_MIN_MARKETABLE_PARCEL_AUD):
            print(f"[trade-plan] {t}: new-position target {_tu}u ≈ ${_pv:,.0f} AUD "
                  f"< ${ASX_MIN_MARKETABLE_PARCEL_AUD:,.0f} ASX min marketable parcel — dropped "
                  f"(will accumulate at a future rebalance)")
            tgt_units_int[t] = 0
    return tgt_units_int.reindex(w_target.index).fillna(0).astype(int)


def generate_targets_mvp_centric(mu_vec, Sigma, span_vol: float = 0.20, n_points: int = 20):
    """
    Compute MVP and return targets around MVP Â± span_vol * MVP_vol.
    """
    mu_arr = np.asarray(mu_vec, dtype=float)
    S = np.asarray(Sigma, dtype=float)
    n = len(mu_arr)

    x0 = np.full(n, 1.0 / n)
    cons = ({"type": "eq", "fun": lambda w: np.sum(w) - 1.0},)
    bnds = [(0.0, 1.0)] * n

    def vol_objective(w):
        return float(np.sqrt(max(w @ S @ w, 0.0)))

    sol = minimize(vol_objective, x0, bounds=bnds, constraints=cons, method="SLSQP")
    if not sol.success:
        raise RuntimeError("Could not compute MVP")

    w_mvp = np.asarray(sol.x, dtype=float)
    mu_mvp = float(w_mvp @ mu_arr)
    vol_mvp = float(np.sqrt(max(w_mvp @ S @ w_mvp, 0.0)))

    lo = mu_mvp - span_vol * vol_mvp
    hi = mu_mvp + span_vol * vol_mvp
    targets = np.linspace(lo, hi, n_points)
    return targets, mu_mvp, vol_mvp


# DIALOG-FIRST: Sigma_opt / mu_vec_opt diagnostic prints moved into
# _run_ff5_and_frontier_setup() (step 10). No module-level execution here.


# =====================================================================
# BLOCK 6 Transaction costs
# =====================================================================
# NOTE: BROKERAGE, MIN_TRADE_VALUE, _market_of, suppress_small_trades_by_value,
# compute_brokerage all moved to brokerage.py (Phase 4 split, 2026-06-29).
# Imported at top of file.


# NOTE: _is_long_term_au, _allocate_sale_to_lots, compute_cgt_tax moved to
# cgt.py (Phase 4 split, 2026-06-29).


def evaluate_transaction_costs(
    trade_df: pd.DataFrame,
    lots_df: pd.DataFrame,
    sale_date: pd.Timestamp,
    marginal_rate: float,
    carry_forward_loss: float = 0.0,
    method: str = "HIFO",
) -> dict:
    brokerage_total, brokerage_per_row = compute_brokerage(trade_df)
    tax_total, tax_bkd = compute_cgt_tax(
        trade_df,
        lots_df,
        sale_date,
        marginal_rate=float(marginal_rate),
        carry_forward_loss=float(carry_forward_loss),
        method=str(method),
    )
    return {
        "brokerage": brokerage_total,
        "cgt_tax": tax_total,
        "total_cost": brokerage_total + tax_total,
        "breakdown": tax_bkd,
        "per_row_brokerage": brokerage_per_row,
    }


def _update_lots_after_trades(
    lots_df: pd.DataFrame,
    trade_df: pd.DataFrame,
    sale_date: pd.Timestamp,
    fx_map: pd.Series | dict,
):
    """
    DEPRECATED — kept only because some test paths still call it.
    The previous trade-plan-driven write was the source of the SMH/SOXX
    corruption (3.4M phantom lots). Production code now rebuilds the
    Lots sheet each run via _build_lots_from_fills_log (or
    _build_lots_from_holdings when LOTS_REBUILD_MODE='holdings').

    Apply executed trades to lots table.
      - Sells: decrement matched lots using LOT_MATCH_METHOD.
      - Buys: append a new lot at current Last Px (AUD).
    Returns a new lots DataFrame.
    """
    base_cols = ["Security", "AcqDate", "Units", "CostBaseAUD"]
    out = lots_df.copy() if lots_df is not None else pd.DataFrame(columns=base_cols)

    for c in base_cols:
        if c not in out.columns:
            out[c] = np.nan

    out["AcqDate"] = pd.to_datetime(out["AcqDate"], errors="coerce")
    out["Units"] = pd.to_numeric(out["Units"], errors="coerce").fillna(0.0)
    out["CostBaseAUD"] = pd.to_numeric(out["CostBaseAUD"], errors="coerce").fillna(0.0)

    if trade_df is None or trade_df.empty:
        return out[base_cols].copy()

    delta_col = _trade_delta_col(trade_df)
    if delta_col is None:
        return out[base_cols].copy()

    for i, tr in trade_df.iterrows():
        sec = _security_from_row(i, tr)
        dU = int(pd.to_numeric(tr.get(delta_col, 0), errors="coerce") or 0)
        px_aud = float(pd.to_numeric(tr.get("Last Px (AUD)", 0.0), errors="coerce") or 0.0)

        if dU < 0:
            lot_block = out[out["Security"] == sec].copy()
            if lot_block.empty:
                continue

            if str(LOT_MATCH_METHOD).upper() == "HIFO":
                lot_block = lot_block.sort_values(by=["CostBaseAUD", "AcqDate"], ascending=[False, True])
            else:
                lot_block = lot_block.sort_values(by=["AcqDate"], ascending=True)

            remaining = abs(dU)
            for lot_idx in lot_block.index:
                if remaining <= 0:
                    break
                have = float(out.at[lot_idx, "Units"])
                take = min(remaining, have)
                out.at[lot_idx, "Units"] = have - take
                remaining -= take

            out = out[out["Units"] > 0.0].copy()

        elif dU > 0:
            new_lot = pd.DataFrame(
                [
                    {
                        "Security": sec,
                        "AcqDate": pd.Timestamp(sale_date),
                        "Units": int(dU),
                        "CostBaseAUD": px_aud,
                    }
                ]
            )
            out = pd.concat([out, new_lot], ignore_index=True)

    return out[base_cols].copy()


# ------------------------------
# Parcel-matching helper (used by trade-plan + CGT audit)
# ------------------------------
# Load parcels once (if sheet missing, returns empty table)
lots_df = _read_lots_from_path(filename, "Lots")


# =====================================================================
# BLOCK 7: Writing into Excel (workbook builder)
# =====================================================================
def ensure_workbook(path):
    if os.path.exists(path):
        return
    with xw.App(visible=False, add_book=True) as app:
        wb = app.books.add()
        for nm in ["Holdings","Tilts","OPT","Input","Cov","FF5F","Lots"]:
            try: wb.sheets[nm]
            except: wb.sheets.add(nm)
        # Minimal headers
        wb.sheets["Holdings"].range("A1").value = [["Security","Units","Last Price","FX to AUD","Market Value","Weight","Include?"]]
        wb.sheets["Tilts"].range("A1").value = [["Factor","Target","Band","Use?"]]
        wb.sheets["Tilts"].range("A2").value = [[f, (1.0 if i==0 else 0.0), 0.20, (i==0)] for i,f in enumerate(TILT_FACTORS)]
        wb.sheets["Lots"].range("A1").value = [["Security","AcqDate","Units","CostBaseAUD"]]
        wb.save(path); wb.close()

def warn_if_workbook_locked(path):
    """Print a clear warning if the workbook has an Office lock file (~$Name.xlsm) sibling.
    A lock means Excel (or another process) has the file open — writes via xlwings will either
    fall back to a read-only copy or trigger a 'File in Use' dialog. Non-fatal: we just surface it."""
    try:
        d, f = os.path.split(path)
        lock = os.path.join(d, "~$" + f)
        if os.path.exists(lock):
            print(f"[warn] Workbook lock file detected: {lock}")
            print("[warn] Close any open Excel windows for 'Stock Analysis.xlsm' (including stray "
                  "EXCEL.EXE processes in Task Manager) before continuing — otherwise the script will "
                  "either save to an _AUTO copy or you'll see a 'File in Use' dialog.")
    except Exception:
        pass

# Call it right before Block 7 seed reads:
ensure_workbook(filename)
warn_if_workbook_locked(filename)
print("[cfg] excel_path:", filename)

# ----- xlwings sheet/format helpers (shared across the Excel builder) -----
# Define path for saving portfolio state if not already defined
state_path = os.path.join(os.path.dirname(filename), "portfolio_state.json")
global results

OPEN_EXCEL_AFTER_SAVE = bool(globals().get("OPEN_EXCEL_AFTER_SAVE", CFG.get("open_excel_after_save", True)))
OPEN_PPT_AFTER_SAVE = bool(globals().get("OPEN_PPT_AFTER_SAVE", CFG.get("open_ppt_after_save", True)))

# -------------------------------
# Writers (used by Block 7)
# -------------------------------
def update_efficient_frontier_chart(
    opt_sheet, stats_df, start_s_row, rf_annual,
    tan_ret, tan_vol, current_point,
    title_text, target_point=None,
    previous_point=None, factor_point=None,
    no_tilt_point=None,
    tilt_point=None
):
    """Safe no-crash chart updater. Creates/refreshes a single scatter chart named 'Efficient Frontier' on OPT."""
    try:
        co = opt_sheet.api.ChartObjects()

        # stats_df was written at A{start_s_row+1} with a header row,
        # so first numeric data row is start_s_row+2
        first_row = start_s_row + 2
        last_row  = first_row + len(stats_df) - 1

        # Build x/y ranges up front (every time)
        x_rng = opt_sheet.range(f"C{first_row}:C{last_row}").api  # Volatility (ann.)
        y_rng = opt_sheet.range(f"B{first_row}:B{last_row}").api  # Achieved Return

        # Find existing chart by title
        the_chart = None
        for i in range(1, co.Count + 1):
            ch_i = co.Item(i).Chart
            if ch_i.HasTitle and "Efficient Frontier" in str(getattr(ch_i.ChartTitle, "Text", "")):
                the_chart = co.Item(i)
                break

        # Create if missing
        if the_chart is None:
            the_chart = opt_sheet.api.ChartObjects().Add(10, 10, 600, 360)  # x,y,w,h
            the_chart.Chart.ChartType = 74  # xlXYScatterSmoothNoMarkers
            the_chart.Chart.HasTitle = True
            the_chart.Chart.ChartTitle.Text = "Efficient Frontier"

        ch = the_chart.Chart

        # Clear series
        try:
            while ch.SeriesCollection().Count > 0:
                ch.SeriesCollection(1).Delete()
        except Exception:
            pass

        # Efficient Frontier (smooth line, no markers)
        s1 = ch.SeriesCollection().NewSeries()
        s1.Name = "Efficient Frontier"
        s1.XValues = x_rng
        s1.Values  = y_rng
        try:
            s1.ChartType = 74          # xlXYScatterSmoothNoMarkers
            s1.MarkerStyle = -4142     # xlMarkerStyleNone
            s1.Smooth = True
        except Exception:
            pass

        # Current portfolio marker
        if current_point:
            s3 = ch.SeriesCollection().NewSeries()
            s3.Name = "Current"
            s3.XValues = [float(current_point[0])]
            s3.Values  = [float(current_point[1])]
            try:
                s3.ChartType = -4169
                s3.MarkerStyle = 8
                s3.MarkerSize = 8
            except Exception:
                pass

        # Previous portfolio marker
        if previous_point:
            sp = ch.SeriesCollection().NewSeries()
            sp.Name = "Previous"
            sp.XValues = [float(previous_point[0])]
            sp.Values  = [float(previous_point[1])]
            try:
                sp.ChartType = -4169
                sp.MarkerStyle = 9
                sp.MarkerSize = 9
            except Exception:
                pass

        # Factor-effected marker
        if factor_point:
            sf = ch.SeriesCollection().NewSeries()
            sf.Name = "Factor-effected"
            sf.XValues = [float(factor_point[0])]
            sf.Values  = [float(factor_point[1])]
            try:
                sf.ChartType = -4169
                sf.MarkerStyle = 4
                sf.MarkerSize = 10
            except Exception:
                pass

        # Target portfolio marker
        if target_point:
            s4 = ch.SeriesCollection().NewSeries()
            s4.Name = "Target"
            s4.XValues = [float(target_point[0])]
            s4.Values  = [float(target_point[1])]
            try:
                s4.ChartType = -4169
                s4.MarkerStyle = 2
                s4.MarkerSize = 10
            except Exception:
                pass

        # Target (No Tilts) point (if available)
        if no_tilt_point:
            s_nt = ch.SeriesCollection().NewSeries()
            s_nt.Name = "Target (No Tilts)"
            s_nt.XValues = [float(no_tilt_point[0])]
            s_nt.Values  = [float(no_tilt_point[1])]
            try:
                s_nt.ChartType = -4169      # markers only
                s_nt.MarkerStyle = 3        # triangle
                s_nt.MarkerSize = 9
            except Exception:
                pass
        # Target (With Tilts) point (if available)
        if tilt_point:
            s_tilt = ch.SeriesCollection().NewSeries()
            s_tilt.Name = "Target (With Tilts)"
            s_tilt.XValues = [float(tilt_point[0])]
            s_tilt.Values  = [float(tilt_point[1])]
            try:
                s_tilt.ChartType = -4169
                s_tilt.MarkerStyle = 8   # circle
                s_tilt.MarkerSize = 9
            except Exception:
                pass

        
        # Title
        ch.ChartTitle.Text = title_text if title_text else "Efficient Frontier"

    except Exception as e:
        print(f"[chart] Skipping chart update (safe wrapper): {e}")


# ---- 10A) Read seeds (no COM; avoids UsedRange issues) ----

seed_units, seed_include = _read_holdings_seed_from_path(filename, sheet_name="Holdings")
tilt_seed = _read_tilts_seed_from_path(filename, sheet_name="Tilts")

# Holdings staleness check (2026-06-27). If ibkr_fills_log.jsonl has
# fills newer than the last engine run that reconciled Holdings, refuse
# to start — the engine would otherwise compute trades against stale
# positions. User must run triage_reset_*.py or update Holdings.Units
# manually before re-running. The check is bypassable via env var
# HOLDINGS_FRESHNESS_BYPASS=1 for triage runs and dev work.
try:
    if not bool(int(os.environ.get("HOLDINGS_FRESHNESS_BYPASS", "0") or "0")):
        _fills_path_check = APP_DIR / "ibkr_fills_log.jsonl"
        _holdings_path_check = APP_DIR / "Stock Analysis.xlsm"
        if _fills_path_check.exists() and _holdings_path_check.exists():
            # Latest FILLED row's exec_timestamp vs Holdings file mtime.
            # Holdings is rewritten by every engine run regardless of
            # reconciliation, so mtime is an UPPER BOUND on last touch
            # — using it means false-negatives (skipped checks) are
            # possible but no false-positives (spurious blocks).
            _latest_fill_ts = None
            try:
                with _fills_path_check.open("r", encoding="utf-8") as _f:
                    for _line in _f:
                        try:
                            _r = json.loads(_line.replace("NaN", "null"))
                            if int(_r.get("qty_filled") or 0) <= 0:
                                continue
                            _ts = _r.get("exec_timestamp")
                            if not _ts:
                                continue
                            _ts_parsed = pd.Timestamp(_ts)
                            if _latest_fill_ts is None or _ts_parsed > _latest_fill_ts:
                                _latest_fill_ts = _ts_parsed
                        except Exception:
                            continue
            except Exception:
                pass
            if _latest_fill_ts is not None:
                _holdings_mtime = pd.Timestamp(_holdings_path_check.stat().st_mtime,
                                                unit="s", tz=None)
                # Permit some clock skew. If holdings is at least 5 minutes
                # older than the latest fill, that's a stale signal.
                if _latest_fill_ts - _holdings_mtime > pd.Timedelta(minutes=5):
                    print()
                    print("=" * 88)
                    print("[HOLDINGS STALE] Engine refusing to run.")
                    print("=" * 88)
                    print(f"  Latest IBKR fill: {_latest_fill_ts}")
                    print(f"  Holdings mtime:   {_holdings_mtime}")
                    print(f"  Holdings is older than the latest broker fill.")
                    print(f"  Run engine against stale positions = trade plan against fiction.")
                    print()
                    print("  RESOLUTION:")
                    print("    1. Reconcile Holdings.Units against your actual IBKR positions")
                    print("       (open the sheet, update Units column, save)")
                    print("    2. OR re-run a triage script (triage_reset_*.py) to seed")
                    print("       Holdings from broker truth")
                    print("    3. Bypass for triage/dev only: HOLDINGS_FRESHNESS_BYPASS=1")
                    print("=" * 88)
                    raise SystemExit(2)
                else:
                    print(f"[holdings] freshness OK — latest fill {_latest_fill_ts.date()}, "
                          f"holdings touched {_holdings_mtime.date()}")
except SystemExit:
    raise
except Exception as _e_fresh:
    print(f"[holdings] freshness check skipped: {_e_fresh}")

# Ensure MOM exists in the seed and rows are in the canonical order
if not isinstance(tilt_seed, pd.DataFrame) or tilt_seed.empty:
    tilt_seed = pd.DataFrame(
        {"Target":[1.0] + [0.0]*(len(TILT_FACTORS)-1),
         "Band":[0.20]*len(TILT_FACTORS),
         "Use?":[True] + [False]*(len(TILT_FACTORS)-1)},
        index=TILT_FACTORS
    )
else:
    for f in TILT_FACTORS:
        if f not in tilt_seed.index:
            tilt_seed.loc[f] = {"Target":0.0, "Band":0.20, "Use?":False}
    tilt_seed = tilt_seed.reindex(TILT_FACTORS)

# Pre-dialog FF5 + frontier setup. Populates B, Sigma_opt, mu_vec_opt, etc.
# so that BOTH the smart tilt seed AND the dialog's "Auto-recommend tilts"
# button work correctly. The canonical setup ALSO runs post-dialog on the
# final universe — running it twice costs ~1 second and is harmless.
#
# Without this pre-dialog call:
#   - The tilt seed degrades to sheet/zero defaults (B is empty placeholder)
#   - The "Auto-recommend tilts" button silently returns zeros because
#     recommended_tilts_for_universe() reads mu_vec_opt/Sigma_opt and finds
#     empty placeholders → optimal_portfolio_tilts returns the zero series
try:
    if _SKIP_LIVE_PIPELINE:
        raise RuntimeError("skip-pipeline mode: skipping pre-dialog FF5 setup")
    _run_ff5_and_frontier_setup(prices)
    # Sync module-level locals from globals (mirror the post-dialog sync block).
    Sigma_opt = globals().get("Sigma_opt", Sigma_opt)
    mu_vec_opt = globals().get("mu_vec_opt", mu_vec_opt)
    B = globals().get("B", B)
    Fcov_daily = globals().get("Fcov_daily", Fcov_daily)
    f_mean_ann = globals().get("f_mean_ann", f_mean_ann)
except Exception as _e_pre_ff5_seed:
    print(f"[tilts] Pre-dialog FF5 setup skipped: {_e_pre_ff5_seed}")

# Seed the tilt TARGETS from the current portfolio's own factor exposures, so the
# dialog opens at "where you are now" rather than arbitrary hardcoded defaults.
try:
    if "B" in globals() and isinstance(B, pd.DataFrame) and not B.empty:
        _last_px_seed = (
            prices.ffill().iloc[-1]
            if isinstance(prices, pd.DataFrame) and not prices.empty
            else pd.Series(dtype=float)
        )
        _w_cur_seed = current_holdings_weights(
            seed_units,
            _last_px_seed,
            list(B.index),
            fx_map_all if "fx_map_all" in globals() else 1.0,
        )
        _cur_tilts = compute_achieved_tilts(B, _w_cur_seed, factors=TILT_FACTORS)
        if _cur_tilts is not None and float(pd.Series(_w_cur_seed).abs().sum()) > 0 and not _cur_tilts.dropna().empty:
            _new_targets = (
                pd.to_numeric(_cur_tilts.reindex(TILT_FACTORS), errors="coerce").fillna(0.0).round(3)
            )
            tilt_seed["Target"] = _new_targets
            print(f"[tilts] Seeded dialog targets from current portfolio factor exposures: "
                  f"{_new_targets.to_dict()}")
            # Diagnostic: also print the top-5 holding weights so we can see
            # what's driving the seed (if all targets round to 0.000 the
            # current portfolio probably has factor-neutral diversification).
            _top_w = pd.Series(_w_cur_seed).sort_values(ascending=False).head(5)
            print(f"[tilts] Seed weights (top 5): {_top_w.to_dict()}")
except Exception as _e_seed:
    print(f"[tilts] Could not seed from current portfolio; using sheet/defaults: {_e_seed}")

# ---- 10B) Combined dialog (holdings + tilts) ----
if _SKIP_LIVE_PIPELINE or _AUTO_PIPELINE_MODE:
    # Diagnostic modes bypass the modal dialog so the run reaches the
    # downstream gate. --auto-pipeline ALSO bypasses (fixed 2026-07-08):
    # its docstring always promised this but the flag was never wired in,
    # so unattended scheduled runs opened an INVISIBLE Tk dialog in the
    # non-interactive session and hung forever — the original 9:30 orphan
    # factory. Sheet seeds (same as user cancelling); TRADE_PLAN_MODE is
    # already forced to 'ensemble' for auto mode at the flag block.
    _why = "auto-pipeline non-interactive" if _AUTO_PIPELINE_MODE else "skip-pipeline"
    print(f"[{_why}] bypassing holdings dialog (using sheet seeds)")
    res = None
else:
    _wire_dialogs_module()
    res = edit_holdings_and_tilts_dialog(
        prices=prices,
        exclude=EXCLUDE_FROM_OPT,
        seed_units=current_holdings_units if 'current_holdings_units' in globals() and current_holdings_units is not None else seed_units,
        seed_include=seed_include,
        seed_tilts=tilt_seed
    )
    # Sync back the globals the dialog's Save handlers write — they now
    # land in the dialogs module's namespace, not this one.
    TRADE_PLAN_MODE = _dialogs.TRADE_PLAN_MODE
    OPEN_EXCEL_AFTER_SAVE = _dialogs.OPEN_EXCEL_AFTER_SAVE
    OPEN_PPT_AFTER_SAVE = _dialogs.OPEN_PPT_AFTER_SAVE
portfolio_value_override = None
if res is None:
    units = seed_units.copy()
    include_flags = seed_include.copy()
    last_px_hold = prices.ffill().iloc[-1].reindex(units.index)
    tilt_df = tilt_seed.copy()
else:
    if len(res) == 6:
        units, last_px_hold, prices, include_flags, tilt_df, portfolio_value_override = res
    else:
        units, last_px_hold, prices, include_flags, tilt_df = res

# Fallback: in --auto-pipeline / skip-live-pipeline modes the dialog never
# fires, so portfolio_value_override stays None. Read portfolio_state.json
# so the OOS backtest runs at the investor's actual NAV (matters because
# IBKR's $5 min binds tighter at $100k than at $1M — friction scales).
# The Roadshow chart's base is also driven from this value so the
# "$X invested" label matches what was actually simulated.
if portfolio_value_override is None:
    try:
        if os.path.exists(state_path):
            with open(state_path, "r") as _f_state:
                _state = json.load(_f_state)
            _pv = float(_state.get("portfolio_value", 0) or 0)
            if np.isfinite(_pv) and _pv > 0:
                portfolio_value_override = _pv
                print(f"[oos-nav] using portfolio_state.json value "
                      f"${_pv:,.0f} as starting_nav_aud for backtest + "
                      f"Roadshow base")
    except Exception as _e_pv_state:
        print(f"[oos-nav] portfolio_state.json read skipped: {_e_pv_state}")

    current_holdings_units = units.copy()

# DIALOG-FIRST: regardless of whether the dialog was saved or cancelled, run
# the canonical FF5 + frontier setup on the FINAL prices DataFrame. This is
# the ONLY FF5/Sigma_opt build path now — the module-level initial pass was
# removed in favour of this single call. Sigma_opt / mu_vec_opt / B /
# W / stats_df / tan_ret / tan_vol / cov_plus / exp_ret_df all come from
# here.
if _SKIP_LIVE_PIPELINE:
    print("[skip-pipeline] skipping live FF5/frontier setup")
try:
    if _SKIP_LIVE_PIPELINE:
        raise RuntimeError("skip-pipeline mode: skipping FF5 setup")
    _run_ff5_and_frontier_setup(prices)
    # Sync module-level locals with the globals the function just wrote.
    Sigma_opt = globals()["Sigma_opt"]
    mu_vec_opt = globals()["mu_vec_opt"]
    B = globals()["B"]
    alpha_daily = globals().get("alpha_daily", alpha_daily)
    resid_var = globals().get("resid_var", resid_var)
    ff5_regression_stats = globals().get("ff5_regression_stats", ff5_regression_stats)
    f_mean_ann = globals().get("f_mean_ann", f_mean_ann)
    Fcov_daily = globals().get("Fcov_daily", Fcov_daily)
    ff_aud = globals().get("ff_aud", ff_aud)
    Sigma_ff_daily = globals().get("Sigma_ff_daily", Sigma_ff_daily)
    Sigma_frontier = globals().get("Sigma_frontier", Sigma_frontier)
    mu_frontier = globals().get("mu_frontier", mu_frontier)
    mu_plus = globals().get("mu_plus", mu_plus)
    cov_plus = globals().get("cov_plus", cov_plus)
    exp_ret_df = globals().get("exp_ret_df", exp_ret_df)
    exp_ret_label = globals().get("exp_ret_label", exp_ret_label)
    W = globals().get("W", W)
    stats_df = globals().get("stats_df", stats_df)
    tan_ret = globals().get("tan_ret", tan_ret)
    tan_vol = globals().get("tan_vol", tan_vol)
    tilt_reco_achievable = globals().get("tilt_reco_achievable", tilt_reco_achievable)
    tilt_reco = globals().get("tilt_reco", tilt_reco)
except Exception as _e_ff5_setup:
    print(f"[ff5-setup] FAILED to build FF5/frontier universe: {_e_ff5_setup}. "
          f"Downstream code will likely crash — check the dialog's price fetch.")

# ---- Make optimiser globals available ----
current_holdings_units = units
securities_opt = list(units.index)
lots_df = lots_df  # already loaded earlier in block 7
gamma_cgt = 0.005        # soft penalty weight for CGT (tune as desired)
beta_brokerage = 0.25    # soft penalty weight for brokerage (tune as desired)

# --- helper: rebuild analytics from (possibly updated) prices ---
def _rebuild_core_from_prices(prices, fx_ticker="USDAUD=X", period="5y"):
    fx_raw = yf.download(fx_ticker, period=period, interval="1d",
                         auto_adjust=True, threads=False, progress=False)
    fx = fx_raw["Close"] if isinstance(fx_raw, pd.DataFrame) else fx_raw
    if isinstance(fx, pd.DataFrame):
        fx = fx.iloc[:, 0]
    fx = pd.to_numeric(fx, errors="coerce").reindex(prices.index).ffill()

    usd_cols = [c for c in prices.columns
                if not str(c).endswith(".AX") and not str(c).startswith("^")]

    prices_aud = prices.copy()
    prices_aud = prices_aud.drop(columns=[c for c in ["PortfolioValue"] if c in prices_aud.columns], errors="ignore")

    if usd_cols:
        prices_aud.update(prices.loc[:, usd_cols].mul(fx, axis=0))

    # FIX: fill missing values AFTER FX conversion but BEFORE returns
    prices_aud = prices_aud.ffill().bfill()

    # Melt into long format
    d = (prices_aud.reset_index()
         .melt(id_vars="Date", var_name="Security", value_name="Close")
         .sort_values(["Security", "Date"]))

    d["Return"] = d.groupby("Security", sort=False)["Close"].pct_change(fill_method=None)
    # Same outlier guard as the canonical df_returns build (see top-of-file helper).
    # Silent here to avoid double-logging — the canonical path already printed.
    d = _drop_return_outliers(d, verbose=False)
    d = d.dropna()

    df_cov_wide = d.pivot(index="Date", columns="Security", values="Return").sort_index()
    rets_opt = df_cov_wide.dropna(how="any")
    Sigma_daily = df_cov_wide.cov()

    d["LogRet"] = np.log1p(d["Return"])
    mu_log_ann = d.groupby("Security")["LogRet"].mean() * 252.0
    mu_ann_geo = np.expm1(mu_log_ann)

    return prices_aud, d, df_cov_wide, Sigma_daily, mu_ann_geo


def run_oos_walk_forward(
    prices_aud: pd.DataFrame,
    train_window_months: int = 24,
    rebalance: str = "MS",
    oos_start=None,
    oos_end=None,
    objective: str = "beat_benchmark",
    benchmark_ticker: str = "SPY",
    beat_premium: float = 0.02,
) -> tuple[pd.Series, pd.DataFrame]:
    """Walk-forward out-of-sample backtest of the long-only max-Sharpe strategy.

    At each rebalance date t we fit mu (annualised geometric) and Sigma (daily
    cov) on the trailing `train_window_months` of AUD-adjusted prices, solve the
    same long-only tangency portfolio used live (max_sharpe_long_only), then hold
    those weights through the next rebalance, accumulating realised daily returns.

    Mirrors the live recipe (geometric mu, daily Sigma, outlier guard) so the
    OOS curve is comparable to the live optimisation rather than a different
    statistical estimator.

    Returns (daily_strategy_returns, weights_history) where weights_history is
    indexed by rebalance date and columns are tickers.
    """
    px = prices_aud.copy()
    px.index = pd.to_datetime(px.index).tz_localize(None)
    px = px.sort_index().ffill().bfill()
    # PortfolioValue is a derived column — exclude from the asset universe.
    px = px.drop(columns=[c for c in ["PortfolioValue"] if c in px.columns], errors="ignore")

    if oos_end is None:
        oos_end = px.index.max()
    oos_end = pd.Timestamp(oos_end)
    lead = pd.DateOffset(months=train_window_months)
    if oos_start is None:
        oos_start = px.index.min() + lead
    else:
        oos_start = max(pd.Timestamp(oos_start), px.index.min() + lead)

    # Pre-compute daily returns once, with the same outlier guard as the live pipeline.
    daily_rets = px.pct_change()
    daily_rets = daily_rets.where(daily_rets.abs() <= RETURN_OUTLIER_THRESHOLD)

    # Rebalance schedule: first trading day on/after each calendar month-start.
    cal_dates = pd.date_range(start=oos_start, end=oos_end, freq=rebalance)
    rebal_dates = []
    for d in cal_dates:
        loc = px.index.searchsorted(d, side="left")
        if loc < len(px.index):
            rebal_dates.append(px.index[loc])
    rebal_dates = pd.DatetimeIndex(sorted(set(rebal_dates)))
    if len(rebal_dates) == 0:
        return pd.Series(dtype=float), pd.DataFrame()

    weights_history: dict[pd.Timestamp, pd.Series] = {}
    segments: list[pd.Series] = []

    _n_rebals_total = len(rebal_dates)
    _progress_every = max(10, _n_rebals_total // 10)  # ~10 progress prints across run
    _t_oos_start = time.perf_counter()
    for i, t in enumerate(rebal_dates):
        # Progress beacon — every ~10% of rebalances.
        if i > 0 and i % _progress_every == 0:
            _elapsed = time.perf_counter() - _t_oos_start
            _pct = i / _n_rebals_total * 100
            _eta = (_elapsed / i) * (_n_rebals_total - i)
            print(f"  [oos-progress] rebal {i}/{_n_rebals_total} ({_pct:.0f}%) "
                  f"@ {t.date()}  elapsed={_elapsed:.1f}s  ETA={_eta:.0f}s")
        train_px = px.loc[t - lead : t]
        if len(train_px) < 60:
            continue
        train_rets = train_px.pct_change()
        train_rets = train_rets.where(train_rets.abs() <= RETURN_OUTLIER_THRESHOLD)

        # Coverage filter: keep tickers with >= 80% non-NaN obs in the window.
        coverage = train_rets.notna().sum() / max(len(train_rets), 1)
        good_cols = coverage[coverage >= 0.8].index.tolist()
        if len(good_cols) < 3:
            continue
        train_rets = train_rets[good_cols].dropna(how="any")
        if len(train_rets) < 60:
            continue

        log_ret = np.log1p(train_rets)
        mu = pd.Series(np.expm1(log_ret.mean() * 252.0), index=train_rets.columns)
        Sigma = train_rets.cov()

        w = pd.Series(dtype=float)
        if objective == "beat_benchmark" and benchmark_ticker in mu.index:
            # Min-vol s.t. E[r] >= benchmark_mu + premium. This delivers SPY-style
            # return targets while clipping variance — the actual fund pitch.
            target_ret = float(mu[benchmark_ticker]) + float(beat_premium)
            # Never target below the max-Sharpe expected return — otherwise the
            # constraint binds nowhere and the optimiser collapses to min-variance.
            tangency_w = max_sharpe_long_only(mu, Sigma, rf=0.0)
            if tangency_w is not None and not tangency_w.empty:
                tangency_mu = float((mu.reindex(tangency_w.index).fillna(0.0) * tangency_w).sum())
                target_ret = max(target_ret, tangency_mu)
            try:
                w_arr, ok, _note = solve_frontier_point_cvxpy(mu, Sigma, target_ret,
                                                              use_inequality=True)
                if ok and w_arr is not None and len(w_arr) > 0 and np.isfinite(w_arr).all():
                    w = pd.Series(w_arr, index=Sigma.index)
            except Exception:
                pass

        if w.empty:
            # Fallback (or default objective == "max_sharpe"): pure tangency
            try:
                w = max_sharpe_long_only(mu, Sigma, rf=0.0)
            except Exception:
                continue

        if w is None or w.empty:
            continue
        w = w[w > 1e-6]
        if w.empty:
            continue
        # Renormalise after clipping
        w = w / w.sum()
        weights_history[t] = w

        # Realised window: held from day after t until the day of the next rebalance.
        if i + 1 < len(rebal_dates):
            seg_end = rebal_dates[i + 1]
        else:
            seg_end = oos_end + pd.Timedelta(days=1)
        held = daily_rets.loc[t:seg_end, w.index].fillna(0.0)
        if len(held) > 0 and held.index[0] == t:
            held = held.iloc[1:]
        if held.empty:
            continue
        seg = (held * w.reindex(held.columns).fillna(0.0)).sum(axis=1)
        segments.append(seg)

    if not segments:
        return pd.Series(dtype=float), pd.DataFrame()

    oos_returns = pd.concat(segments).sort_index()
    # If rebalance dates produce a duplicate boundary day, keep the later (post-rebalance) value.
    oos_returns = oos_returns[~oos_returns.index.duplicated(keep="last")]
    oos_weights = pd.DataFrame.from_dict(weights_history, orient="index").fillna(0.0)
    return oos_returns, oos_weights


# --- Strategy ensemble (regime-aware multi-objective blend) ---

# Canonical 5-slot menu. Each slot specifies a return-floor premium over the
# benchmark's training-window mu (None = pure tangency, no floor).
# Slots span an upside-tilted aggression range — the goal is to BEAT SPY, not
# to defend versus it. Inverse-ETF / hedge ballast (BBUS/BEAR/GOLD/AGVT) is
# still available to the solver inside every slot when the optimiser deems
# it worthwhile; it just isn't structurally anchored by a Defensive bucket
# whose return target undershoots SPY.
ENSEMBLE_SLOTS: tuple[tuple[str, float | None], ...] = (
    # Modest now serves as the low-end fallback — solver hits SPY's expected
    # return at lower-vol weights when the regime signal is risk-off.
    ("Modest (SPY+0%)",       0.00),
    ("Aggressive (SPY+5%)",   0.05),
    ("Bold (SPY+10%)",        0.10),
    ("Maximum (SPY+15%)",     0.15),
    # Stretch at SPY+25%: forces concentration into the top-mu assets (SMH,
    # VLUE, VVLU, IOO, etc.) when feasible. Falls back to Maximum if infeasible.
    # This is the slot that actually competes with SPY's tech-driven runs.
    ("Stretch (SPY+25%)",     0.25),
)
ENSEMBLE_SLOT_NAMES: tuple[str, ...] = tuple(name for name, _ in ENSEMBLE_SLOTS)

# Sync engine-canonical config into the solvers module (module split #18). The
# solvers read caps via globals().get() and resolve default slots from their own
# module globals; the engine is the source of truth, so push the finalised values
# (caps already env-overridden above) once, before any solve runs. Same-object
# for the caps dicts so later in-place mutations propagate.
_solvers.PER_ASSET_WEIGHT_CAPS = PER_ASSET_WEIGHT_CAPS
_solvers.SECTOR_GROUP_CAPS = SECTOR_GROUP_CAPS
_solvers.ENSEMBLE_SLOTS = ENSEMBLE_SLOTS
_solvers.ENSEMBLE_SLOT_NAMES = ENSEMBLE_SLOT_NAMES

# Sync the ~25 config values the OOS engine reads (incl. ENSEMBLE_SLOT_NAMES,
# defined just above). Placed here — after all config is defined and before any
# pipeline execution — so both run_oos and the live-pipeline helper calls see it.
_sync_oos_engine()


# NOTE: softmax_ensemble_weights moved to ensemble.py (Phase 4 split).


# ============================================================================
# OOS disk cache — Phase 3a, 2026-06-29
# ----------------------------------------------------------------------------
# Backtest runs are expensive (~60-100s each). Cache results keyed by a
# deterministic hash of the inputs that actually affect the output:
#   - starting_nav_aud (drives brokerage scaling)
#   - data fingerprint (shape + max date + last row sum → catches data updates)
#   - config knobs (rebalance, lambda_temp, slot_weights, crash_hedge)
#   - engine version (BUILD_GIT_SHA)
# Cache invalidates whenever any of these change. Cache hits load in
# ~50ms vs ~90s recompute → ~1800× speedup when warm.
#
# Cache file format: pickle of the full ensemble_out dict. Pickle is
# acceptable here because the cache is in-repo, never shared.
# ============================================================================

_OOS_CACHE_DIR = APP_DIR / ".cache" / "oos"


def _oos_cache_fingerprint(prices_aud: pd.DataFrame,
                            starting_nav_aud: float,
                            **kwargs) -> str:
    """Compute deterministic 16-char hex key for an OOS run.

    Data fingerprint is intentionally coarse: shape + column set +
    date range only. The first cache attempt also hashed the last-row
    price sum to invalidate on intraday data updates, but yfinance
    returns floats that vary by sub-cent amounts between calls (FX
    recalc, intraday quotes), so that key was always changing → 0%
    cache hit rate. The current fingerprint is stable across runs
    that fetch the same daily-resolution window, which is what we
    actually care about. Data refreshes that add a new date cleanly
    invalidate via the `dates:` hash component.
    """
    import hashlib
    h = hashlib.sha256()
    try:
        h.update(f"shape:{prices_aud.shape}".encode())
        h.update(f"cols:{','.join(sorted(str(c) for c in prices_aud.columns))}".encode())
        h.update(f"dates:{prices_aud.index.min()}_{prices_aud.index.max()}".encode())
    except Exception:
        h.update(b"data_hash_failed")
    # NAV
    h.update(f"nav:{float(starting_nav_aud):.2f}".encode())
    # All other kwargs in sorted order for determinism
    for k in sorted(kwargs.keys()):
        v = kwargs[k]
        # Serialize dict / list / scalar to a deterministic string
        try:
            h.update(f"{k}:{json.dumps(v, sort_keys=True, default=str)}".encode())
        except Exception:
            h.update(f"{k}:{repr(v)}".encode())
    # Engine version — invalidates cache on any code change
    try:
        gsha = str(globals().get("BUILD_GIT_SHA", "unknown"))
        h.update(f"git:{gsha}".encode())
    except Exception:
        pass
    # Behavioral config that changes the backtest but is NOT passed as
    # kwargs — read from globals/files at solve/TLH time. Without these,
    # editing caps or tlh_pairs.json between runs on identical data
    # silently returns stale results (nearly invalidated the 2026-07-02
    # TLH-pairs A/B: lockboxed data → same key → HIT on the old config).
    try:
        h.update(f"caps:{json.dumps(globals().get('PER_ASSET_WEIGHT_CAPS', {}), sort_keys=True)}".encode())
        h.update(f"gcaps:{json.dumps(globals().get('SECTOR_GROUP_CAPS', {}), sort_keys=True)}".encode())
        h.update(f"mu_shrink:{float(globals().get('MU_SHRINKAGE_LAMBDA', 0.0) or 0.0)}".encode())
        h.update(f"lt_defer:{int(globals().get('LT_DEFER_WINDOW_DAYS', 0) or 0)}".encode())
        h.update(f"lt_defer_cond:{int(bool(globals().get('LT_DEFER_DD_CONDITIONAL', False)))}".encode())
        h.update(f"lt_defer_reldd:{float(globals().get('LT_DEFER_RELEASE_DD', 0.0) or 0.0)}".encode())
        h.update(f"skip_calm:{float(globals().get('SKIP_REBAL_DELTA_CALM', 0.0) or 0.0)}".encode())
        h.update(f"stretch_floor:{float(globals().get('STRETCH_FLOOR_CALM', 0.0) or 0.0)}".encode())
        h.update(f"stretch_pred:{int(bool(globals().get('STRETCH_FLOOR_PREDICTIVE', False)))}".encode())
        h.update(f"trend_sleeve:{float(globals().get('TREND_SLEEVE_WEIGHT', 0.0) or 0.0)}".encode())
        h.update(f"cov_shrink:{int(bool(globals().get('COV_SHRINKAGE', False)))}".encode())
        h.update(f"vol_target:{float(globals().get('VOL_TARGET_ANNUAL', 0.0) or 0.0)}".encode())
        h.update(f"crisis_hedge:{float(globals().get('CRISIS_HEDGE_WEIGHT', 0.0) or 0.0)}".encode())
        h.update(f"crisis_ma:{int(globals().get('CRISIS_HEDGE_MA_DAYS', 200))}".encode())
        h.update(f"crisis_band:{float(globals().get('CRISIS_HEDGE_BAND_SD', 0.0) or 0.0)}".encode())
    except Exception:
        h.update(b"caps_hash_failed")
    try:
        _pairs_p = APP_DIR / "tlh_pairs.json"
        if _pairs_p.exists():
            h.update(b"tlh_pairs:" + _pairs_p.read_bytes())
    except Exception:
        h.update(b"tlh_pairs_hash_failed")
    return h.hexdigest()[:16]


def _oos_cache_load(key: str) -> "Optional[dict]":
    """Try to load cached result. Returns None on miss or read error."""
    import pickle
    cache_file = _OOS_CACHE_DIR / f"oos_{key}.pkl"
    if not cache_file.exists():
        return None
    try:
        with open(cache_file, "rb") as f:
            return pickle.load(f)
    except Exception as e:
        print(f"[oos-cache] read failed for {key}: {e}")
        return None


def _oos_cache_save(key: str, value: dict) -> None:
    """Persist OOS result to cache. Best-effort — failure is non-fatal."""
    import pickle
    try:
        _OOS_CACHE_DIR.mkdir(parents=True, exist_ok=True)
        cache_file = _OOS_CACHE_DIR / f"oos_{key}.pkl"
        with open(cache_file, "wb") as f:
            pickle.dump(value, f, protocol=pickle.HIGHEST_PROTOCOL)
    except Exception as e:
        print(f"[oos-cache] write failed for {key}: {e}")


def run_oos_ensemble_walk_forward_cached(prices_aud: pd.DataFrame,
                                           **kwargs) -> dict:
    """Cached wrapper around run_oos_ensemble_walk_forward.

    Cache hit on identical (NAV, data fingerprint, config, engine version)
    returns the prior result instantly. Cache miss runs the full backtest
    and stores the result for next time.

    Useful for:
      - Scale-sensitivity sweeps when same NAVs appear repeatedly
      - Iterative dev work — change UI, rerun, OOS hits cache
      - Re-running after a partial failure (e.g. sanity halt → fix
        Holdings → rerun, OOS still cached)
    """
    starting_nav_aud = float(kwargs.get("starting_nav_aud", 1_000_000.0))
    # Pass `kwargs` as kwargs WITHOUT extracting starting_nav_aud — the
    # fingerprint expects nav as a positional arg AND other config as
    # kwargs. Extracting starting_nav_aud out of kwargs avoids the
    # "multiple values for argument 'starting_nav_aud'" error.
    _fp_kwargs = {k: v for k, v in kwargs.items() if k != "starting_nav_aud"}
    key = _oos_cache_fingerprint(prices_aud, starting_nav_aud, **_fp_kwargs)
    cached = _oos_cache_load(key)
    if cached is not None:
        print(f"[oos-cache] HIT key={key} NAV=${starting_nav_aud:,.0f} "
              f"(saved ~60-100s)")
        return cached
    print(f"[oos-cache] MISS key={key} NAV=${starting_nav_aud:,.0f} "
          f"— computing...")
    result = run_oos_ensemble_walk_forward(prices_aud, **kwargs)
    _oos_cache_save(key, result)
    return result


# --- OOS kernel mode early-exit (Phase 3b, 2026-06-29) ---
# Subprocess workers spawned by the parallel scale-sensitivity loop set
# OOS_KERNEL_MODE=1 before importing the engine. By the time execution
# reaches here, all the heavy work needed by run_oos_ensemble_walk_forward
# is in scope: imports, constants, helper functions (_run_tlh_pass,
# LotBook, etc.), the OOS function itself, and the cache wrapper. The
# parent process catches the SystemExit via importlib.exec_module and
# inspects the partially-imported module namespace to grab the function.
# Skips: dialog box, live OOS execution, PPT/Excel write — all the
# "main pipeline" that workers don't need.
if OOS_KERNEL_MODE:
    print(f"[oos-kernel] OOS_KERNEL_MODE=1 detected; "
          f"exiting after function defs to act as a library for worker subprocess")
    sys.exit(0)


# --- OOS metrics helpers (Phase 2) ---

# NOTE: _series_metrics, _ir_vs_bench, _capm_alpha_beta, _ff5_alpha moved to
# metrics.py (Phase 4 split).


def _annual_turnover(weights_history: pd.DataFrame, start_dt, end_dt) -> float:
    if weights_history is None or weights_history.empty:
        return np.nan
    w = weights_history.copy()
    w.index = pd.to_datetime(w.index).tz_localize(None)
    w = w.sort_index()
    mask = (w.index >= pd.Timestamp(start_dt)) & (w.index <= pd.Timestamp(end_dt))
    w = w[mask]
    if len(w) < 2:
        return np.nan
    # |w_t - w_{t-1}|, summed across tickers per rebalance, averaged, scaled to
    # actual rebalance frequency (12/yr monthly, 4/yr quarterly, etc.)
    dw = w.diff().abs().sum(axis=1).iloc[1:]
    return float(dw.mean() * float(globals().get("REBALANCES_PER_YEAR", 12)))


def compute_oos_metrics(strat_returns: pd.Series,
                        spy_returns: pd.Series,
                        aord_returns: pd.Series,
                        ff5_factors: pd.DataFrame | None = None,
                        weights_history: pd.DataFrame | None = None,
                        horizons_years: tuple[int, ...] = (3, 5, 10)) -> pd.DataFrame:
    """Build a (metric × [horizon, series]) MultiIndex DataFrame of OOS stats."""
    if strat_returns is None or strat_returns.empty:
        return pd.DataFrame()
    # Defensively normalise the indices (sort + tz-strip) so .loc slicing and
    # reindexing don't blow up on duplicated or non-monotonic data.
    strat_returns = strat_returns.copy()
    strat_returns.index = pd.to_datetime(strat_returns.index).tz_localize(None)
    strat_returns = strat_returns.sort_index()
    strat_returns = strat_returns[~strat_returns.index.duplicated(keep="last")]
    spy_returns = spy_returns.copy() if spy_returns is not None else pd.Series(dtype=float)
    if not spy_returns.empty:
        spy_returns.index = pd.to_datetime(spy_returns.index).tz_localize(None)
        spy_returns = spy_returns.sort_index()
        spy_returns = spy_returns[~spy_returns.index.duplicated(keep="last")]
    aord_returns = aord_returns.copy() if aord_returns is not None else pd.Series(dtype=float)
    if not aord_returns.empty:
        aord_returns.index = pd.to_datetime(aord_returns.index).tz_localize(None)
        aord_returns = aord_returns.sort_index()
        aord_returns = aord_returns[~aord_returns.index.duplicated(keep="last")]

    end_dt = strat_returns.index.max()

    metric_order = ["Cumulative Return", "Annualised Return", "Annualised Volatility",
                    "Sharpe Ratio", "Sortino Ratio", "Max Drawdown", "IR vs ^AORD",
                    "Beta vs SPY", "Alpha vs SPY (ann)", "Alpha vs FF5 (ann)",
                    "Annual Turnover"]

    blocks = []
    for h in horizons_years:
        start_dt = end_dt - pd.DateOffset(years=h)
        # Boolean mask avoids label-existence/monotonicity errors when
        # start_dt or end_dt happens to land on a non-trading day.
        s = strat_returns[(strat_returns.index >= start_dt) & (strat_returns.index <= end_dt)]
        if s.empty:
            continue
        sp = spy_returns.reindex(s.index).fillna(0.0)
        ao = aord_returns.reindex(s.index).fillna(0.0)

        m_s = _series_metrics(s)
        m_sp = _series_metrics(sp)
        m_ao = _series_metrics(ao)

        ir = _ir_vs_bench(s, ao)
        alpha_spy, beta_spy = _capm_alpha_beta(s, sp)
        alpha_ff5 = _ff5_alpha(s, ff5_factors) if ff5_factors is not None else np.nan
        turn = _annual_turnover(weights_history, start_dt, end_dt)

        strat_col = [m_s["Cumulative Return"], m_s["Annualised Return"],
                     m_s["Annualised Volatility"], m_s["Sharpe Ratio"],
                     m_s["Sortino Ratio"], m_s["Max Drawdown"], ir, beta_spy,
                     alpha_spy, alpha_ff5, turn]
        spy_col = [m_sp["Cumulative Return"], m_sp["Annualised Return"],
                   m_sp["Annualised Volatility"], m_sp["Sharpe Ratio"],
                   m_sp["Sortino Ratio"], m_sp["Max Drawdown"],
                   np.nan, np.nan, np.nan, np.nan, np.nan]
        aord_col = [m_ao["Cumulative Return"], m_ao["Annualised Return"],
                    m_ao["Annualised Volatility"], m_ao["Sharpe Ratio"],
                    m_ao["Sortino Ratio"], m_ao["Max Drawdown"],
                    np.nan, np.nan, np.nan, np.nan, np.nan]

        block = pd.DataFrame({"Strategy": strat_col, "SPY (AUD)": spy_col,
                              "^AORD": aord_col}, index=metric_order)
        block.columns = pd.MultiIndex.from_tuples(
            [(f"{h}Y", c) for c in block.columns]
        )
        blocks.append(block)

    if not blocks:
        return pd.DataFrame()
    return pd.concat(blocks, axis=1)


# === GFC stress test (Tier-1 #2) =============================================
# Standalone walk-forward over a pre-2006 minimal universe to ask: does the
# ensemble engine survive a -56% SPY drawdown? Triggered by `--stress-test`
# on the CLI; exits before the live pipeline runs.
#
# Stretch-only comparison: if --stretch-only is also passed, runs BOTH
# 5-slot blend and Stretch-only allocations side by side, reporting deltas
# so we can quantify how much of the engine's GFC defence comes from slot
# blending vs concentration in high-mu assets.

# Sync shared helpers + config into research_modes before any research dispatch
# below. All 15 injected names are defined by here (config, ENSEMBLE_SLOT_NAMES,
# and the 6 helper fns incl. conditional _apply_data_lockbox).
_sync_research_modes()

if "--stress-test" in sys.argv:
    _exit_code = _run_gfc_stress_test()
    sys.exit(_exit_code)


# === Scale analysis (--scale-analysis) ======================================
# Runs the OOS walk-forward at 6 starting NAVs ($10k → $10M) to surface the
# fee-drag curve. Answers: "at what AUM does this strategy actually work?"
# Same walk-forward engine, same universe, same ~10-year window — only the
# starting_nav changes, which scales the flat-fee component appropriately.
if "--scale-analysis" in sys.argv:
    _exit_code = _run_scale_analysis()
    sys.exit(_exit_code)


# === Dev/validation split (--dev-validation) ================================
# Re-run the same engine on two disjoint OOS windows to expose meta-parameter
# overfitting. The ensemble design (slot menu, lambda_temp, gaussian_width,
# halflife, weight caps, scoring rule) was selected on 2016-2026 data — so
# good 2016-2026 OOS Sharpe is not honest evidence of generalisation. The
# discipline going forward: tune on the DEV window only, then open the
# VALIDATION lock-box at most once per change. A big dev→validation gap means
# the change was overfit.
#
#   Dev OOS:        2015-01-01 → 2020-02-19  (SPY pre-COVID ATH)
#   Validation OOS: 2020-02-20 → today       (LOCK BOX)
#
# The engine derives oos_start = data_start + 24mo, so we pre-roll each data
# slice by 24 months to make the OOS evaluation window match exactly.
if "--dev-validation" in sys.argv:
    _exit_code = _run_dev_validation()
    sys.exit(_exit_code)


# === SKIP_REBAL_DELTA sweep (--rebal-skip-sweep) ============================
# Dev-only sweep of the skip-rebalance threshold to find the value that
# maximises net Sharpe on the dev window. The hypothesis is that a higher
# threshold (5-7% vs the current 3%) lets more positions cross the 12-month
# LT-discount line before being trimmed, reducing CGT drag. After picking
# the winner on DEV, we open the VALIDATION lock-box ONCE on the winner to
# confirm the choice generalises (per the dev/validation discipline in
# ARCHITECTURE.md §10).
if "--rebal-skip-sweep" in sys.argv:
    _exit_code = _run_rebal_skip_sweep()
    sys.exit(_exit_code)


# === Turnover-penalty sweep (--turnover-penalty-sweep) ======================
# Dev-only sweep of the cost-aware solver penalty. The penalty
# turnover_penalty * ||w - w_prev||_1 is added to the cvxpy objective in
# solve_frontier_point_cvxpy → solve_candidate_portfolios → walk-forward.
# Hypothesis: regime-switch turnover (the Modest ↔ Stretch slot flips)
# drives most of the engine's CGT drag. Penalising weight changes inside
# the optimiser should reduce that turnover without hurting risk
# adjustment. Tested values bracket the typical magnitude of w'Σw (~1e-4
# to 1e-3 daily variance for long-only weights).
if "--turnover-penalty-sweep" in sys.argv:
    _exit_code = _run_turnover_penalty_sweep()
    sys.exit(_exit_code)


# === Walk-forward CV (--walk-forward-cv) ===================================
# Multi-fold OOS evaluation: slice the walk-forward output into N non-
# overlapping calendar-year folds and report per-fold Sharpe + mean ± std.
# Designed as the preferred statistical hygiene tool for parameter
# selection — does NOT touch the validation lock-box, so it doesn't burn
# the peek budget documented in ARCHITECTURE.md §10.
#
# Each fold gives 1 year of OOS observations (~252 trading days). For
# 10-year history this yields 7-8 independent Sharpe estimates. Standard
# error of mean Sharpe scales as 1/sqrt(N), so 8 folds ≈ 35% noise
# reduction vs single-window evaluation.
#
# Future extension: wrap this in a parameter-sweep loop (gamma, lambda_temp,
# etc.) to do robust parameter selection — winner = max mean Sharpe across
# folds with reasonable std, NOT max single-window Sharpe.
if "--walk-forward-cv" in sys.argv:
    _exit_code = _run_walk_forward_cv()
    sys.exit(_exit_code)


# === Performance attribution (--attribution) ================================
# Decomposes the engine's OOS returns to answer:
#   1) Per-slot:   which of the 5 candidate slots adds Sharpe / α?
#   2) Per-asset:  which tickers are the biggest cumulative contributors?
#   3) Per-regime: bull (SPY 20d>50d) vs bear (SPY 20d<50d) split
# Purely descriptive — no parameters to tune, no validation budget consumed.
# Modern-universe filter (post-2016) matches walk-forward-cv default so
# results are directly comparable.
if "--attribution" in sys.argv:
    _exit_code = _run_attribution()
    sys.exit(_exit_code)


# === Crash-hedge A/B test (--crash-hedge-test) ==============================
# Walk-forward CV with crash hedge ON vs OFF. Per-fold uplift table +
# aggregate verdict. Used to test the asymmetric defensive overlay against
# the unguarded engine baseline.
if "--crash-hedge-test" in sys.argv:
    _exit_code = _run_crash_hedge_test()
    sys.exit(_exit_code)


# === Crash-hedge release-threshold sweep (--crash-hedge-release-sweep) =====
# Sweeps the hedge release threshold across {-3, -5, -8, -10, -12}% with
# trigger fixed at -15%. For each value: walk-forward CV with hedge ON,
# per-fold metrics, aggregate mean Sharpe + alpha vs baseline (no hedge).
# Picks winner by mean Sharpe across modern folds. Designed to fix the
# V-shape recovery miss observed in 2020 / 2025 with the -5% release.
if "--crash-hedge-release-sweep" in sys.argv:
    _exit_code = _run_crash_hedge_release_sweep()
    sys.exit(_exit_code)


# === Stretch-only A/B test (--stretch-only-test) ===========================
# Walk-forward CV comparing the 5-slot ensemble against a Stretch-only
# allocation (slot_weights_override forcing 100% on Stretch). Tests
# whether the defensive slots (Modest/Aggressive) are worth their alpha
# cost — per attribution, Modest has alpha -5.5%/yr and Stretch +4.8%/yr
# so it's plausible that Stretch alone beats the blend. Settles the
# question definitively before any further defensive-layer redesign.
if "--stretch-only-test" in sys.argv:
    _exit_code = _run_stretch_only_test()
    sys.exit(_exit_code)


# === Stretch + hedge synthesis sweep (--stretch-hedge-sweep) ================
# THE synthesis test: Stretch-only base (proven +3.6% alpha in modern era)
# combined with crash hedge overlay at multiple release thresholds. The
# Stretch base lacks intrinsic defense — so the hedge has more headroom
# to add value here than it did on the 5-slot blend (where slot blending
# was already providing defense).
if "--stretch-hedge-sweep" in sys.argv:
    _exit_code = _run_stretch_hedge_sweep()
    sys.exit(_exit_code)


# === Auto factor-tilt A/B test (--tilted-ensemble-test) =====================
# Walk-forward CV comparing 5-slot ensemble WITH auto-recommended factor
# tilts (recomputed per rebal from trailing-3M factor Sharpes) vs the
# baseline ensemble without tilts. Per-fold + aggregate metrics plus
# full-period peak-to-trough MaxDD (per the 2026-06-19 measurement lesson).
if "--tilted-ensemble-test" in sys.argv:
    _exit_code = _run_tilted_ensemble_test()
    sys.exit(_exit_code)


# === Rebuild core analytics ===
prices_aud_for_returns, df_melt, df_cov_wide, Sigma_daily, mu_ann_geo = _rebuild_core_from_prices(prices)
globals()["returns_wide_df"] = df_cov_wide.copy()

# === OOS walk-forward validation (Task #7 Part B) ===
# Default ON via config.oos_validation. Skipped silently when False.
# NOTE: the live pipeline downloads only 2y of prices (PRICE_DOWNLOAD_PERIOD),
# which is too short to start a 24mo-trained walk-forward backtest. We fetch
# a separate long-history (12y) view here so live optimisation behaviour is
# unchanged. Task #11 will consolidate this into a single data store.
oos_returns_daily = pd.Series(dtype=float)
oos_weights_history = pd.DataFrame()
oos_prices_aud_long = pd.DataFrame()
if bool(CFG.get("oos_validation", True)):
    try:
        _oos_t0 = time.perf_counter()
        _oos_tickers = [c for c in prices.columns if c != "PortfolioValue"]
        _oos_long_raw = yf.download(
            _oos_tickers,
            period="12y",
            interval="1d",
            auto_adjust=True,
            threads=False,
            progress=False,
        )
        _oos_long_px = _normalize_yfinance_close(_oos_long_raw)
        _oos_long_px.index = pd.to_datetime(_oos_long_px.index).tz_localize(None)
        _oos_long_px = _oos_long_px.sort_index().ffill().bfill()
        # FX-adjust USD tickers into AUD (same recipe as _rebuild_core_from_prices).
        _fx_raw = yf.download("USDAUD=X", period="12y", interval="1d",
                              auto_adjust=True, threads=False, progress=False)
        _fx = _fx_raw["Close"] if isinstance(_fx_raw, pd.DataFrame) else _fx_raw
        if isinstance(_fx, pd.DataFrame):
            _fx = _fx.iloc[:, 0]
        _fx = pd.to_numeric(_fx, errors="coerce").reindex(_oos_long_px.index).ffill()
        _usd_cols = [c for c in _oos_long_px.columns
                     if not str(c).endswith(".AX") and not str(c).startswith("^")]
        oos_prices_aud_long = _oos_long_px.copy()
        if _usd_cols:
            oos_prices_aud_long.update(_oos_long_px.loc[:, _usd_cols].mul(_fx, axis=0))
        oos_prices_aud_long = oos_prices_aud_long.ffill().bfill()

        # Use the investor's actual NAV so brokerage friction (IBKR $5 min)
        # scales correctly — at $100k the min binds, at $1M the rate
        # dominates. Without this the Roadshow slide would compound
        # $1M-scale returns onto a $100k chart base and overstate
        # achievable net returns at small scale.
        _oos_nav = (float(portfolio_value_override)
                    if portfolio_value_override is not None
                       and np.isfinite(portfolio_value_override)
                       and portfolio_value_override > 0
                    else 1_000_000.0)
        ensemble_out = run_oos_ensemble_walk_forward_cached(
            oos_prices_aud_long,
            train_window_months=24,
            rebalance=REBALANCE_FREQ,
            benchmark_ticker="SPY",
            score_lookback_days=252,
            lambda_temp=3.0,
            starting_nav_aud=_oos_nav,
            # PRODUCTION CONFIG: the metrics shown in PPT/Excel must match
            # what the live trade plan actually trades. See PRODUCTION_*
            # constants at the top of the file for the empirical rationale.
            slot_weights_override=PRODUCTION_SLOT_OVERRIDE,
            crash_hedge=PRODUCTION_CRASH_HEDGE,
        )
        globals()["_oos_starting_nav_aud"] = float(_oos_nav)
        oos_returns_daily = ensemble_out["blended_returns"]
        oos_weights_history = ensemble_out["blended_weights"]

        # Optional second OOS backtest at ROADSHOW_BASE_NAV — for the PDS /
        # Roadshow chart, which compares the fund at the user's NAV vs at
        # the assumed wholesale entry size. Both lines appear on the
        # chart and three extra rows in the metrics table. Off by default
        # so daily ops aren't slowed.
        oos_returns_daily_roadshow = pd.Series(dtype=float)
        oos_weights_history_roadshow = pd.DataFrame()
        if ROADSHOW_DUAL_NAV and abs(ROADSHOW_BASE_NAV - _oos_nav) > 1.0:
            try:
                print(f"[oos-roadshow] running second backtest at "
                      f"${ROADSHOW_BASE_NAV:,.0f} for dual-NAV chart...")
                ensemble_out_rs = run_oos_ensemble_walk_forward_cached(
                    oos_prices_aud_long,
                    train_window_months=24,
                    rebalance=REBALANCE_FREQ,
                    benchmark_ticker="SPY",
                    score_lookback_days=252,
                    lambda_temp=3.0,
                    starting_nav_aud=float(ROADSHOW_BASE_NAV),
                    slot_weights_override=PRODUCTION_SLOT_OVERRIDE,
                    crash_hedge=PRODUCTION_CRASH_HEDGE,
                )
                oos_returns_daily_roadshow = ensemble_out_rs["blended_returns"]
                oos_weights_history_roadshow = ensemble_out_rs["blended_weights"]
                print(f"[oos-roadshow] second backtest complete, "
                      f"{len(oos_returns_daily_roadshow)} days")
            except Exception as _e_rs:
                print(f"[oos-roadshow] second backtest failed: {_e_rs}")
        globals()["oos_returns_daily_roadshow"] = oos_returns_daily_roadshow
        globals()["oos_weights_history_roadshow"] = oos_weights_history_roadshow
        globals()["_roadshow_nav_aud"] = float(ROADSHOW_BASE_NAV)

        # Scale-sensitivity sweep — N additional OOS backtests at the
        # NAVs in SCALE_SENSITIVITY_NAVS. Used by the dedicated scale
        # slide that follows Roadshow in the PDS deck. Reuses already-
        # computed results where the NAVs match (within $1 tolerance)
        # so we don't waste 100s per duplicate.
        oos_scale_results: dict[float, dict] = {}
        if SCALE_SENSITIVITY and SCALE_SENSITIVITY_NAVS:
            # Index the primary + roadshow runs by NAV for reuse.
            _existing: dict[float, pd.Series] = {}
            if not oos_returns_daily.empty:
                _existing[float(_oos_nav)] = oos_returns_daily
            if (ROADSHOW_DUAL_NAV
                and not oos_returns_daily_roadshow.empty
                and abs(ROADSHOW_BASE_NAV - _oos_nav) > 1.0):
                _existing[float(ROADSHOW_BASE_NAV)] = oos_returns_daily_roadshow

            # Phase 3b parallelism: bucket the NAVs into reuses (instant)
            # vs cache misses (compute-bound). Cache-miss NAVs are dispatched
            # to a ProcessPoolExecutor when SCALE_PARALLEL is on (default
            # on). Each worker subprocess imports the engine in
            # OOS_KERNEL_MODE=1 (skips main pipeline), grabs the OOS
            # function, runs it, returns the result via pickle.
            #
            # Why we still check cache: cache hits are ~50ms each, no
            # point spawning a 30s subprocess for them. Only true cache
            # misses get parallelised.
            #
            # Disable via SCALE_PARALLEL=0 (env var) if the worker
            # subprocess path proves flaky in the frozen exe environment.
            _scale_parallel = bool(int(os.environ.get(
                "SCALE_PARALLEL", "1") or "0"))
            _common_kwargs = dict(
                train_window_months=24,
                rebalance=REBALANCE_FREQ,
                benchmark_ticker="SPY",
                score_lookback_days=252,
                lambda_temp=3.0,
                slot_weights_override=PRODUCTION_SLOT_OVERRIDE,
                crash_hedge=PRODUCTION_CRASH_HEDGE,
            )
            _navs_to_compute: list[float] = []
            for _nav in SCALE_SENSITIVITY_NAVS:
                _tol = max(1.0, 0.01 * _nav)
                _match = next((k for k in _existing
                               if abs(k - _nav) <= _tol), None)
                if _match is not None:
                    oos_scale_results[float(_nav)] = {
                        "returns": _existing[_match],
                        "nav_aud": float(_nav),
                        "reused": True,
                    }
                    print(f"[scale] @ ${_nav:,.0f}: reused existing backtest")
                    continue
                _navs_to_compute.append(float(_nav))

            if _navs_to_compute and _scale_parallel and len(_navs_to_compute) > 1:
                # Parallel path via direct subprocess.Popen (NOT
                # ProcessPoolExecutor — the latter auto-reimports the
                # parent script in each worker, which then hits the
                # OOS_KERNEL_MODE sys.exit and dies before the worker
                # function can be called).
                #
                # Each worker is a standalone `python oos_worker.py
                # in.pkl out.pkl` invocation. The worker imports
                # Portfolio_Optimiser in kernel mode, gets the OOS
                # function, runs it, writes the pickled result. Parent
                # waits for all subprocesses and reads results.
                try:
                    import subprocess as _subp
                    import pickle as _pickle

                    _max_workers = min(len(_navs_to_compute), 4)
                    print(f"[scale] dispatching {len(_navs_to_compute)} "
                          f"OOS workers (subprocess pool, "
                          f"{_max_workers} concurrent target)...")

                    _stage_dir = APP_DIR / ".cache" / "scale_workers"
                    _stage_dir.mkdir(parents=True, exist_ok=True)

                    # Build the env for workers — kernel mode on, bypass
                    # the freshness check (workers don't trade), strip
                    # SCALE_SENSITIVITY and ROADSHOW_DUAL_NAV so the
                    # worker doesn't try to re-run those code paths.
                    _worker_env = os.environ.copy()
                    _worker_env["OOS_KERNEL_MODE"] = "1"
                    _worker_env["HOLDINGS_FRESHNESS_BYPASS"] = "1"
                    _worker_env.pop("SCALE_SENSITIVITY", None)
                    _worker_env.pop("ROADSHOW_DUAL_NAV", None)

                    _worker_script = str(APP_DIR / "oos_worker.py")
                    _python_exe = sys.executable

                    # Stage all inputs first
                    _jobs = []   # (nav, in_path, out_path)
                    for _nav in _navs_to_compute:
                        _in_path = _stage_dir / f"in_{int(_nav)}.pkl"
                        _out_path = _stage_dir / f"out_{int(_nav)}.pkl"
                        _kw = dict(_common_kwargs)
                        _kw["starting_nav_aud"] = float(_nav)
                        with open(_in_path, "wb") as _fp:
                            _pickle.dump((oos_prices_aud_long, _kw), _fp,
                                         protocol=_pickle.HIGHEST_PROTOCOL)
                        _jobs.append((_nav, _in_path, _out_path))

                    # Spawn workers up to _max_workers concurrent. As
                    # one finishes, start the next. Tracks via a list
                    # of (Popen, nav, out_path).
                    _active: list = []
                    _pending = list(_jobs)
                    _done = []  # list of (nav, out_path, returncode)
                    while _pending or _active:
                        # Fill up to _max_workers
                        while _pending and len(_active) < _max_workers:
                            _nav, _in_path, _out_path = _pending.pop(0)
                            _proc = _subp.Popen(
                                [_python_exe, _worker_script,
                                 str(_in_path), str(_out_path)],
                                env=_worker_env,
                                stdout=_subp.PIPE,
                                stderr=_subp.PIPE,
                            )
                            _active.append((_proc, _nav, _out_path))
                        # Poll for finishes — small sleep prevents
                        # tight loop CPU burn.
                        import time as _time_pool
                        _time_pool.sleep(0.5)
                        _still_active = []
                        for _proc, _nav, _out_path in _active:
                            if _proc.poll() is None:
                                _still_active.append((_proc, _nav, _out_path))
                            else:
                                _rc = _proc.returncode
                                _stderr_out = _proc.stderr.read().decode(
                                    "utf-8", errors="replace")[:500]
                                _done.append((_nav, _out_path, _rc, _stderr_out))
                        _active = _still_active

                    for _nav, _out_path, _rc, _stderr in _done:
                        if _rc != 0:
                            print(f"[scale] @ ${_nav:,.0f}: worker exit {_rc}")
                            if _stderr.strip():
                                print(f"  stderr: {_stderr.strip()[:200]}")
                            continue
                        try:
                            with open(_out_path, "rb") as _fp:
                                _scale_out = _pickle.load(_fp)
                            oos_scale_results[float(_nav)] = {
                                "returns": _scale_out["blended_returns"],
                                "nav_aud": float(_nav),
                                "reused": False,
                            }
                            print(f"[scale] @ ${_nav:,.0f}: complete via worker, "
                                  f"{len(_scale_out['blended_returns'])} days")
                        except Exception as _e_load:
                            print(f"[scale] @ ${_nav:,.0f}: result load failed — {_e_load}")

                    for _p in _stage_dir.glob("*.pkl"):
                        try: _p.unlink()
                        except Exception: pass
                except Exception as _e_pool:
                    import traceback as _tb_pool
                    print(f"[scale] subprocess pool failed — "
                          f"falling back to sequential: {type(_e_pool).__name__}: {_e_pool}")
                    print(_tb_pool.format_exc())
                    _navs_to_compute_seq = list(_navs_to_compute)
                    _navs_to_compute = []
                    for _nav in _navs_to_compute_seq:
                        try:
                            print(f"[scale] running OOS at ${_nav:,.0f} (sequential)...")
                            _scale_out = run_oos_ensemble_walk_forward_cached(
                                oos_prices_aud_long,
                                starting_nav_aud=float(_nav),
                                **_common_kwargs,
                            )
                            oos_scale_results[float(_nav)] = {
                                "returns": _scale_out["blended_returns"],
                                "nav_aud": float(_nav),
                                "reused": False,
                            }
                        except Exception as _e_seq:
                            print(f"[scale] @ ${_nav:,.0f} (seq) failed: {_e_seq}")
            elif _navs_to_compute:
                # Sequential path: either parallel disabled, or only one
                # NAV to compute (no point in pool overhead).
                for _nav in _navs_to_compute:
                    try:
                        print(f"[scale] running OOS at ${_nav:,.0f} (sequential)...")
                        _scale_out = run_oos_ensemble_walk_forward_cached(
                            oos_prices_aud_long,
                            starting_nav_aud=float(_nav),
                            **_common_kwargs,
                        )
                        oos_scale_results[float(_nav)] = {
                            "returns": _scale_out["blended_returns"],
                            "nav_aud": float(_nav),
                            "reused": False,
                        }
                        print(f"[scale] @ ${_nav:,.0f}: complete, "
                              f"{len(_scale_out['blended_returns'])} days")
                    except Exception as _e_scale:
                        print(f"[scale] @ ${_nav:,.0f}: failed — {_e_scale}")
        globals()["oos_scale_results"] = oos_scale_results
        globals()["_scale_sensitivity_navs"] = list(SCALE_SENSITIVITY_NAVS)
        # Per-candidate returns + softmax history available for downstream
        # ensemble-mix displays (roadshow stacked area, trade-plan regime row).
        globals()["oos_per_candidate_returns"] = ensemble_out["per_candidate_returns"]
        globals()["oos_softmax_history"] = ensemble_out["softmax_history"]
        globals()["oos_per_candidate_weights"] = ensemble_out["per_candidate_weights"]
        globals()["oos_rebalance_costs"] = ensemble_out.get("rebalance_costs", pd.Series(dtype=float))
        globals()["oos_rebalance_taxes"] = ensemble_out.get("rebalance_taxes", pd.Series(dtype=float))
        # TLH events from the OOS engine — consumed by the Excel TLH_Log writer
        # and the PPT scorecard row. Empty list if TLH_ENABLED is False.
        globals()["oos_tlh_events"] = ensemble_out.get("tlh_events", []) or []
        # Report annualised drag from brokerage + CGT separately (informative).
        _cost_ser = ensemble_out.get("rebalance_costs", pd.Series(dtype=float))
        _tax_ser = ensemble_out.get("rebalance_taxes", pd.Series(dtype=float))
        _years = max(len(oos_returns_daily) / ANNUAL_TRADING_DAYS, 1e-6)
        if not _cost_ser.empty:
            _ann_cost_bps = float(_cost_ser.sum()) / _years * 10_000
            print(f"[oos] brokerage ({BROKER_CONFIG['name']}): "
                  f"avg {float(_cost_ser.mean())*10000:.1f} bps/rebal, "
                  f"~{_ann_cost_bps:.0f} bps/year")
        if not _tax_ser.empty and _tax_ser.sum() > 0:
            _ann_tax_bps = float(_tax_ser.sum()) / _years * 10_000
            _mtr_pct = (float(CGT_CONFIG['marginal_tax_rate']) +
                        float(CGT_CONFIG.get('medicare_levy', 0.0)) *
                        (1.0 if CGT_CONFIG.get('include_medicare', True) else 0.0)) * 100
            print(f"[oos] CGT (MTR {_mtr_pct:.0f}%, LT discount "
                  f"{int(float(CGT_CONFIG['lt_discount_rate'])*100)}%): "
                  f"avg {float(_tax_ser.mean())*10000:.1f} bps/rebal, "
                  f"~{_ann_tax_bps:.0f} bps/year")
            print(f"[oos] TOTAL drag (brokerage + CGT): "
                  f"~{_ann_cost_bps + _ann_tax_bps:.0f} bps/year (NET in metrics)")
        # TLH summary for the live OOS run: quantifies harvesting activity over
        # the full backtested window. Gross headline tax-saved figure assumes
        # 100% utilisation against future gains — see dev-validation output for
        # the realistic net impact on Sharpe.
        _tlh_events_live = ensemble_out.get("tlh_events", []) or []
        if _tlh_events_live:
            _tlh_loss = float(sum(e.get("loss_aud", 0.0) for e in _tlh_events_live))
            _eff_st = _effective_cgt_rate(short_term=True)
            _eff_lt = _effective_cgt_rate(short_term=False)
            _tlh_tax_est = _tlh_loss * (_eff_st + _eff_lt) / 2.0
            # bps vs the ACTUAL backtest starting NAV — a hardcoded $1M
            # denominator understated TLH's real scale ~4x at a $251k NAV,
            # which mis-sized the 2026-07-02 pairs-expansion decision.
            _tlh_nav_denom = float(globals().get("_oos_starting_nav_aud") or 1_000_000.0)
            _tlh_bps = _tlh_tax_est / _tlh_nav_denom / _years * 10_000
            print(f"[tlh] {len(_tlh_events_live)} events over window, "
                  f"${_tlh_loss:,.0f} loss realised, "
                  f"~${_tlh_tax_est:,.0f} gross tax saved est "
                  f"(~{_tlh_bps:.0f} bps/yr; gross — actual net depends on FY-end netting)")
        elif TLH_ENABLED:
            print(f"[tlh] enabled but 0 events triggered over window "
                  f"(threshold={TLH_MIN_LOSS_PCT*100:+.0f}%, ${TLH_MIN_LOSS_AUD:.0f} min)")
        _oos_t1 = time.perf_counter()
        if not oos_returns_daily.empty:
            print(
                f"[oos] ensemble walk-forward ({REBALANCE_FREQ}, "
                f"~{REBALANCES_PER_YEAR:.1f}/yr): "
                f"{oos_returns_daily.index.min().date()} → "
                f"{oos_returns_daily.index.max().date()} "
                f"({len(oos_returns_daily)} days, {len(oos_weights_history)} rebalances, "
                f"{_oos_t1 - _oos_t0:.1f}s)"
            )
            # Conditional rebalancing diagnostics
            _ns = int(ensemble_out.get("n_scheduled", 0))
            _net = int(ensemble_out.get("n_early_triggered", 0))
            _nex = int(ensemble_out.get("n_executed", 0))
            _nsk = int(ensemble_out.get("n_skipped", 0))
            print(
                f"[oos] rebal mix → scheduled={_ns}, early-triggered={_net}, "
                f"executed={_nex}, skipped={_nsk} "
                f"(skip<{SKIP_REBAL_DELTA*100:.0f}% Δw, DD-trigger>{EARLY_TRIGGER_DD_DEEPEN*100:.0f}%)"
            )
            # Latest regime mix — useful for the trade plan slide.
            if not ensemble_out["softmax_history"].empty:
                _latest = ensemble_out["softmax_history"].iloc[-1]
                _mix = ", ".join(f"{n.split(' ')[0]}={float(_latest.get(n,0))*100:.0f}%"
                                  for n in ENSEMBLE_SLOT_NAMES)
                print(f"[oos] latest regime mix → {_mix}")
        else:
            print("[oos] ensemble walk-forward returned empty series — check training-window coverage")
    except Exception as _e:
        print(f"[oos] walk-forward failed: {_e}")
        oos_returns_daily = pd.Series(dtype=float)
        oos_weights_history = pd.DataFrame()
globals()["oos_returns_daily"] = oos_returns_daily
globals()["oos_weights_history"] = oos_weights_history
globals()["oos_prices_aud_long"] = oos_prices_aud_long

# === Live ensemble recommendation (Task #10) ===
# Use TODAY's mu/Sigma (live 2y window) for candidate solving, but score the
# candidates with OOS HISTORICAL per-candidate returns (10y of evidence). This
# is the unified engine: the roadshow shows what the live engine would have
# done, and the live engine produces the same thing today.
w_ensemble_live = pd.Series(dtype=float)
ensemble_mix_live = pd.Series(dtype=float)
try:
    _mu_live = _apply_mu_shrinkage(pd.Series(mu_ann_geo).astype(float).dropna())
    _Sigma_live = Sigma_daily.copy()
    # Ship-consistency: the live solve must use the SAME covariance estimator
    # as the shipped backtest. Overlay Ledoit-Wolf on the good-coverage recent
    # window (matches the walk-forward), keep sample cov for ragged-edge
    # tickers. Fallback to sample cov on any failure (never break the plan).
    if bool(globals().get("COV_SHRINKAGE", False)) and "df_cov_wide" in globals():
        try:
            _win = df_cov_wide.tail(504)
            _win = _win.loc[:, _win.notna().mean() >= 0.8].dropna(how="any")
            if _win.shape[0] >= 60 and _win.shape[1] >= 3:
                _lw_live, _lwd = _ledoit_wolf_cc(_win)
                _Sigma_live.loc[_lw_live.index, _lw_live.columns] = _lw_live.values
                print(f"[cov-shrink] live Σ shrunk (Ledoit-Wolf δ={_lwd:.3f}, "
                      f"{_win.shape[1]} tickers)")
        except Exception as _e_lwlive:
            print(f"[cov-shrink] live LW failed ({_e_lwlive}); using sample cov")
            _Sigma_live = Sigma_daily.copy()
    _spy_mu_live = float(_mu_live["SPY"]) if "SPY" in _mu_live.index else None
    _cand_live = solve_candidate_portfolios(_mu_live, _Sigma_live, _spy_mu_live)
    # Publish the live slot weights so downstream consumers (EF chart, TLH,
    # attribution) can read per-slot composition without re-solving.
    globals()["LIVE_SLOT_WEIGHTS"] = {n: w.copy() for n, w in _cand_live.items()}
    _oos_cand_rets = globals().get("oos_per_candidate_returns", pd.DataFrame())
    # For IR scoring, pass SPY daily returns aligned to OOS candidate returns.
    _spy_bench_for_score = None
    if "SPY" in oos_prices_aud_long.columns:
        _spy_bench_for_score = oos_prices_aud_long["SPY"].pct_change().dropna()
    ensemble_mix_live = softmax_ensemble_weights(
        _oos_cand_rets, lookback_days=252, lambda_temp=3.0, halflife_days=60,
        benchmark_returns=_spy_bench_for_score,
    )
    # Blend with forward-looking SPY regime signal (same as OOS engine).
    if "SPY" in oos_prices_aud_long.columns and not oos_prices_aud_long.empty:
        _fwd_live = compute_forward_regime_signal(
            benchmark_prices=oos_prices_aud_long["SPY"],
            as_of_date=oos_prices_aud_long.index.max(),
        )
        ensemble_mix_live = blend_ensemble_signals(
            backward_weights=ensemble_mix_live,
            forward_weights=_fwd_live,
            backward_alpha=0.5,
        )
    # PRODUCTION CONFIG override: force ensemble_mix_live to the shipped
    # allocation (e.g. 100% Stretch) instead of the softmax+forward blend.
    # See PRODUCTION_* constants for the empirical rationale.
    if PRODUCTION_SLOT_OVERRIDE:
        _override = pd.Series(PRODUCTION_SLOT_OVERRIDE).reindex(
            ENSEMBLE_SLOT_NAMES, fill_value=0.0).astype(float)
        _ovs = float(_override.sum())
        if _ovs > 0:
            ensemble_mix_live = _override / _ovs
            print(f"[ensemble] PRODUCTION override applied: "
                  f"{dict(PRODUCTION_SLOT_OVERRIDE)}")
    # Blend
    _ticker_idx = sorted(set().union(*[set(c.index) for c in _cand_live.values() if not c.empty]))
    if _ticker_idx and not ensemble_mix_live.empty:
        w_ensemble_live = pd.Series(0.0, index=_ticker_idx)
        for n in ENSEMBLE_SLOT_NAMES:
            cand_w = _cand_live.get(n, pd.Series(dtype=float))
            if cand_w.empty or ensemble_mix_live.get(n, 0.0) <= 0:
                continue
            w_ensemble_live = w_ensemble_live.add(
                cand_w.reindex(_ticker_idx).fillna(0.0) * float(ensemble_mix_live[n]),
                fill_value=0.0,
            )
        w_ensemble_live = w_ensemble_live[w_ensemble_live > 1e-6]
        if not w_ensemble_live.empty and w_ensemble_live.sum() > 0:
            w_ensemble_live = w_ensemble_live / w_ensemble_live.sum()
        # Volatility targeting (ship-consistency with the backtest): cap the
        # live ex-ante vol at VOL_TARGET_ANNUAL by scaling toward cash. Uses
        # the same LW live Σ; long-only (de-risk only). Mirrors the walk-forward.
        _vt_live = float(globals().get("VOL_TARGET_ANNUAL", 0.0) or 0.0)
        if _vt_live > 0 and not w_ensemble_live.empty:
            _cvl = [c for c in w_ensemble_live.index if c in _Sigma_live.index]
            if len(_cvl) >= 2:
                _wvl = w_ensemble_live.reindex(_cvl).fillna(0.0).values
                _Sl = _Sigma_live.reindex(index=_cvl, columns=_cvl).fillna(0.0).values
                _sig_l = float(np.sqrt(max(float(_wvl @ _Sl @ _wvl), 0.0) * 252.0))
                if _sig_l > _vt_live > 0:
                    _sc = _vt_live / _sig_l
                    w_ensemble_live = w_ensemble_live * _sc
                    print(f"[vol-target] live ex-ante vol {_sig_l*100:.1f}% > "
                          f"{_vt_live*100:.0f}% target → scaled to {_sc:.2f} (rest cash)")
        # Market-timed inverse-ETF crisis hedge (ship-consistency with the
        # backtest): if SPY is below its 200d SMA, carve _ch_live into BEAR.AX
        # and scale the long book to (1-w). Mirrors the walk-forward overlay.
        _ch_live = float(globals().get("CRISIS_HEDGE_WEIGHT", 0.0) or 0.0)
        _ch_tkr = str(globals().get("CRISIS_HEDGE_TICKER", "BEAR.AX"))
        _ch_ma = int(globals().get("CRISIS_HEDGE_MA_DAYS", 200))
        _ch_band = float(globals().get("CRISIS_HEDGE_BAND_SD", 0.0) or 0.0)
        if (_ch_live > 0 and not w_ensemble_live.empty
                and "SPY" in oos_prices_aud_long.columns):
            _spy_l = oos_prices_aud_long["SPY"].sort_index()
            _mp_l = max(2, _ch_ma // 2)
            _ma_l = _spy_l.rolling(window=_ch_ma, min_periods=_mp_l).mean()
            _sd_l = _spy_l.rolling(window=_ch_ma, min_periods=_mp_l).std()
            _lb_l = _ma_l - _ch_band * _sd_l
            _spy_last = _spy_l.iloc[-1] if not _spy_l.empty else np.nan
            _lb_last = _lb_l.iloc[-1] if not _lb_l.empty else np.nan
            _band_lbl = (f"{_ch_ma}dMA" if _ch_band <= 0 else f"{_ch_ma}dMA-{_ch_band:g}σ")
            if pd.notna(_spy_last) and pd.notna(_lb_last):
                if float(_spy_last) < float(_lb_last) and _ch_tkr in oos_prices_aud_long.columns:
                    w_ensemble_live = w_ensemble_live * (1.0 - _ch_live)
                    w_ensemble_live.loc[_ch_tkr] = (
                        w_ensemble_live.get(_ch_tkr, 0.0) + _ch_live)
                    w_ensemble_live = w_ensemble_live[w_ensemble_live > 1e-6]
                    print(f"[crisis-hedge] live SPY {float(_spy_last):.2f} < "
                          f"{_band_lbl} {float(_lb_last):.2f} → {_ch_live*100:.0f}% "
                          f"into {_ch_tkr} (long book scaled to {(1-_ch_live)*100:.0f}%)")
                else:
                    print(f"[crisis-hedge] live armed, SPY {float(_spy_last):.2f} "
                          f">= {_band_lbl} {float(_lb_last):.2f} → not triggered")
    # PRODUCTION CONFIG crash-hedge: check trigger NOW; if active, replace
    # w_ensemble_live with hedge basket. Mirrors the engine's per-rebalance
    # hedge check so the live trade plan reflects current crash-hedge status.
    if PRODUCTION_CRASH_HEDGE and "SPY" in oos_prices_aud_long.columns:
        try:
            _live_hedge_state = {"active": False}
            _is_crashing = _check_crash_trigger(
                spy_history=oos_prices_aud_long["SPY"],
                as_of=oos_prices_aud_long.index.max(),
                state=_live_hedge_state,
            )
            if _is_crashing:
                _avail = oos_prices_aud_long.columns
                w_with_basket = w_ensemble_live.reindex(
                    w_ensemble_live.index.union(CRASH_HEDGE_BASKET.keys())
                ).fillna(0.0)
                w_ensemble_live = _apply_crash_hedge(
                    weights=w_with_basket,
                    basket=CRASH_HEDGE_BASKET,
                    available_tickers=_avail,
                )
                w_ensemble_live = w_ensemble_live[w_ensemble_live > 1e-6]
                if not w_ensemble_live.empty and w_ensemble_live.sum() > 0:
                    w_ensemble_live = w_ensemble_live / w_ensemble_live.sum()
                print(f"[ensemble] PRODUCTION crash hedge ACTIVE — "
                      f"portfolio overridden to hedge basket "
                      f"(SPY DD {_live_hedge_state.get('last_dd', 0.0)*100:+.1f}%)")
            else:
                print(f"[ensemble] PRODUCTION crash hedge armed, not triggered "
                      f"(SPY DD {_live_hedge_state.get('last_dd', 0.0)*100:+.1f}%)")
        except Exception as _eh:
            print(f"[ensemble] PRODUCTION crash hedge check failed: {_eh}")
    _mix_str = ", ".join(
        f"{n.split(' ')[0]}={float(ensemble_mix_live.get(n,0))*100:.0f}%"
        for n in ENSEMBLE_SLOT_NAMES
    )
    print(f"[ensemble] Live regime mix → {_mix_str}")
    print(f"[ensemble] Live recommendation: {len(w_ensemble_live)} positions, "
          f"top: {w_ensemble_live.nlargest(5).to_dict() if not w_ensemble_live.empty else '{}'}")
except Exception as _e:
    print(f"[ensemble] Live ensemble blend failed: {_e}")
globals()["W_ENSEMBLE_SER"] = w_ensemble_live
globals()["ensemble_mix_live"] = ensemble_mix_live

# Metrics table (Phase 2): horizons × series, written to Excel + Slide 2 later.
oos_metrics_table = pd.DataFrame()
if not oos_returns_daily.empty and not oos_prices_aud_long.empty:
    try:
        _spy_aud = oos_prices_aud_long.get("SPY") if "SPY" in oos_prices_aud_long.columns else None
        _aord = oos_prices_aud_long.get("^AORD") if "^AORD" in oos_prices_aud_long.columns else None
        _spy_ret = _spy_aud.pct_change().dropna() if _spy_aud is not None else pd.Series(dtype=float)
        _aord_ret = _aord.pct_change().dropna() if _aord is not None else pd.Series(dtype=float)
        _ff5 = globals().get("ff5_raw", None)
        oos_metrics_table = compute_oos_metrics(
            strat_returns=oos_returns_daily,
            spy_returns=_spy_ret,
            aord_returns=_aord_ret,
            ff5_factors=_ff5,
            weights_history=oos_weights_history,
            horizons_years=(3, 5, 10),
        )
        if not oos_metrics_table.empty:
            print(f"[oos] metrics computed: {oos_metrics_table.shape[0]} metrics × {oos_metrics_table.shape[1]} cols")
    except Exception as _e:
        print(f"[oos] metrics computation failed: {_e}")
globals()["oos_metrics_table"] = oos_metrics_table

# Optional roadshow-NAV metrics table — only the "Strategy" columns are
# kept (benchmarks are identical) and renamed so the Roadshow slide can
# discriminate. Quiet no-op when dual mode is off.
oos_metrics_table_roadshow = pd.DataFrame()
_oos_rets_rs = globals().get("oos_returns_daily_roadshow", pd.Series(dtype=float))
if isinstance(_oos_rets_rs, pd.Series) and not _oos_rets_rs.empty:
    try:
        _spy_aud_rs = oos_prices_aud_long.get("SPY") if "SPY" in oos_prices_aud_long.columns else None
        _aord_rs = oos_prices_aud_long.get("^AORD") if "^AORD" in oos_prices_aud_long.columns else None
        _spy_ret_rs = _spy_aud_rs.pct_change().dropna() if _spy_aud_rs is not None else pd.Series(dtype=float)
        _aord_ret_rs = _aord_rs.pct_change().dropna() if _aord_rs is not None else pd.Series(dtype=float)
        _ff5_rs = globals().get("ff5_raw", None)
        oos_metrics_table_roadshow = compute_oos_metrics(
            strat_returns=_oos_rets_rs,
            spy_returns=_spy_ret_rs,
            aord_returns=_aord_ret_rs,
            ff5_factors=_ff5_rs,
            weights_history=globals().get("oos_weights_history_roadshow", pd.DataFrame()),
            horizons_years=(3, 5, 10),
        )
        if not oos_metrics_table_roadshow.empty:
            print(f"[oos-roadshow] metrics computed: "
                  f"{oos_metrics_table_roadshow.shape[1]} cols")
    except Exception as _e_rs_mtx:
        print(f"[oos-roadshow] metrics computation failed: {_e_rs_mtx}")
globals()["oos_metrics_table_roadshow"] = oos_metrics_table_roadshow

# Per-NAV metrics for the scale-sensitivity slide. Same compute_oos_metrics
# call as Roadshow, applied to each scale backtest.
oos_scale_metrics: dict[float, pd.DataFrame] = {}
_scale_results_local = globals().get("oos_scale_results", {})
if _scale_results_local:
    try:
        _spy_aud_sc = oos_prices_aud_long.get("SPY") if "SPY" in oos_prices_aud_long.columns else None
        _aord_sc = oos_prices_aud_long.get("^AORD") if "^AORD" in oos_prices_aud_long.columns else None
        _spy_ret_sc = _spy_aud_sc.pct_change().dropna() if _spy_aud_sc is not None else pd.Series(dtype=float)
        _aord_ret_sc = _aord_sc.pct_change().dropna() if _aord_sc is not None else pd.Series(dtype=float)
        _ff5_sc = globals().get("ff5_raw", None)
        for _nav, _payload in _scale_results_local.items():
            _rets_for_nav = _payload.get("returns", pd.Series(dtype=float))
            if not isinstance(_rets_for_nav, pd.Series) or _rets_for_nav.empty:
                continue
            try:
                _mtx = compute_oos_metrics(
                    strat_returns=_rets_for_nav,
                    spy_returns=_spy_ret_sc,
                    aord_returns=_aord_ret_sc,
                    ff5_factors=_ff5_sc,
                    weights_history=pd.DataFrame(),
                    horizons_years=(3, 5, 10),
                )
                if not _mtx.empty:
                    oos_scale_metrics[float(_nav)] = _mtx
            except Exception as _e_sc:
                print(f"[scale] metrics @ ${_nav:,.0f}: {_e_sc}")
        if oos_scale_metrics:
            print(f"[scale] metrics computed for {len(oos_scale_metrics)} NAVs")
    except Exception as _e_sc_all:
        print(f"[scale] metrics computation failed: {_e_sc_all}")
globals()["oos_scale_metrics"] = oos_scale_metrics

# Persist a metrics snapshot to metrics_history.jsonl + warn on regressions.
# Compares 10Y Sharpe / MaxDD / α-vs-SPY to the prior run. Non-fatal —
# diagnostic only. See _append_metrics_snapshot() for thresholds.
try:
    _tlh_events_for_log = globals().get("oos_tlh_events", []) or []
    _tlh_loss_for_log = float(sum(e.get("loss_aud", 0.0) for e in _tlh_events_for_log))
    _ens_out_local = globals().get("ensemble_out", {})
    _append_metrics_snapshot(
        metrics_table=oos_metrics_table,
        ensemble_mix_live=globals().get("ensemble_mix_live", pd.Series(dtype=float)),
        w_ensemble_live=globals().get("W_ENSEMBLE_SER", pd.Series(dtype=float)),
        tlh_events_n=len(_tlh_events_for_log),
        tlh_loss_aud=_tlh_loss_for_log,
        n_executed=int(_ens_out_local.get("n_executed", 0)) if isinstance(_ens_out_local, dict) else 0,
        n_skipped=int(_ens_out_local.get("n_skipped", 0)) if isinstance(_ens_out_local, dict) else 0,
        # When SCALE_SENSITIVITY=1 (default in daily_auto.ps1), the per-NAV
        # metrics are populated and get logged so metrics_history.jsonl
        # accumulates a continuous multi-scale evidence track over time.
        scale_metrics=globals().get("oos_scale_metrics", {}) or None,
    )
except Exception as _e_metrics:
    print(f"[metrics] snapshot call failed: {_e_metrics}")

# Tables used later
n_opt = len(securities_opt)
cov_plus = Sigma_opt.copy()
cov_plus.loc[:, 'w'] = 0.0
cov_plus.loc['w', :] = 0.0
cov_plus.loc['w', 'w'] = 0.0
exp_ret_df = mu_vec_opt.rename(exp_ret_label).to_frame()

# FX map used by Holdings + trade plan
usd_aud    = get_usd_aud_fx()
fx_map_all = fx_to_aud_for_tickers(prices.columns, usd_aud)

# ---- 10D) Reopen Excel and WRITE everything, then close ----
if USE_XLWINGS:
    try:
        with xw.App(visible=False, add_book=False) as app:
            filename = os.path.abspath(filename)
            wb = app.books.open(filename, update_links=False, read_only=False)      
            
            if bool(wb.api.ReadOnly):
                # If Excel forces read-only (usually file is already open/locked), write to a new file instead of CSV fallback
                base, ext = os.path.splitext(filename)
                alt = base + "_AUTO" + ext
                shutil.copy2(filename, alt)
                print(f"[warn] Workbook opened read-only. Will write to: {alt}")
                wb.close()
                wb = app.books.open(alt, update_links=False, read_only=False)
            
            wb.activate()
            app.display_alerts = False
            app.screen_updating = False
            try: app.api.EnableEvents = False
            except Exception: pass
            time.sleep(0.2)

            # Pick the max-Sharpe portfolio column once for reuse
            sh = pd.to_numeric(stats_df['Sharpe'], errors='coerce').fillna(-1)
            best_idx = int(sh.values.argmax()) if len(sh) else 0
            w_star = W.iloc[:, best_idx].reindex(W.index).fillna(0.0)

            w_star_no_tilts = pd.to_numeric(w_star, errors="coerce").reindex(Sigma_opt.index).fillna(0.0)
            if float(w_star_no_tilts.sum()) != 0:
                w_star_no_tilts = w_star_no_tilts / float(w_star_no_tilts.sum())
            
            R_star = float(stats_df.loc[best_idx, "Achieved Return"])

            if "w_tilt" in locals():
                print("[debug] len(Sigma_opt.index) =", len(Sigma_opt.index), "| type(w_tilt) =", type(w_tilt), "| len(w_tilt) =", (len(w_tilt) if w_tilt is not None else None))
            else:
                print("[debug] len(Sigma_opt.index) =", len(Sigma_opt.index), "| w_tilt not in locals()")

            # --- Target WITHOUT tilts ---
            use_mask_no_tilts = {f: False for f in tilt_df.index}
            w_nt_raw, _, _ = solve_frontier_point_cvxpy_with_tilts(
                mu_vec_opt,
                Sigma_opt,
                R_star,
                B,
                tilt_df["Target"],
                tilt_df["Band"],
                use_mask_no_tilts
            )
            
            # Convert solver output back to ticker-indexed Series
            w_nt_raw = np.asarray(w_nt_raw, dtype=float).reshape(-1)
            w_nt_raw = w_nt_raw[:len(Sigma_opt.index)]  # safety if solver returns an extra element
            w_star_no_tilts = pd.Series(w_nt_raw, index=Sigma_opt.index).fillna(0.0)
            
            # Normalise
            s = float(w_star_no_tilts.sum())
            if s != 0.0:
                w_star_no_tilts = w_star_no_tilts / s
            
            # --- Target WITH tilts ---
            use_mask_with_tilts = tilt_df["Use?"].astype(bool).to_dict()
            w_wt_raw, _, _ = solve_frontier_point_cvxpy_with_tilts(
                mu_vec_opt,
                Sigma_opt,
                R_star,
                B,
                tilt_df["Target"],
                tilt_df["Band"],
                use_mask_with_tilts
            )
            
            # Convert solver output back to ticker-indexed Series
            w_wt_raw = np.asarray(w_wt_raw, dtype=float).reshape(-1)
            w_wt_raw = w_wt_raw[:len(Sigma_opt.index)]  # safety if solver returns an extra element
            w_star_with_tilts = pd.Series(w_wt_raw, index=Sigma_opt.index).fillna(0.0)
            
            # Normalise (use the WITH-TILTS sum, not the no-tilts sum)
            s_wt = float(w_star_with_tilts.sum())
            if s_wt != 0.0:
                w_star_with_tilts = w_star_with_tilts / s_wt

            # Always publish the series so the PPT performance chart can plot it
            globals()["W_WITH_TILTS_SER"] = w_star_with_tilts.copy()


            # 1) Cov sheet
            cov = get_or_clear_sheet(wb, 'Cov')
            cov.range('A1').options(pd.DataFrame, index=True, header=True).value = Sigma_opt

            # 2) Input sheet
            inp = get_or_clear_sheet(wb, 'Input')
            inp.range('A1').options(pd.DataFrame, index=False, header=True).value = df_melt

            # 3) OPT sheet
            opt = get_or_clear_sheet(wb, 'OPT')

            # Header
            opt.range('A1').value = 'Optimal Portfolio Theory (long-only where possible)'
            opt.range('A2').value = f"Generated: {datetime.now():%Y-%m-%d %H:%M:%S}"
            opt.range('A3').value = 'Expected returns use geometric (log-based) annualisation.'
            opt.range('A4').value = 'Variance is daily; annual vol = sqrt(252) * stdev.'
            try:
                opt.range('A1').api.Font.Bold = True; opt.range('A1').api.Font.Size = 14
            except Exception:
                pass

            # Expected returns
            opt.range('A6').value = exp_ret_label
            opt.range('A7').options(pd.DataFrame, index=True, header=True).value = exp_ret_df
            n_rows = exp_ret_df.shape[0] + 1
            set_number_formats(opt, {f"B8:B{7+n_rows}": "0.00%"})

            # Covariance (+ weight row/col)
            start_cov_row = 9 + n_rows
            opt.range(f"A{start_cov_row}").value = 'Covariance Matrix (daily, model) with weight row/column'
            opt.range(f"A{start_cov_row+1}").options(pd.DataFrame, index=True, header=True).value = cov_plus.fillna(0.0)

            # Weights grid
            start_w_row = start_cov_row + cov_plus.shape[0] + 4
            opt.range(f"A{start_w_row}").value = 'Optimised Weights by Target Return'
            opt.range(f"A{start_w_row+1}").options(pd.DataFrame, index=True, header=True).value = W
            # --- Dynamic format for W (optimised weights table) ---
            w_first = start_w_row + 1               # header row
            w_data_first = w_first + 1              # first data row
            w_rows = W.shape[0]
            w_cols = W.shape[1]
            
            # Percent format for all weight cells
            rng_w = opt.range(
                f"B{w_data_first}:{chr(ord('A')+w_cols)}{w_data_first + w_rows - 1}"
            )
            try:
                rng_w.api.NumberFormat = "0.00%"
            except:
                pass

            # Portfolio Statistics
            start_s_row = start_w_row + W.shape[0] + 4
            opt.range(f"A{start_s_row}").value = 'Portfolio Statistics'
            opt.range(f"A{start_s_row+1}").options(pd.DataFrame, index=False, header=True).value = stats_df
            # ==========================================================
            #  Efficient Frontier Chart
            # ==========================================================
            co_old = opt.api.ChartObjects()
            to_delete = []
            for i in range(1, co_old.Count + 1):
                o = co_old.Item(i)
                try:
                    title_text = o.Chart.ChartTitle.Text
                    if "Efficient Frontier" in str(title_text):
                        to_delete.append(o)
                except:
                    pass
            
            for o in to_delete:
                o.Delete()

            co = opt.api.ChartObjects()
            chart_obj = None
            
            # Find existing chart by the *dynamic* title
            for i in range(1, co.Count + 1):
                o = co.Item(i)
                try:
                    if o.Chart.HasTitle and "Efficient Frontier" in str(o.Chart.ChartTitle.Text):
                        chart_obj = o
                        break
                except:
                    pass
            
            # If not found, create it
            if chart_obj is None:
                left = opt.range("I1").api.Left       # right of the stats table
                top  = opt.range(f"A{start_s_row+1}").api.Top
                width = 480
                height = 245
            
                chart_obj = co.Add(left, top, width, height)
                ch = chart_obj.Chart
                ch.ChartType = -4169       # XY scatter
                ch.HasTitle = True
                ch.ChartTitle.Text = chart_title
            else:
                ch = chart_obj.Chart
            
            # Reposition every run
            chart_obj.Left   = opt.range("I1").api.Left
            chart_obj.Top    = opt.range(f"A{start_s_row+1}").api.Top
            chart_obj.Width  = 480
            chart_obj.Height = 245

            # ----------------------------------------------------------
            # Format Portfolio Statistics
            # ----------------------------------------------------------
            
            stat_rows = stats_df.shape[0]
            if stat_rows > 0:
                header_row = start_s_row + 1
                data_first = header_row + 1
            
                for col_name, fmt in {
                    "Achieved Return": "0.00%",
                    "Volatility (ann.)": "0.00%",
                    "Sharpe": "0.00"
                }.items():
                    if col_name in stats_df.columns:
                        col_idx = list(stats_df.columns).index(col_name)
                        col_letter = chr(ord("A") + col_idx)
                        try:
                            opt.range(
                                f"{col_letter}{data_first}:{col_letter}{data_first + stat_rows - 1}"
                            ).api.NumberFormat = fmt
                        except:
                            pass

            # ================= Efficient Frontier chart updater =================
            def _col_letter(idx0: int) -> str:
                n = idx0 + 1  # A=1
                letters = ""
                while n:
                    n, rem = divmod(n - 1, 26)
                    letters = chr(65 + rem) + letters
                return letters
            
            def _get_chart_by_title(opt_sheet, title_text: str):
                # --- Hard reset: delete all previous Efficient Frontier charts ---
                co = opt_sheet.api.ChartObjects()
                delete_list = []
                for i in range(1, co.Count + 1):
                    o = co.Item(i)
                    try:
                        t = o.Chart.ChartTitle.Text
                        if "Efficient Frontier" in str(t):
                            delete_list.append(o)
                    except Exception:
                        pass
                
                for o in delete_list:
                    o.Delete()

                """Return the COM Chart object whose Title text equals title_text (case/space-insensitive)."""
                def _norm(s): return " ".join(str(s).split()).casefold()
                co = opt_sheet.api.ChartObjects()
                want = _norm(title_text)
                for i in range(1, co.Count + 1):
                    o = co.Item(i)
                    try:
                        ch = o.Chart
                        if ch.HasTitle and _norm(ch.ChartTitle.Text) == want:
                            return ch
                    except Exception:
                        pass
                return None
                            
           
            # -------- Per-slot frontier points (replaces Current/Previous/
            # Optimised/With Tilts markers as of 2026-06-22 user request). The
            # 5 ENSEMBLE_SLOTS are themselves frontier-optimal by construction
            # (each solves solve_frontier_point_cvxpy at a specific target),
            # so projecting them onto Sigma_opt / mu_vec_opt shows where each
            # regime preference lands on the curve. Ensemble (separate marker
            # below) is a softmax blend of these slots — by concavity it sits
            # strictly inside the frontier unless one slot has 100% weight.
            slot_points: dict[str, tuple[float, float]] = {}
            ensemble_point = None

            # Project everything in the LIVE basis (mu_ann_geo / Sigma_daily)
            # — the same basis the slots and ensemble were solved in. This is
            # the basis-consistent view: slot markers land ON the live frontier
            # by construction, and Ensemble's distance below the curve is the
            # pure concavity / regime tax. (Prior version used Sigma_opt which
            # produced an idealised FF5-projected curve the engine can't
            # actually achieve — visually misleading. Migrated 2026-06-22.)
            try:
                _idx_live = list(Sigma_daily.index)
                mu_use = pd.Series(mu_ann_geo).reindex(_idx_live).fillna(0.0).values
                S_use = Sigma_daily.values

                _slot_w_map = globals().get("LIVE_SLOT_WEIGHTS", {}) or {}
                for _slot_name, _slot_w in _slot_w_map.items():
                    try:
                        if not isinstance(_slot_w, pd.Series) or _slot_w.empty:
                            continue
                        _ws = pd.Series(_slot_w, dtype=float).reindex(
                            _idx_live).fillna(0.0)
                        if float(_ws.sum()) <= 0:
                            continue
                        _ws = _ws / float(_ws.sum())
                        _wv = _ws.values
                        _sret = float(mu_use @ _wv)
                        _svol = float(np.sqrt(_wv @ S_use @ _wv) * np.sqrt(252.0))
                        slot_points[_slot_name] = (_svol, _sret)
                    except Exception as _e_slot:
                        print(f"[chart] Slot point error for {_slot_name}: {_e_slot}")

                # --- Ensemble (live regime-blend portfolio) ---
                # Convex combination of the 5 slots — by concavity of the
                # frontier this sits strictly inside the curve unless one
                # slot has 100% softmax weight.
                try:
                    _w_ens = globals().get("W_ENSEMBLE_SER", None)
                    if isinstance(_w_ens, pd.Series) and not _w_ens.empty:
                        w_ens_ser = pd.Series(_w_ens, dtype=float).reindex(_idx_live).fillna(0.0)
                        if float(w_ens_ser.sum()) > 0:
                            w_ens_ser = w_ens_ser / float(w_ens_ser.sum())
                            wv_ens = w_ens_ser.values
                            ens_ret = float(mu_use @ wv_ens)
                            ens_vol = float(np.sqrt(wv_ens @ S_use @ wv_ens) * np.sqrt(252.0))
                            ensemble_point = (ens_vol, ens_ret)
                except Exception as e:
                    print(f"[chart] Ensemble point error: {e}")

            except Exception as e:
                print(f"[chart] Point compute error: {e}")
                slot_points = {}
                ensemble_point = None


            # --- Build Efficient Frontier PNG for PowerPoint (optional) ---
            charts = globals().get("charts", {}) or {}

            try:
                # Build the LIVE-basis frontier (mu_ann_geo / Sigma_daily) so
                # the curve is the actual achievable frontier the engine
                # optimises on. Slots solved by solve_candidate_portfolios are
                # frontier points in this exact basis → they will land ON the
                # curve. Falls back to the old stats_df (FF5-projected) curve
                # if the live sweep fails for any reason.
                try:
                    _W_live, _stats_live, _, _ = _build_frontier(
                        pd.Series(mu_ann_geo), Sigma_daily, n_points=24,
                    )
                    _x = pd.to_numeric(_stats_live["Volatility (ann.)"], errors="coerce")
                    _y = pd.to_numeric(_stats_live["Achieved Return"], errors="coerce")
                    print(f"[ef-chart] live frontier built: {len(_x)} points, "
                          f"vol [{_x.min():.3f}, {_x.max():.3f}], "
                          f"ret [{_y.min():.3f}, {_y.max():.3f}]")
                except Exception as _e_live_fr:
                    print(f"[ef-chart] live frontier build failed, "
                          f"falling back to stats_df: {_e_live_fr}")
                    _x = pd.to_numeric(stats_df["Volatility (ann.)"], errors="coerce")
                    _y = pd.to_numeric(stats_df["Achieved Return"], errors="coerce")

                fig, ax = plt.subplots(figsize=(7.5, 4.8))
                ax.plot(_x, _y, linewidth=2.0, label="Frontier (live)")
                ax.set_title(chart_title)
                ax.set_xlabel("Volatility (ann.)")
                ax.set_ylabel("Return (ann.)")
                ax.xaxis.set_major_formatter(mtick.PercentFormatter(1.0))
                ax.yaxis.set_major_formatter(mtick.PercentFormatter(1.0))
                
                # Per-slot markers (Modest / Aggressive / Bold / Maximum /
                # Stretch). Each slot is a frontier point by construction, so
                # if Σ_opt ≈ the live Σ used at slot-solve time these markers
                # should land on (or very near) the curve. Visible drift = the
                # basis difference between mu_ann_geo/Sigma_daily (slot solve)
                # and mu_vec_opt/Sigma_opt (chart axis).
                _slot_style = {
                    "Modest (SPY+0%)":      ("o", "#1f77b4", 70),
                    "Aggressive (SPY+5%)":  ("s", "#2ca02c", 70),
                    "Bold (SPY+10%)":       ("D", "#ff7f0e", 70),
                    "Maximum (SPY+15%)":    ("^", "#9467bd", 90),
                    "Stretch (SPY+25%)":    ("v", "#d62728", 90),
                }
                for _sn, _spt in slot_points.items():
                    _mk, _co, _sz = _slot_style.get(_sn, ("o", "#555555", 70))
                    _short = _sn.split(" (")[0]
                    ax.scatter(
                        [float(_spt[0])], [float(_spt[1])],
                        s=_sz, marker=_mk, color=_co,
                        edgecolors="black", linewidths=0.5,
                        label=_short, zorder=5,
                    )
                    ax.annotate(
                        _short,
                        (float(_spt[0]), float(_spt[1])),
                        xytext=(6, 6), textcoords="offset points",
                        fontsize=8, color=_co,
                    )

                if ensemble_point:
                    # Star out the regime-blend portfolio so it stands out
                    # against the 5 slot points. By concavity of the frontier
                    # this convex combination sits strictly inside the curve
                    # unless one slot has 100% softmax weight.
                    ax.scatter(
                        [float(ensemble_point[0])], [float(ensemble_point[1])],
                        s=170, marker="*", color="#c00000",
                        edgecolors="black", linewidths=0.8,
                        label="Ensemble (blend)", zorder=7,
                    )
                    ax.annotate(
                        "Ensemble",
                        (float(ensemble_point[0]), float(ensemble_point[1])),
                        xytext=(8, -12),
                        textcoords="offset points",
                        fontsize=9, fontweight="bold", color="#c00000",
                    )

                ax.legend(loc="best", fontsize=8)
                _eff_buf = io.BytesIO()
                fig.savefig(_eff_buf, format="png", bbox_inches="tight")
                plt.close(fig)
                _eff_buf.seek(0)

                charts["efficient_frontier_image"] = _eff_buf
                _fp = {sn: pt for sn, pt in slot_points.items()}
                _fp["Ensemble"] = ensemble_point
                charts["frontier_points"] = _fp
                globals()["charts"] = charts
            except Exception as _e_eff_png:
                print(f"[pptx] Efficient frontier PNG build skipped: {_e_eff_png}")
            
            # ---- Store achieved tilts for PPT Slide 5 (With Tilts + Without Tilts) ----
            try:
                if (B is None) or B.empty:
                    raise ValueError("B is None or empty")
            
                factor_order = ["Mkt-RF", "SMB", "HML", "RMW", "CMA", "MOM"]
            
                def _norm_w(w, idx):
                    s = pd.Series(np.asarray(w, dtype=float).reshape(-1), index=idx).fillna(0.0)
                    tot = float(s.sum())
                    return (s / tot) if tot != 0 else s
            
                # --- With Tilts achieved betas (use optimiser w_star_with_tilts) ---
                w_with = _norm_w(w_star_with_tilts, Sigma_opt.index).reindex(B.index).fillna(0.0)
                with_beta = (B.T @ w_with)
            
                # --- Without Tilts achieved betas (use optimiser w_star) ---
                w_without = _norm_w(w_star, Sigma_opt.index).reindex(B.index).fillna(0.0)
                without_beta = (B.T @ w_without)
            
                out = pd.DataFrame(index=[f for f in factor_order if f in with_beta.index])
                out["With Tilts"] = with_beta.reindex(out.index).astype(float)
                out["Without Tilts"] = without_beta.reindex(out.index).astype(float)
            
                # Targets (from Tilts sheet)
                if isinstance(tilt_df, pd.DataFrame) and (not tilt_df.empty) and ("Target" in tilt_df.columns):
                    tgt = tilt_df.reindex(out.index)
                    out["Target"] = pd.to_numeric(tgt["Target"], errors="coerce")
                    print("[debug] tilt targets used:", out["Target"].to_dict())

                    # Optional: filter to Use? if present
                    if "Use?" in tgt.columns:
                        use_mask = tgt["Use?"].astype(bool)
                        out = out.loc[use_mask.reindex(out.index).fillna(False)]
            
                charts["tilts_comparison_rows"] = (
                    out.reset_index()
                    .rename(columns={"index": "Factor"})
                    .to_dict("records")
                )
                
            except Exception as _e_ppt_front:
                print(f"[pptx] Tilt table storage skipped: {_e_ppt_front}")
            
            globals()["charts"] = charts
            
            print("[debug] tilts_comparison_rows sample:", (charts.get("tilts_comparison_rows") or [])[:2])

            
            # Finally, update the existing chart on 'OPT'
            # --- Efficient Frontier Chart Update (safe version) ---
            # PPT chart was migrated to slot-point markers 2026-06-22. The
            # Excel chart still accepts the legacy current/previous/etc
            # markers but those values are no longer computed, so pass None.
            # The Excel chart will show only the frontier curve + tangency.
            # Migrate the Excel updater to slot points in a follow-up if you
            # want them on the OPT sheet too.
            try:
                update_efficient_frontier_chart(
                    opt_sheet=opt,
                    stats_df=stats_df,
                    start_s_row=start_s_row,
                    rf_annual=float(rf_annual),
                    tan_ret=float(tan_ret),
                    tan_vol=float(tan_vol),
                    current_point=None,
                    title_text=chart_title,
                    target_point=None,
                    previous_point=None,
                    factor_point=None,
                    no_tilt_point=None,
                    tilt_point=None,
                )

            except Exception as e:
                print(f"[chart] Skipping chart update: {e}")
 
            co = opt.api.ChartObjects()
            for i in range(1, co.Count + 1):
                o = co.Item(i)
                title = ""
                try:
                    if o.Chart.HasTitle:
                        title = o.Chart.ChartTitle.Text
                except Exception:
                    pass
                print(i, "name:", o.Name, "| title:", title)

          
            # ---- Build trade plan & costs - writing Trade Plan/Costs/Tilts ----
            _tp_mode = str(globals().get("TRADE_PLAN_MODE", "ask")).lower().strip()
            print(f"[tradeplan] resolved mode at entry: {_tp_mode!r}")

            # Decide which portfolio drives the ACTIVE trade plan
            if _tp_mode == "ask":
                _tp_mode = ask_tradeplan_portfolio_choice()

            elif _tp_mode == "auto":
                # Validation-based choice (Sharpe over lookback). Now considers
                # the regime-adaptive ensemble as a 3rd candidate.
                _rwide = globals().get("returns_wide_df", None)
                _w_ensemble_live = globals().get("W_ENSEMBLE_SER", pd.Series(dtype=float))
                print(f"[tradeplan] auto branch entered. returns_wide_df type={type(_rwide).__name__}, "
                      f"shape={getattr(_rwide,'shape',None)}; W_ENSEMBLE_SER len={len(_w_ensemble_live)}")
                if isinstance(_rwide, pd.DataFrame) and not _rwide.empty:
                    try:
                        choice_label, w_chosen, diag = choose_portfolio_for_tradeplan(
                            returns_df=_rwide,
                            w_no_tilts=pd.Series(w_star, index=Sigma_opt.index),
                            w_with_tilts=pd.Series(w_star_with_tilts, index=Sigma_opt.index),
                            rf_annual=float(rf_annual),
                            lookback_days=int(globals().get("VALIDATION_LOOKBACK_DAYS", 252)),
                            w_ensemble=_w_ensemble_live,
                        )
                        _tp_mode = choice_label
                        print(f"[tradeplan] auto-selected: {choice_label} "
                              f"(Sharpe — ens:{diag.get('sharpe_ensemble', np.nan):.2f}, "
                              f"with:{diag.get('sharpe_with_tilts', np.nan):.2f}, "
                              f"no:{diag.get('sharpe_no_tilts', np.nan):.2f})")
                    except Exception as _e_choose:
                        print(f"[tradeplan] choose_portfolio_for_tradeplan FAILED: {_e_choose!r}")
                        _tp_mode = "no_tilts"
                else:
                    print("[tradeplan] returns_wide_df not usable → defaulting to no_tilts")
                    _tp_mode = "no_tilts"

            # Resolve ensemble weights for the active plan (if requested).
            _w_ensemble_live = globals().get("W_ENSEMBLE_SER", pd.Series(dtype=float))

            # === Hybrid pricing: IBKR delayed override (Tier-1 #1) ===
            # Replace yfinance last-prices with IBKR's where available so the
            # trade plan's cash flow + brokerage estimates match what we'd
            # actually fill at. Silently falls back to yfinance if TWS isn't
            # running.
            if USE_IBKR_LIVE_PRICES:
                try:
                    _t_ibkr = time.perf_counter()
                    _ibkr_px = fetch_ibkr_live_prices_native(list(last_px_hold.index))
                    if _ibkr_px:
                        last_px_hold, _diag = apply_ibkr_price_override(last_px_hold, _ibkr_px)
                        print(f"[ibkr-price] applied to {_diag['n_overridden']} tickers "
                              f"in {time.perf_counter()-_t_ibkr:.1f}s "
                              f"(max divergence: {_diag['max_bps']:.0f} bps "
                              f"on {_diag.get('max_bps_ticker') or 'n/a'}; "
                              f"{_diag['n_warn']} >{IBKR_DIVERGENCE_WARN_BPS}bps warned)")
                    else:
                        print(f"[ibkr-price] no IBKR prices returned; using yfinance")
                except Exception as _e_px:
                    print(f"[ibkr-price] override skipped: {_e_px}")

            # === Live TLH injection (Phase 4) ====================================
            # Run TLH against the CURRENT lot book BEFORE building trade plans so
            # the resulting rebalance delta already includes the harvest swaps.
            # Result: the user executes one atomic batch via Phase 3 — rebalance
            # + TLH baked into the same recommended_trades list. The swap events
            # are also recorded as `tlh_swaps` in the rec log so Phase 3 can
            # annotate them, the PPT slide can show a TLH-specific summary, etc.
            #
            # Cooldown state is persisted across runs in tlh_cooldown_state.json
            # so the engine doesn't recommend re-buying a substitute we recently
            # harvested into (wash-swap protection under TR 2008/1).
            live_tlh_events: list[dict] = []
            if TLH_ENABLED:
                try:
                    _tlh_cooldown_path = APP_DIR / "tlh_cooldown_state.json"
                    _live_cooldown_state = _load_tlh_cooldown_state(_tlh_cooldown_path)
                    _live_lot_book = _build_lot_book_from_df(lots_df)
                    # Build price snapshot dict (ticker -> AUD price) from
                    # last_px_hold; substitute tickers may be outside the
                    # current holdings universe so widen with all known prices.
                    _px_snap = {}
                    try:
                        for _tk, _v in last_px_hold.items():
                            if pd.notna(_v) and float(_v) > 0:
                                _px_snap[str(_tk)] = float(_v)
                    except Exception:
                        pass
                    # Also accept any tickers in the wider prices DataFrame so
                    # substitutes outside last_px_hold can still be priced.
                    try:
                        _last_row = prices.iloc[-1]
                        for _tk, _v in _last_row.items():
                            if pd.notna(_v) and float(_v) > 0 and str(_tk) not in _px_snap:
                                _px_snap[str(_tk)] = float(_v)
                    except Exception:
                        pass
                    _tlh_today = pd.Timestamp(prices.index[-1])
                    _live_nav = (float(portfolio_value_override)
                                  if (portfolio_value_override is not None
                                      and np.isfinite(portfolio_value_override)
                                      and portfolio_value_override > 0)
                                  else None)
                    _tlh_out = _run_tlh_pass(
                        lot_book=_live_lot_book,
                        price_snapshot=_px_snap,
                        as_of=_tlh_today,
                        cooldown_state=_live_cooldown_state,
                        pairs=TLH_PAIRS,
                        nav_aud=_live_nav,
                    )
                    live_tlh_events = _tlh_out.get("events", []) or []
                    if live_tlh_events:
                        # Apply swap deltas to `units` so the rebalance delta
                        # below already incorporates them. Net effect: the
                        # user executes one batch via Phase 3 and the harvest
                        # swap + rebalance both clear together.
                        units = pd.Series(units, dtype=float).copy()
                        for _ev in live_tlh_events:
                            _stk = str(_ev["ticker_sold"])
                            _btk = str(_ev["ticker_bought"])
                            _su = float(_ev["units_sold"])
                            _bu = float(_ev["units_bought"])
                            units[_stk] = float(units.get(_stk, 0.0)) - _su
                            units[_btk] = float(units.get(_btk, 0.0)) + _bu
                        _save_tlh_cooldown_state(_tlh_cooldown_path,
                                                   _live_cooldown_state)
                        _loss_total = sum(float(ev.get("loss_aud", 0))
                                            for ev in live_tlh_events)
                        print(f"[tlh-live] {len(live_tlh_events)} swap(s) "
                              f"recommended, ${_loss_total:,.0f} loss to "
                              f"realise. units adjusted; cooldown saved.")
                    else:
                        print(f"[tlh-live] no qualifying loss lots "
                              f"(threshold {TLH_MIN_LOSS_PCT*100:+.0f}%, "
                              f"min ${TLH_MIN_LOSS_AUD:.0f}, "
                              f"cooldown {TLH_COOLDOWN_DAYS}d, "
                              f"{len(lots_df) if lots_df is not None else 0} "
                              f"lots scanned)")
                except Exception as _e_live_tlh:
                    print(f"[tlh-live] skipped: {_e_live_tlh}")
            globals()["LIVE_TLH_EVENTS"] = live_tlh_events

            # Cash-fit: size all three plans to the cash actually on hand so
            # the net buys are always fundable (live IBKR → last snapshot →
            # None=NAV sizing). Fetched ONCE and shared across the variants.
            _avail_cash_aud = _get_available_cash_aud()
            if _avail_cash_aud is not None:
                print(f"[cash-fit] sizing live plans to holdings + "
                      f"${_avail_cash_aud:,.0f} cash − reserve "
                      f"(was NAV ${(portfolio_value_override or 0):,.0f})")

            # Build ALL three trade plans up-front; downstream code can compare
            # what each one implies, and the "active" one is picked below.
            trade_no, resid_no = make_trade_plan(
                units, last_px_hold, fx_map_all, w_star,
                include_zero_lines=True, include_flags=include_flags,
                portfolio_value_override=portfolio_value_override,
                available_cash_aud=_avail_cash_aud
            )

            trade_with, resid_with = make_trade_plan(
                units, last_px_hold, fx_map_all, w_star_with_tilts,
                include_zero_lines=True, include_flags=include_flags,
                portfolio_value_override=portfolio_value_override,
                available_cash_aud=_avail_cash_aud
            )

            trade_ens, resid_ens = None, None
            if isinstance(_w_ensemble_live, pd.Series) and not _w_ensemble_live.empty:
                try:
                    # The ensemble can recommend tickers that aren't in Sigma_opt
                    # (e.g. GOLD.AX, GOVT.AX — added to prices but missing FF5
                    # regional betas). Use the UNION of Sigma_opt and the
                    # ensemble's own tickers so nothing gets dropped.
                    _target_idx = _w_ensemble_live.index.union(Sigma_opt.index)
                    _target_idx = _target_idx.difference({"^AORD", "PortfolioValue"})
                    _w_ens_full = pd.Series(0.0, index=_target_idx)
                    _w_ens_full.loc[_w_ensemble_live.index.intersection(_target_idx)] = \
                        _w_ensemble_live.reindex(_target_idx.intersection(_w_ensemble_live.index))
                    if _w_ens_full.sum() > 0:
                        _w_ens_full = _w_ens_full / _w_ens_full.sum()
                    # Diagnostic: surface tickers the ensemble picked that
                    # weren't in Sigma_opt (FF5 filtered them out).
                    _extras = sorted(set(_w_ensemble_live.index) - set(Sigma_opt.index) - {"^AORD"})
                    if _extras:
                        print(f"[tradeplan] ensemble picks NOT in Sigma_opt "
                              f"(FF5 filter): {_extras}")
                    trade_ens, resid_ens = make_trade_plan(
                        units, last_px_hold, fx_map_all, _w_ens_full,
                        include_zero_lines=True, include_flags=include_flags,
                        portfolio_value_override=portfolio_value_override,
                        available_cash_aud=_avail_cash_aud
                    )
                except Exception as _e_ens:
                    print(f"[tradeplan] ensemble plan build failed: {_e_ens}")
                    trade_ens, resid_ens = None, None

            # Map _tp_mode → (w_tradeplan, trade_rec, resid_rec)
            # The ensemble path can include tickers outside Sigma_opt (e.g.
            # GOLD.AX or dialog-added NDQ.AX that lack FF5 regional betas).
            # Capture the actual weight index per-branch so we don't force a
            # length mismatch when the ensemble's universe is wider.
            if _tp_mode == "ensemble" and trade_ens is not None:
                w_tradeplan_vals = _w_ens_full.values
                w_tradeplan_idx = _w_ens_full.index
                trade_rec, resid_rec = trade_ens, resid_ens
            elif _tp_mode == "with_tilts":
                w_tradeplan_vals = w_star_with_tilts
                w_tradeplan_idx = Sigma_opt.index
                trade_rec, resid_rec = trade_with, resid_with
            else:
                _tp_mode = "no_tilts"
                w_tradeplan_vals = w_star
                w_tradeplan_idx = Sigma_opt.index
                trade_rec, resid_rec = trade_no, resid_no
            w_tradeplan = pd.Series(np.asarray(w_tradeplan_vals, dtype=float),
                                    index=w_tradeplan_idx)

            # === Rebalance trigger verdict ============================================
            # Computes Σ|Δw| from the trade_rec DataFrame directly:
            #     Σ|Δw| = Σ|Δ_units · last_px| / portfolio_value
            # This is the actual trade volume normalised by NAV — exactly the
            # L1 distance between current and target weight vectors, computed
            # WITHOUT the weight-vector arithmetic that has produced impossible
            # values twice now (39.87 → 78.53 → 155.91 across 3 runs, despite
            # raw sums of both inputs being 1.0). The trade_rec path uses the
            # engine's already-validated make_trade_plan() output and inherits
            # its rounding for partial-share trades.
            #
            # Rationale recap: a previous diff-of-weight-vectors approach
            # silently degraded when one side had signed/cancelling values.
            # The trade-volume approach can't degrade that way — it's a
            # straight cash-flow sum.
            try:
                # Use the existing helper so we match whatever encoding the
                # delta column has (Δ Units, Delta Units, or the mojibake
                # variant from legacy Excel round-trips).
                _delta_col_v = _trade_delta_col(trade_rec)
                _px_col_v = ("Last Px (AUD)"
                               if "Last Px (AUD)" in trade_rec.columns
                               else None)
                # Portfolio value: the same value the rec log will record so
                # the verdict and the log agree.
                if (portfolio_value_override is not None
                        and np.isfinite(portfolio_value_override)
                        and portfolio_value_override > 0):
                    _portfolio_val_for_verdict = float(portfolio_value_override)
                else:
                    _cu_v = pd.to_numeric(trade_rec.get("Curr Units", 0),
                                            errors="coerce").fillna(0)
                    _lp_v = pd.to_numeric(trade_rec.get("Last Px (AUD)", 0),
                                            errors="coerce").fillna(0)
                    _portfolio_val_for_verdict = float((_cu_v * _lp_v).sum())

                if _delta_col_v is None or _px_col_v is None:
                    raise ValueError(
                        f"trade_rec missing delta or price column "
                        f"(delta={_delta_col_v}, px={_px_col_v})"
                    )
                _deltas = pd.to_numeric(trade_rec[_delta_col_v],
                                          errors="coerce").fillna(0)
                _pxs = pd.to_numeric(trade_rec[_px_col_v],
                                       errors="coerce").fillna(0)
                _trade_volume_aud = float(np.abs(_deltas * _pxs).sum())
                _n_trades = int((_deltas != 0).sum())
                if _portfolio_val_for_verdict > 0:
                    _summed_abs_dw = _trade_volume_aud / _portfolio_val_for_verdict
                else:
                    _summed_abs_dw = float("nan")

                _verdict = ("SKIP" if (np.isfinite(_summed_abs_dw)
                                         and _summed_abs_dw < float(SKIP_REBAL_DELTA))
                              else "RUN")
                print(f"[rebal-trigger] summed_|Δw|={_summed_abs_dw:.4f}  "
                      f"threshold={float(SKIP_REBAL_DELTA):.4f}  "
                      f"verdict={_verdict}  "
                      f"mode={_tp_mode}  "
                      f"portfolio_aud={_portfolio_val_for_verdict:,.0f}  "
                      f"trade_volume_aud={_trade_volume_aud:,.0f}  "
                      f"n_trades={_n_trades}")
                globals()["REBAL_TRIGGER_VERDICT"] = _verdict
                globals()["REBAL_TRIGGER_SUMMED_DW"] = _summed_abs_dw
                # Broker-vs-engine mark reconciliation (READ-ONLY, informs
                # nothing downstream — the verdict above is already final).
                # Uses the latest --snapshot-nav row if it's fresh (<3 days):
                # a persistent gap means fees/marks the engine isn't seeing.
                try:
                    _bnav = _load_broker_nav_series()
                    if (not _bnav.empty
                            and (pd.Timestamp.now().normalize() - _bnav.index[-1]).days <= 3
                            and _portfolio_val_for_verdict > 0):
                        _gap = (_portfolio_val_for_verdict - float(_bnav.iloc[-1])) / float(_bnav.iloc[-1])
                        _sev = "[drift][WARN]" if abs(_gap) > 0.01 else "[drift]"
                        print(f"{_sev} engine mark ${_portfolio_val_for_verdict:,.0f} vs "
                              f"broker NetLiq ${float(_bnav.iloc[-1]):,.0f} "
                              f"({_gap*100:+.2f}%, snapshot {_bnav.index[-1].date()})")
                except Exception:
                    pass
            except Exception as _e_rebal_trig:
                print(f"[rebal-trigger] verdict computation failed: {_e_rebal_trig}")
                globals()["REBAL_TRIGGER_VERDICT"] = "UNKNOWN"
                globals()["REBAL_TRIGGER_SUMMED_DW"] = float("nan")

            # Persist labels/weights for PPT + achieved-tilts table
            globals()["TRADEPLAN_LABEL"] = _tp_mode
            globals()["TRADEPLAN_WEIGHTS_SER"] = w_tradeplan.copy()

            # Keep the others available for Excel writing / comparison
            globals()["TRADEPLAN_DF_NO_TILTS"] = trade_no.copy()
            globals()["TRADEPLAN_DF_WITH_TILTS"] = trade_with.copy()
            if trade_ens is not None:
                globals()["TRADEPLAN_DF_ENSEMBLE"] = trade_ens.copy()
            
            # --- Ensure 'Security' is a proper column BEFORE any downstream functions ---
            trade_rec = trade_rec.copy()
            trade_rec.columns = [str(c).strip() for c in trade_rec.columns]
            trade_rec.index.name = "Security"
            if "Security" not in trade_rec.columns:
                trade_rec = trade_rec.reset_index()
            
            # --- Now it is safe to compute costs (some code expects trade_rec["Security"]) ---
            costs_rec = evaluate_transaction_costs(
                trade_rec, lots_df, pd.Timestamp(prices.index[-1]), MARGINAL_TAX_RATE
            )
            
            # --- Add per-row brokerage (keep your existing logic) ---
            row_b = costs_rec.get("per_row_brokerage", pd.Series(0.0, index=trade_rec.index))
            row_b = pd.to_numeric(row_b, errors="coerce").reindex(trade_rec.index).fillna(0.0)
            
            # No brokerage where trade delta is zero (support both Delta Units and ÃŽâ€ Units)
            _delta_col = (_trade_delta_col(trade_rec) if "_trade_delta_col" in globals() else ("Delta Units" if "Delta Units" in trade_rec.columns else "ÃŽâ€ Units"))
            _delta_vals = pd.to_numeric(trade_rec.get(_delta_col, 0), errors="coerce").fillna(0).astype(int)
            row_b = np.where(_delta_vals == 0, 0.0, row_b)
            trade_rec["Brokerage (AUD)"] = pd.Series(row_b, index=trade_rec.index).round(2)

            trade_rec.drop(
                columns=[c for c in trade_rec.columns if str(c).lower().startswith("promo")],
                errors="ignore",
                inplace=True
            )

            # === Sanity layer (2026-06-27) =====================================
            # Halt the engine on structurally absurd trade plans BEFORE any
            # side effects (recommendation log, PPT, state file, paper exec).
            # Defends against silent state-corruption bugs like the
            # 2026-06-26 SMH→SOXX phantom-lots incident — see
            # _validate_trade_plan_sanity() docstring for thresholds and
            # the SanityViolation class for the contract.
            #
            # Resolve NAV the same way the drift-log block does immediately
            # below, so the check uses the same number the rest of the run
            # operates on.
            if (portfolio_value_override is not None
                    and np.isfinite(portfolio_value_override)
                    and portfolio_value_override > 0):
                _sanity_nav = float(portfolio_value_override)
            else:
                _cu_s = pd.to_numeric(trade_rec.get("Curr Units", 0), errors="coerce").fillna(0)
                _lp_s = pd.to_numeric(trade_rec.get("Last Px (AUD)", 0), errors="coerce").fillna(0)
                _sanity_nav = float((_cu_s * _lp_s).sum())
            # Let SanityViolation propagate — top-level run loop will catch
            # and exit with a clear error. We intentionally do NOT wrap this
            # in try/except — that would defeat the entire purpose.
            _validate_trade_plan_sanity(trade_rec, _sanity_nav)

            # === Drift tracker (Tier-1 #3): recommendation log =================
            # One JSONL line per run. Foundation for later fill/slippage compare
            # once IBKR API or manual fill sheet is wired up.
            try:
                _drift_log_path = APP_DIR / "trade_recommendation_log.jsonl"
                if (portfolio_value_override is not None
                        and np.isfinite(portfolio_value_override)
                        and portfolio_value_override > 0):
                    _portfolio_val_for_log = float(portfolio_value_override)
                else:
                    _cu = pd.to_numeric(trade_rec.get("Curr Units", 0), errors="coerce").fillna(0)
                    _lp = pd.to_numeric(trade_rec.get("Last Px (AUD)", 0), errors="coerce").fillna(0)
                    _portfolio_val_for_log = float((_cu * _lp).sum())
                append_trade_recommendation_log(
                    _drift_log_path,
                    selected_mode=_tp_mode,
                    trade_df=trade_rec,
                    w_target=w_tradeplan,
                    current_units=pd.Series(units),
                    portfolio_value_aud=_portfolio_val_for_log,
                    regime_mix=globals().get("ensemble_mix_live", pd.Series(dtype=float)),
                    expected_brokerage_aud=float(costs_rec.get("brokerage", 0.0)),
                    expected_cgt_aud=float(costs_rec.get("cgt_tax", 0.0)),
                    broker_name=str(BROKER_CONFIG.get("name", "unknown")),
                    cgt_mtr=float(CGT_CONFIG.get("marginal_tax_rate", 0.30)),
                    universe_size=int(len(w_tradeplan)),
                    tlh_events=globals().get("LIVE_TLH_EVENTS", []) or [],
                )
            except Exception as _e_drift:
                print(f"[drift] recommendation log skipped: {_e_drift}")

            # === Drift tracker v2/v3: fills + NAV history + warnings ===========
            # Append live NAV; ensure Actual_Fills sheet; join fills against
            # recommendation log; compute monthly drift if live trading is on.
            try:
                _nav_path = APP_DIR / "live_nav_history.jsonl"
                append_live_nav_history(_nav_path, _portfolio_val_for_log)
                _live_nav = _load_live_nav_series(_nav_path)
                _live_dd = compute_live_max_drawdown(_live_nav)
                _ensure_actual_fills_sheet(wb)
                _fills_df = _read_actual_fills(wb)
                _fills_drift = compute_fill_drift(_fills_df, _drift_log_path)
                _oos_ret = globals().get("oos_returns_daily", pd.Series(dtype=float))
                _nav_drift = compute_monthly_nav_drift(
                    _live_nav, _oos_ret, LIVE_TRADING_START_DATE,
                )
                _write_drift_sheets(wb, _fills_drift, _nav_drift, _live_nav, _live_dd)
                _n_warn = _print_drift_warnings(_fills_drift, _nav_drift, _live_dd)
                _adherent = 0 if _fills_drift.empty else int(_fills_drift["Recommended"].sum())
                _total_fills = 0 if _fills_drift.empty else len(_fills_drift)
                print(f"[drift] tracker: NAV samples={int(_live_nav.size)}, "
                      f"current DD {_live_dd*100:+.2f}%, "
                      f"fills {_adherent}/{_total_fills} adherent, "
                      f"warnings={_n_warn}")
            except Exception as _e_drift_v2:
                print(f"[drift] v2/v3 tracker skipped: {_e_drift_v2}")

            # === Cash ledger (persistent) =====================================
            # Snapshot every run so the user can see where money is going:
            # cumulative brokerage + CGT, drift vs start, drift vs $1M target,
            # unexplained delta (if non-zero, brokerage/CGT/market doesn't add up).
            try:
                _cash_ledger_path = APP_DIR / "cash_ledger.jsonl"
                # Compute trade-plan-summary values inline (the live block that
                # defines `total_portfolio` etc. lives below this point — we
                # mirror its logic here so the ledger runs in the right scope).
                _cash_total_brokerage = float(costs_rec.get("brokerage", 0.0))
                _cash_net_invested = 0.0
                _cash_balance_local = 0.0
                _cash_total_portfolio = 0.0
                if trade_rec is not None and not trade_rec.empty:
                    _tgt_units = pd.to_numeric(trade_rec.get("Target Units"), errors="coerce").fillna(0.0)
                    _last_px = pd.to_numeric(trade_rec.get("Last Px (AUD)"), errors="coerce").fillna(0.0)
                    _cash_net_invested = float((_tgt_units * _last_px).sum())
                    if (portfolio_value_override is not None
                            and np.isfinite(portfolio_value_override)
                            and float(portfolio_value_override) > 0):
                        _cash_total_portfolio = float(portfolio_value_override)
                        _cash_balance_local = (_cash_total_portfolio
                                               - _cash_net_invested
                                               - _cash_total_brokerage)
                    else:
                        _cash_balance_local = float(
                            pd.to_numeric(trade_rec.get("Cash Flow (AUD)"),
                                          errors="coerce").fillna(0.0).sum()
                        )
                        _cash_total_portfolio = _cash_net_invested + _cash_balance_local
                append_cash_ledger(
                    _cash_ledger_path,
                    portfolio_value_aud=_cash_total_portfolio,
                    net_invested_aud=_cash_net_invested,
                    cash_balance_aud=_cash_balance_local,
                    brokerage_this_run_aud=_cash_total_brokerage,
                    cgt_this_run_aud=float(costs_rec.get("cgt_tax", 0.0)),
                    loss_cf_tax_aud=float(
                        costs_rec.get("breakdown", {}).get("loss_carry_forward", 0.0)
                    ) * float(CGT_CONFIG.get("marginal_tax_rate", 0.30)),
                    selected_mode=_tp_mode,
                    broker_name=str(BROKER_CONFIG.get("name", "unknown")),
                )
                _ledger_df = _load_cash_ledger(_cash_ledger_path)
                _write_cash_ledger_sheet(wb, _ledger_df)
                if not _ledger_df.empty:
                    _latest = _ledger_df.iloc[-1]
                    _runs = len(_ledger_df)
                    _unex = _latest.get("unexplained_delta_aud")
                    _unex_str = (f"${float(_unex):,.0f}"
                                 if _unex is not None and not pd.isna(_unex)
                                 else "(first run — no prior to compare)")
                    print(f"[cash] ledger: {_runs} run(s) recorded. "
                          f"Drift vs ${TARGET_PORTFOLIO_VALUE_AUD:,.0f}: "
                          f"${float(_latest['drift_vs_target_aud']):,.0f} | "
                          f"Cum. brokerage ${float(_latest['cum_brokerage_aud']):,.0f} | "
                          f"Cum. CGT ${float(_latest['cum_cgt_aud']):,.0f} | "
                          f"Unexplained Δ {_unex_str}")
            except Exception as _e_cash:
                print(f"[cash] ledger skipped: {_e_cash}")

            # --- Lot expansion (safe now that 'Security' exists as a column) ---
            lot_expanded = expand_with_lots(
                trade_rec,
                lots_df,
                sale_date=pd.Timestamp(prices.index[-1]),
                method="FIFO"
            )
            print("\n=== LOT-EXPANDED TABLE ===")
            print(lot_expanded.head(20))


            if "Security" not in trade_rec.columns and trade_rec.index.name == "Security":
                trade_rec = trade_rec.reset_index()
            if "Security" not in trade_rec.columns:
                trade_rec.insert(0, "Security", trade_rec.index.astype(str))
            
            if isinstance(lot_expanded, pd.DataFrame):
                if "Security" not in lot_expanded.columns and lot_expanded.index.name == "Security":
                    lot_expanded = lot_expanded.reset_index()

            
            # === Build CGT audit table (parcel-level) ===
            try:
                tax_bkd = costs_rec.get("breakdown", {})
                audit_df = tax_bkd.get("audit", pd.DataFrame()).copy()

                if not audit_df.empty:
                    # Ensure proper dtypes
                    audit_df["AcqDate"] = pd.to_datetime(audit_df["AcqDate"], errors="coerce")
                    audit_df["SaleDate"] = pd.to_datetime(audit_df["SaleDate"], errors="coerce")
                    audit_df["Qty"]      = pd.to_numeric(audit_df["Qty"], errors="coerce")
                    audit_df["Proceeds"] = pd.to_numeric(audit_df["Proceeds"], errors="coerce")
                    audit_df["CostBase"] = pd.to_numeric(audit_df["CostBase"], errors="coerce")
                    audit_df["Gain"]     = pd.to_numeric(audit_df["Gain"], errors="coerce")

                    # Holding period & discount flag (12-month rule)
                    audit_df["HoldingDays"] = (audit_df["SaleDate"] - audit_df["AcqDate"]).dt.days
                    audit_df["LongTermEligible"] = audit_df["LongTermEligible"].astype(bool)

                    # 50% discount only for positive gains that are LT eligible
                    audit_df["DiscountRate"] = 0.0
                    audit_df.loc[(audit_df["Gain"] > 0) & (audit_df["LongTermEligible"]), "DiscountRate"] = 0.5

                    audit_df["DiscountedGainIllustrative"] = audit_df["Gain"]
                    mask_disc = (audit_df["Gain"] > 0) & (audit_df["LongTermEligible"])
                    audit_df.loc[mask_disc, "DiscountedGainIllustrative"] = (
                        audit_df.loc[mask_disc, "Gain"] * 0.5
                    )

                    # === Write parcel-level audit sheet ===
                    try:
                        sht_cgt = get_or_clear_sheet(wb, "CGT_Audit")
                        sht_cgt.range("A1").value = [[
                            "Security",
                            "Qty",
                            "AcqDate",
                            "SaleDate",
                            "Proceeds",
                            "CostBase",
                            "Gain",
                            "LongTermEligible",
                            "HoldingDays",
                            "DiscountRate",
                            "DiscountedGainIllustrative",
                        ]]
                        sht_cgt.range("A2").options(index=False, header=False).value = audit_df[
                            [
                                "Security",
                                "Qty",
                                "AcqDate",
                                "SaleDate",
                                "Proceeds",
                                "CostBase",
                                "Gain",
                                "LongTermEligible",
                                "HoldingDays",
                                "DiscountRate",
                                "DiscountedGainIllustrative",
                            ]
                        ]
                    except Exception as e_cgt_sheet:
                        print(f"[cgt] could not write CGT_Audit sheet: {e_cgt_sheet}")

                    # === Optional security-level summary ===
                    try:
                        sec_grp = audit_df.groupby("Security", as_index=False).agg(
                            ProceedsTotal=("Proceeds", "sum"),
                            CostBaseTotal=("CostBase", "sum"),
                            GainTotal=("Gain", "sum"),
                        )

                        lt_mask = audit_df["LongTermEligible"]
                        st_mask = ~audit_df["LongTermEligible"]

                        lt_sum = (
                            audit_df.loc[lt_mask]
                            .groupby("Security")["Gain"]
                            .sum()
                            .rename("LongTermGain")
                        )
                        st_sum = (
                            audit_df.loc[st_mask]
                            .groupby("Security")["Gain"]
                            .sum()
                            .rename("ShortTermGain")
                        )

                        sec_summary = (
                            sec_grp
                            .merge(lt_sum, on="Security", how="left")
                            .merge(st_sum, on="Security", how="left")
                            .fillna(0.0)
                        )

                        sht_cgt.range("L1").value = [[
                            "Security",
                            "ProceedsTotal",
                            "CostBaseTotal",
                            "GainTotal",
                            "LongTermGain",
                            "ShortTermGain",
                        ]]
                        sht_cgt.range("L2").options(index=False, header=False).value = sec_summary[
                            [
                                "Security",
                                "ProceedsTotal",
                                "CostBaseTotal",
                                "GainTotal",
                                "LongTermGain",
                                "ShortTermGain",
                            ]
                        ]
                    except Exception as e_cgt_summary:
                        print(f"[cgt] could not write CGT summary: {e_cgt_summary}")

                else:
                    print("[cgt] audit_df is empty (no CGT-relevant sells).")

            except Exception as e_cgt:
                print(f"[cgt] error building CGT audit table: {e_cgt}")

            
            # ---- Achieved factor tilts table (use implemented target holdings if available) ----
            tilts_out = None
            if (B is not None) and (not B.empty):
                factor_order = ["Mkt-RF","SMB","HML","RMW","CMA","MOM"]
            
                achieved_series = None
            
                # 1) Prefer implemented portfolio from trade plan target units
                try:
                    if isinstance(trade_rec, pd.DataFrame) and (not trade_rec.empty):
                        tr = trade_rec.copy()
                        if "Security" not in tr.columns and tr.index.name == "Security":
                            tr = tr.reset_index()
            
                        if ("Security" in tr.columns) and ("Target Units" in tr.columns) and ("Last Px (AUD)" in tr.columns):
                            tgt_u = pd.to_numeric(tr["Target Units"], errors="coerce").fillna(0.0)
                            px_aud = pd.to_numeric(tr["Last Px (AUD)"], errors="coerce").fillna(0.0)
                            val = tgt_u * px_aud
            
                            w_impl = pd.Series(val.values, index=tr["Security"].astype(str))
                            w_impl = w_impl.reindex(B.index).fillna(0.0)
            
                            s = float(w_impl.sum())
                            if s > 0:
                                w_impl = w_impl / s
                                achieved_series = (B.T @ w_impl).reindex(factor_order)
                except Exception as _e_tilt_impl:
                    achieved_series = None
            
                # 2) Fallback to model weights if needed
                if achieved_series is None:
                    w_use = None
                
                    # Prefer w_tilt if available (and align safely)
                    if ("w_tilt" in locals()) and (w_tilt is not None):
                        if isinstance(w_tilt, pd.Series):
                            w_use = pd.to_numeric(w_tilt, errors="coerce").reindex(Sigma_opt.index).fillna(0.0)
                        else:
                            wt = np.asarray(w_tilt, dtype=float).reshape(-1)
                            if wt.shape[0] == len(Sigma_opt.index):
                                w_use = pd.Series(wt, index=Sigma_opt.index).fillna(0.0)
                            elif wt.shape[0] == len(Sigma_opt.index) + 1:
                                w_use = pd.Series(wt[:len(Sigma_opt.index)], index=Sigma_opt.index).fillna(0.0)
                            else:
                                raise ValueError(f"w_tilt length {wt.shape[0]} does not match Sigma_opt universe {len(Sigma_opt.index)}")
                
                    # If no w_tilt, fall back to w_star
                    if w_use is None:
                        w_use = pd.Series(w_star, index=Sigma_opt.index).reindex(Sigma_opt.index).fillna(0.0)
                
                    # If trade plan weights were stored, use those; otherwise fall back to the with-tilts optimiser weights
                    w_use = globals().get("TRADEPLAN_WEIGHTS_SER", None)
                    if w_use is None:
                        w_use = w_star_with_tilts
                    
                    w_use = pd.Series(w_use).reindex(B.index).fillna(0.0)
                    
                    s = float(w_use.sum())
                    if s > 0:
                        w_use = w_use / s
                    
                    achieved_series = (B.T @ w_use).reindex(factor_order)

            
                # Build output table
                if isinstance(tilt_df, pd.DataFrame) and not tilt_df.empty:
                    tgt = tilt_df.reindex(factor_order)
                    tilts_out = pd.DataFrame({
                        "Use?": tgt["Use?"].astype(str).str.upper().isin(["TRUE","1","Y","YES","T"]).map({True:"Yes", False:"No"}),
                        "Target Beta": pd.to_numeric(tgt["Target"], errors="coerce"),
                        "Band": pd.to_numeric(tgt["Band"], errors="coerce"),
                        "Achieved Beta": achieved_series,
                    })
                    tilts_out["Diff"] = tilts_out["Achieved Beta"] - tilts_out["Target Beta"]
                    tilts_out["Within Band?"] = (tilts_out["Diff"].abs() <= tilts_out["Band"]).map({True: "Yes", False: "No"})
                else:
                    tilts_out = achieved_series.to_frame()

            # ---------- Layout anchors (avoid overlaps) ----------
            anchor_row = start_s_row + stats_df.shape[0] + 4
            TP_COL, COST_COL, TILT_COL = "A", "J", "M"
            # Pre-compute summary_row so the alt-plan block (which anchors below it) can reference it.
            summary_row = anchor_row + trade_rec.shape[0] + 4

            # ---------- LEFT: Trade Plan ----------
            opt.range(f"{TP_COL}{anchor_row}").value = "Trade Plan (rounded units)"

            # Write the main trade plan body. Clear formatting on destination first to avoid
            # bold/currency carryover from prior runs (Excel preserves cell formats when only contents are cleared).
            if isinstance(trade_rec, pd.DataFrame) and not trade_rec.empty:
                tp_header = anchor_row + 1
                tp_data_first = tp_header + 1
                tp_data_last = tp_header + trade_rec.shape[0]
                try:
                    opt.range(f"{TP_COL}{tp_header}:G{tp_data_last+4}").api.ClearFormats()
                except Exception:
                    pass
                opt.range(f"{TP_COL}{tp_header}").options(pd.DataFrame, index=False, header=True).value = trade_rec
                set_number_formats(opt, {
                    f"B{tp_data_first}:D{tp_data_last}": "0",
                    f"E{tp_data_first}:G{tp_data_last}": "$0.00",
                })

            # --- Write an Alternative Trade Plan block (full + aligned + with summaries) ---
            try:
                # Identify which DF is the alternative one
                alt_df = globals().get("TRADEPLAN_DF_NO_TILTS", None)
                if str(globals().get("TRADEPLAN_LABEL", "")).lower().strip() == "no_tilts":
                    alt_df = globals().get("TRADEPLAN_DF_WITH_TILTS", None)
            
                if isinstance(alt_df, pd.DataFrame) and not alt_df.empty:
                    # Force same columns as the main trade plan (prevents "missing columns")
                    alt_df = alt_df.copy()
                    alt_df = alt_df.copy()
                    if "Security" not in alt_df.columns:
                        alt_df.insert(0, "Security", alt_df.index.astype(str))
                    
                    alt_df = alt_df.reindex(columns=trade_rec.columns)
                    # Place BELOW the main trade plan summary (so nothing gets overwritten)
                    alt_anchor = summary_row + 4
                    alt_header = alt_anchor + 1
                    alt_data_first = alt_header + 1
                    alt_data_last = alt_header + alt_df.shape[0]

                    # Clear formatting on destination to avoid carryover from prior runs (rows that
                    # previously held the main plan inherit its $/bold formats otherwise).
                    try:
                        opt.range(f"{TP_COL}{alt_anchor}:G{alt_data_last+4}").api.ClearFormats()
                    except Exception:
                        pass
                    opt.range(f"{TP_COL}{alt_anchor}").value = "Alternative Trade Plan (rounded units)"
                    opt.range(f"{TP_COL}{alt_header}").options(pd.DataFrame, index=False, header=True).value = alt_df
                    set_number_formats(opt, {
                        f"B{alt_data_first}:D{alt_data_last}": "0",
                        f"E{alt_data_first}:G{alt_data_last}": "$0.00",
                    })
            
                    # Compute alt costs (so you can show brokerage/CGT/total for the alt plan too)
                    alt_for_costs = alt_df.copy()
                    alt_for_costs.columns = [str(c).strip() for c in alt_for_costs.columns]
                    alt_for_costs.index.name = "Security"
                    if "Security" not in alt_for_costs.columns:
                        alt_for_costs = alt_for_costs.reset_index()
            
                    alt_costs = evaluate_transaction_costs(
                        alt_for_costs, lots_df, pd.Timestamp(prices.index[-1]), MARGINAL_TAX_RATE
                    )
                    alt_total_brokerage = float(alt_costs.get("brokerage", 0.0))
            
                    # Alt portfolio value + cash summary
                    alt_net_invested = 0.0
                    alt_cash_balance = 0.0
                    alt_total_portfolio = 0.0
            
                    if not alt_df.empty:
                        alt_tgt_units = pd.to_numeric(alt_df["Target Units"], errors="coerce").fillna(0.0)
                        alt_last_px = pd.to_numeric(alt_df["Last Px (AUD)"], errors="coerce").fillna(0.0)
                        alt_net_invested = float((alt_tgt_units * alt_last_px).sum())
            
                        if portfolio_value_override is not None and np.isfinite(portfolio_value_override) and float(portfolio_value_override) > 0:
                            alt_total_portfolio = float(portfolio_value_override)
                            alt_cash_balance = alt_total_portfolio - alt_net_invested - alt_total_brokerage
                        else:
                            alt_cash_balance = float(pd.to_numeric(alt_df["Cash Flow (AUD)"], errors="coerce").fillna(0.0).sum())
                            alt_total_portfolio = alt_net_invested + alt_cash_balance
            
                    alt_summary_row = alt_anchor + alt_df.shape[0] + 4
                    opt.range(f"{TP_COL}{alt_summary_row}").value = [
                        ["Portfolio Value (Holdings)", alt_net_invested],
                        ["Cash", alt_cash_balance],
                        ["Total Portfolio", alt_total_portfolio],
                    ]
                    try:
                        rng_labels2 = opt.range(f"{TP_COL}{alt_summary_row}:{TP_COL}{alt_summary_row+2}").api
                        rng_labels2.Font.Bold = True
                        rng_vals2 = opt.range(f"{TP_COL}{alt_summary_row}:{TP_COL}{alt_summary_row+2}").offset(0, 1).api
                        rng_vals2.NumberFormat = "$0.00"
                    except Exception:
                        pass
            
                    # Alt costs summary (middle column block)
                    opt.range(f"{COST_COL}{alt_anchor}").value = "Transaction Costs (AUD) - Alternative"
                    opt.range(f"{COST_COL}{alt_anchor+1}").value = [
                        ["Brokerage", "CGT Tax", "Total"],
                        [alt_costs.get("brokerage", 0.0), alt_costs.get("cgt_tax", 0.0), alt_costs.get("total_cost", 0.0)],
                    ]
                    try:
                        opt.range(f"{COST_COL}{alt_anchor+2}").api.NumberFormat = "0.00"
                        opt.range(f"{COST_COL}{alt_anchor+2}").offset(0,1).api.NumberFormat = "0.00"
                        opt.range(f"{COST_COL}{alt_anchor+2}").offset(0,2).api.NumberFormat = "0.00"
                    except Exception:
                        pass
            
            except Exception as _e_alt_plan:
                print(f"[excel] Alternative trade plan write skipped: {_e_alt_plan}")


            # --- Add portfolio value & cash summary underneath the Trade Plan ---
            # Hoisted from below so the override branch at line ~4680 can use it.
            total_brokerage = float(costs_rec.get("brokerage", 0.0))
            net_invested = 0.0
            cash_balance = 0.0

            if not trade_rec.empty:
                # Value of target holdings
                tgt_units = pd.to_numeric(trade_rec["Target Units"], errors="coerce").fillna(0.0)
                last_px_aud = pd.to_numeric(trade_rec["Last Px (AUD)"], errors="coerce").fillna(0.0)
                net_invested = float((tgt_units * last_px_aud).sum())
        
                # Cash handling:
                # - If portfolio_value_override is provided, treat it as TOTAL portfolio value (holdings + cash),
                #   and compute cash as the residual after funding the TARGET holdings.
                # - Otherwise, fall back to deriving cash from the net trade cashflows.
                if portfolio_value_override is not None and np.isfinite(portfolio_value_override) and float(portfolio_value_override) > 0:
                    total_portfolio = float(portfolio_value_override)
                    cash_balance = total_portfolio - net_invested - total_brokerage
                else:
                    # Net cash after trades (positive = cash released, negative = extra cash needed)
                    cash_balance = float(pd.to_numeric(trade_rec["Cash Flow (AUD)"], errors="coerce").fillna(0.0).sum())
                    total_portfolio = net_invested + cash_balance

        
            # summary_row was pre-computed above the trade-plan writes
            opt.range(f"{TP_COL}{summary_row}").value = [
                ["Portfolio Value (Holdings)", net_invested],
                ["Cash",                      cash_balance],
                ["Total Portfolio",           total_portfolio],
            ]
            try:
                # Bold the three summary labels and format the numbers as currency
                rng_labels = opt.range(f"{TP_COL}{summary_row}:{TP_COL}{summary_row+2}").api
                rng_labels.Font.Bold = True
                rng_vals = opt.range(f"{TP_COL}{summary_row}:{TP_COL}{summary_row+2}").offset(0, 1).api
                rng_vals.NumberFormat = "$0.00"
            except Exception:
                pass

            # ---------- MIDDLE: Transaction Costs summary ----------
            opt.range(f"{COST_COL}{anchor_row}").value = "Transaction Costs (AUD)"
            opt.range(f"{COST_COL}{anchor_row+1}").value = [
                ["Brokerage", "CGT Tax", "Total"],
                [costs_rec["brokerage"], costs_rec["cgt_tax"], costs_rec["total_cost"]],
            ]
            try:
                opt.range(f"{COST_COL}{anchor_row+2}").api.NumberFormat = "0.00"
                opt.range(f"{COST_COL}{anchor_row+2}").offset(0,1).api.NumberFormat = "0.00"
                opt.range(f"{COST_COL}{anchor_row+2}").offset(0,2).api.NumberFormat = "0.00"
            except Exception:
                pass

            # ---------- RIGHT: Achieved Factor Tilts ----------
            if tilts_out is not None:
                opt.range(f"{TILT_COL}{anchor_row}").value = "Achieved Factor Tilts vs Targets"
                opt.range(f"{TILT_COL}{anchor_row+1}").options(pd.DataFrame, index=True, header=True).value = tilts_out
                t_rows = tilts_out.shape[0] + 1
                t_first = anchor_row + 1
                t_data_first = t_first + 1
                try:
                    for col_name in ["Target Beta","Band","Achieved Beta","Diff"]:
                        if col_name in tilts_out.columns:
                            idx = list(tilts_out.columns).index(col_name)
                            col_letter = chr(ord(TILT_COL) + 1 + idx)  # after index column
                            opt.range(f"{col_letter}{t_data_first}:{col_letter}{t_first+t_rows}").api.NumberFormat = "0.000"
                except Exception:
                    pass
            # ---------- BELOW RIGHT: Factor Feasible Ranges (long-only, sum=1) ----------
            if (B is not None) and (not B.empty):
                factor_order = ["Mkt-RF","SMB","HML","RMW","CMA","MOM"]
                rng_df = compute_factor_feasible_ranges(B, include_flags=include_flags, factor_order=factor_order)
            
                # Optional: show your target & achieved alongside the ranges
                if isinstance(tilts_out, pd.DataFrame):
                    # pull Target and Achieved columns safely
                    tgt = pd.to_numeric(tilts_out.get("Target Beta", np.nan), errors="coerce")
                    ach = pd.to_numeric(tilts_out.get("Achieved Beta", np.nan), errors="coerce")
                    rng_df = rng_df.join(tgt.rename("Target Beta")).join(ach.rename("Achieved Beta"))
                    min_col = "Min Beta" if "Min Beta" in rng_df.columns else ("Min Î²" if "Min Î²" in rng_df.columns else "Min beta")
                    max_col = "Max Beta" if "Max Beta" in rng_df.columns else ("Max Î²" if "Max Î²" in rng_df.columns else "Max beta")
                    if ("Target Beta" in rng_df.columns) and (min_col in rng_df.columns) and (max_col in rng_df.columns):
                        rng_df["Within Range?"] = (rng_df["Target Beta"] >= rng_df[min_col]) & (rng_df["Target Beta"] <= rng_df[max_col])
                    else:
                        rng_df["Within Range?"] = np.nan
            
                # place a few rows *below* the achieved-tilts table to avoid overlap
                tilt_rows = (tilts_out.shape[0] + 2) if isinstance(tilts_out, pd.DataFrame) else 3
                ranges_anchor = anchor_row + tilt_rows + 2
            
                opt.range(f"{TILT_COL}{ranges_anchor}").value = "Factor Feasible Ranges (long-only, sum=1)"
                opt.range(f"{TILT_COL}{ranges_anchor+1}").options(pd.DataFrame, index=True, header=True).value = rng_df
            
                # number formats
                rr = ranges_anchor + 1
                rr_rows = rng_df.shape[0] + 1
                try:
                    # format numeric columns to 3 decimals if present
                    for col_name in ["Min Beta","Max Beta","Min beta","Max beta","Target Beta","Achieved Beta"]:
                        if col_name in rng_df.columns:
                            idx = list(rng_df.columns).index(col_name)
                            # first data column is one to the right of TILT_COL
                            col_letter = chr(ord(TILT_COL) + 1 + idx)
                            opt.range(f"{col_letter}{rr+1}:{col_letter}{rr+rr_rows}").api.NumberFormat = "0.000"
                except Exception:
                    pass

            # Final tidy
            try: opt.autofit()
            except Exception: pass

            # 4) FF5F sheet (optional transparency) — AUD-adjusted aggregate (US factors)
            # used downstream for factor moments / tilt-target μ + Σ.
            ff5s = get_or_clear_sheet(wb, 'FF5F')
            ff5s.range('A1').options(pd.DataFrame, index=True, header=True).value = ff_aud

            # FF5F_Regional sheet (audit) — the actual factor matrices each
            # security was regressed against. Stacked long-format with a Region
            # marker column so you can sort/filter to see only one region.
            try:
                _ff5_regional = globals().get("ff5_regional_windows", None)
                if isinstance(_ff5_regional, dict) and _ff5_regional:
                    _parts = []
                    for _region, _df_r in _ff5_regional.items():
                        if _df_r is None or _df_r.empty:
                            continue
                        _tmp = _df_r.copy().reset_index().rename(columns={"index": "Date"})
                        _tmp.insert(0, "Region", _region)
                        _parts.append(_tmp)
                    if _parts:
                        _audit = pd.concat(_parts, ignore_index=True)
                        ff5r = get_or_clear_sheet(wb, 'FF5F_Regional')
                        ff5r.range('A1').options(pd.DataFrame, index=False, header=True).value = _audit
            except Exception as _e_ff5_reg:
                print(f"[excel] FF5F_Regional write skipped: {_e_ff5_reg}")

            # Regression_Diagnostics sheet — per-security FF5+MOM regression stats
            # for the user to sanity-check fit quality. Annualised alpha and resid σ
            # are added for human-readable comparison against the asset μ values.
            try:
                _stats = globals().get("ff5_regression_stats", None)
                if isinstance(_stats, pd.DataFrame) and not _stats.empty:
                    _stats_out = _stats.copy()
                    # Annualise the daily figures (compounding ignored — small-return regime).
                    if "alpha_daily" in _stats_out.columns:
                        _stats_out["alpha_annual"] = _stats_out["alpha_daily"] * TRADING_DAYS
                    if "resid_std_daily" in _stats_out.columns:
                        _stats_out["resid_std_annual"] = _stats_out["resid_std_daily"] * np.sqrt(TRADING_DAYS)
                    # Friendly column order: bookkeeping first, fit quality, then per-factor t-stats,
                    # then alpha block, then residual block.
                    _factor_t_cols = [c for c in _stats_out.columns if c.endswith("_t") and c != "alpha_t"]
                    _front = [c for c in ["Region", "Standardised", "N obs", "R^2", "R^2 adj"] if c in _stats_out.columns]
                    _alpha_block = [c for c in ["alpha_daily", "alpha_annual", "alpha_t"] if c in _stats_out.columns]
                    _resid_block = [c for c in ["resid_std_daily", "resid_std_annual"] if c in _stats_out.columns]
                    _ordered = _front + _factor_t_cols + _alpha_block + _resid_block
                    _stats_out = _stats_out[[c for c in _ordered if c in _stats_out.columns]]

                    diag = get_or_clear_sheet(wb, 'Regression_Diagnostics')
                    diag.range('A1').options(pd.DataFrame, index=True, header=True).value = _stats_out
                    # Number formats: ratios 3dp, t-stats 2dp, alphas 2dp/2dp%.
                    n_rows = _stats_out.shape[0] + 1  # +1 for header row
                    _fmts = {}
                    for col_name in ["R^2", "R^2 adj"]:
                        if col_name in _stats_out.columns:
                            col_letter = chr(ord("B") + list(_stats_out.columns).index(col_name))
                            _fmts[f"{col_letter}2:{col_letter}{n_rows}"] = "0.000"
                    for col_name in _factor_t_cols + (["alpha_t"] if "alpha_t" in _stats_out.columns else []):
                        col_letter = chr(ord("B") + list(_stats_out.columns).index(col_name))
                        _fmts[f"{col_letter}2:{col_letter}{n_rows}"] = "0.00"
                    if "alpha_annual" in _stats_out.columns:
                        col_letter = chr(ord("B") + list(_stats_out.columns).index("alpha_annual"))
                        _fmts[f"{col_letter}2:{col_letter}{n_rows}"] = "0.00%"
                    if "resid_std_annual" in _stats_out.columns:
                        col_letter = chr(ord("B") + list(_stats_out.columns).index("resid_std_annual"))
                        _fmts[f"{col_letter}2:{col_letter}{n_rows}"] = "0.00%"
                    set_number_formats(diag, _fmts)
                    diag.autofit()
                    print(f"[excel] Regression_Diagnostics: {len(_stats_out)} securities written")
            except Exception as _e_diag:
                print(f"[excel] Regression_Diagnostics write skipped: {_e_diag}")

            # OOS_Validation sheet — walk-forward backtest metrics across 3y/5y/10y.
            try:
                _oos_tbl = globals().get("oos_metrics_table", None)
                if isinstance(_oos_tbl, pd.DataFrame) and not _oos_tbl.empty:
                    # Flatten MultiIndex columns to "(horizon, series)" strings for xlwings.
                    _oos_flat = _oos_tbl.copy()
                    _oos_flat.columns = [f"{h} {s}" for h, s in _oos_tbl.columns]
                    oos_sht = get_or_clear_sheet(wb, 'OOS_Validation')
                    oos_sht.range('A1').options(pd.DataFrame, index=True, header=True).value = _oos_flat

                    # Format metrics by row (percentages vs ratios).
                    pct_rows = {"Cumulative Return", "Annualised Return", "Annualised Volatility",
                                "Max Drawdown", "Alpha vs SPY (ann)", "Alpha vs FF5 (ann)",
                                "Annual Turnover"}
                    ratio_rows = {"Sharpe Ratio", "Sortino Ratio", "IR vs ^AORD", "Beta vs SPY"}
                    n_cols = len(_oos_flat.columns)
                    _fmts_oos = {}
                    for i, metric in enumerate(_oos_flat.index):
                        row_excel = i + 2  # +1 for header, +1 for 1-indexed
                        end_letter = chr(ord("A") + n_cols)  # last data column letter
                        rng = f"B{row_excel}:{end_letter}{row_excel}"
                        if metric in pct_rows:
                            _fmts_oos[rng] = "0.00%"
                        elif metric in ratio_rows:
                            _fmts_oos[rng] = "0.00"
                    set_number_formats(oos_sht, _fmts_oos)
                    oos_sht.autofit()
                    print(f"[excel] OOS_Validation: {_oos_flat.shape[0]} metrics × {_oos_flat.shape[1]} cols written")
            except Exception as _e_oos:
                print(f"[excel] OOS_Validation write skipped: {_e_oos}")

            # TLH_Log sheet — every harvesting swap fired by the walk-forward
            # engine across the OOS window. Top rows summarise; below is the
            # per-event ledger. Empty sheet if TLH disabled or no events.
            try:
                _tlh_events = globals().get("oos_tlh_events", []) or []
                _tlh_sht = get_or_clear_sheet(wb, 'TLH_Log')
                if _tlh_events:
                    _tlh_df = pd.DataFrame(_tlh_events)
                    # Coerce dates to date-only strings for clean Excel display.
                    for _dc in ("date", "lot_date"):
                        if _dc in _tlh_df.columns:
                            _tlh_df[_dc] = pd.to_datetime(_tlh_df[_dc]).dt.strftime("%Y-%m-%d")
                    _tlh_loss_total = float(_tlh_df["loss_aud"].sum())
                    _eff_st = _effective_cgt_rate(short_term=True)
                    _eff_lt = _effective_cgt_rate(short_term=False)
                    _tax_saved_est = _tlh_loss_total * (_eff_st + _eff_lt) / 2.0
                    _years_tlh = max(len(oos_returns_daily) / ANNUAL_TRADING_DAYS, 1e-6) \
                                 if isinstance(oos_returns_daily, pd.Series) else 1.0
                    _bps_yr = (_tax_saved_est
                               / float(globals().get("_oos_starting_nav_aud") or 1_000_000.0)
                               / _years_tlh * 10_000)
                    # Header summary rows
                    _tlh_sht.range("A1").value = "Tax-Loss Harvesting — backtest engine activity"
                    _tlh_sht.range("A2").value = "TLH events"
                    _tlh_sht.range("B2").value = int(len(_tlh_events))
                    _tlh_sht.range("A3").value = "Total loss realised (AUD)"
                    _tlh_sht.range("B3").value = round(_tlh_loss_total, 2)
                    _tlh_sht.range("A4").value = "Tax saved (gross est, AUD)"
                    _tlh_sht.range("B4").value = round(_tax_saved_est, 2)
                    _tlh_sht.range("A5").value = "Drag offset (bps/yr, gross)"
                    _tlh_sht.range("B5").value = round(_bps_yr, 1)
                    _tlh_sht.range("A6").value = ("Note: gross figure assumes 100% loss utilisation. "
                                                  "Net Sharpe impact is typically smaller because FY-end "
                                                  "netting already offsets most realised gains.")
                    # Per-event detail starting row 8
                    _tlh_sht.range("A8").options(pd.DataFrame, index=False, header=True).value = _tlh_df
                    _tlh_sht.autofit()
                    print(f"[excel] TLH_Log: {len(_tlh_events)} events written "
                          f"(${_tlh_loss_total:,.0f} loss, ~${_tax_saved_est:,.0f} gross tax saved)")
                else:
                    _tlh_sht.range("A1").value = "Tax-Loss Harvesting — backtest engine activity"
                    _tlh_sht.range("A2").value = ("No TLH events triggered over the OOS window "
                                                  if TLH_ENABLED else "TLH disabled in config ")
                    _tlh_sht.range("A3").value = (f"(threshold={TLH_MIN_LOSS_PCT*100:+.0f}%, "
                                                  f"min=${TLH_MIN_LOSS_AUD:.0f}, "
                                                  f"cooldown={TLH_COOLDOWN_DAYS}d, "
                                                  f"pairs={len(TLH_PAIRS)})")
                    _tlh_sht.autofit()
                    print(f"[excel] TLH_Log: 0 events (TLH_ENABLED={TLH_ENABLED})")
            except Exception as _e_tlh:
                print(f"[excel] TLH_Log write skipped: {_e_tlh}")

            # ---- FY Tax Ledger sheet (ACTUAL fills, per financial year) ----
            # Reconciles live NAV against the backtest's assumption that CGT
            # is paid (and TLH savings reinvested) INSIDE the portfolio. In
            # reality tax settles at lodgement in the user's bank account.
            # This ledger shows, per AU FY: realised ST/LT gains, losses,
            # carry-forward chain, and the estimated CGT at lodgement — so
            # the user knows what to contribute back (refund) or expect to
            # owe, keeping live NAV comparable to the simulation. Derived
            # fresh from ibkr_fills_log.jsonl + lots_seed.json every run.
            try:
                _fy_ledger_df = compute_fy_tax_ledger(
                    APP_DIR / "ibkr_fills_log.jsonl",
                    seed_path=APP_DIR / "lots_seed.json",
                    fx_map=fx_map_all,
                    lot_match_method=LOT_MATCH_METHOD,
                )
                globals()["FY_TAX_LEDGER_DF"] = _fy_ledger_df
                _tax_sht = get_or_clear_sheet(wb, 'Tax_Ledger')
                _tax_sht.range("A1").value = ("FY Tax Ledger — ACTUAL fills (broker truth), "
                                              f"profile={ACTIVE_CGT_PROFILE}")
                _tax_sht.range("A2").value = ("At lodgement: contribute any refund back into the "
                                              "portfolio (or fund the bill externally) so live NAV "
                                              "stays comparable to the backtest, which models tax "
                                              "flows inside the portfolio.")
                if not _fy_ledger_df.empty:
                    _tax_sht.range("A4").options(pd.DataFrame, index=False, header=True).value = _fy_ledger_df
                    print(f"[excel] Tax_Ledger: {len(_fy_ledger_df)} FY row(s) written "
                          f"(current: {_fy_ledger_df.iloc[-1]['FY']}, "
                          f"CGT ${float(_fy_ledger_df.iloc[-1]['CGT at Lodgement (AUD)']):,.0f}, "
                          f"c/f ${float(_fy_ledger_df.iloc[-1]['Carry-Fwd Out']):,.0f})")
                else:
                    _tax_sht.range("A4").value = "No sell fills recorded yet — ledger starts with the first executed sell."
                    print("[excel] Tax_Ledger: no sell fills yet (sheet stubbed)")
                _tax_sht.autofit()
            except Exception as _e_taxled:
                print(f"[excel] Tax_Ledger write skipped: {_e_taxled}")

            # ---- IBKR Actual Fills sheet (from ibkr_fills_log.jsonl) ----
            # Surfaces the Phase 3 paper-execution log inside the workbook so
            # the user can reconcile what was actually filled (broker truth)
            # against what the engine recommended. Refreshes every engine run;
            # safe to be missing (shows "no log entries yet").
            try:
                _fills_path = APP_DIR / "ibkr_fills_log.jsonl"
                _fills_sht = get_or_clear_sheet(wb, 'Actual_Fills')
                _fills_rows = []
                if _fills_path.exists():
                    with open(_fills_path, "r", encoding="utf-8") as _fh:
                        for _line in _fh:
                            _line = _line.strip()
                            if not _line:
                                continue
                            try:
                                _fills_rows.append(json.loads(_line))
                            except json.JSONDecodeError:
                                continue
                if _fills_rows:
                    _fills_rows.sort(key=lambda r: r.get("exec_timestamp", ""),
                                       reverse=True)
                    _latest_ts = _fills_rows[0].get("exec_timestamp", "?")
                    _latest_batch = [r for r in _fills_rows
                                       if r.get("exec_timestamp") == _latest_ts]
                    _n_filled = sum(1 for r in _latest_batch
                                       if r.get("status_final") == "Filled"
                                       or r.get("status") == "Filled")
                    _n_cancelled = sum(1 for r in _latest_batch
                                         if r.get("status_final") == "Cancelled"
                                         or r.get("status") == "Cancelled")
                    _n_pending = sum(1 for r in _latest_batch
                                       if not r.get("is_done", False))
                    _fills_sht.range("A1").value = (
                        "IBKR Actual Fills — Phase 3 paper-trade execution log"
                    )
                    _fills_sht.range("A2").value = (
                        f"Source: ibkr_fills_log.jsonl  ·  Total rows: {len(_fills_rows)}"
                    )
                    _fills_sht.range("A3").value = (
                        f"Most recent batch: {_latest_ts}"
                    )
                    _fills_sht.range("A4").value = (
                        f"  Submitted: {len(_latest_batch)}  ·  "
                        f"Filled: {_n_filled}  ·  "
                        f"Cancelled: {_n_cancelled}  ·  "
                        f"Pending: {_n_pending}"
                    )
                    _fills_sht.range("A5").value = (
                        "Note: fills_log captures script-side state at write time. "
                        "For broker truth use: ibkr_paper_exec.py --check-fills"
                    )
                    _fills_df = pd.DataFrame([{
                        "Exec TS":       r.get("exec_timestamp", ""),
                        "Rec Run TS":    r.get("rec_log_run_at", ""),
                        "Ticker":        r.get("ticker", ""),
                        "Side":          r.get("side", ""),
                        "Qty Req":       r.get("qty_requested", 0),
                        "Qty Filled":    r.get("qty_filled", 0),
                        "Qty Remaining": r.get("qty_remaining", 0),
                        "Avg Fill Px":   r.get("avg_fill_price_local", None),
                        "Rec Px (AUD)":  r.get("rec_px_aud", None),
                        "Status":        (r.get("status_final")
                                            or r.get("status") or "?"),
                        "Done":          r.get("is_done", False),
                        "OrderId":       r.get("order_id", 0),
                        "PermId":        r.get("ibkr_perm_id", 0),
                        "N Fills":       r.get("n_fills", 0),
                    } for r in _fills_rows])
                    _fills_sht.range("A7").options(pd.DataFrame, index=False,
                                                     header=True).value = _fills_df
                    _fills_sht.autofit()
                    print(f"[excel] Actual_Fills: {len(_fills_rows)} rows "
                          f"(latest batch {_latest_ts}: "
                          f"{_n_filled}F/{_n_cancelled}C/{_n_pending}P)")
                elif _fills_path.exists():
                    _fills_sht.range("A1").value = "IBKR Actual Fills — log empty"
                    _fills_sht.range("A2").value = (
                        "ibkr_fills_log.jsonl exists but has no rows. Run "
                        "ibkr_paper_exec.py --execute to populate."
                    )
                    print("[excel] Actual_Fills: log file exists but empty")
                else:
                    _fills_sht.range("A1").value = "IBKR Actual Fills — log not found"
                    _fills_sht.range("A2").value = (
                        "ibkr_fills_log.jsonl does not exist yet. Run "
                        "ibkr_paper_exec.py --execute to start populating; "
                        "this sheet refreshes on every engine run."
                    )
                    print("[excel] Actual_Fills: log file does not exist yet")
            except Exception as _e_fills:
                print(f"[excel] Actual_Fills write skipped: {_e_fills}")

            # ---- Rebuild Lots from authoritative source ---------------------
            # 'fills' (default): only count confirmed IBKR fills. Empty until
            #   the first real fill — exactly the right baseline for a paper
            #   account whose past 12 orders all Cancelled.
            # 'holdings': single lot per held ticker at today's AUD price.
            #   Plumbed for brokers without per-fill export. Loses CGT
            #   history — use only when no fills log is available.
            #
            # Replaces the legacy `_update_lots_after_trades(lots_df, trade_rec, ...)`
            # path which wrote the engine's *recommended* trades as if they'd
            # filled. After many runs that inflated SMH lots to 3.4M units
            # and triggered the -1.7M/+1.7M trade plan corruption (2026-06).
            _mode = str(globals().get("LOTS_REBUILD_MODE", "fills")).lower()
            _fills_path = APP_DIR / "ibkr_fills_log.jsonl"
            _seed_path = APP_DIR / "lots_seed.json"
            if _mode == "holdings":
                UPDATED_LOTS = _build_lots_from_holdings(
                    units, last_px_hold, today=pd.Timestamp(prices.index[-1]),
                )
                print(f"[lots] rebuilt from Holdings ({len(UPDATED_LOTS)} positions, "
                      f"CGT-naive baseline)")
            else:
                UPDATED_LOTS = _build_lots_from_fills_log(
                    _fills_path, fx_map=fx_map_all,
                    lot_match_method=LOT_MATCH_METHOD,
                    seed_path=_seed_path,
                )
                _seed_note = " (with seed)" if _seed_path.exists() else ""
                print(f"[lots] rebuilt from {_fills_path.name}{_seed_note} "
                      f"({len(UPDATED_LOTS)} lots)")

            sht_lots = get_or_clear_sheet(wb, 'Lots')
            sht_lots.range("A1").value = [["Security","AcqDate","Units","CostBaseAUD"]]
            if not UPDATED_LOTS.empty:
                sht_lots.range("A2").options(index=False, header=False).value = UPDATED_LOTS
            
            tgt_units_full = compute_target_units_for_holdings(
                units, last_px_hold, fx_map_all, w_star, include_flags,
                portfolio_value_override=portfolio_value_override,
                available_cash_aud=globals().get("_avail_cash_aud"))

            # Holdings reconciliation fix (2026-06-27): pass the engine's
            # CURRENT units, NOT the new target. Writing target back to
            # Holdings.Units created a self-referential loop where the
            # engine read its own previous target as next run's "current".
            # The HBRD=17,122 anomaly on 2026-06-27 originated here:
            # user's actual paper position was HBRD=2,034 (per 2026-06-26
            # triage), but successive engine runs each wrote larger
            # target_units to Holdings until "current" no longer matched
            # broker reality.
            #
            # Holdings.Units is now immutable from the engine's perspective.
            # User updates it via:
            #   - triage_reset_*.py for initial seeding
            #   - manual edit after broker fills
            #   - (future) ibkr_paper_exec.py --reconcile flag
            #
            # Prices / Market Value / Weight still get refreshed here
            # because those are derivable from current units + current
            # market data; only Units stays sticky.
            _write_holdings_sheet(wb, prices, units, include_flags, sheet_name="Holdings", fx_to_aud_map=fx_map_all)

            # --- Step 1: Compute current portfolio values ---
            if not trade_rec.empty:
                trade_rec["Target Units"] = pd.to_numeric(trade_rec["Target Units"], errors="coerce").fillna(0.0)
                trade_rec["Last Px (AUD)"] = pd.to_numeric(trade_rec["Last Px (AUD)"], errors="coerce").fillna(0.0)
                trade_rec["Value"] = trade_rec["Target Units"] * trade_rec["Last Px (AUD)"]
                net_invested = float(trade_rec["Value"].sum())
                # Net cash after trades (already net of brokerage)
                cash_balance = float(pd.to_numeric(trade_rec["Cash Flow (AUD)"], errors="coerce").fillna(0.0).sum())
            else:
                net_invested = 0.0
                cash_balance = 0.0
        
            # Brokerage (for reporting)
            total_brokerage = float(costs_rec.get("brokerage", 0.0))
            
            # Total portfolio + cash:
            # If the user provided a portfolio value override, treat it as TOTAL portfolio value.
            # Cash becomes the residual AFTER funding target holdings AND paying brokerage.
            pvo = None
            try:
                pvo = float(portfolio_value_override) if portfolio_value_override is not None else None
            except Exception:
                pvo = None
            
            if pvo is not None and np.isfinite(pvo) and pvo > 0:
                total_portfolio = float(pvo)
                cash_balance = total_portfolio - float(net_invested) - float(total_brokerage)
            else:
                # Fall back to "cash from trade cashflows" (your current behaviour)
                # NOTE: if your Cash Flow already includes brokerage, then total_portfolio should be:
                # holdings + cash
                total_portfolio = float(net_invested) + float(cash_balance)
            
            print(f"[debug] Current totals â†’ Portfolio: {total_portfolio:.2f}, Net Invested: {net_invested:.2f}")
            
            # --- Step 2: Load previous run data (AFTER calculating current totals) ---
            if os.path.exists(state_path):
                with open(state_path, "r") as f:
                    prev_state = json.load(f)
                previous_portfolio = prev_state.get("portfolio_value", 0.0)
                previous_invested = prev_state.get("net_invested", 0.0)
                print(f"[debug] Previous totals â†’ Portfolio: {previous_portfolio:.2f}, Net Invested: {previous_invested:.2f}")
            else:
                previous_portfolio = 0.0
                previous_invested = 0.0
                print("[info] No previous state file found â€” starting fresh deltas at 0.")
            
            # --- Step 3: Compute deltas for PowerPoint ---
            _bkd_for_cgt = costs_rec.get("breakdown", {}) if "costs_rec" in globals() else {}
            _mtr_for_cgt = float(CGT_CONFIG.get("marginal_tax_rate", 0.30))
            results = {
                "total_brokerage": total_brokerage,
                "net_invested": net_invested,
                "total_portfolio_value": total_portfolio,
                "portfolio_change": total_portfolio - previous_portfolio,
                "net_invested_change": net_invested - previous_invested,
                "cash_balance": cash_balance,
                # CGT this rebalance (positive = tax owed on realised gains).
                "total_cgt": float(costs_rec.get("cgt_tax", 0.0)),
                # Unused losses available to carry forward × MTR = tax saved.
                "loss_carry_forward_tax_aud": (
                    float(_bkd_for_cgt.get("loss_carry_forward", 0.0)) * _mtr_for_cgt
                ),
            }
            
            # --- Step 4: Save current state for next comparison ---
            with open(state_path, "w") as f:
                json.dump(
                    {"portfolio_value": total_portfolio, "net_invested": net_invested},
                    f,
                    indent=2
                )
            print(f"[debug] Saved new state â†’ Portfolio: {total_portfolio:.2f}, Net Invested: {net_invested:.2f}")
            
            # --- Step 5: Generate PowerPoint summary ---
            trades = trade_rec.copy()
            # --- Label which portfolio this trade plan represents (used later by PPT) ---
            globals()["TRADE_PLAN_PORTFOLIO_LABEL"] = str(globals().get("choice_label", globals().get("TRADE_PLAN_MODE", "unknown")))
            globals()["TRADE_PLAN_SOURCE"] = "trade_rec"
           
            charts = dict(globals().get("charts", {}) or {})
            if ("tilts_comparison_rows" not in charts) or (not charts.get("tilts_comparison_rows")):
                try:
                    if "out" in locals() and isinstance(out, pd.DataFrame) and (not out.empty):
                        charts["tilts_comparison_rows"] = (
                            out.reset_index()
                            .rename(columns={"index": "Factor"})
                            .to_dict("records")
                        )
                except Exception:
                    pass
            charts.pop("tilts_comparison_rows", None)
            charts.pop("with_tilts_achieved_tilts", None)     

            # Persist report payload for downstream launcher cells
            globals()["results"] = results
            globals()["trades"] = trades
            globals()["charts"] = charts

            # --- Step 6: Compute PortfolioValue for PowerPoint charts (no Excel readback) ---
            try:
                # Use the in-memory target units you already computed
                units_ser = pd.to_numeric(pd.Series(tgt_units_full), errors="coerce").fillna(0.0)
                valid_tickers = [t for t in units_ser.index.astype(str) if t in prices.columns]
            
                if not valid_tickers:
                    raise ValueError("No valid tickers found in prices for target holdings.")
            
                port_prices = prices[valid_tickers].copy().ffill().bfill()
                u = units_ser.reindex(valid_tickers).astype(float).values
                portfolio_value_series = (port_prices * u).sum(axis=1).ffill().bfill()
                portfolio_value_series = portfolio_value_series.copy()
                
                print(f"[pptx prep] PortfolioValue series computed for {len(valid_tickers)} securities.")
                # Diagnostic: confirm the date range that downstream slides will anchor to.
                # If this end date is much earlier than today, something upstream is
                # truncating `prices` (yfinance cache, dropna, reindex, etc.).
                try:
                    _pv_idx = portfolio_value_series.dropna().index
                    print(f"[pptx prep] portfolio_value_series spans "
                          f"{_pv_idx.min().date()} -> {_pv_idx.max().date()} "
                          f"({len(_pv_idx)} rows); prices.index.max()={prices.index.max().date()}")
                except Exception:
                    pass
            except Exception as e:
                print(f"[pptx prep] Could not compute PortfolioValue: {e}")
         
                # Rebuild tilts rows if missing (prevents Slide 5 table disappearing)
                if ("tilts_comparison_rows" not in charts) or (not charts.get("tilts_comparison_rows")):
                    charts["tilts_comparison_rows"] = (globals().get("charts", {}) or {}).get("tilts_comparison_rows", [])
                    print("[pptx prep] tilts_comparison_rows length:", len(charts.get("tilts_comparison_rows") or []))

                try:
                    ppt_path = export_to_ppt(results, trades, charts)
                except Exception as e:
                    print(f"[pptx] Skipped PowerPoint generation: {e}")

            wb.save()
            wb.close()

    except SanityViolation:
        # NEVER swallow a sanity violation here. The whole point of the
        # safety layer is to halt the engine BEFORE side effects ship.
        # If this catch were broad enough to consume SanityViolation,
        # the engine would continue past the violation and generate
        # the PPT / write state / submit orders — exactly the failure
        # mode the layer exists to prevent. Bug surfaced 2026-06-27:
        # an Excel COM exception handler was catching SanityViolation
        # and treating it as "Excel had a hiccup, fall back to CSVs".
        # The engine then completed the run as if nothing was wrong.
        # Re-raising here propagates to the top-level handler, which
        # exits the run cleanly with a clear error.
        raise
    except Exception as e:
        import traceback as _tb_xl
        print(f"[Excel fallback] xlwings/COM error â†’ exporting CSVs instead: {e}")
        print(f"[Excel fallback] full traceback:\n{_tb_xl.format_exc()}")
        export_dir = os.path.join(os.path.dirname(filename), "Exports")
        try: os.makedirs(export_dir, exist_ok=True)
        except Exception: pass
        try: exp_ret_df.to_csv(os.path.join(export_dir, "expected_returns.csv"))
        except Exception as ee: print(f"[export] expected_returns.csv: {ee}")
        try: cov_plus.to_csv(os.path.join(export_dir, "covariance_plus.csv"))
        except Exception as ee: print(f"[export] covariance_plus.csv: {ee}")
        try: W.to_csv(os.path.join(export_dir, "weights_grid.csv"))
        except Exception as ee: print(f"[export] weights_grid.csv: {ee}")
        try: stats_df.to_csv(os.path.join(export_dir, "portfolio_stats.csv"), index=False)
        except Exception as ee: print(f"[export] portfolio_stats.csv: {ee}")
        try: tilt_df.to_csv(os.path.join(export_dir, "tilts.csv"))
        except Exception as ee: print(f"[export] tilts.csv: {ee}")
        try: df_melt.to_csv(os.path.join(export_dir, "returns_long.csv"), index=False)
        except Exception as ee: print(f"[export] returns_long.csv: {ee}")
else:
    # ---------- Headless fallback: write key outputs as CSVs ----------
    export_dir = os.path.join(os.path.dirname(filename), "Exports")
    try: os.makedirs(export_dir, exist_ok=True)
    except Exception: pass
    try: exp_ret_df.to_csv(os.path.join(export_dir, "expected_returns.csv"))
    except Exception as e: print(f"[export] expected_returns.csv: {e}")
    try: cov_plus.to_csv(os.path.join(export_dir, "covariance_plus.csv"))
    except Exception as e: print(f"[export] covariance_plus.csv: {e}")
    try: W.to_csv(os.path.join(export_dir, "weights_grid.csv"))
    except Exception as e: print(f"[export] weights_grid.csv: {e}")
    try: stats_df.to_csv(os.path.join(export_dir, "portfolio_stats.csv"), index=False)
    except Exception as e: print(f"[export] portfolio_stats.csv: {e}")
    try: tilt_df.to_csv(os.path.join(export_dir, "tilts.csv"))
    except Exception as e: print(f"[export] tilts.csv: {e}")
    try: df_melt.to_csv(os.path.join(export_dir, "returns_long.csv"), index=False)
    except Exception as e: print(f"[export] returns_long.csv: {e}")

# Ensure report payload exists even if Excel/COM path was skipped or failed
if "results" not in globals():
    _net_invested = 0.0
    _cash_balance = 0.0
    if "trade_rec" in globals() and isinstance(trade_rec, pd.DataFrame) and not trade_rec.empty:
        _tgt = pd.to_numeric(trade_rec.get("Target Units"), errors="coerce").fillna(0.0)
        _px = pd.to_numeric(trade_rec.get("Last Px (AUD)"), errors="coerce").fillna(0.0)
        _net_invested = float((_tgt * _px).sum())
        _cash_balance = float(pd.to_numeric(trade_rec.get("Cash Flow (AUD)"), errors="coerce").fillna(0.0).sum())
    _total_brokerage = float(costs_rec.get("brokerage", 0.0)) if "costs_rec" in globals() else 0.0
    _total_portfolio = float(_net_invested + _cash_balance)
    _bkd_fb = costs_rec.get("breakdown", {}) if "costs_rec" in globals() else {}
    _mtr_fb = float(CGT_CONFIG.get("marginal_tax_rate", 0.30))
    results = {
        "total_brokerage": _total_brokerage,
        "net_invested": _net_invested,
        "total_portfolio_value": _total_portfolio,
        "portfolio_change": 0.0,
        "net_invested_change": 0.0,
        "cash_balance": _cash_balance,
        "total_cgt": float(costs_rec.get("cgt_tax", 0.0)) if "costs_rec" in globals() else 0.0,
        "loss_carry_forward_tax_aud": (
            float(_bkd_fb.get("loss_carry_forward", 0.0)) * _mtr_fb
        ),
    }
if "trades" not in globals():
    trades = trade_rec.copy() if "trade_rec" in globals() and isinstance(trade_rec, pd.DataFrame) else pd.DataFrame()
if "charts" not in globals() or not isinstance(charts, dict):
    charts = dict(globals().get("charts", {}) or {})
if isinstance(charts, dict) and ("portfolio_value_series" not in charts or charts.get("portfolio_value_series") is None):
    try:
        _pv_series = None
        if "prices" in globals() and isinstance(prices, pd.DataFrame) and not prices.empty and isinstance(trades, pd.DataFrame) and not trades.empty:
            _sec = trades.get("Security")
            _tgt = pd.to_numeric(trades.get("Target Units"), errors="coerce").fillna(0.0)
            if _sec is not None:
                _u = pd.Series(_tgt.values, index=_sec.astype(str)).groupby(level=0).sum()
                _valid = [t for t in _u.index if t in prices.columns]
                if _valid:
                    _px = prices[_valid].copy().ffill().bfill()
                    _vals = _u.reindex(_valid).astype(float).values
                    _pv_series = (_px * _vals).sum(axis=1).ffill().bfill()
        if _pv_series is not None and len(_pv_series) > 0:
            charts["portfolio_value_series"] = _pv_series
            globals()["portfolio_value_series"] = _pv_series
    except Exception:
        pass
globals()["results"] = results
globals()["trades"] = trades
globals()["charts"] = charts

print("Workbook Successfully Updated")

# --- Create a Desktop shortcut (optional, safe in any context) ---
try:
    if HAS_WIN32COM:
        shortcut_path = str(Path.home() / "Desktop" / "Portfolio Optimiser.lnk")

        # Prefer the exe if it exists; otherwise point at the script weâ€™re running.
        # Works when frozen, when run as .py, and in Jupyter (falls back to .py name in APP_DIR).
        if getattr(sys, "frozen", False):
            target = Path(sys.executable)
        else:
            # Try the current file if available; else fall back to a known script name in this folder
            if "__file__" in globals():
                target = Path(__file__).resolve()
            else:
                # Adjust the name if your launcher script is 'Main.py' instead
                # (You have both Main.py and Portfolio_Optimiser3110.py in your screenshot.)
                candidate = APP_DIR / "Portfolio_Optimiser1411.py"
                target = candidate if candidate.exists() else (APP_DIR / "Main.py")

        shell = win32.Dispatch("WScript.Shell")
        sc = shell.CreateShortCut(shortcut_path)
        sc.WindowStyle = 1  # normal window
        sc.Arguments = ""   # no extra args      
        sc.Targetpath = str(target)
        sc.WorkingDirectory = str(target.parent)
        # Use icon.ico if present; otherwise the target itself
        icon_path = APP_DIR / "icon.ico"
        sc.IconLocation = str(icon_path if icon_path.exists() else target)
        sc.save()
    else:
        print("[shortcut] pywin32 not available; skipping Desktop shortcut.")
except Exception as e:
    print(f"[shortcut] skipped due to error: {e}")
print("=== MU VEC (sorted) ===")
print(mu_vec_opt.sort_values())
print("\nMin mu:", mu_vec_opt.min())
print("Max mu:", mu_vec_opt.max())
print("Mean mu:", mu_vec_opt.mean())
print("Top 10 assets by expected return:")
print(mu_vec_opt.sort_values().tail(10))
print(ff5_raw.head())
# Post-write OPT annotation + alternative plan placement + shortcut repair.
from datetime import datetime
from pathlib import Path

try:
    _xl = str(globals().get("filename", "")).strip()
    _diag = dict(globals().get("TRADEPLAN_VALIDATION_DIAG", {}) or {})
    _mode = str(_diag.get("mode", globals().get("TRADE_PLAN_MODE", ""))).strip().lower()
    _selected = str(_diag.get("selected", globals().get("TRADEPLAN_LABEL", ""))).strip()
    _lookback = _diag.get("lookback_days", globals().get("VALIDATION_LOOKBACK_DAYS", 252))
    _sh0 = _diag.get("sharpe_no_tilts", np.nan)
    _sh1 = _diag.get("sharpe_with_tilts", np.nan)

    if not _xl or not os.path.exists(_xl):
        print("[post] Skipped OPT post-write fixes: workbook path not available.")
    else:
        _rows = [
            ("Mode", _mode),
            ("Lookback (days)", _lookback),
            ("Selected Portfolio", _selected),
            ("Sharpe (Optimised)", _sh0),
            ("Sharpe (With Tilts)", _sh1),
        ]

        _written = False

        # Use a dedicated hidden xlwings App so Excel.exe terminates when the `with` exits.
        # (Bare xw.Book(path) attaches to the default app, which then stays alive as a blank
        # window after _book.close() — caused the spurious second Excel window.)
        try:
            with xw.App(visible=False, add_book=False) as _app:
                _book = _app.books.open(_xl, update_links=False, read_only=False)
                _ws = _book.sheets["OPT"] if "OPT" in [s.name for s in _book.sheets] else _book.sheets.add("OPT")

                _ws.range("W2").value = "Trade Plan Validation"
                _ws.range("W3").value = _rows
                _ws.range("X6:X7").api.NumberFormat = "0.000"

                _book.save()
                _book.close()
            _written = True
            print(f"[post] Wrote OPT validation + layout fixes to: {_xl}")
        except Exception as _e_xlw:
            print(f"[post] xlwings post-write fallback triggered: {_e_xlw}")

        # Fallback to openpyxl; if locked, save a timestamped copy.
        if not _written:
            _wbx = load_workbook(_xl, keep_vba=True)
            _ws = _wbx["OPT"] if "OPT" in _wbx.sheetnames else _wbx.create_sheet("OPT")
            _r0, _c0 = 2, 23  # W2
            _ws.cell(_r0, _c0, "Trade Plan Validation")
            for i, (k, v) in enumerate(_rows, start=1):
                _ws.cell(_r0 + i, _c0, k)
                _ws.cell(_r0 + i, _c0 + 1, v)

            try:
                _wbx.save(_xl)
                print(f"[post] Wrote Trade Plan Validation block to OPT sheet in: {_xl}")
            except PermissionError:
                _stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                _copy = _xl.replace(".xlsm", f"_validated_{_stamp}.xlsm")
                _wbx.save(_copy)
                print(f"[post] Workbook locked. Saved validated copy instead: {_copy}")

    # Repair desktop shortcut safely.
    try:
        if globals().get("HAS_WIN32COM", False):
            _shortcut = str(Path.home() / "Desktop" / "Portfolio Optimiser.lnk")
            _app_dir = Path(globals().get("APP_DIR", os.getcwd()))
            if getattr(sys, "frozen", False):
                _target = Path(sys.executable)
            elif "__file__" in globals():
                _target = Path(__file__).resolve()
            else:
                _cand = _app_dir / "Portfolio_Optimiser1411.py"
                _target = _cand if _cand.exists() else (_app_dir / "Main.py")

            _shell = win32.Dispatch("WScript.Shell")
            _sc = _shell.CreateShortCut(_shortcut)
            _sc.WindowStyle = 1
            _sc.Arguments = ""
            _sc.Targetpath = str(_target)
            _sc.WorkingDirectory = str(_target.parent)
            _icon = _app_dir / "icon.ico"
            _sc.IconLocation = str(_icon if _icon.exists() else _target)
            _sc.save()
            print(f"[post] Desktop shortcut refreshed: {_shortcut}")
    except Exception as _e_sc:
        print(f"[post] Shortcut refresh skipped: {_e_sc}")

except Exception as _e_post_val:
    print(f"[post] OPT validation annotation skipped: {_e_post_val}")


# =====================================================================
# --- Block 8: Finishers / Launchers Excel and Code ---
# =====================================================================
# --- Optional: auto-open outputs after save ---
OPEN_EXCEL_AFTER_SAVE = bool(globals().get("OPEN_EXCEL_AFTER_SAVE", CFG.get("open_excel_after_save", True)))
OPEN_PPT_AFTER_SAVE = bool(globals().get("OPEN_PPT_AFTER_SAVE", CFG.get("open_ppt_after_save", True)))

if bool(globals().get("OPEN_AFTER_SAVE", True)) and OPEN_EXCEL_AFTER_SAVE:
    _excel_path = str(filename) if "filename" in globals() else ""
    if _excel_path and os.path.exists(_excel_path):
        try:
            _os_open(_excel_path)
        except Exception as exc:
            print(f"[open] Could not open Excel workbook: {exc}")
    else:
        print(f"[open] Workbook not found, skipping open: {_excel_path}")


# =====================================================================
# --- BLOCK 9: PowerPoint Report Generator ---
# =====================================================================
# --- BLOCK 9: PowerPoint Report Generator ---

# -----------------------------------------------------------------
# PPTX helpers (module level so they're not redefined per export call)
# -----------------------------------------------------------------
# =====================================================================
# --- Block 10: Finishers / Launchers PPTX ---
# =====================================================================
# --- Block 10: Finishers / Launchers PPTX ---
ppt_path = None

def _wait_for_pptx_ready(path, timeout_s=10.0, stable_s=1.0, poll_s=0.2):
    """Wait for a valid, stable PPTX file on disk before opening."""
    import os
    import time
    import zipfile

    if not path:
        return False

    t0 = time.time()
    last_size = -1
    last_change = t0

    while (time.time() - t0) < timeout_s:
        try:
            if not os.path.exists(path):
                time.sleep(poll_s)
                continue

            size = os.path.getsize(path)
            if size != last_size:
                last_size = size
                last_change = time.time()

            if (time.time() - last_change) >= stable_s:
                with zipfile.ZipFile(path, "r") as zf:
                    if zf.testzip() is None:
                        return True
        except Exception:
            pass

        time.sleep(poll_s)

    return False

if CFG.get("generate_report", True):
    _results = globals().get("results")
    _trades = globals().get("trades")
    _charts = globals().get("charts")

    _missing = []
    if not isinstance(_results, dict):
        _missing.append("results")
    if _trades is None:
        _missing.append("trades")
    if _charts is None:
        _missing.append("charts")

    if _missing:
        print(f"[pptx] Skipping report generation. Missing data: {', '.join(_missing)}. Run Cell 15 first.")
        ppt_path = None
    else:
        if not isinstance(_charts, dict):
            _charts = dict(_charts or {})
            globals()["charts"] = _charts

        if isinstance(_trades, pd.DataFrame) and ("Brokerage" not in _trades.columns):
            _trades = _trades.copy()
            if "Brokerage (AUD)" in _trades.columns:
                _trades["Brokerage"] = pd.to_numeric(_trades["Brokerage (AUD)"], errors="coerce").fillna(0.0)
            else:
                _trades["Brokerage"] = 0.0
            globals()["trades"] = _trades

        try:
            _tilt_rows = _charts.get("tilts_comparison_rows") if isinstance(_charts, dict) else None

            # Rebuild from optimisation globals when charts payload is narrowed.
            if not _tilt_rows:
                try:
                    _factor_order = ["Mkt-RF", "SMB", "HML", "RMW", "CMA", "MOM"]
                    _B = globals().get("B")
                    _Sigma = globals().get("Sigma_opt")
                    _w_with = globals().get("w_star_with_tilts")
                    _w_without = globals().get("w_star")
                    _tilt_df = globals().get("tilt_df", None)

                    if isinstance(_B, pd.DataFrame) and (not _B.empty) and hasattr(_Sigma, "index") and (_w_with is not None) and (_w_without is not None):
                        def _norm_w_local(w, idx):
                            s = pd.Series(np.asarray(w, dtype=float).reshape(-1), index=idx).fillna(0.0)
                            tot = float(s.sum())
                            return (s / tot) if tot != 0 else s

                        _w_with_s = _norm_w_local(_w_with, _Sigma.index).reindex(_B.index).fillna(0.0)
                        _w_without_s = _norm_w_local(_w_without, _Sigma.index).reindex(_B.index).fillna(0.0)
                        _with_beta = (_B.T @ _w_with_s)
                        _without_beta = (_B.T @ _w_without_s)

                        _out = pd.DataFrame(index=[f for f in _factor_order if f in _with_beta.index])
                        _out["With Tilts"] = _with_beta.reindex(_out.index).astype(float)
                        _out["Without Tilts"] = _without_beta.reindex(_out.index).astype(float)

                        if isinstance(_tilt_df, pd.DataFrame) and (not _tilt_df.empty) and ("Target" in _tilt_df.columns):
                            _tgt = _tilt_df.reindex(_out.index)
                            _out["Target"] = pd.to_numeric(_tgt["Target"], errors="coerce")
                            if "Use?" in _tgt.columns:
                                _use_mask = _tgt["Use?"].astype(bool)
                                _out = _out.loc[_use_mask.reindex(_out.index).fillna(False)]

                        _tilt_rows = _out.reset_index().rename(columns={"index": "Factor"}).to_dict("records")
                        if _tilt_rows:
                            _charts["tilts_comparison_rows"] = _tilt_rows
                except Exception as _e_tilt_rebuild:
                    print(f"[pptx] Tilt-row rebuild skipped: {_e_tilt_rebuild}")

            if not _tilt_rows and isinstance(globals().get("charts"), dict):
                _tilt_rows = globals().get("charts", {}).get("tilts_comparison_rows")
                if _tilt_rows:
                    _charts["tilts_comparison_rows"] = _tilt_rows
            if not _tilt_rows:
                _keys = list(_charts.keys()) if isinstance(_charts, dict) else []
                print(f"[pptx] Warning: missing charts['tilts_comparison_rows']; Slide 5 tilt table may be skipped. Available keys: {_keys}")
            ppt_path = export_to_ppt(_results, _trades, _charts)
        except Exception as exc:
            print(f"[pptx] Report generation failed: {exc}")
            ppt_path = None

if OPEN_PPT_AFTER_SAVE and ppt_path:
    if _wait_for_pptx_ready(ppt_path):
        open_ppt_if_enabled(ppt_path)
    else:
        print(f"[pptx] Saved but not opened automatically (file not stable yet): {ppt_path}")


# === RUN HEALTH SUMMARY ============================================
# One-shot block at the very end of the live pipeline. Designed so a
# silent failure (Excel skipped a sheet, metrics regressed, PPT slide
# missing, etc.) can't slip through unnoticed. Scans the run.log for
# [WARN] / [ERROR] / metrics-warn counts and reports them alongside
# the in-memory state.
def _print_run_health_summary():
    import os
    print()
    print("=" * 88)
    print("[health] === RUN HEALTH SUMMARY ===")
    print("=" * 88)
    # Runtime
    try:
        elapsed = _time_for_health.perf_counter() - _SCRIPT_START_TIME
        mins = int(elapsed // 60)
        secs = int(elapsed % 60)
        print(f"  Runtime:              {mins}m {secs:02d}s")
    except Exception:
        pass
    print(f"  Build:                {_BUILD_GIT_SHA} at {_BUILD_TIME}")
    _prod_label = (next(iter(PRODUCTION_SLOT_OVERRIDE.keys()))
                    if PRODUCTION_SLOT_OVERRIDE else "5-slot blend")
    print(f"  Production config:    slot={_prod_label}  "
          f"crash_hedge={'ON' if PRODUCTION_CRASH_HEDGE else 'off'}")

    # PPT
    if ppt_path and os.path.exists(ppt_path):
        try:
            from pptx import Presentation as _P
            _p = _P(ppt_path)
            print(f"  PPT generated:        OK ({len(_p.slides)} slides)")
        except Exception:
            print(f"  PPT generated:        OK")
    else:
        print(f"  PPT generated:        FAILED / not saved")

    # Excel
    try:
        _xlpath = globals().get("filename")
        if _xlpath and os.path.exists(_xlpath):
            mtime = os.path.getmtime(_xlpath)
            age_s = _time_for_health.time() - mtime
            if age_s < 120:
                print(f"  Excel workbook:       OK (updated {int(age_s)}s ago)")
            else:
                print(f"  Excel workbook:       WARNING (last modified {int(age_s/60)}m ago)")
    except Exception:
        pass

    # Metrics history
    try:
        _hist_path = APP_DIR / "metrics_history.jsonl"
        if _hist_path.exists():
            with _hist_path.open(encoding="utf-8") as f:
                _n_hist = sum(1 for line in f if line.strip())
            print(f"  Metrics snapshot:     OK ({_n_hist} total runs logged)")
    except Exception:
        pass

    # Live recommendation
    _w_live = globals().get("W_ENSEMBLE_SER", pd.Series(dtype=float))
    if isinstance(_w_live, pd.Series) and not _w_live.empty:
        print(f"  Live recommendation:  {len(_w_live)} positions")

    # TLH events
    _tlh = globals().get("oos_tlh_events", []) or []
    if _tlh:
        _loss = float(sum(e.get("loss_aud", 0.0) for e in _tlh))
        print(f"  TLH (backtest):       {len(_tlh)} events  (${_loss:,.0f} loss realised)")

    # Drift tracker last state
    try:
        _drift_summary = globals().get("DRIFT_LAST_SUMMARY", None)
        if isinstance(_drift_summary, str):
            print(f"  Drift tracker:        {_drift_summary}")
    except Exception:
        pass

    # Scan run.log for warnings + errors + metrics regressions
    try:
        base_dir = os.path.dirname(sys.executable) if getattr(sys, "frozen", False) \
            else os.path.dirname(os.path.abspath(__file__))
        log_path = os.path.join(base_dir, "run.log")
        if os.path.exists(log_path):
            with open(log_path, encoding="utf-8", errors="ignore") as f:
                log_text = f.read()
            n_warn = (log_text.count("[WARN") + log_text.count("[warn")
                      + log_text.count("WARNING") + log_text.count("Warning"))
            # subtract the count of the word "Warning" inside the health block we're emitting now
            # (rough — best-effort).
            n_err = (log_text.count("[ERROR") + log_text.count("Traceback")
                     + log_text.count("FAILED") + log_text.count("Exception:"))
            n_regress = log_text.count("[metrics-warn]")
            _warn_str = f"  Warnings in log:      {n_warn}"
            if n_warn > 0:
                _warn_str += "  (grep run.log for [WARN])"
            print(_warn_str)
            _err_str = f"  Errors in log:        {n_err}"
            if n_err > 0:
                _err_str += "  ← INVESTIGATE (grep run.log for [ERROR] / Traceback)"
            print(_err_str)
            if n_regress:
                print(f"  Metrics regressions:  {n_regress}  ← REVIEW (grep run.log for [metrics-warn])")
            print(f"  Run log:              {log_path}")
    except Exception as _e_log:
        print(f"  (log scan failed: {_e_log})")
    print("=" * 88)


try:
    _print_run_health_summary()
except Exception as _e_health:
    print(f"[health] summary print failed: {_e_health}")

# === Sentinel for the daily_auto wrapper =====================================
# The wrapper's Start-Process -Wait can hang on PyInstaller --noconsole
# child processes (Tk, matplotlib, Excel COM workers) even after the
# engine PID has exited. By dropping a flag file with the finish
# timestamp here, the wrapper can poll for a definitive "engine done"
# signal instead of relying solely on Process.WaitForExit().
#
# File is small (one JSON line) and overwritten each run. Wrapper
# compares mtime > its own start time to filter out stale sentinels.
try:
    import json as _json_done
    from datetime import datetime as _dt_done
    _flag_path = APP_DIR / "engine_done.flag"
    _flag_path.write_text(_json_done.dumps({
        "finished_at": _dt_done.now().isoformat(timespec="seconds"),
        "build_stamp": str(globals().get("BUILD_STAMP", "?")),
    }))
except Exception as _e_flag:
    print(f"[sentinel] engine_done.flag write failed: {_e_flag}")
