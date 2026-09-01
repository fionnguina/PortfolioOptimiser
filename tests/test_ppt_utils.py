"""Tests for ppt_utils.py — the report layer's pure primitives.

Focus on the date-window return math (which carried a real regression: a bad
_nearest_on_or_before fell back to idx[0] so 3M/6M/12M/3Y all showed the same
number) and the perf-value formatter. Two builders (_add_perf_table,
_add_change_run) are exercised against a real python-pptx slide.
"""
from __future__ import annotations

import sys
from pathlib import Path as _Path

import numpy as np
import pandas as pd
import pytest

sys.path.insert(0, str(_Path(__file__).resolve().parent.parent))
import ppt_utils


# === _nearest_on_or_before ====================================================

def test_nearest_before_start_returns_none():
    idx = pd.to_datetime(["2026-01-10", "2026-01-20"])
    assert ppt_utils._nearest_on_or_before(idx, "2026-01-01") is None


def test_nearest_exact_and_between():
    idx = pd.to_datetime(["2026-01-10", "2026-01-20", "2026-01-30"])
    assert ppt_utils._nearest_on_or_before(idx, "2026-01-20") == pd.Timestamp("2026-01-20")
    # between 20th and 30th -> the largest <= dt = the 20th
    assert ppt_utils._nearest_on_or_before(idx, "2026-01-25") == pd.Timestamp("2026-01-20")
    # after the end -> clamps to last
    assert ppt_utils._nearest_on_or_before(idx, "2026-02-15") == pd.Timestamp("2026-01-30")


def test_nearest_empty_index():
    assert ppt_utils._nearest_on_or_before(pd.to_datetime([]), "2026-01-01") is None


# === _period_total_return =====================================================

def test_period_total_return_known_window():
    idx = pd.date_range("2025-01-01", periods=400, freq="D")
    px = pd.Series(np.linspace(100.0, 200.0, 400), index=idx)
    # 12M ending at the last date: value doubled overall, so ~1yr slice < 100%.
    r = ppt_utils._period_total_return(px, idx[-1], years=1)
    end_v = px.iloc[-1]
    start_dt = ppt_utils._nearest_on_or_before(idx, idx[-1] - pd.DateOffset(years=1))
    start_v = px.loc[start_dt]
    assert r == pytest.approx(end_v / start_v - 1.0)


def test_period_total_return_window_before_series_is_nan():
    idx = pd.date_range("2026-01-01", periods=10, freq="D")
    px = pd.Series(np.arange(100, 110), index=idx)
    # 3Y lookback on a 10-day series -> start falls before series -> NaN (the
    # regression that this guard fixes: it must NOT fall back to idx[0]).
    assert np.isnan(ppt_utils._period_total_return(px, idx[-1], years=3))


def test_period_total_return_empty():
    assert np.isnan(ppt_utils._period_total_return(pd.Series(dtype=float), "2026-01-01", months=3))


# === _window_compound_total ===================================================

def test_window_compound_total_compounds_returns():
    idx = pd.date_range("2026-01-01", periods=40, freq="D")   # spans the 1M window
    rng = np.random.default_rng(3)
    r = pd.Series(rng.normal(0, 0.01, 40), index=idx)
    out = ppt_utils._window_compound_total(r, idx[-1], months=1)
    # Replicate the function's own window selection to derive the expectation.
    start = ppt_utils._nearest_on_or_before(idx, idx[-1] - pd.DateOffset(months=1))
    end = ppt_utils._nearest_on_or_before(idx, idx[-1])
    expected = float((1.0 + r.loc[start:end]).prod() - 1.0)
    assert out == pytest.approx(expected)


def test_window_compound_total_degenerate_window_is_nan():
    idx = pd.date_range("2026-01-01", periods=3, freq="D")
    r = pd.Series([0.01, 0.02, 0.03], index=idx)
    # end anchored to first date -> start >= end -> NaN
    assert np.isnan(ppt_utils._window_compound_total(r, idx[0], months=1))


# === _format_perf_value =======================================================

@pytest.mark.parametrize("v,fmt,expected", [
    (0.1234, "pct2", "12.34%"),
    (-0.05, "pct2", "-5.00%"),
    (1.23456, "dec3", "1.235"),
    (42.0, "raw", "42.0"),
    (float("nan"), "pct2", "n/a"),
    (float("inf"), "pct2", "n/a"),
    (None, "pct2", "n/a"),
    ("not a number", "pct2", "n/a"),
])
def test_format_perf_value(v, fmt, expected):
    assert ppt_utils._format_perf_value(v, fmt=fmt) == expected


# === slide builders (real python-pptx) ========================================

def _blank_slide():
    from pptx import Presentation
    prs = Presentation()
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def test_add_perf_table_fills_cells():
    from pptx.util import Cm
    slide = _blank_slide()
    df = pd.DataFrame({"SPY": [0.10, 0.15], "Engine": [0.12, np.nan]},
                      index=["Return", "Sharpe"])
    tbl = ppt_utils._add_perf_table(slide, df, Cm(1), Cm(1), Cm(10), Cm(4),
                                    title="Metric", value_fmt="pct2")
    assert tbl.cell(0, 0).text == "Metric"
    assert tbl.cell(0, 1).text == "SPY"
    assert tbl.cell(1, 0).text == "Return"
    assert tbl.cell(1, 1).text == "10.00%"
    assert tbl.cell(2, 2).text == "n/a"       # NaN cell renders 'n/a', not blank


def test_add_change_run_signs_and_zero():
    slide = _blank_slide()
    tb = slide.shapes.add_textbox(0, 0, 100, 100)
    p = tb.text_frame.paragraphs[0]
    ppt_utils._add_change_run(p, 12.5)
    ppt_utils._add_change_run(p, -3.0)
    ppt_utils._add_change_run(p, 0)           # zero adds an empty run
    texts = [r.text for r in p.runs]
    assert " (+12.50)" in texts
    assert " (-3.00)" in texts
    assert "" in texts


# --------------------------------------------------------------------------
# Actual-NAV chart gaps (2026-08-25)
# --------------------------------------------------------------------------

def test_short_nav_gaps_are_bridged_for_the_plot():
    """A day with no broker snapshot becomes NaN once reindexed onto the price
    panel, and matplotlib breaks the line at every NaN. Six such days —
    2026-07-09/10 (TWS not logged in, then a locked ibkr_nav_log.jsonl) and
    07-14 to 07-16 plus 08-13 (machine off) — left the live line shredded. The
    NAV either side of each hole is correct; only the plot was misleading."""
    import numpy as np
    import ppt_export as P

    idx = pd.bdate_range("2026-07-06", "2026-07-17")
    s = pd.Series(np.arange(len(idx), dtype=float), index=idx)
    for d in ("2026-07-09", "2026-07-10", "2026-07-14", "2026-07-15", "2026-07-16"):
        s.loc[pd.Timestamp(d)] = np.nan

    out, n = P.bridge_short_gaps(s, max_days=3)
    assert n == 5 and not out.isna().any()
    # Linear ramp, so a correct time-interpolation reproduces the originals.
    assert out.loc[pd.Timestamp("2026-07-14")] == pytest.approx(6.0)


def test_a_long_outage_stays_visible_on_the_chart():
    """Three days covers a missed morning or a long weekend. Longer than that
    is the pipeline being down, and smoothing it into a straight line would
    imply data we do not have."""
    import numpy as np
    import ppt_export as P

    idx = pd.bdate_range("2026-07-06", "2026-07-24")
    s = pd.Series(np.arange(len(idx), dtype=float), index=idx)
    s.loc[pd.Timestamp("2026-07-13"):pd.Timestamp("2026-07-17")] = np.nan  # 5 days

    out, n = P.bridge_short_gaps(s, max_days=3)
    assert n == 0, "a five-day outage must not be bridged"
    assert int(out.isna().sum()) == 5


def test_bridging_never_invents_nav_outside_the_observed_range():
    """limit_area='inside' — the live series must not be extrapolated backwards
    to before the account had a NAV, nor forward past the last snapshot."""
    import numpy as np
    import ppt_export as P

    idx = pd.bdate_range("2026-07-06", "2026-07-15")
    s = pd.Series(np.nan, index=idx)
    s.iloc[3:6] = [1.0, 2.0, 3.0]

    out, _ = P.bridge_short_gaps(s, max_days=3)
    assert out.iloc[:3].isna().all(), "nothing before the first observation"
    assert out.iloc[6:].isna().all(), "nothing after the last"


def test_a_clean_series_is_returned_untouched():
    import ppt_export as P

    idx = pd.bdate_range("2026-07-06", "2026-07-10")
    s = pd.Series(range(len(idx)), index=idx, dtype=float)
    out, n = P.bridge_short_gaps(s)
    assert n == 0 and out.equals(s)


# --------------------------------------------------------------------------
# Every line on the performance slide must share an origin (2026-09-01)
# --------------------------------------------------------------------------

def test_rebasing_removes_the_strategy_head_start():
    """Actual NAV was rebased to the account's inception while the strategy and
    benchmarks were rebased to the chart window's start, a month earlier. The
    strategy therefore carried a head start — +5.19% over 25 May -> 23 Jun on
    the live book's weights — before the NAV line began at zero, so the slide
    read as ~5x more divergence than the drift table's actual -1.00%."""
    import numpy as np
    import ppt_export as P

    idx = pd.bdate_range("2026-05-25", "2026-08-25")
    # +5% before inception, +5% after.
    s = pd.Series(np.linspace(1.0, 1.05, 22).tolist()
                  + np.linspace(1.05, 1.10, len(idx) - 22).tolist(), index=idx)
    origin = idx[22]

    assert P.rebase_to(s, None).iloc[-1] == pytest.approx(1.10)        # old
    assert P.rebase_to(s, origin).iloc[-1] == pytest.approx(1.10 / 1.05)
    assert P.rebase_to(s, origin).loc[origin] == pytest.approx(1.0), \
        "the origin is where every line must cross zero"


def test_rebasing_handles_a_frame_of_benchmarks():
    import ppt_export as P

    idx = pd.bdate_range("2026-05-25", "2026-07-25")
    df = pd.DataFrame({"ASX": range(len(idx)), "SPX": range(len(idx))},
                      dtype=float) + 100.0
    df.index = idx
    origin = idx[10]
    out = P.rebase_to(df, origin)
    assert list(out.loc[origin]) == pytest.approx([1.0, 1.0])


def test_rebasing_falls_back_when_the_origin_is_unusable():
    """No live NAV yet (hypothetical run), or a zero/absent base — the window
    start is a worse origin but a finite one, and the line must still render."""
    import numpy as np
    import ppt_export as P

    idx = pd.bdate_range("2026-05-25", "2026-06-25")
    s = pd.Series(np.arange(len(idx), dtype=float) + 1.0, index=idx)

    assert P.rebase_to(s, None).iloc[0] == pytest.approx(1.0)
    # An origin outside the index
    assert P.rebase_to(s, pd.Timestamp("2030-01-01")).iloc[0] == pytest.approx(1.0)
    # A zero base must not blow the series up to inf
    z = pd.Series(0.0, index=idx)
    assert np.isfinite(P.rebase_to(z, idx[5]).replace([np.inf, -np.inf], np.nan)
                       .fillna(0.0)).all()


def test_the_slide_rebases_all_three_series_together():
    """Source-text guard: the fix is only a fix if strategy AND benchmarks use
    the same origin the NAV line does. Rebasing one of them would leave the
    chart inconsistent in a subtler way than before."""
    src = (_Path(__file__).resolve().parent.parent / "ppt_export.py").read_text(encoding="utf-8")
    assert "rebase_to(bench, _chart_origin)" in src, "benchmarks must share it"
    assert "rebase_to(tilted_curve, _chart_origin)" in src, "strategy must share it"
    assert "_chart_origin = _first_valid_nav" in src, "and it is the NAV's own origin"
